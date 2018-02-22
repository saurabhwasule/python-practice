Advanced OOP features in Python
    * Python OO architectural overview
    * Python OO paradigms and metaprogramming concepts
    * An overview on Python special methods:
        - __new__(), __init__(), __del__(),
        - __str__(), __repr__(), __len__(), __nonzero__()
        - __cmp__(), __eq__(), __ne__() and family
        - Operator overloading
        - Dictionary and List emulation
        - Accessor methods - __getattr__(), __setattr__(), 
                             __delattr__()
                             __delattr__()
        - Callable methods - __call__()
        - Iterator methods - __iter__()

Common design patterns with examples and exercises
    * Singleton/Borg patterns
    * Factory method
    * Factory pattern
    * Accessor pattern
    * Command pattern
    * Event-Delegate pattern
    * Decorator patterns and built-in decorators
    
Database connectivity using DBAPI 2.0 
[Examples using sqlite3 module]
    * Connecting to database
    * Database connection object methods
    * Creating cursor objects for DML queries
        - Preparing and executing queries
        - Fetching rows from cursor
    * Managing transactions
    * Best practices and idioms 

Network programming
    * An overview of the socket module
    * Creating network servers and clients using socket module
    * Creating network applications using multiprocessing.Listener, multiprocessing.Client and multiprocessing.Connection objects
    * Creating asynchronous Socket server and clientsusing asyncore module
    * Using the SocketServer framework for creating scalable network servers
    * Introductory overview on Twisted Matrix framework

Hands On
    * Creating a simple multi-tasking fileserver that can serve files to client, supporting multiple concurrent client connections.
    * Simple examples on using ftplib and telnetlib
    * Simple examples for SSH automation
    * Sending and receiving e-mails using python script

Day 2
Automating processes and CLI
    * Using the subprocess module
    * Using pexpect module

Network automation modules
    * FTP automation using ftplib
    * Telnet automation using telnetlib
    * Automating HTTP using urllib, urllib2 and httplib
    * Mail handling using smtplib and poplib
    * Using paramiko for SSH automation
    * Using Exscript and fabric modules

Web programming and automation using Python
    * User-Agent automation using urllib and urllib2
    * An overview on CGI programming using cgi module
    * An overview on WSGI 
    * Using the wsgi module
        - Scheme handlers
        - HTTP status, request and response setup
        - Request and Response headers
        - CGIHandler, BaseHandler, SimpleHandler
    * Using Bottle and cherrypy for basic web application development.
    * Parsing HTML using HTMLParser
    * Parsing XML using ElementTree XML API

Hands On
    * A simple web crawler program that can crawl a website and report broken links on the same by email.
    * A simple web application and collects and reports latest news updates using RSS feed of a popular news site

	
	
-----------------------------------------------------------------------------------------------------------------------------------------
###Making Python2.x or py3.x equivalent to Py3.x
#(applicable from py2.1) 

#Use import from __future__
#Example, do below and in Py2.7, you get- print function:  print(*objects, sep=' ', end='\n', file=sys.stdout)
from __future__ import print_function

#List of features , Note in python2.7, all Py3.x syntax is possible with 
from __future__ import print_function, division, unicode_literals

#Name 			 Introduced in 	Description
nested_scopes 	 2.2			introducing E in LEGB resolution
generators 		 2.3 			Generators ie ( .. for x in ..)
with_statement 	 2.6 			Introducing with Statement 

division 		 3.0 			Introducing  // for int div, / for true div
print_function   3.0 			Make print a function 
unicode_literals 3.0 			Introducing Bytes literals  , b''
absolute_import  3.0 			Introducing () in import, 
                                from Tkinter import (Tk, Frame....) #can be multiline
								'import package' is always absolute import 
								Introducing . as relative import 'from' syntax

##Relative import 
#A single leading dot indicates a relative import, 
#starting with the current package. 

#Two or more leading dots give a relative import to the parent(s) of the current package,

#Example 
package/
    __init__.py
    subpackage1/
        __init__.py
        moduleX.py
        moduleY.py
    subpackage2/
        __init__.py
        moduleZ.py
    moduleA.py


#Current file is moduleX.py :
 
from .moduleY import spam
from .moduleY import spam as ham
from . import moduleY
from ..subpackage1 import moduleY
from ..subpackage2.moduleZ import eggs
from ..moduleA import foo
from ...package import bar



###Meaning of __init__.py

#Ex:
package/
    __init__.py
    file.py
    file2.py
    subpackage/
        __init__.py
        submodule1.py
        submodule2.py


##USAGE-1: Normally, use import like
from package.file import File

#if package/__init__.py contains below 
from file import File

#Can use like below in user file
from package import File

##USAGE-2: when you do below
from package import *

#if  package/__init__.py contains below 
__all__ = ['file', 'file2']
#and subpackage/__init__.py
__all__ = ['submodule1', 'submodule2']   # list of module names(ie file name) that should be imported 

#it imports all modules from __all__

#Example 
>>>from package import * 
>>>file.function_X()
>>>from package.subpackage import * 
>>>submodule1.function_X()


##USAGE-3:Any class, methods defined in __init__.py would be automatically available 
#when user import the package 
#An example
database/
    __init__.py
    schema.py
   ...

#__init__.py
import os

from sqlalchemy.orm import sessionmaker
from sqlalchemy import create_engine

engine = create_engine(os.environ['DATABASE_URL'])
Session = sessionmaker(bind=engine)

def method():
    pass
    
class C:
    pass

#User can do below
from database import Session, method, C
session = Session()

#Usage-4:__init__.py can  contain anything as it itself is a py file, 
#But below are generally used for special purpose(not so standardised)
__version__ = '0.1'
__author__ = 'Cardinal Biggles'

>>> import email
>>> email.__version__
'5.1.0'







###Namespace package (>Py3.3)
#Usage is exactly same as normal package , but can span on multiple dirs 


#While looking for 'import foo' 
#for each directory, <directory> in the parent path(from sys.path): 
1. If  <directory>/foo/__init__.py   is found, a regular package is imported and returned. 
2. If not, but  <directory>/foo.{py,pyc,so,pyd}   is found, a module is imported and returned. 
3. If not, but  <directory>/foo   is found and is a directory, it is recorded 
   and the scan continues with the next directory in the parent path. 
4. Otherwise the scan continues with the next directory in the parent path. 

#If the scan completes without returning a module or package, 
#and at least one directory was recorded, 
#then a namespace package is created with 
1. __path__  attribute set to an iterable of the path strings that were found above 
2. without  __file__  attribute. 

#Namespace packages and regular packages are very similar. 
#The differences are: 
1. Portions of namespace packages need not all come from the same directory structure
2. Namespace packages have no  __file__  attribute. 
3. Namespace packages'  __path__  attribute is a read-only iterable of strings, 
   which is automatically updated when the parent path is modified. 
4. Namespace packages have no  __init__.py  module. 
5. Namespace packages have a different type of object for their  __loader__  attribute. 


#Example - Nested namespace packages  
#directory structure: Note no __init__.py
Lib/test/namespace_pkgs
                project1
                    parent
                        child
                            one.py
                project2
                    parent
                        child
                            two.py


import sys
sys.path += ['Lib/test/namespace_pkgs/project1', 'Lib/test/namespace_pkgs/project2']

import parent.child.one  #namespace package 
parent.__path__       #_NamespacePath(['Lib/test/namespace_pkgs/project1/parent', 'Lib/test/namespace_pkgs/project2/parent'])
parent.child.__path__ #_NamespacePath(['Lib/test/namespace_pkgs/project1/parent/child', 'Lib/test/namespace_pkgs/project2/parent/child'])
import parent.child.two
#...

## Dynamic path computation  
#directory structure: Note no __init__.py
Lib/test/namespace_pkgs
                project1
                    parent
                        child
                            one.py
                project2
                    parent
                        child
                            two.py
                project3
                    parent
                        child
                            three.py


# add the first two parent paths to sys.path
import sys
sys.path += ['Lib/test/namespace_pkgs/project1', 'Lib/test/namespace_pkgs/project2']

# parent.child.one can be imported, because project1 was added to sys.path:
import parent.child.one
parent.__path__  #_NamespacePath(['Lib/test/namespace_pkgs/project1/parent', 'Lib/test/namespace_pkgs/project2/parent'])

# parent.child.__path__ contains project1/parent/child and project2/parent/child, but not project3/parent/child:
parent.child.__path__ #_NamespacePath(['Lib/test/namespace_pkgs/project1/parent/child', 'Lib/test/namespace_pkgs/project2/parent/child'])

# parent.child.two can be imported, because project2 was added to sys.path:
import parent.child.two

# we cannot import parent.child.three, because project3 is not in the path:
import parent.child.three   #ImportError: No module named 'parent.child.three'

# now add project3 to sys.path:
sys.path.append('Lib/test/namespace_pkgs/project3')

# and now parent.child.three can be imported:
import parent.child.three

# project3/parent has been added to parent.__path__:
parent.__path__  #_NamespacePath(['Lib/test/namespace_pkgs/project1/parent', 'Lib/test/namespace_pkgs/project2/parent', 'Lib/test/namespace_pkgs/project3/parent'])

# and project3/parent/child has been added to parent.child.__path__
parent.child.__path__  #_NamespacePath(['Lib/test/namespace_pkgs/project1/parent/child', 'Lib/test/namespace_pkgs/project2/parent/child', 'Lib/test/namespace_pkgs/project3/parent/child'])
























###Executing modules as scripts

def fib(n):    # write Fibonacci series up to n
    a, b = 0, 1
    while b < n:
        print(b, end=' ')
        a, b = b, a+b
    print()

if __name__ == "__main__":       #this is true when module is executed as script 
    import sys
    fib(int(sys.argv[1]))

	

###Advanced OOP features in Python

##Class- Name mangling 
#all attributes of class are stored in instance.__dict__
#hence only one variable name as key is possible 

class C1:
	def meth1(self): self.X = 88 	
	def meth2(self): print(self.X)

class C2:
	def metha(self): self.X = 99 	
	def methb(self): print(self.X)

class C3(C1, C2): pass

>>> c = C3()
>>> c.meth1()
>>> c.__dict__
{'X': 88}
>>> c.metha()
>>> c.__dict__
{'X': 99}			# replaces original X


##Solution - Name X as __X, becomes _CLASS__X

class C1:
	def meth1(self): self.__X = 88 	
	def meth2(self): print(self.__X)

class C2:
	def metha(self): self.__X = 99 	
	def methb(self): print(self.__X)

class C3(C1, C2): pass

>>> c = C3()
>>> c.meth1()
>>> c.__dict__
{'_C1__X': 88}
>>> c.metha()
>>> c.__dict__
{'_C2__X': 99, '_C1__X': 88}
>>> c.X
Traceback (most recent call last):
  File "<stdin>", line 1, in <module>
AttributeError: 'C3' object has no attribute 'X'
>>> c._C2__X
99
>>> c.meth2()
88
>>> c.methb()
99



## Dictionary and List emulation
'''
Comparisons in Python 2.X use specific methods such as __lt__ for "less than" if
present, or else the general __cmp__. 
Python 3.X uses only specific methods, not __cmp__
'''


object.__len__(self)   
#For bool(obj), if __bool__(self) is not defined 
#and len(self) gives zero, it is False
# for py2.x , __bool__ is named as __nonzero__

#slicing
a[1:2] = b
#is translated to
a[slice(1, 2, None)] = b

#Internal Methods
object.__getitem__(self, key)
object.__setitem__(self, key, value)
object.__delitem__(self, key)
'''
for self[key], raises KeyError for non existing key
for List type , key can be int or slice object, can raise IndexError, TypeError
slice has following methods - indices, start , step, stop
'for' loops expect that an IndexError will be raised for illegal indexes 
'''

object.__missing__(self, key)
#Called by dict.__getitem__() to missing  key 


object.__iter__(self)
#For iterator, must return a  new iterator object , 
#for map, return keys iterator

object.__contains__(self, item)
#automatic implementation is provided for 'in' and 'not in'
#or provide this method


#Py3.x- index could be int or slice

class Indexer:
	def __init__(self):
		self.data = [5, 6, 7, 8, 9]
	def __getitem__(self, index):
		if isinstance(index, int): 
			print('indexing', index)
			return self.data[index]
		else:
			print('slicing', index.start, index.stop, index.step)
			return self.data[index]

	def __setitem__(self, index, value): # Intercept index or slice assignment
		self.data[index] = value 		 # Assign index or slice
			
			
#2.x
# use __getitem__ and __setitem__ 
#allow for both indexes and slice objects as arguments
#Can use __getslice__ for slice(no 'step'), but  is removed in 3.x

class Slicer:
	def __getitem__(self, index): print index
	def __getslice__(self, i, j): print i, j
	def __setslice__(self, i, j,seq): print i, j,seq
	
	
>>> Slicer()[1] # Runs __getitem__ with int, like 3.X
1
>>> Slicer()[1:9] # Runs __getslice__ if present, else __getitem__
1 9
>>> Slicer()[1:9:2] # Runs __getitem__ with slice(), like 3.X!
slice(1, 9, 2)		

#__index__ method in Python 3.X is for int, hex, bin, oct interception
#In order to have a coherent integer type class, when __index__() is defined 
#__int__() should also be defined, and both should return the same value

#__index__ is used to convert the numeric object to an int object 
#(such as in slicing, or in the built-in bin(), hex() and oct() functions). 
# Must return an integer.

class C:
	def __index__(self):
		return 255

>>> X = C()
>>> hex(X) 	
'0xff'
>>> bin(X)
'0b11111111'
>>> oct(X)
'0o377'


#is also used in contexts that require an integer—including indexing
>>> ('C' * 256)[255]
'C'
>>> ('C' * 256)[X] # X is converted to int via X.__index__
'C'
>>> ('C' * 256)[X:] # As index (not X[i:])
'C'

#in py2.x, same usage but  use __hex__ and __oct__ for other



#other uses of __getitem__

#if __iter__ is not defined,
#'for' statement works by repeatedly indexing a sequence from zero to higher indexes,
#until an out-of-bounds IndexError exception is detected

class StepperIndex:
	def __getitem__(self, i):
		return self.data[i]
		
>>> X = StepperIndex() 	# X is a StepperIndex object
>>> X.data = "Spam"
>>>
>>> X[1] 				# Indexing calls __getitem__
'p'
>>> for item in X: 				# for loops call __getitem__
		print(item, end=' ') 	# for indexes items 0..N
S p a m

>>> 'p' in X # All call __getitem__ too
True

>>> [c for c in X] # List comprehension
['S', 'p', 'a', 'm']

>>> list(map(str.upper, X)) 		# map calls (use list() in 3.X)
['S', 'P', 'A', 'M']

>>> (a, b, c, d) = X # Sequence assignments
>>> a, c, d
('S', 'a', 'm')

>>> list(X), tuple(X), ''.join(X) 		# And so on...
(['S', 'p', 'a', 'm'], ('S', 'p', 'a', 'm'), 'Spam')








###collections.abc — Abstract Base Classes for Containers
#Py3.x
#ABC		Inherits from		Abstract Methods				Methods you get automatically
Container   					__contains__   
Hashable   						__hash__   
Iterable   						__iter__   
Iterator 	Iterable 			__next__                        __iter__ 
Reversible  Iterable            __reversed__   
Generator 	Iterator 			send, throw 					close, __iter__, __next__ 


Sized   						__len__   
Callable   						__call__   

Sequence 	Sized, Iterable, 
			Container 			__getitem__, __len__ 			__contains__, __iter__, __reversed__, index, and count 
MutableSequence 
			Sequence 			__getitem__, __setitem__, 
								__delitem__, __len__,insert 	Inherited Sequence methods 
                                                                and append, reverse, extend, pop, remove, and __iadd__ 
Set 		Sized, Iterable, 
			Container 			__contains__, __iter__, 
								__len__ 						__le__, __lt__, __eq__, __ne__, 
                                                                __gt__, __ge__, __and__, __or__, 
                                                                __sub__, __xor__, and isdisjoint 
                                                                
MutableSet 	Set 				__contains__, __iter__, 
								__len__, add, discard 			Inherited Set methods 
                                                                and clear, pop, remove, __ior__, __iand__, __ixor__, and __isub__ 
Mapping 	Sized, Iterable, 
			Container 			__getitem__, __iter__,
								__len__ 						__contains__, keys, items, 
                                                                values, get, __eq__, and __ne__ 
MutableMapping 
			Mapping 			__getitem__, __setitem__, 
								__delitem__, __iter__, 
								__len__ 						Inherited Mapping methods 
                                                                and pop, popitem, clear, update, and setdefault 


#ABC		Inherits from		Abstract Methods				Methods you get automatically
MappingView
            Sized                                               __len__ 
ItemsView 
            MappingView, Set                                    __contains__, __iter__ 
KeysView 
            MappingView, Set                                    __contains__, __iter__ 
ValuesView 
            MappingView                                         __contains__, __iter__ 
Awaitable                        __await__   
Coroutine 
            Awaitable           send, throw                     close 
AsyncIterable                   __aiter__   
AsyncIterator 
            AsyncIterable      __anext__                        __aiter__ 
AsyncGenerator 
            AsyncIterator       asend, athrow                   aclose, __aiter__, __anext__ 
	


##Usage  
#These ABCs allow us to ask classes or instances 
#if they provide any particular functionality, 

#For example if myvar implements __len__ ,ie collections.abc.Sized
size = None
if isinstance(myvar, collections.abc.Sized):
    size = len(myvar)

    
    
##Usage as  as mixins 
#eg implement  __contains__(), __iter__(), and __len__()
#and get Set functionality 

#ABC		Inherits from		Abstract Methods				Methods you get automatically
Set 		Sized, Iterable, 
			Container 			__contains__, __iter__, 
								__len__ 						__le__, __lt__, __eq__, __ne__, 
                                                                __gt__, __ge__, __and__, __or__, 
                                                                __sub__, __xor__, and isdisjoint 
#Example 
import collections.abc

class ListBasedSet(collections.abc.Set):
     ''' Alternate set implementation favoring space over speed
         and not requiring the set elements to be hashable. '''
     def __init__(self, iterable):
         self.elements = lst = []
         for value in iterable:
             if value not in lst:
                 lst.append(value)
     def __iter__(self):
         return iter(self.elements)
     def __contains__(self, value):
         return value in self.elements
     def __len__(self):
         return len(self.elements)

s1 = ListBasedSet('abcdef')
s2 = ListBasedSet('defghi')
overlap = s1 & s2            # The __and__() method is supported automatically

#To add set hashability using mixins, 
#inherit from both Set() and Hashable(), 
#then define __hash__ = Set._hash
#(3.The Set mixin provides a _hash() method to compute a hash value for the set; )


##Example -Set that remembers original insertion order.
import collections

class OrderedSet(collections.MutableSet):

    def __init__(self, iterable=None):
        self.end = end = [] 
        end += [None, end, end]         # sentinel node for doubly linked list
        self.map = {}                   # key --> [key, prev, next]
        if iterable is not None:
            self |= iterable

    def __len__(self):
        return len(self.map)

    def __contains__(self, key):
        return key in self.map

    def add(self, key):
        if key not in self.map:
            end = self.end
            curr = end[1]
            curr[2] = end[1] = self.map[key] = [key, curr, end]

    def discard(self, key):
        if key in self.map:        
            key, prev, next = self.map.pop(key)
            prev[2] = next
            next[1] = prev

    def __iter__(self):
        end = self.end
        curr = end[2]
        while curr is not end:
            yield curr[0]
            curr = curr[2]

    def __reversed__(self):
        end = self.end
        curr = end[1]
        while curr is not end:
            yield curr[0]
            curr = curr[1]

    def pop(self, last=True):
        if not self:
            raise KeyError('set is empty')
        key = self.end[1][0] if last else self.end[2][0]
        self.discard(key)
        return key

    def __repr__(self):
        if not self:
            return '%s()' % (self.__class__.__name__,)
        return '%s(%r)' % (self.__class__.__name__, list(self))

    def __eq__(self, other):
        if isinstance(other, OrderedSet):
            return len(self) == len(other) and list(self) == list(other)
        return set(self) == set(other)

            
if __name__ == '__main__':
    s = OrderedSet('abracadaba')
    t = OrderedSet('simsalabim')
    print(s | t)
    print(s & t)
    print(s - t)



###Module - Types
#Py2 and Py3 
#Types module contains various types 
#Typical use of these names is for isinstance() or issubclass() checks.

#Example checking for int or list 
def delete(mylist, item):
    if isinstance(item, int):
       del mylist[item]
    else:
       mylist.remove(item)

       

#It contains below types
#CodeType, FrameType, FunctionType,GeneratorType, 
#LambdaType, ModuleType
#types.CoroutineType, types.AsyncGeneratorType, types.TracebackType
#types.BuiltinFunctionType, types.BuiltinMethodType
#types.MethodType
    The type of methods of user-defined class instances.


#Usage 
import types
type(attribute) is types.FunctionType

#In general, for built-in types, we use below way
isinstance(attribute, int)



###Callable class instance 
#define  - def __call__(self, [args ...])

class A:
	def __init__(self):
		print "init"
	def __call__(self):
		print "call"


>>> A()
init
>>> A()()
init
call




###Meta Class 

#when Class is defined, 
#Python calls Meta Class's(default is type) __new__(meta as first arg)  for allocation 
#(default implemention)type's __new__ calls Meta Class's __init__(class as first arg) to init the class

#When a class instance is created, Py calls metaclass's(default is type) __call__
#(default implemention)type's __call__ calls class's __new__(class as first arg)
#(default implemention)class's __new__ is calls object's __new__(class as first arg)
#(default implemention)object's __new__ calls class's __init__(instance as first arg)


#m is metaclass as inherited from type 
class m(type):
	def __new__(meta, c, s, cd):  #meta, classname, supers, classdict #at class definition
		print("meta.__new__")
		print(meta, c, s, cd)
		return type.__new__(meta,c,s,cd)  #calls meta.__init__
	def __call__(*args, **kargs):    #calls at instance creation
		print("meta.__call__")
		return type.__call__(*args, **kargs)
	def __init__ (c, cn, s, cd):   #class, classname, supers, classdict #class definition
		print("meta.__init__")
		print(c,cn,s,cd)
		return type.__init__(c,cn,s,cd)
	
	
class A(metaclass=m):   #object py2.x
	#__metaclass__ = m  #py2.x
	def __new__(cls):           #called at instance creation
		print("A's new")
		return object.__new__(cls)  #calls self.__init__
	def __init__(self):
		print("A's init")
	def __call__(self, *args, **kargs):
		print("A's call")


#above instantly prints below
m.__new__
<class '__main__.m'> A () {'__module__': '__main__', '__qualname__': 'A'}
m.__init__
<class '__main__.A'> A () {'__module__': '__main__', '__qualname__': 'A'}

a = A()
#prints below
meta.__call__
A's new
A's init

a()
#prints below
A's call

##Creation of Meta Class inheritance

class mm(type):
	def __new__(meta, c, s, cd):
		print("mm.__new__")
		print(meta, c, s, cd)
		return type.__new__(meta,c,s,cd)
	def __call__(*args, **kargs):
		print("mm.__call__")
		print(*args, **kargs)
		return type.__call__(*args, **kargs)
	def __init__ (c, cn, s, cd):
		print("mm.__init__")
		print(c,cn,s,cd)
		return type.__init__(c,cn,s,cd)


#when m is defined, mm.__new__ and mm.__init__ are called 		

		
class m(type, metaclass = mm):
	def __new__(meta, c, s, cd):
		print("m.__new__")
		print(meta, c, s, cd)
		return type.__new__(meta,c,s,cd)
	def __call__(*args, **kargs):
		print("m.__call__")
		return type.__call__(*args, **kargs)
	def __init__ (c, cn, s, cd):
		print("m.__init__")
		print(c,cn,s,cd)
		return type.__init__(c,cn,s,cd)
		

		
#prints below
mm.__new__
<class '__main__.mm'> m (<class 'type'>,) {'__init__': <function m.__init__ at 0x0225C228>, '__new__
': <function m.__new__ at 0x0225C978>, '__module__': '__main__', '__call__': <function m.__call__ at
 0x0225C858>, '__qualname__': 'm'}
 
mm.__init__
<class '__main__.m'> m (<class 'type'>,) {'__init__': <function m.__init__ at 0x0225C228>, '__new__'
: <function m.__new__ at 0x0225C978>, '__module__': '__main__', '__call__': <function m.__call__ at
0x0225C858>, '__qualname__': 'm'}


#Now for class definition 'class A(..)', calls  mm's __call__  
#type's __call__  calls m's __new__ and __init__
	
class A(metaclass=m):pass

#prints 
mm.__call__
<class '__main__.m'> A () {'__module__': '__main__', '__qualname__': 'A'}
m.__new__
<class '__main__.m'> A () {'__module__': '__main__', '__qualname__': 'A'}
m.__init__
<class '__main__.A'> A () {'__module__': '__main__', '__qualname__': 'A'}


class B(A): pass

mm.__call__
<class '__main__.m'> B (<class '__main__.A'>,) {'__module__': '__main__', '__qualname__': 'B'}
m.__new__
<class '__main__.m'> B (<class '__main__.A'>,) {'__module__': '__main__', '__qualname__': 'B'}
m.__init__
<class '__main__.B'> B (<class '__main__.A'>,) {'__module__': '__main__', '__qualname__': 'B'}

#Now 'B()' instance creation, calls m's __call__
#type's __call__  calls Class's __new__ and __init__

b = B()
#prints
m.__call__



##MetaClass with Attribute accessing
#class variable's MRO - class and then class's __class__.__mro__ ie metaclasses hierarchy
#instance variable's MRO - instance and then instance's __class__.__mro__ ie class hierarchy

class M(type):pass
class A(metaclass = M):pass
class B(A): pass

b = B()

B.a1 = 100
A.a2 = 200
M.a3 = 300

b.a1, b.a2, b.a3  #100, 200, Error  # does not  search Metaclass for instance's 
B.a1, B.a2, B.a3  #100,200,300      # searches metaclass for class's


# with instance variable

class M(type):
	def __new__(m,cn,s,cd):
		m.m1 = 10
		return type.__new__(m,c,s,cd)
	def __init__(c,cn,s,d):
		c.ma1 = 20
		

class A(metaclass = M):
	def __init__ (self):
		self.ia1 = 30

	

class B(A): 
	def __init__ (self):
		self.ia2 = 40
		super().__init__()
	

b = B()

B.a1 = 100
A.a2 = 200
M.a3 = 300

b.a1, b.a2, b.a3  			# 100, 200, Error # does not  search Metaclass for instance's 

b.ia2, b.ia1, b.ma1, b.m1  	# 40, 30, 20, Error # does not  search Metaclass for instance's 

B.a1, B.a2, B.a3  			#100,200,300

B.ia2, B.ia1, B.ma1, B.m1  	# Error, Error, 20, 10



##Usage-Introducing a new function for all classes 

def func4(self): return self.value * 4


class Extender(type):
	def __new__(meta, classname, supers, classdict):
		classdict['func4'] = func4					#adding func4 method
		return type.__new__(meta, classname, supers, classdict)



class A(metaclass=Extender):
	def __init__(self, value): self.value = value
	def func2(self): return self.value * 2

	
a = A(1)
a.func4()


##Usage- Tracing
def trace(f):
	def inner(*args, **kargs):
		print("-->" + f.__name__)
		res = f(*args, **kargs)
		print("<--" + f.__name__)
		return res
	return inner

	

#Meta Class 

class TracingO(type):
	def __new__(meta, classname, supers, classdict):
		classdict['func2'] = trace(classdict['func2'])
		return type.__new__(meta, classname, supers, classdict)

	
	

class B(metaclass=TracingO):
	def __init__(self, value): self.value = value
	def func2(self): return self.value * 2


	
a = B(2)
a.func2()



##Usage - Adding tracing to all functions

import types

class Tracing(type):
	def __new__(meta, classname, supers, classdict):
		classdict = {k:(trace(v) if type(v) is types.FunctionType and v.__name__.startswith("fun") else v ) for k,v in classdict.items()}
		return type.__new__(meta, classname, supers, classdict)



class B(metaclass=Tracing):
	def __init__(self, value): self.value = value
	def func(self): return self.value * 2



	



		
		
		



@@@
###Advanced-  Accessor methods 
__getattr__(), __setattr__(),   __delattr__()
'''
-In Python 3.X, all classes are automatically 'new style'
whether they explicitly inherit from object or not.

-In Python 2.x, you need to inheritate from object to get "New Style", eg 
class MyClass(object) :

Properties of 'new style'

	object as root class

	Attribute/special method Fetch for an instance of class and for class

case 1. Instance attribute access – Search below
a. The __dict__ of the instance I
b. The __dict__ of all classes on the __mro__ found at I’s __class__, from left to right  

case 2. Class attribute access – Search Below
a. The __dict__ of all classes on the __mro__ found at C itself, from left to right
b. The __dict__ of all metaclasses on the __mro__ found at C’s __class__, from left to right

4. In both rule 1 and 2, skip step (a) 
for built-in  implicit operations ( eg __X__ for  str, len , [], in )
ie searches only class (case 1) or metaclass(case 2)


5. Built-in  implicit operations ( eg __X__ for  str, len , [], in )(Py3.x):
 __getattr__ (if __X__ undefined)and __getattribute__ (for all read access)
of class(case 1) or metaclass(case 2 ) are not called  
but not skipped for explicit method(via .__X__(..)) invocation

Hence required operator overloading methods for builtins 
must be implemented at class(case 1) and metaclass(case 2)
'''

def __getattr__(self, name): 		# On undefined attribute fetch [obj.name], for old and new style
def __getattribute__(self, name): 	# On all attribute fetch [obj.name], new style
def __setattr__(self, name, value): # On all attribute assignment [obj.name=value]
def __delattr__(self, name): 		# On all attribute deletion [del obj.name]

##builtins method to access above 
getattr(object, name[, default]) -> value  # object.name 
setattr(object, name, value)               #object.name = value
delattr(object, name)					   # del object.name


class Meta(type):
	def __getattribute__(*args):
		print("Metaclass getattribute invoked")
		return type.__getattribute__(*args)


class C(object, metaclass=Meta):
	def __len__(self):
		return 10
	def __getattribute__(*args):
		print("Class getattribute invoked")
		return object.__getattribute__(*args)

c = C()
>>> c.__len__()                 # Explicit lookup via instance, class's
Class getattribute invoked      #__getattribute__ called 
10

>>> type(c).__len__(c)          # Explicit lookup via type, metaclass's getttribute
Metaclass getattribute invoked
10
>>> len(c)                      # Implicit lookup , __getattribute__ of class skipped
10

>>> len(C)						# Implicit lookup , __getattribute__ of metaclass skipped
Traceback (most recent call last):
  File "<stdin>", line 1, in <module>
TypeError: object of type 'Meta' has no len()

>>> bool(c)               #from object's
True

>>> c < C()                        #Implicit lookup , __getattribute__ of class skipped
Traceback (most recent call last):
  File "<stdin>", line 1, in <module>
TypeError: unorderable types: C() < C()

>>> c.__lt__(C())                 # Explicit lookup via instance, class's getttribute
Class getattribute invoked
NotImplemented

>>> C.__lt__(C,C)                 # Explicit lookup via type, metaclass's getttribute
Metaclass getattribute invoked
NotImplemented

>>> c.__len__ = lambda self : 3    #instance's

>>> len(c)						   # only class's one called, instance skipped
10

>>> c.__len__()                    # if you add from outside, you have to pass 'self' explicitly  
Class getattribute invoked
Traceback (most recent call last):
  File "<stdin>", line 1, in <module>
TypeError: <lambda>() missing 1 required positional argument: 'self'

>>> c.__len__(c)                  #class's getattribute gets the calls, but self has to be explicit
Class getattribute invoked
3

>>> getattr(c, '__len__')       #Explicit lookup via instance, class's getttribute
Class getattribute invoked
<function <lambda> at 0x0239C4B0>

>>> getattr(c, '__lt__')  	#Explicit lookup via instance, class's getttribute
Class getattribute invoked
<method-wrapper '__lt__' of C object at 0x02361F90>

>>> getattr(c, '__lt__')(c)
Class getattribute invoked
NotImplemented




###pickle — Python object serialization
#a Python object hierarchy is converted into a byte stream, 
#and 'unpickling' is the inverse operation

#Important methods
#fix_imports must be true and protocol =2  for compatibility with python2
#Protocol can be 3 (default one) for new proptocol of Python3

pickle.dump(obj, file, protocol=None, *, fix_imports=True)  #dumps obj to file which has .write() method eg file or io.BytesIO 
pickle.dumps(obj, protocol=None, *, fix_imports=True)       #dumps obj to bytes string
class pickle.Pickler(file, protocol=None, *, fix_imports=True)
    dump(obj)
        Write a pickled representation of obj to the open file object given in the constructor.


#unpickling
pickle.load(file, *, fix_imports=True, encoding="ASCII", errors="strict") #file object 
pickle.loads(bytes_object, *, fix_imports=True, encoding="ASCII", errors="strict")  #reads from bytes 
class pickle.Unpickler(file, *, fix_imports=True, encoding="ASCII", errors="strict")
    load()
        Read a pickled object representation from the open file object given in the constructor, and return the reconstituted object 


>>> a = 100.2
>>> b = pickle.loads(pickle.dumps(a))
>>> b == a
True


#The following types can be pickled:
#When functions, classes are pickled, and then unpickled, 
#the module defining class/functions must be be imported to get their definition

#While class is pickled, only instance data is pickled, 
#neither class variables nor class method definition

1.None, True, and False
2.integers, floating point numbers, complex numbers
3.strings, bytes, bytearrays
4.tuples, lists, sets, and dictionaries containing only picklable objects
5.functions defined at the top level of a module (using def, not lambda)
6.built-in functions defined at the top level of a module
7.classes that are defined at the top level of a module
8.instances of such classes whose __dict__ 
  or the result of calling __getstate__() is picklable 



#Default behaviour of class instance for pickle and unpickle are  Ok for most of the cases
#but can be customized by 
object.__getstate__()		
    if the class defines the method __getstate__(), it is called 
    and the returned object is pickled as the contents for the instance,
    instead of the contents of the instance’s dictionary

object.__setstate__(state)
    Upon unpickling, if the class defines __setstate__(), 
    it is called with the unpickled state

#Example:
class TextReader:
    """Print and number lines in a text file."""
    def __init__(self, filename):
        self.filename = filename
        self.file = open(filename)
        self.lineno = 0
    def readline(self):
        self.lineno += 1
        line = self.file.readline()
        if not line:
            return None
        if line.endswith('\n'):
            line = line[:-1]
        return "%i: %s" % (self.lineno, line)
    def __getstate__(self):
        # Copy the object's state from self.__dict__ which contains
        # all our instance attributes. Always use the dict.copy()
        # method to avoid modifying the original state.
        state = self.__dict__.copy()
        # Remove the unpicklable entries.
        del state['file']
        return state
    def __setstate__(self, state):
        # Restore instance attributes (i.e., filename and lineno).
        self.__dict__.update(state)
        # Restore the previously opened file's state. To do so, we need to
        # reopen it and read from it until the line count is restored.
        file = open(self.filename)
        for _ in range(self.lineno):
            file.readline()
        # Finally, save the file.
        self.file = file


#usage 
>>> reader = TextReader("hello.txt")
>>> reader.readline()
'1: Hello world!'
>>> reader.readline()
'2: I am line number two.'
>>> new_reader = pickle.loads(pickle.dumps(reader))
>>> new_reader.readline()
'3: Goodbye!'



###Context Management Pattern - to be used with 'with'
#same in py2.7 and Py3.x


from __future__ import print_function  #py2.x
from __future__ import with_statement  #py2.x

class A(object):
	def __init__(self, name):
		self.name = name
	def hello(self):
		print('hello %s!' % (self.name,))
	def __enter__(self):
		print('Enter the function')
		return self   #must return self or a class which implements __exit__
	def __exit__(self, exc_type, exc_value, traceback):
		print('Exit the function')


with A("das") as a:
	print(a.hello())
	




###Decorator patterns - same in py2.7 and py3.x

def ConvertString(org):
	def _inner(*args, **kargs):
		res = org(*args, **kargs)
		return str(res)
	return _inner

	
@ConvertString
def f(x,y):
	return x+y


##stacking of decorator
def  mockReturnConstant(fun):
	def _mydec(*args, **kargs):		
		res = fun(*args, **kargs)		
		return 10
	return _mydec
	
@ConvertString        #last application's decorator must be at first
@mockReturnConstant
def f(x):
	return x*x

#tracing
def  mydec(fun):
	def _mydec(*args, **kargs):
		print("calling " + fun.__name__)
		res = fun(*args, **kargs)
		print("ending " + fun.__name__)
		return res
	return _mydec

@mydec
def f(x):
	return x*x
	
f(2)

##Decorator taking argument
def makeDefault100(f):
	def org(*args):
		val = args[1] if len(args) > 1 else 100
		return f(*args) if len(args) > 1 else f(*args,y=val)
	return org

@makeDefault100
def f1(x,y):
	return x+y

	
f1(2)
f1(2,3)	

# Example with Decorator argument
def makeDefault(d):
	def aDec(f):
		def org(*args):
			val = args[1] if len(args) > 1 else d
			return f(*args) if len(args) > 1 else f(*args,y=val)
		return org
	return aDec

@makeDefault(50)
def f1(x,y):
	return x+y

	
f1(2)
f1(2,3)	

##Creating a well defined decorator , 
#such that all attributes of original are preserved
from functools import wraps

def iter_f(func):
	@wraps(func)
    def newf(*args, **kwargs):
        for i in range(10):
            res = func(*args, **kwargs)
			print(res)
    return newf

@iter_f
def f(x):
	return x*x
	
	
def iter_f(d):
	def inner(f):
		@wraps(f)
		def org(*args, **kwargs):
			for i in range(10):
				res = f(*args, **kwargs)
				print(res)
		return org
	return inner
	
	
@iter_f(10)
def f(x):
	return x*x

##Recursive decorator
	
def  mea(fun):
	def _inner(*args, **kargs):
		import time
		now = time.time()
		res = fun(*args,**kargs)
		print(time.time() - now)
		return res
	return _inner
	
	
@mea
def f(x,y):
	return x+y
	

#for recursion it calls many times
@mea
def fib(n):
	return n if n<2 else  fib(n-1) + fib(n-2)
	
#to fix

def  mea(fun):
	def _inner(*args, _first = False, **kargs):
		import time
		if _inner.first :
			_inner.first = False
			_first = True
		now = time.time()
		res = fun(*args,**kargs)
		if _first : 
			_inner.first = True;
			print(time.time() - now)
		return res
	_inner.first = True
	return _inner
	
@mea
def fib(n):
	return n if n<2 else  fib(n-1) + fib(n-2)
	
	
##class decorator
#use below, 
#but for builtins( ie len(x) etc), __getattr__/__getattribute__ is not called
def __getattr__(self, name): 			# On undefined attribute fetch [obj.name], for old and new style
def __getattribute__(self, name): 		# On all attribute fetch [obj.name], new style
def __setattr__(self, name, value): 	# On all attribute assignment [obj.name=value]
def __delattr__(self, name): 			# On all attribute deletion [del obj.name]




def decorator(cls): 					# On @ decoration
	class Wrapper:
		def __init__(self, *args): 		# On instance creation
			self.wrapped = cls(*args)
		def __getattr__(self, name): 	# On attribute fetch
			print("Getting " + name + "...")
			return getattr(self.wrapped, name)
	return Wrapper

@decorator
class C: 							# C = decorator(C)
	def __init__(self, x, y): 		# Run by Wrapper.__init__
		self.attr = 'spam'

x = C(6, 7) 		# Really calls Wrapper(6, 7)
print(x.attr) 		# Runs Wrapper.__getattr__, prints "spam"





###Iterator Pattern
		
# Examples of Iterators - Single Pass
class Squares0: 				
	def __init__(self, start, stop): 	
		self.start = start
		self.stop = stop
		self.c = start
	def __iter__(self):
		return self				#must return self or a class instance which implements __next__
	def __next__(self):			# for py3.3 , for py2.7 it is next()
		if self.c >= self.stop : raise StopIteration
		self.c += 1
		return (self.c-1)**2
s = Squares0(1,10)
I1 = iter(s); I2=iter(s)
next(I1), next(I2)  #(1, 4)
next(I1), next(I2) #(9, 16)

#Alternate version- MultiPass- use Yield
class Squares: 				
	def __init__(self, start, stop): 	
		self.start = start
		self.stop = stop
	def __iter__(self):
		for value in range(self.start, self.stop + 1):
			yield value ** 2

			
>>> i = iter(Squares(1,5))
>>> next(i)
1


##Generator  pattern along with Iterator


def generate_ints(N):
    for i in range(N):
        yield i*i

def counter (maximum):
    i = 0
    while i < maximum:
        val = (yield i)
        # If value provided, change counter
        if val is not None:
            i = val
        else:
            i += 1
			
s = iter(counter(1000))
next(s)
s.send(50)
next(s)


##generator object for both version

a =  (x*x for x in range(100))
next(a)
next(a)


#Upto some n
def fib(n):
    a, b = 0, 1
    for _ in range(n):
        yield a
        a, b = b, a + b

#Infinite
def fib():
    a, b = 0, 1
    while 1:
        yield a
        a, b = b, a + b

a = iter(fib())
next(a)
0

import itertools as it

for x in it.islice(next(a), 20):
    print x


class Fib: 				
	def __init__(self, a=0, b=1): 	
		self.a, self.b = a,b 
	def __iter__(self):
		while 1:
			yield self.a
			self.a, self.b = self.b, self.a + self.b

	
#with functions			
def fib1():
	a,b = 0,1
	def g():
		nonlocal a,b  #else UnboundLocalError: local variable 'b' referenced before assignment
		a,b = b, a+b
		return a
	return g

	
###comprehensions  Pattern
#- eager - Note lazy is not possible as there should be stop in range else goes infinite
#to find pythogoran triplets
[(x,y,z) for x in range(1,100) for y in range(x,100) for z in range(x,100) if x*x + y*y == z*z ]


###Lazy computations via yield from
	
#https://docs.python.org/3/library/itertools.html

#itertools.count(start,[step])   -> start, start+step...
#itertools.repeat(element, [n])

#dropwhile(predicate, iterable)
#takewhile(predicate, iterable)
#islice(iterable, start, stop[, step]) or islice(iterable, stop)
#tee(iterable, n=2) --> tuple of n independent iterators. (memoized version of it, hence very fast)


# Few lazy definition
from itertools import *

def iterate(x,f):
	yield f(x)
	yield from iterate(f(x),f)
	
#eg #itertools.count(start,[step])  can be implemented as 

def count(st):
	import functools
	import operator
	f = functools.partial(operator.add, 1)
	yield st
	yield from iterate(st, f)
	
	
#Few more functions
def tail(iterable):  
	return islice(iterable, 1, None)
	

def take(iterable, n):  
	return list(islice(iterable, 0, n))

def head(iterable):  
	return list(islice(iterable, 1))

#or if islice is not wanted

def tail(it):
	it = iter(it)
	next(it)
	return it


def take(it, n):
	it = iter(it)
	return [ next(it) for i in range(n) ]


def head(it):
	it = iter(it);
	return next(it)
	
#few fibs implementations

def fibfrom(a = 0, b = 1):
	yield a
	yield from fibfrom(b, a+b)
	

>>> take(fibfrom(), 100)

def fib_inf():
	yield from iterate( (0,1) , lambda t: ( t[1], t[0] + t[1]) )

>>> take(map(lambda x: x[0], fib_inf()), 10)

#inefficient as double recursion via fibs() and  tail(fibs()) and no memoization
def fibs():
	import operator
	yield 0
	yield 1
	yield from map(operator.add, fibs(), tail(fibs()))

>>> take(fibs(), 10)

#efficient, using tee to copy instead of double recursion
def fibs():
	import operator
	import itertools
	yield 1
	yield 1
	fibs1, fibs2 = itertools.tee(fibs())
	yield from map(operator.add, fibs1, tail(fibs2))

	
	
###Recursion pattern - accumulator pattern
#P3.3
def mysum1(L):
	first, *rest = L
	return first if not rest else first + mysum1(rest)


#P2.7
def mysum1(L):
	first, rest = L[0], L[1:]
	return first if not rest else first + mysum1(rest)
	
	
def s(lst):
	return 0 if not lst else lst[0] + s(lst[1:])
	
def mysum(l):
	f, *rest = l + ([0] if not l else [])
	return f if not rest else f+sum(rest)

>>> mysum([1,2,3])
6
>>> mysum([])
0

import sys
sys.getrecursionlimit()


def mysum3(L, csum):  #csum is accumulator 
	while True:                     
		if not L: return csum
		L, csum = L[1:], csum + L[0]   

		
#fibs
def fibs(n, s = [0,1]):
	return s if n == 0 else fibs(n-1, s + [ s[-2] + s[-1] ] )

#TCO way
def fibs(n, s = [0,1]):
	while True:
		if not n : return s
		n, s = n-1, s + [ s[-2] + s[-1] ]


#mymap
def mymap(f, lst):
	return [ ] if not lst else [ f(lst[0]) ] + mymap (f, lst[1:])
	
#TCO way

def mymap(f, lst, acc = []):
	while True:
		if not lst : return acc
		lst, acc = lst[1:], acc + [ f(lst[0]) ] 


##Recursion pattern - Memoization

def fib1(n):
	if n == 0 : return 0
	if n ==  1 : return 1
	return fib1(n-1) + fib1(n-2)

	
#3.2
import functools
@functools.lru_cache(maxsize=None)   # by default 128
def fib_m(n):
	return n if n<2 else  fib_m(n-1) + fib_m(n-2)
	
#Recurisve
@functools.lru_cache(maxsize=None)
def fib_r(n, prev=0, next=1):
	while True:
		if n == 0 : return prev
		if n ==  1 : return next
		n, prev, next = n-1, next, prev+next


		
#2.7

from repoze.lru import lru_cache

@lru_cache(maxsize=5000)
def fib_r(n, prev, next):
	while True:
		if n == 0 : return prev
		if n ==  1 : return next
		n, prev, next = n-1, next, prev+next
		
fib_r(100001,0,1)



	
###Singleton Pattern

#simple by using __new__ .
# __new__ gets called when instantiation with 'cls' , must call object.__new__ to get a instance, __new__ calls __init__ automatically 
class SingleTone(object):
    __instance = None             #__ means names are mangled 
    def __new__(cls, val):
        if SingleTone.__instance is None:
            SingleTone.__instance = object.__new__(cls)
        SingleTone.__instance.val = val
        return SingleTone.__instance

	
#complex 	way 
class OnlyOne(object):
    class __OnlyOne:
        def __init__(self):
            self.val = None
        def __str__(self):
            return str(self.val)
    instance = None
    def __new__(cls): 				# __new__ always a classmethod
        if not OnlyOne.instance:
            OnlyOne.instance = OnlyOne.__OnlyOne()
        return OnlyOne.instance
    def __getattr__(self, name):     #name is str ,  for any attributes, but not for builtins eh len() etc 
        return getattr(self.instance, name)
    def __setattr__(self, name):
        return setattr(self.instance, name)

x = OnlyOne()
x.val = 'sausage'
print(x)
y = OnlyOne()
y.val = 'eggs'
print(y)
z = OnlyOne()
z.val = 'spam'
print(z)
print(x)
print(y)
#<hr>
output = '''
<__main__.__OnlyOne instance at 0x00798900>sausage
<__main__.__OnlyOne instance at 0x00798900>eggs
<__main__.__OnlyOne instance at 0x00798900>spam
<__main__.__OnlyOne instance at 0x00798900>spam
<__main__.__OnlyOne instance at 0x00798900>spam
'''


#Singleton using Borg patterns - can be implemented by using inheritance of class Borg
#is to have a single set of state data for all objects
#by setting all the __dict__  to the same static piece of storage

class Borg:
	_shared_state = {}
	def __init__(self):
		self.__dict__ = self._shared_state

		

class Singleton(Borg):
	def __init__(self, arg):
		Borg.__init__(self)
		self.val = arg
	def __str__(self): return self.val

x = Singleton('sausage')
print(x)
y = Singleton('eggs')
print(y)
z = Singleton('spam')
print(z)
print(x)
print(y)


output = '''
sausage
eggs
spam
spam
spam
'''


#using class decorator
class SingletonDecorator:
	def __init__(self,klass):
		self.klass = klass
		self.instance = None
	def __call__(self,*args,**kwds):
		if self.instance == None:
			self.instance = self.klass(*args,**kwds)
		return self.instance


@SingletonDecorator		#not for py2.x
class foo: pass         #foo = SingletonDecorator(foo)  for py2.x


x=foo()
y=foo()
z=foo()
x.val = 'sausage'
y.val = 'eggs'
z.val = 'spam'
print(x.val)

#Using metaclass - Example here for no change after  once creation

class Singleton(type):
    _instances = {}
    def __call__(cls, *args, **kwargs): #called at instance creation
        if cls not in cls._instances:
            cls._instances[cls] = type.__call__(*args, **kwargs)
        return cls._instances[cls]

#Python2
class MyClass(object):
    __metaclass__ = Singleton

#Python3
class MyClass(metaclass=Singleton):
	def __init__(self,val):
		self.val = val
	def __str__(self):
		return str(self.val)



x = MyClass('sausage')
y=MyClass('eggs')
z=MyClass('spam')
print(x)
print(y)
print(z)
print(x is y is z)




###Factory Pattern 


from __future__ import generators
import random

class Shape(object):
    @staticmethod
    def factory(type):        
        if type == "Circle": return Circle()
        if type == "Square": return Square()
        assert 0, "Bad shape creation: " + type


class Circle(Shape):
    def draw(self): print("Circle.draw")
    def erase(self): print("Circle.erase")

class Square(Shape):
    def draw(self): print("Square.draw")
    def erase(self): print("Square.erase")

# Generate shape name strings:
def shapeNameGen(n):
    import random 
    types = Shape.__subclasses__()
    for i in range(n):
        yield random.choice(types).__name__

shapes = [ Shape.factory(i) for i in shapeNameGen(7)]

for shape in shapes:
    shape.draw()
    shape.erase()


# Preventing direct creation

import random

class Shape(object):pass


def factory(name):
	class Circle(Shape):
		def draw(self): print("Circle.draw")
		def erase(self): print("Circle.erase")
	class Square(Shape):
		def draw(self): print("Square.draw")
		def erase(self): print("Square.erase")
	if name == "Circle": return Circle()
	if name == "Square": return Square()
	assert 0, "Bad shape creation: " + name
def shapeNameGen(n):
    for i in range(n):
        yield factory(random.choice(["Circle", "Square"]))

# Circle() # Not defined

for shape in shapeNameGen(7):
    shape.draw()
    shape.erase()

	
###Polymorphic Factories-Factory pattern
#Every class has create method 

from __future__ import generators
import random

class ShapeFactory:
	factories = {}	
	def createShape(id):
		if id not in ShapeFactory.factories:
			ShapeFactory.factories[id] =  eval(id + '.Factory()')
		return ShapeFactory.factories[id].create()
	createShape = staticmethod(createShape)

class Shape(object): pass

class Circle(Shape):
    def draw(self): print("Circle.draw")
    def erase(self): print("Circle.erase")
    class Factory:
        def create(self): return Circle()

class Square(Shape):
    def draw(self):
        print("Square.draw")
    def erase(self):
        print("Square.erase")
    class Factory:
        def create(self): return Square()
		
def shapeNameGen(n):
    types = Shape.__subclasses__()
    for i in range(n):
        yield random.choice(types).__name__

shapes = [ ShapeFactory.createShape(i) for i in shapeNameGen(7)]

for shape in shapes:
    shape.draw()
    shape.erase()

	
	
	
	
###Abstract class pattern

#In Python 3.X, 

from abc import ABCMeta, abstractmethod
class Super(metaclass=ABCMeta):
		@abstractmethod
		def method(self, arg):
			pass

			
#In Python 2.6 and 2.7, 
class Super:
		__metaclass__ = ABCMeta
		@abstractmethod
		def method(self, arg):
			pass

			

>>> X = Super()
TypeError: Can't instantiate abstract class Super with abstract methods action

>>> class Sub(Super): pass
>>> X = Sub()
TypeError: Can't instantiate abstract class Sub with abstract methods action

class A(Super):
	def method(self, arg):
		print(arg)


>>> a = A()
>>> a.method(2)
2






### Accessors pattern - Various techniques



##Descriptor - Py3.x and Py2.7 works same way if Py2.7 class are new style(derives from object)

#object.__get__(self, instance, owner)
#self = descriptor class instance
#instance = instance of owner class (for instance.attr), or None for class.attr
#owner = the owner class always


#Note: readonly attribute, don't omit __set__, but use below
#def __set__(*args): raise AttributeError('cannot set')

class AgeDesc(object):
	def __get__(self, instance, owner): return 40 if instance else 50 #return 40 if instance else return 50 for class variable 
	def __set__(self, instance, value): instance._age = value
	def __delete__(self, instance): pass

class descriptors(object):
		age = AgeDesc()

>>> x = descriptors()
>>> x.age 			# Runs AgeDesc.__get__
40
>>> x.age = 42 		# Runs AgeDesc.__set__
>>> x._age 			# Normal fetch: no AgeDesc call
42



##property - attribute -  Py3.x and Py2.7 works same way if Py2.7 class are new style(derives from object)
#property(fget, fset, fdel, doc)
#for readonly - make property(fget, None, None, doc)

class D(object):	
    def __init__(self, p):
        self._prop = p
    def p_get(self):
        return self._prop
    def p_set(self, val):
        self._prop = val
    def p_del(self):
        del self._prop	
    prop = property(p_get, p_set, p_del,"This is property")   #instance.prop

#Or using a decorator, for read-only, keep only @property section

class C:
    def __init__(self):
        self._x = None

    @property			#get 
    def x(self):
        """I'm the 'x' property."""
        return self._x

    @x.setter				#set
    def x(self, value):
        self._x = value

    @x.deleter				#del 
    def x(self):
        del self._x


		
##Slots: Attribute Declarations(Py2.7 and Py3.x) - limiting only few attributes in a class
#prevents the automatic creation of __dict__

#When inheriting from a class without __slots__, the __dict__ attribute of that class will always be accessible, 
#so a __slots__ definition in the subclass is meaningless

#The action of a __slots__ declaration is limited to the class where it is defined. 
#As a result, subclasses will have a __dict__ unless they also define __slots__ 
#(which must only contain names of any additional slots).


class limiter(object):
	__slots__ = ['a', ]

>>> x = limiter()
>>> x.a			
AttributeError: a
>>> x.a = 40 			
>>> x.a
40
>>> x.b = 1000 		
AttributeError: 'limiter' object has no attribute 'b‘
>>> limiter.__dict__

>>> x.__dict__
Traceback (most recent call last):
  File "<stdin>", line 1, in <module>
AttributeError: 'limiter' object has no attribute '__dict__'
		

		
		
		
###proxy pattern

#Proxy- in Py3.x(only for user defined attributes, not for overloaded operator)
#Note in new style(Py3, Py2.7 when deriving from object), you must implement for all overloading operator explicitly 
#because - __getattr__ and __getattribute__ are not called for builtins eg [], len etc 
def __getattr__(self, name): 		# On undefined attribute fetch [obj.name], for old and new style
def __getattribute__(self, name): 	# On all attribute fetch [obj.name], new style
def __setattr__(self, name, value): # On all attribute assignment [obj.name=value]
def __delattr__(self, name): 		# On all attribute deletion [del obj.name]

#Class -Attribute fetching - details

def __getattr__(self, attrname): 	
		if attrname == 'age':
			return 40
		else:
			raise AttributeError(attrname)

def __setattr__(self, attr, value):  			# for all setting
		if attr == 'age':
			self.__dict__[attr] = value + 10 	# does not include __slot__
		else:
			raise AttributeError(attr + ' not allowed')

def __getattribute__(self, name):
		x = object.__getattribute__(self, name)   # Must use object. else infine loop


#Note following would loops infinitely in __setattr__
self.age = value + 10 			# Loops
setattr(self, attr, value + 10) # Loops (attr is 'age')

#use below in __setattr__
self.__dict__[attr] = value + 10 			# OK: doesn't loop
object.__setattr__(self, attr, value + 10) 	# OK: doesn't loop (new-style only)
		
#To be inclusive of slot and properties, use always object.__setattr__


#proxy - pattern
class C:
  data = 'spam'
  def __getattr__(self, name):
    print('getattr: ' + name)
    return getattr(self.data, name)
	
	
>>> X = C()
>>> X.__getitem__(1) 		# Traditional mapping works but new-style's does not
getattr: __getitem__
'p'
>>> X[1]
getattr: __getitem__
'p'
>>> getattr(X,'data')
'spam'

>>> getattr(X,'__getitem__')
getattr: __getitem__
<method-wrapper '__getitem__' of str object at 0x0084DF60>

>>> getattr(X,'__getitem__')(1)
getattr: __getitem__
'p'



#Using Proxy decorator

class A(object):
	def __init__(self, name):
		self.name = name
	def hello(self):
		print('hello %s!' % (self.name,))
	def __enter__(self):
		print('Enter the function')
		return self   #must return self or a class which implements __exit__
	def __exit__(self, exc_type, exc_value, traceback):
		print('Exit the function')

class Proxy(object):
	def __init__(self, object_a):
		self._object_a = object_a
	def decorateEnterExit(self, obj, f):
		print("inside decorateEnterExit")
		def inner(*args, **kwargs):
			print("inside decorateEnterExit.inner")
			with obj as _:
				return f(*args, **kwargs)
		return inner
	def __getattribute__(self, name):
		obj = object.__getattribute__(self, '_object_a')
		dee = object.__getattribute__(self, 'decorateEnterExit')
		print("inside __getattribute__")
		return dee(obj, getattr(obj, name))
		
>>> Proxy(A('Ax')).hello()
inside __getattribute__
inside decorateEnterExit
inside decorateEnterExit.inner
Enter the function
hello Ax!
Exit the function






### Command pattern - Abstract Base contains execute command

from abc import ABCMeta
from abc import abstractmethod
import os

class Command(object):   #in py3.x,  class Command(metaclass=ABCMeta):
	"""
	Abstract / Interface base class for commands.
	"""
	__metaclass__ = ABCMeta    #In Py2.x, this is commented 
	@abstractmethod
	def execute(self):  pass
	@abstractmethod
	def undo(self): pass


class CreateCommand(Command):
	"""
	Create command implementation.
	 """
	def __init__(self, name):
		self.file_name = name
	def execute(self, name):
		fp = open(self.file_name, 'w')
		fp.close()
	def undo(self):
		os.remove(self.file_name)
        


class MoveCommand(Command):
    """
    Move command implementation.
    """
    def __init__(self, src, dest):
        self.src = src
        self.dest = dest

    def execute(self, src, dest):
        os.rename(self.src, self.dest)
    def undo(self):
        os.rename(self.dest, self.src)
        


class Invoker(object):
    def __init__(self, command):
        self.command = command
    def do(self):
        self.command.execute()
    def undo(self):
        self.command.undo() 


# Client for the command pattern
if __name__ == '__main__':
    create_cmd = CreateCommand('/tmp/foo.txt')
    move_cmd = MoveCommand('/tmp/foo.txt', '/tmp/bar.txt')
    create_invoker = Invoker(create_cmd)
    move_invoker = Invoker(move_cmd)
    create_invoker.do()
    move_invoker.do()
    move_invoker.undo()
    create_invoker.undo()

	
	

###StateMachine Pattern 
#imposes a structure to automatically change the implementation from one object to the next

#code pattern
class SuperState( object ):             #can be abstract
    def do(self, input) ):
        raise NotImplementedError()
    def transitionRule( self, input ):
        raise NotImplementedError()

class SomeState( SuperState ):
    def do( self, input ):
        #actually do something() after retrieving few chunks input 
		return self.transitionRule(rest_input)
    def transitionRule( self, input ):
		#check input and depending on call some other state
        return NextState()      #nextState or all states are derived from SuperState


#One Example
		
class State( object ):
	def transitionRule( self, input ):
		cls = eval(self.map[input])()  
		print(cls.__class__)
		return cls

class S1( State ): 
	map = { "input": "S2", "other": "S3" }
	pass # Overrides to state-specific methods


class S2( State ):
	map = { "foo": "S1", "bar": "S2" }


class S3( State ):
	map = { "quux": "S1" }


#Usage 
s = S1()
i = ["input", "bar", "bar", "foo", "other","quux"]
for e in i:
	s = s.transitionRule(e)

	
#Above example with lambda

class State( object ):
	def transitionRule( self, input ):
		next_states = [ s for f,s in self.map if f(input)  ]
		assert len(next_states) >= 1, "faulty transition rule"
		return eval(next_states[0])()

class S1( State ):
    map = [ (lambda x: x == "input", "S2"), (lambda x: x == "other", "S3" ) ]

class S2( State ):
    map = [ (lambda x: "bar" <= x <= "foo", "S3"), (lambda x: True, "S1") ]



#table driven State machines

#All State derive from State 

class State:
    def __init__(self, name): self.name = name
    def __str__(self): return str(self.name)

	
#Input class to tag input , 
#Derive many Input for each input
class Input: pass


#Condition validates the input to decide whether this row in the table is the correct transition
#Derive many Condition for one Condition with one Input
class Condition:
    boolean condition(input) : #Input instance 
        assert 0, "condition() not implemented"


#Transition Actions
#If the Condition returns true, then this state's action happens in transition method 
#Derive many transitions for each transition with Input 
class Transition:
    def transition(self, input):  #Input instance 
        assert 0, "transition() not implemented"

#The Table, Given current state with input, if Condition true, transition is executed with next state 

{(CurrentState, InputA) : (ConditionA, TransitionA, NextStateA),
 (CurrentState, InputB) : (ConditionB, TransitionB, NextStateB),
 (CurrentState, InputC) : (ConditionC, TransitionC, NextStateC),
 #...
}

# A table-driven state machine

class StateMachine:
    def __init__(self, initialState, tranTable):
		self.state = initialState
		self.transitionTable = tranTable   #above table
		self.end = False
		self.exception = None
    def nextState(self, input):
		try:
			cond, trans, nState = self.transitionTable.get((self.state, input), (None, None, None))
			if cond and cond(input):
				trans.transition(input)
				self.state = nState
			else:
				self.end = True
		except Exception as e:
			self.exception = e
			self.end = True
		finally:
			print("closing the machine")
	
		
ins = [InputA, InputB, InputC,.....]
s = StateMachine(initialState, tranTable)
for i in ins:
	s.nextState(i)
	
if s.exception :
	print("error", s.exception)

	
	
	
	
	
	


###Observer pattern
#Used when a group of objects needs to update themselves when some object changes state.
#Observer is an 'interface' class with method update( ), called by Observable objects  


#Derive your class from Observer and override update 
#your update logic tells what to be done when observable calls update
class Observer:
	def update(self, observable, arg):
		'''Called when the observed object is
		modified. You call an Observable object's
		notifyObservers method to notify all the
		object's observers of the change.'''
		print(self, observable,arg)

		
#In multi threading, self.obj must be synchronized
#Basic flow is - Create one observable, add all observers here by addObserver
#Any observer can call notifyObservers, all observer's update would be called

class Observable:
	def __init__(self):
		self.obs = []
		self.changed = 0		
	def addObserver(self, observer):
		if observer not in self.obs:
			self.obs.append(observer)
	def deleteObserver(self, observer):
		self.obs.remove(observer)
	def notifyObservers(self, arg = None):
		'''If 'changed' indicates that this object
		has changed, notify all its observers, then
		call clearChanged(). Each observer has its
		update() called with two arguments: this
		observable object and the generic 'arg'.'''
		try:
			if not self.changed: return
			# Make a local copy in case of synchronous
			# additions of observers:
			localArray = self.obs[:]
			self.clearChanged()
		finally:
			pass
        # Updating is not required to be synchronized:
		for observer in localArray:
			observer.update(self, arg)
	def deleteObservers(self): self.obs = []
	def setChanged(self): self.changed = 1
	def clearChanged(self): self.changed = 0
	def hasChanged(self): return self.changed
	def countObservers(self): return len(self.obs)

  
 
ol = Observable()
[ ol.addObserver(Observer()) for i in range(5) ]

ol.setChanged()
ol.notifyObservers()

 
 
 
###Visitor pattern - kind of multiple dispatching
#A class which wants to be visited,  'accepts' a Visitor and then calls Visitor's visit method passing ownself

#User can create Various kind of visitor's logic 

# Demonstration of "visitor" pattern.
from __future__ import generators
import random


class Flower(object):
    def accept(self, visitor):
        visitor.visit(self)
    def pollinate(self, pollinator):
        print(self, "pollinated by", pollinator)
    def eat(self, eater):
        print(self, "eaten by", eater)
    def __str__(self):
        return self.__class__.__name__


		
class Gladiolus(Flower): pass

class Runuculus(Flower): pass

class Chrysanthemum(Flower): pass


#Various kind of visitors
class Visitor:
	def __str__(self):
		return self.__class__.__name__

class Bug(Visitor): pass

class Pollinator(Bug): pass

# Add the ability to do "Bee" activities:
class Bee(Pollinator):
	def visit(self, flower):
		flower.pollinate(self)


#Another Visitor
class Predator(Bug): pass

class Worm(Predator):
	def visit(self, flower):
		flower.eat(self)

		

def flowerGen(n):
	flwrs = Flower.__subclasses__()
	for i in range(n):
		yield random.choice(flwrs)()

bee = Bee()
worm = Worm()
for flower in flowerGen(10):
	flower.accept(bee)
	flower.accept(worm)





#Single dispatch Pattern - functools has functionality -  @functools.singledispatch(default)


#dispatch based on first argument
from functools import singledispatch
@singledispatch
def fun(arg, verbose=False):       #default version if no specific is found
	if verbose:
		print("Let me just say,", end=" ")
	print(arg)


#to add overloaded instance

@fun.register(int)   #'fun' comes from above
def _(arg, verbose=False):
	if verbose:
		print("Strength in numbers, eh?", end=" ")
	print(arg)
	

@fun.register(list)
def _(arg, verbose=False):
	if verbose:
		print("Enumerate this:")
	for i, elem in enumerate(arg):
		print(i, elem)


#for lambda and preexisting functions
def nothing(arg, verbose=False):
	print("Nothing.")

fun.register(type(None), nothing)


#for multiple registrtion

@fun.register(float)
@fun.register(Decimal)
def fun_num(arg, verbose=False):
	if verbose:
		print("Half of your number:", end=" ")
	print(arg / 2)


#usage 

>>>>>> fun("Hello, world.")
Hello, world.
>>> fun("test.", verbose=True)
Let me just say, test.
>>> fun(42, verbose=True)
Strength in numbers, eh? 42
>>> fun(['spam', 'spam', 'eggs', 'spam'], verbose=True)
Enumerate this:
0 spam
1 spam
2 eggs
3 spam
>>> fun(None)
Nothing.
>>> fun(1.23)
0.615

#to know which function is registered
>>> fun.dispatch(float)
<function fun_num at 0x1035a2840>
>>> fun.dispatch(dict)    # note: default implementation
<function fun at 0x103fe0000>




# Multiple dispatch Pattern -  using a table
from __future__ import generators  #for P2.x
import random

class Outcome:
	def __init__(self, value, name):
		self.value = value
		self.name = name
	def __str__(self): return self.name
	def __eq__(self, other):
		return self.value == other.value

		
#class variable 	

Outcome.WIN = Outcome(0, "win")
Outcome.LOSE = Outcome(1, "lose")
Outcome.DRAW = Outcome(2, "draw")

#compete is multiple dispatched method
class Item(object):
	def compete(self, item):
		# Use a tuple for table lookup:
		return outcome[self.__class__, item.__class__]
	def __str__(self):
		return self.__class__.__name__


class Paper(Item): pass

class Scissors(Item): pass

class Rock(Item): pass


outcome = {
  (Paper, Rock): Outcome.WIN,
  (Paper, Scissors): Outcome.LOSE,
  (Paper, Paper): Outcome.DRAW,
  (Scissors, Paper): Outcome.WIN,
  (Scissors, Rock): Outcome.LOSE,
  (Scissors, Scissors): Outcome.DRAW,
  (Rock, Scissors): Outcome.WIN,
  (Rock, Paper): Outcome.LOSE,
  (Rock, Rock): Outcome.DRAW,
}

def match(item1, item2):
	print("%s <--> %s : %s" % (item1, item2, item1.compete(item2)))
	  
	  

# Generate the items:
def itemPairGen(n):
    # Create a list of instances of all Items:
    Items = Item.__subclasses__()
    for i in range(n):
        yield (random.choice(Items)(), random.choice(Items)())


for item1, item2 in itemPairGen(20):
    match(item1, item2)

	
	
	
	
	
	
	
	
	
	
	
	
	

###Database connectivity using DBAPI 2.0 

#mysql -running daemon in cygwin
/usr/bin/mysqld_safe &
#or run it services.msc
#shutting down
mysqladmin.exe -h 127.0.0.1 -u root   --connect-timeout=5 shutdown
#mysql admin #  default port 3306, 
mysql -u root    -h 127.0.0.1 
#few commands
show databases;
create database python;
use python;
show tables;
create table employes ( id INT, first_name VARCHAR(20), last_name VARCHAR(20), hire_date  DATE);
desc employes;
insert into employes values (3, "das", "das", '1999-03-30');
select * from employes; 


#Module - mysql connector -  installation of mysql connector- Py3.x or Py2.x
1. get python connector from mysql site
2. unzip to a dir
3. cd dir
4. python setup.py install

import mysql.connector
cnx = mysql.connector.connect(user='root', password='', host='127.0.0.1', database='python')

#fetching 
cursor = cnx.cursor()
cursor.execute("select * from employes")

results = cursor.fetchall()
results
for result in results:
		print(result[0], result[1], result[2])

cursor.close()
cnx.close()

#other methods
cursor.column_names
cursor.rowcount



#	More example of fetches
head_rows = cursor.fetchmany(size=2)
remaining_rows = cursor.fetchall()
cursor.fetchone()  # get one row


#Inserting row

import datetime

cursor = cnx.cursor()

insert_stmt = (
  "INSERT INTO employes (id, first_name, last_name, hire_date) "
  "VALUES (%s, %s, %s, %s)"
)   # Note we use only %s, backend converts to real SQL type

data = (5, 'N', 'Das', datetime.date(2012, 3, 23))
cursor.execute(insert_stmt, data)
cnx.commit()

#	To insert multiple rows
c.executemany(
      """INSERT INTO employes (id, first_name, last_name, hire_date) 
  	VALUES (%s, %s, %s, %s)""",
      [
      (2, 'N', 'Das', datetime.date(2012, 3, 23)),
      (3, 'N', 'Das', datetime.date(2012, 3, 23)),
      (4, 'N', 'Das', datetime.date(2012, 3, 23))
      ] )

#	Selecting with variable

select_stmt = "SELECT * FROM employes WHERE id = %(emp_no)s"
cursor.execute(select_stmt, { 'emp_no': 2 })

cursor.fetchall()


cursor.execute("SELECT last_name, first_name, hire_date "
               "FROM employes WHERE id = %s", (2,))
cursor.fetchall()


#Getting column name and printing
cursor.execute("SELECT last_name, first_name, hire_date "
               "FROM employes WHERE id = %s", (2,))
row = dict(zip(cursor.column_names, cursor.fetchone()))
print("{0[last_name]}, {0[first_name]}: {0[hire_date]}".format(row))

cursor.close()
cnx.close()

#mysqlDB – another python lib, Py2.x , for Py3.x - use Mysqlclient
#MySQL-4.1 through 5.5  supported
#check your python xbit at  http://www.lfd.uci.edu/~gohlke/pythonlibs/
#and use pip3 install some-package.whl
#or check http://victorjabur.com/2011/06/05/compiling-python-2-7-modules-on-windows-32-and-64-using-msvc-2008-express/
#Install VC for Python2.7 http://www.microsoft.com/en-us/download/details.aspx?id=44266
#Add to path,  C:\Users\das\AppData\Local\Programs\Common\Microsoft\Visual C++ for Python\9.0\VC\Bin\
#for py2.7, uses VC to compile c module
#check https://blogs.msdn.microsoft.com/pythonengineering/2016/04/11/unable-to-find-vcvarsall-bat/
#Python Version 	    You will need
3.5 and later 	    Visual C++ Build Tools 2015 or Visual Studio 2015
3.3 and 3.4 	    Windows SDK for Windows 7 and .NET 4.0
                    (Alternatively, Visual Studio 2010 if you have access to it)
2.6 to 3.2 	        Microsoft Visual C++ Compiler for Python 2.7




pip2 install mysql-python



import MySQLdb
cnx=MySQLdb.connect(user='root', passwd='', host='127.0.0.1', db='books')

#Then usage is same as above
#https://mysqlclient.readthedocs.org/en/latest/user_guide.html#mysqldb


#Module-sqllit , standard package
from sqlite3 import connect

conn = connect(r'D:/temp.db')
curs = conn.cursor()
curs.execute('create table emp (who, job, pay)')

prefix = 'insert into emp values '
curs.execute(prefix + "('Bob', 'dev', 100)")
curs.execute(prefix + "('Sue', 'dev', 120)")

curs.execute("select * from emp where pay > 100")
for (who, job, pay) in curs.fetchall():
		print(who, job, pay)

result = curs.execute("select who, pay from emp")
result.fetchall()

query = "select * from emp where job = ?"
curs.execute(query, ('dev',)).fetchall()











#Network programming
#Python's standard library releases the GIL around each blocking i/o call
#Hence threads can be created with send/recv which runs in multicore

#check socket
#low level library for TCP, UDP
import socket
dir(socket)

#server - create, bind, listen and accept
#client - create, bind, send  and rx

#high level library, Only override handle()
import socketserver  # contains TCPServer, UDPServer and as well as asynchronous 

#Http server
import http.server


#using multiprocessing.Listener, multiprocessing.Client and multiprocessing.Connection objects
#multiprocessing lib has Process, Q, Pipe, sharing data and Pool object

# multiprocessing.connection has Listener for server side and Client for client side 
#Message based API with a message boundary, Must be used with pickleable data

import multiprocesing


# Creating asynchronous Socket server and clients using asyncore module (standard module)
#asyncore.dispatcher contains below asynchronous event handlers
handle_read()  Called data to be read and readable returns true
handle_write() called when data to be written ie when writable() return true
handle_expt()  Called when there is out of band (OOB) data 
handle_connect() Called when the active opener’s socket actually makes a connection. 
handle_close()  Called when the socket is closed.
handle_error()  Called when an exception is raised and not otherwise handled.
handle_accept() Called on listening channels to accept
readable()   Called each time to get bool value for interested to read
writable() Called each time to get bool value for interested to write





#Twisted Matrix framework - Py2.7 (Py3.x work ongoing)
#https://twistedmatrix.com/documents/current/core/howto/index.html
#for windows
#require pywin32 , install from http://sourceforge.net/projects/pywin32/files/. 
#install VC http://aka.ms/vcpython27
#pip2 install twisted
#for other
#pip2 install twisted 

#examples are for cygwin


Twisted includes an event-driven (no GIL) web server 
supports SMTP, POP3, IMAP, SSHv2, and DNS


#Steps for writing TCP Server

0. twisted.internet.endpoints # Network address is called endpoints

1. Subclass from twisted.internet.protocol.Protocol
Or Subclass twisted.protocols.basic.LineReceiver for Line based Server

2. For any incoming connection, implement handler as def connectionMade(self) 

3. Handle conntaion close by def connectionLost(self, reason)

4. For Data Recieved , implement handler as def dataReceived(self, data) 
or Implement line based server as def lineReceived(self, line)

5. Use self.transport.write  or self.transport.sendLine to send data. 
You need to give control back to reactor to actually send the data

6. Subclass twisted.internet.protocol.Factory and implement buildProtocol  to instantiate the Server
The factory is used to share state that exists beyond the lifetime of any given connection

7. create the server
twisted.internet.endpoints.TCP4ServerEndpoint(reactor, port)
OR
twisted.internet.endpoints.serverFromString(reactor, b"tcp:1234:interface=127.0.0.1")

8. endpoint.listen(ServerFactory()) tells the reactor to handle connections to the endpoint’s address 

9. Use twisted.internet.reactor.run to run the Factory

10. can call loseConnection() and abortConnection() for graceful close or abort

11. stop the reactor by hitting Control-C in a terminal or calling reactor.stop().

#Steps for writing TCP Client
1.  Steps are same as of Server

2. Subclass from ReconnectingClientFactory to create Reconnecting ClientFactory
Handle def startedConnecting(self, connector) for connection starting
handle def clientConnectionLost(self, connector, reason) when connection is broken
handle def clientConnectionFailed(self, connector, reason) when connection is failed after retry

3. create client by reactor.connectTCP(ip_string, port, ClientFactory())





#asyncio  Module - standard module - Py2.x, Py3.x
#event driven  concurrent code using coroutine 
#coroutines are functions whose execution you can pause
#an event loop "is a programming construct that waits for and dispatches events or messages in a program". 
#uses the selectors module

#Example 
import asyncio

# Borrowed from http://curio.readthedocs.org/en/latest/tutorial.html.
@asyncio.coroutine
def countdown(number, n):  
    while n > 0:
        print('T-minus', n, '({})'.format(number))
        yield from asyncio.sleep(1)
        n -= 1

loop = asyncio.get_event_loop()  
tasks = [  
    asyncio.ensure_future(countdown("A", 2)),
    asyncio.ensure_future(countdown("B", 3))]
loop.run_until_complete(asyncio.wait(tasks))  
loop.close()  

'''
The event loop starts one of the countdown() coroutine call, executing until it hits yield from 
and the asyncio.sleep() function. That returns an asyncio.Future object
While waiting for sleep() in one call, it starts other countdown call 
It monitors all "yield from"/asyncio.Future object for readiness (like select ) 
and when one is ready, result of that is passed to paused coroutine via return from "yield from"
and continues to end of that routine
'''
##With Py3.5 
# This also works in Python 3.5.
@asyncio.coroutine
def py34_coro():  
    yield from stuff()
    
#In Python 3.5, the types.coroutine decorator has been added to replace asyncio.coroutine
#you can use "yield from" inside those 
#only Py3.5 
#only await and reurn are allowed inside 

async def py35_coro():  
    await stuff()

##await takes a coroutine (defined with asyncio.coroutine or types.coroutine) OR 
# an awaitable object: an object that defines an __await__() method which returns an iterator which is not a coroutine itself . 


###*** Details 
#Run an event loop by asyncio.get_event_loop()
BaseEventLoop.run_forever()  Runs till stop() is called
BaseEventLoop.run_until_complete(future)  Returns future
BaseEventLoop.is_running()
BaseEventLoop.stop()
BaseEventLoop.is_closed()
BaseEventLoop.close()

#Calls
BaseEventLoop.call_soon(callback, *args) calls as soon as possible via FIFO Q
callback can not take  keywords based arg, use functools.partial 
loop.call_soon(functools.partial(print, "Hello", flush=True)) will call print("Hello", flush=True)
returns asyncio.Handle 

#example
import asyncio

def hello_world(loop):
    print('Hello World')
    loop.stop()

loop = asyncio.get_event_loop()

# Schedule a call to hello_world()
loop.call_soon(hello_world, loop)

# Blocking call interrupted by loop.stop()
loop.run_forever()
loop.close()



#Delayed calls
BaseEventLoop.call_later(delay, callback, *args)
BaseEventLoop.call_at(when, callback, *args)
BaseEventLoop.time() Return the current time, as a float value,
asyncio.sleep() 
 
#example:
import asyncio
import datetime

def display_date(end_time, loop):
    print(datetime.datetime.now())
    if (loop.time() + 1.0) < end_time:
        loop.call_later(1, display_date, end_time, loop)
    else:
        loop.stop()

loop = asyncio.get_event_loop()

# Schedule the first call to display_date()
end_time = loop.time() + 5.0
loop.call_soon(display_date, end_time, loop)

# Blocking call interrupted by loop.stop()
loop.run_forever()
loop.close()




#Coroutines - chained by 'yield from' syntax . You should only do below 
•result = yield from future_or_task
suspends the coroutine until the future is done, then returns the future’s result, 
or raises an exception, which will be propagated. 
(If the future is cancelled, it will raise a CancelledError exception.) 
Note that tasks are also futures

•result = yield from coroutine 
wait for another coroutine to produce a result (or raise an exception, which will be propagated). 
The coroutine expression must be a call to another coroutine.
Note asyncio has many standard co-rouines , you must get result via 'yield from'

•return expression – produce a result to the coroutine that is waiting for this one using yield from.
•raise exception – raise an exception in the coroutine that is waiting for this one using yield from.

#How to run coroutine
call yield from coroutine from another coroutine 
or schedule its execution using the async() function or the BaseEventLoop.create_task() method
or call loop.run_until_complete(coroutine)

#Coroutines (and tasks) can only run when the event loop is running

#example
import asyncio

@asyncio.coroutine
def hello_world():
    print("Hello World!")

loop = asyncio.get_event_loop()
# Blocking call which returns when the hello_world() coroutine is done
loop.run_until_complete(hello_world())
loop.close()

#example
Example of coroutine displaying the current date every second during 5 seconds using the sleep() function:


import asyncio
import datetime

@asyncio.coroutine
def display_date(loop):
    end_time = loop.time() + 5.0
    while True:
        print(datetime.datetime.now())
        if (loop.time() + 1.0) >= end_time:
            break
        yield from asyncio.sleep(1)

loop = asyncio.get_event_loop()
# Blocking call which returns when the display_date() coroutine is done
loop.run_until_complete(display_date(loop))
loop.close()


#Example chaining coroutines:


import asyncio

@asyncio.coroutine
def compute(x, y):
    print("Compute %s + %s ..." % (x, y))
    yield from asyncio.sleep(1.0)
    return x + y

@asyncio.coroutine
def print_sum(x, y):
    result = yield from compute(x, y)
    print("%s + %s = %s" % (x, y, result))

loop = asyncio.get_event_loop()
loop.run_until_complete(print_sum(1, 2))
loop.close()


#Example combining a Future and a coroutine function:
#Future has result(), done(), cancel(), add_done_callback(fn),use partial, set_result()

import asyncio

@asyncio.coroutine
def slow_operation(future):
    yield from asyncio.sleep(1)
    future.set_result('Future is done!')

loop = asyncio.get_event_loop()
future = asyncio.Future()
asyncio.ensure_future(slow_operation(future))
loop.run_until_complete(future)
print(future.result())
loop.close()

#The previous example can be written differently using the 
#Future.add_done_callback() method to describe explicitly the control flow:


import asyncio

@asyncio.coroutine
def slow_operation(future):
    yield from asyncio.sleep(1)
    future.set_result('Future is done!')

def got_result(future):
    print(future.result())
    loop.stop()

loop = asyncio.get_event_loop()
future = asyncio.Future()
asyncio.ensure_future(slow_operation(future))
future.add_done_callback(got_result)
try:
    loop.run_forever()
finally:
    loop.close()

	

#Task
A task is responsible for executing a coroutine object in an event loop. 
If the wrapped coroutine yields from a future, the task suspends the execution of the wrapped coroutine 
and waits for the completition of the future. 

When the future is done, the execution of the wrapped coroutine restarts with the result 
or the exception of the future.

use the async() function or the BaseEventLoop.create_task(coroutine) method.
Schedule the execution of a coroutine: wrap it in a future. 
A task is a subclass of Future

#Example executing 3 tasks (A, B, C) in parallel:


import asyncio

@asyncio.coroutine
def factorial(name, number):
    f = 1
    for i in range(2, number+1):
        print("Task %s: Compute factorial(%s)..." % (name, i))
        yield from asyncio.sleep(1)
        f *= i
    print("Task %s: factorial(%s) = %s" % (name, number, f))

loop = asyncio.get_event_loop()
tasks = [
    asyncio.ensure_future(factorial("A", 2)),
    asyncio.ensure_future(factorial("B", 3)),
    asyncio.ensure_future(factorial("C", 4))]
loop.run_until_complete(asyncio.wait(tasks))
loop.close()



#Task functions
asyncio.as_completed(fs, *, loop=None, timeout=None)
for f in as_completed(fs):
    result = yield from f  # The 'yield from' may raise
    # Use result


asyncio.ensure_future(coro_or_future, *, loop=None)
Returns Task

asyncio.wait(futures_or_coros, *, loop=None, timeout=None, return_when=ALL_COMPLETED)
Returns two sets of Future: (done, pending).
done, pending = yield from asyncio.wait(fs)

asyncio.wait_for(single_fut_or_corou, timeout, *, loop=None)
Returns result of the Future or coroutine
result = yield from asyncio.wait_for(fut, 60.0)


#Transports and protocols 
Similar to Twisted library with handler methods
Transports are classes provided by asyncio in order to abstract various kinds of communication channels
Once the communication channel is established, a transport is always paired with a protocol instance

#Streams (coroutine based API)
Network connections are like files

#Create a subprocess
Run subprocesses asynchronously using the subprocess module.
coroutine asyncio.create_subprocess_exec(*args, stdin=None, stdout=None, stderr=None, loop=None, limit=None, **kwds)

#Synchronization
#it has Lock, Event, Condition, Semaphore
#Usages is similar to thread


lock = Lock()
...
yield from lock
try:
    ...
finally:
    lock.release()


Context manager usage:


lock = Lock()
...
with (yield from lock):
     ...


Lock objects can be tested for locking state:


if not lock.locked():
   yield from lock
else:
   # lock is acquired
    ...





#Concurrency and multithreading

An event loop runs in a thread and executes all callbacks and tasks in the same thread. 
While a task is running in the event loop, no other task is running in the same thread. 
But when the task uses yield from, the task is suspended and the event loop executes the next task.

To schedule a callback from a different thread, 
loop.call_soon_threadsafe(callback, *args)


Most asyncio objects are not thread safe.
You should only worry if you access objects outside the event loop. 
For example, to cancel a future, don’t call directly its Future.cancel() method, but:
loop.call_soon_threadsafe(fut.cancel)


To handle signals and to execute subprocesses, the event loop must be run in the main thread.

To schedule a coroutine object from a different thread, 
future = asyncio.run_coroutine_threadsafe(coro_func(), loop)
result = future.result(timeout)  # Wait for the result with a timeout


The BaseEventLoop.run_in_executor() method can be used with a thread pool executor to execute a callback 
in different thread to not block the thread of the event loop.


#Handle blocking functions correctly
Blocking functions should not be called directly. 
For networking and subprocesses, the asyncio module provides high-level APIs like protocols.
OR use Executor 
coroutine BaseEventLoop.run_in_executor(executor, func, *args)
Arrange for a func to be called in the specified concurrent.futures.Executor instance or None for default
Use functools.partial to pass keywords to the *func*.




##########################
#Module -subprocess - replaces os.system and os.spawn

import subprocess

#Arguments meaning
stdin, stdout and stderr specify the executed program’s standard input, 
standard output and standard error file handles, 
values are PIPE, DEVNULL, an existing file descriptor (a positive integer), an existing file object, and None

stderr can be STDOUT
If universal_newlines is False the file objects stdin, stdout and stderr will be opened as binary streams,
and no line ending conversion is done

If shell is True, the specified command will be executed through the shell
accesses to shell pipes, filename wildcards, environment variable expansion
On POSIX with shell=True, the shell defaults to /bin/sh. 
Popen(['/bin/sh', '-c', args[0], args[1], ...])
On windows shell=True to execute is built into the shell (e.g. dir or copy)

Handle  OSError or A ValueError for  invalid arguments.
check_call() and check_output() will raise CalledProcessError(check .returncode) for non-zero return code.


#subprocess.call(args, *, stdin=None, stdout=None, stderr=None, shell=False, timeout=None)
Run the command described by args. Wait for command to complete, then return the returncode attribute
# Do not use stdout=PIPE or stderr=PIPE with this function

#subprocess.check_output(args, *, input=None, stdin=None, stderr=None, shell=False, universal_newlines=False, timeout=None)
#input (which is sent to cmd) must be a byte sequence, or a string if universal_newlines=True.
# Do not use stdout=PIPE or stderr=PIPE with this function

>>> subprocess.check_output(["echo", "Hello World!"])
b'Hello World!\n'

>>> subprocess.check_output(["echo", "Hello World!"]).decode("utf-8")
'Hello World!\n'

>>> subprocess.check_output(["sed", "-e", "s/foo/bar/"], input=b"when in the course of fooman events\n")
b'when in the course of barman events\n'

>>> subprocess.check_output(
     "ls non_existent_file; exit 0",
     stderr=subprocess.STDOUT,
     shell=True)
'ls: non_existent_file: No such file or directory\n'

>>> subprocess.check_output("exit 1", shell=True)
Traceback (most recent call last):
   ...
subprocess.CalledProcessError: Command 'exit 1' returned non-zero exit status 1



#class subprocess.Popen(args, bufsize=0, executable=None, stdin=None, stdout=None, stderr=None, preexec_fn=None, close_fds=False, shell=False, cwd=None, env=None, universal_newlines=False, startupinfo=None, creationflags=0)
#subprocess.check_output(args, *, stdin=None, stderr=None, shell=False, universal_newlines=False)

Popen.poll()
Check if child process has terminated. Set and return returncode attribute.

Popen.wait(timeout=None)
Wait for child process to terminate. Set and return returncode attribute.

Popen.send_signal(signal)
Sends the signal signal to the child.

Popen.terminate()
Stop the child. 

Popen.kill()
Kills the child

Popen.args
Popen.stdin  #if stdin argument was PIPE, it is file object, use write()
Popen.stdout #If the stdout argument was PIPE, it is file object, use read()
Popen.stderr
Use communicate() rather than .stdin.write, .stdout.read or .stderr.read to avoid deadlocks 

Popen.pid
The process ID of the child process.

Popen.returncode

(stdoutdata, stderrdata) = Popen.communicate(input=None, timeout=None) 
Send data to stdin. Read data from stdout and stderr, until end-of-file is reached
The data read is buffered in memory, so do not use this method if the data size is large or unlimited.

proc = subprocess.Popen(...)
try:
    outs, errs = proc.communicate(timeout=15)
except TimeoutExpired:
    proc.kill()
    outs, errs = proc.communicate()

#works with with


with Popen(["ifconfig"], stdout=PIPE) as proc:
    log.write(proc.stdout.read())

	
	
#shlex.split() can be useful when determining the correct tokenization for args, especially in complex cases
 

>>> import shlex, subprocess

>>> command_line = input()
/bin/vikings -input eggs.txt -output "spam spam.txt" -cmd "echo '$MONEY'"

>>> args = shlex.split(command_line)

>>> print(args)
['/bin/vikings', '-input', 'eggs.txt', '-output', 'spam spam.txt', '-cmd', "echo '$MONEY'"]
>>> p = subprocess.Popen(args)



#Replacing /bin/sh shell backquote
output=`mycmd myarg`
# becomes
output = check_output(["mycmd", "myarg"])

# Replacing shell pipeline
#cat regex.py | grep def

p1 = Popen(["cat", "regex.py"], stdout=PIPE)
p2 = Popen(["grep", "def"], stdin=p1.stdout, stdout=PIPE)
p1.stdout.close()  ## Allow p1 to receive a SIGPIPE if p2 exits
output_str = p2.communicate()[0]  # or p2.stdout.read().decode("utf-8")
print(output_str.decode("utf-8"))

#OR

output=`dmesg | grep hda`
# becomes
output=check_output("dmesg | grep hda", shell=True)

#Replacing os.system()

sts = os.system("mycmd" + " myarg")
# becomes
sts = call("mycmd" + " myarg", shell=True)

#or 

try:
    retcode = call("mycmd" + " myarg", shell=True)
    if retcode < 0:
        print("Child was terminated by signal", -retcode, file=sys.stderr)
    else:
        print("Child returned", retcode, file=sys.stderr)
except OSError as e:
    print("Execution failed:", e, file=sys.stderr)



#Replacing the os.spawn family
#P_NOWAIT example:
pid = os.spawnlp(os.P_NOWAIT, "/bin/mycmd", "mycmd", "myarg")
#becomes
pid = Popen(["/bin/mycmd", "myarg"]).pid


#P_WAIT example:
retcode = os.spawnlp(os.P_WAIT, "/bin/mycmd", "mycmd", "myarg")
#becomes
retcode = call(["/bin/mycmd", "myarg"])


#Vector example:
os.spawnvp(os.P_NOWAIT, path, args)
#becomes
Popen([path] + args[1:])


#Environment example:
os.spawnlpe(os.P_NOWAIT, "/bin/mycmd", "mycmd", "myarg", env)
#becomes
Popen(["/bin/mycmd", "myarg"], env={"PATH": "/usr/bin"})



#Replacing os.popen(), os.popen2(), os.popen3()
(child_stdin, child_stdout) = os.popen2(cmd, mode, bufsize)
#becomes
p = Popen(cmd, shell=True, bufsize=bufsize,
          stdin=PIPE, stdout=PIPE, close_fds=True)
(child_stdin, child_stdout) = (p.stdin, p.stdout)



(child_stdin,
 child_stdout,
 child_stderr) = os.popen3(cmd, mode, bufsize)
#becomes
p = Popen(cmd, shell=True, bufsize=bufsize,
          stdin=PIPE, stdout=PIPE, stderr=PIPE, close_fds=True)
(child_stdin,
 child_stdout,
 child_stderr) = (p.stdin, p.stdout, p.stderr)



(child_stdin, child_stdout_and_stderr) = os.popen4(cmd, mode, bufsize)
#becomes
p = Popen(cmd, shell=True, bufsize=bufsize,
          stdin=PIPE, stdout=PIPE, stderr=STDOUT, close_fds=True)
(child_stdin, child_stdout_and_stderr) = (p.stdin, p.stdout)


#Return code handling translates as follows:
pipe = os.popen(cmd, 'w')
...
rc = pipe.close()
if rc is not None and rc >> 8:
    print("There were some errors")
==>
process = Popen(cmd, stdin=PIPE)
...
process.stdin.close()
if process.wait() != 0:
    print("There were some errors")






	
#Module - pexpect for linux for py3.x or p2.7

#for windows-There are some ports with PyWin32 
#for Py2.7, wexpect (downloaded win-pexpect/) cd and start py prompt

#Windows SSH client and server
#http://www.mls-software.com/opensshd.html
#with older SSH server 
#ssh -oKexAlgorithms=+diffie-hellman-group1-sha1 ftpuser@127.0.0.1

#example -cygwin


#start ssh service , from services.msc for service 'CYGWIN sshd'


#pexpect4.0 works on  Windows and POSIX systems. but not  pexpect.spawn and pexpect.run()
#use pexpect.popen_spawn.PopenSpawn, or a file descriptor, using pexpect.fdpexpect.fdspawn (experimental)
#class pexpect.popen_spawn.PopenSpawn(cmd, timeout=30, maxread=2000, searchwindowsize=None, logfile=None, cwd=None, env=None, encoding=None, codec_errors='strict')


#for linux/cygwin
#main methods are 

#pexpect.spawn(command, args=[], timeout=30, maxread=2000, searchwindowsize=None, logfile=None, cwd=None, env=None, ignore_sighup=False, echo=True, preexec_fn=None, encoding=None, codec_errors='strict', dimensions=None)

Pexpect does NOT interpret shell meta characters such as redirect, pipe, or wild cards , hence, use shell
child = pexpect.spawn('/bin/bash -c "ls -l | grep LOG > logs.txt"')
child.expect(pexpect.EOF)

'searchwindowsiz'e is None (default), the full buffer is searched at each iteration of receiving incoming data

'logfile' member turns on or off logging. 
All input and output will be copied to the given file object. 
Set logfile to None to stop logging. This is the default. 
Set logfile to sys.stdout to echo everything to standard output. The logfile is flushed after each write.

child = pexpect.spawn('some_command')
fout = open('mylog.txt','wb')
child.logfile = fout

# In Python 2:
child = pexpect.spawn('some_command')
child.logfile = sys.stdout

# In Python 3, spawnu should be used to give str to stdout:
child = pexpect.spawnu('some_command')
child.logfile = sys.stdout


#You only want to log what the child sends back. for Py3.x , pass an encoding to spawn 
child = pexpect.spawn('some_command')
child.logfile_read = sys.stdout

#To separately log output sent to the child use logfile_send:
child.logfile_send = fout

'delaybeforesend' uses a delay before sending , use 0.05 if some delay is required 

To get the exit status of the child , call close()and check self.exitstatus or self.signalstatus

'echo' attribute may be set to False to disable echoing of input

#expect(pattern, timeout=-1, searchwindowsize=-1, async=False)
pattern can be a string(re pattern), pexpect.EOF, pexpect.TIMEOUT, a compiled re, 
or a list of any of those types for any match 

This returns the index into the pattern list or 0 if only one pattern given 
This may raise exceptions for pexpect.EOF or pexpect.TIMEOUT. 
To avoid the EOF or TIMEOUT exceptions add EOF or TIMEOUT to the pattern list

When a match is found for the given pattern, 
'match' becomes an re.MatchObject result or exception
'before' and 'after' are before and after match

#example
child.expect('password:')
child.sendline(my_secret_password)
# We expect any of these three patterns...
i = child.expect (['Permission denied', 'Terminal type', '[#\$] '])
if i==0:
    print('Permission denied on host. Can\'t login')
    child.kill(0)
elif i==1:
    print('Login OK... need to send terminal type.')
    child.sendline('vt100')
    child.expect('[#\$] ')
elif i==2:
    print('Login OK.')
    print('Shell command prompt', child.after)


On Python3.4, , passing async=True will make this return an asyncio coroutine, 
index = yield from p.expect(patterns, async=True)


#expect_exact(pattern_list, timeout=-1, searchwindowsize=-1, async=False) 
pattern_list is list of literal string

#send(s)
#sendline(s='')
Sends string s to the child process, returning the number of bytes written. 

#writelines(sequence)
#write(s)
This is similar to send() except that there is no return value.


#sendcontrol(char) 
child.sendcontrol('g')

#sendeof() 
This sends an EOF to the child

#sendintr()
This sends a SIGINT to the child. 

#readline(size=-1)
#read(size=-1) 
This reads at most "size" bytes from the file 

#eof()
This returns True if the EOF exception was ever raised.

#interact(escape_character='\x1d', input_filter=None, output_filter=None)
This gives control of the child process to the interactive user (the human at the keyboard).
 
 
#Controlling the child process from spawn
kill(sig)
terminate(force=False)
isalive()
wait()
close

#Handling unicode
By default, spawn is a bytes interface
With 'encoding' strings sent or bytes recieved will be encoded or decoded  using that encoding
In Pexpect 3.x, use spawnu to have unicode interface

#pexpect.run(command, timeout=30, withexitstatus=False, events=None, extra_args=None, logfile=None, cwd=None, env=None, **kwargs)[source] 
This function runs the given command; waits for it to finish; then returns all output as a string. 

run("ssh username@machine.example.com 'ls -l'",  events={'(?i)password':'secret\n'})


#Find the end of line – CR/LF conventions
The $ pattern for end of line match is useless, Use below
child.expect('\r\n')  
child.expect('\w+\r\n')

Pexpect compiles all regular expressions with the re.DOTALL flag. 
With the DOTALL flag, a "." will match a newline

#+ and * at the end of patterns - in pexpect, always non greedy

#match  one character
child.expect ('.+')

#match no characters
child.expect ('.*')


#Debugging
use str(child) where child = spawn(..)

	

#Example
import pexpect

pexpect.run('ls -l') # in py3.x , use .decode("ascii"), each line ends with '\r\n'

#using spawn  , can use spawnu for Pyx as unitcode version 

child = pexpect.spawn('scp open.pl ftpuser@localhost:.') #moves var.txt to cygwin's ftpuser 
child.expect ('password: ')
child.sendline ('ftpuser')

#for windows with wexpect (cd to win-pexpect)
#wexpect.py needs modification:line 2397, make Y coordinate of rect small eg 24

import wexpect
child = wexpect.spawn("ssh.exe", ['ftpuser@localhost', 'ls -al'], timeout=60)
child.expect ('password: ')
child.sendline ('ftpuser')
for line in child:
	print(">>" + line.decode("utf-8").strip())
	
child.close()



#ftplib Module - windows, linux , py2.x, py3.x 
#start the server from computer management
#IIS
#%windir%\system32\compmgmt.msc
#check ipconfig and change binding for correct ip



import ftplib
import os
filename = ".bashrc"
ftp = ftplib.FTP("192.168.1.106")
ftp.login("ftpuser", "ftpuser")
os.chdir(r'/home/das')    #from windows c:/cygwin64/home/das'
ftp.retrlines('LIST')   # list directory contents
ftp.nlst()  			#['.bashrc', 'class-ex.py', 'one.txt']

ftp.cwd("dump")
ftp.storbinary("STOR " + filename , open(filename, 'rb'))    # for uploading binary file, 
ftp.storlines("STOR " + filename + ".t" , open(filename, 'rb'))    # for uploading text file, always 'b' 
ftp.retrbinary('RETR ' + filename, open(filename + ".bak" , 'wb').write) #for downloading binary file
ftp.retrlines('RETR ' + filename, open(filename + ".bak2" , 'w').write) #for downloding text file

data = []
ftp.dir(data.append)
print("\n".join(data))
ftp.quit()


#Telnetlib - py3.x, py2.x, 
#for example start Telnet service from services.msc
#Py3.x - Always  interaction using bytes, hence encode/decode with 'ascii' or use b' '
#ending line by '\r\n'

Telnet.read_until(expected, timeout=None)
Read until a given byte string, expected, is encountered or until timeout seconds have passed.

Telnet.read_all()
Read all data until EOF as bytes; block until connection closed.

Telnet.read_some()
Read at least one byte of cooked data(ie IAC processed) unless EOF is hit. Return b'' if EOF is hit. 

Telnet.read_very_eager()
Read everything that can be without blocking in I/O (eager).

Telnet.read_eager()
Read readily available data.

Telnet.read_lazy()
Process and return data already in the queues (lazy).

Telnet.read_very_lazy()
Return any data available in the cooked queue (very lazy).

Telnet.read_sb_data()
Return the data collected between a SB/SE pair (suboption begin/end).

Telnet.write(buffer)
Write a byte string to the socket, doubling any IAC characters. 

(index_of_match, matchObject, data_till_match) = Telnet.expect(list, timeout=None)
Read until one from a list of a regular expressions matches ie compiled or byte string
If a regular expression ends with a greedy match (such as .*) or if more than one expression can match the same input, 
the results are non-deterministic
When nothing matches, return (-1, None, data)


#details of Telnet 
'''
telnet mode, as described in RFC854
Lines are expected to end with either '\r\n' or '\r\0', or '\n'  and ASCII 255 is used for telnet control codes
IAC in data is doubled , inserts \0 after \r  

Raw mode - basically raw tcpip socket , acts as a transparent bridge, transmitting all bytes across the socket unmodified
lines end with the ASCII NUL character \0, and no control codes are present. IAC in data is not doubled , no insertion of \0 after \r
'''

#Explains nicely options :  https://support.microsoft.com/en-us/kb/231866
  
#Options : http://www.iana.org/assignments/telnet-options/telnet-options.xhtml
'''
Senders wants to do an option  IAC WILL opt   
		receiver responds IAC DO opt or IAC DONT opt
sender asks receiver to do an options IAC DO opt    
		receiver responds IAC WILL opt or IAC WONT opt
'''

#Telnetlib has constants for options (check telnetlib.py)
'''
telnetlib.IAC  
telnetlib.DONT 
telnetlib.DO   
telnetlib.WONT 
telnetlib.WILL 
telnetlib.ECHO    #echoing data characters it receives over the   TELNET connection back to the sender of the data characters
telnetlib.SGA     #Supress GA  for full duplex operation.  by default 
'''
#setting debug level 
Telnet.set_debuglevel(debuglevel)
Set the debug level. >= 1 to get debug output  (on sys.stdout).

#installing callback for options 
Telnet.set_option_negotiation_callback(callback)
Each time a telnet option is read on the input flow, this callback (if set) is called with the following parameters: 
callback(telnet_socket, command (DO/DONT/WILL/WONT), option). 


# Example : How to disable telnet echo in python telnetlib
'''
telnetlib.py automatically responds to IAC commands . (See telnetlib.process_rawq()) if callback is not set 
If telnetlib gets WILL, sends DONT and if gets DO, sends WONT for the same option automatically 
(t.set_debuglevel(1) to get many outputs) 

Hence , telnetlib sends IAC DONT ECHO whenever gets IAC WILL ECHO
However, it might not be enough to turn off echo .
The solution most commonly used is  to say that reciever will do the echoing, which stops other end  doing echoing:
telnetlib.IAC + telnetlib.WILL + telnetlib.ECHO
But in this case, you need to echo back whatever you receive
'''

# set call back
t.set_option_negotiation_callback(callback)

#callback 
def callabck(sock, cmd, opt):
	if cmd == telnetlib.WILL and opt == telnetlib.ECHO:        #Senders  WILL ECHO
		sock.sendall(telnetlib.IAC + telnetlib.DONT + telnetlib.ECHO)
		sock.sendall(telnetlib.IAC + telnetlib.WILL + telnetlib.ECHO) 
	elif opt == telnetlib.ECHO and 	cmd == telnetlib.DO :   # server would respond back with DO for ECHO, ignore that 
		sock.sendall(telnetlib.IAC + telnetlib.WILL + telnetlib.ECHO) 
	else:
		#default handling 
		if cmd in (telnetlib.DO, telnetlib.DONT):
			sock.sendall(telnetlib.IAC + telnetlib.WONT + opt)
		elif cmd in (telnetlib.WILL, telnetlib.WONT):
			sock.sendall(telnetlib.IAC + telnetlib.DONT + opt)


#Write raw sequence
#write - Write a string to the socket, doubling any IAC characters.
#to write raw sequence , get inner socket 

def write_raw_sequence(tn, seq):
	sock = tn.get_socket()
	if sock is not None:
		sock.sendall(seq)

write_raw_sequence(tn, telnetlib.IAC + telnetlib.WILL + telnetlib.ECHO)


#example
import sys
import telnetlib
user = "ftpuser"
passw = "ftpuser"

tn = telnetlib.Telnet("localhost", 23)
print(tn.read_until(b"login: "))
tn.write(user.encode('ascii') + b"\r\n")

print(tn.read_until(b"password: "))
tn.write(passw.encode('ascii') + b"\r\n")
print(tn.read_until(b"ftpuser>"))
tn.write(b"dir\r\n")

print(tn.read_until(b"ftpuser>").decode("ascii"))
tn.write(b"exit\r\n")
tn.close()

#py2.7
import sys
import telnetlib
user = "ftpuser"
passw = "ftpuser"

tn = telnetlib.Telnet("localhost", 23)
print tn.read_until("login: ")

tn.write(user + "\r\n")
print tn.read_until("password: ") 
tn.write(passw + "\r\n")
print tn.read_until("ftpuser>")
tn.write("dir\r\n")

print tn.read_until("ftpuser>")
tn.write("exit\r\n")
tn.close()





#paramiko module - (2.6+, 3.3+) implementation of the SSHv2 protocol , provides server and client
#for windows pycrypto-3.4 use https://github.com/axper/python3-pycrypto-windows-installer
#or compile through VC 
#pip3 install paramiko
#linux/cygwin OK with compilation

# for example start cygwin ssh from services.msc



import paramiko
paramiko.common.logging.basicConfig(level=paramiko.common.DEBUG)
ssh    = paramiko.SSHClient()
ssh.set_missing_host_key_policy(paramiko.AutoAddPolicy())
ssh.connect("127.0.0.1", username="ftpuser", password="ftpuser")
i, o, e = ssh.exec_command("ls -l")
i.flush()
print(" ".join(o.readlines()))
ssh.close()

#for file transfer

import paramiko
ssh = paramiko.SSHClient()
ssh.set_missing_host_key_policy( paramiko.AutoAddPolicy())
ssh.connect("localhost", username="ftpuser", password="ftpuser")
ftp = ssh.open_sftp()
ftp.put('var.txt', 'remotefile.py')
ftp.close()
ftp = ssh.open_sftp() 
ftp.get('remotefile.py', 'localfile.py') 
ftp.close() 


		

#Automating HTTP using urllib, urllib2 and httplib
#Py2.x - urllib and urllib2 , renamed in Python 3 to urllib.request, urllib.parse, and urllib.error. 
#Py3.x urllib.request.urlopen() is equivalent to urllib2.urlopen() and urllib.urlopen() has been removed
#Use apache requests for high level API

•urllib.request for opening and reading URLs for file, ftp and http, raises urllib.error
•urllib.parse for parsing URLs
•urllib.robotparser for parsing robots.txt files

#urllib.request.urlopen(url, data=None, [timeout, ]*, cafile=None, capath=None, cadefault=False, context=None)
provide data (must be in bytes) for http POST, use urllib.parse.urlencode({k:v,..})

#example - read() gives in bytes

import urllib.request
with urllib.request.urlopen('http://www.python.org/') as f:
	print(f.read(100).decode('utf-8'))  #can use read() and readline() or can be used as iterator


#get method

import urllib.request
import urllib.parse
params = urllib.parse.urlencode({'spam': 1, 'eggs': 2, 'bacon': 0})
url = "http://www.musi-cal.com/cgi-bin/query?%s" % params
with urllib.request.urlopen(url) as f:
	print(f.read().decode('utf-8'))



# POST method instead. 

import urllib.request
import urllib.parse
data = urllib.parse.urlencode({'spam': 1, 'eggs': 2, 'bacon': 0})
data = data.encode('ascii')
with urllib.request.urlopen("http://requestb.in/xrbl82xr", data) as f:
	print(f.read().decode('utf-8'))


#parse method
from urllib.parse import urlparse
o = urlparse('http://www.cwi.nl:80/%7Eguido/Python.html')
>> o
ParseResult(scheme='http', netloc='www.cwi.nl:80', path='/%7Eguido/Python.html',
            params='', query='', fragment='')
>>> o.scheme
'http'
>>> o.port
80
>>> o.geturl()
'http://www.cwi.nl:80/%7Eguido/Python.html'

#Other methods 
urllib.parse.quote(string, safe='/', encoding=None, errors=None)
Replace special characters in string using the %xx escape.
Example: quote('/El Niño/') yields '/El%20Ni%C3%B1o/'.

urllib.parse.quote_plus(string, safe='', encoding=None, errors=None)
Like quote(), but also replace spaces by plus signs
as required for quoting HTML form values when building up a query string to go into a URL
Example: quote_plus('/El Niño/') yields '%2FEl+Ni%C3%B1o%2F'.

urllib.parse.unquote(string, encoding='utf-8', errors='replace')
Replace %xx escapes by their single-character equivalent

urllib.parse.unquote_plus(string, encoding='utf-8', errors='replace')
Like unquote(), but also replace plus signs by spaces

urllib.parse.urlencode(query, doseq=False, safe='', encoding=None, errors=None)
Convert a mapping object or a sequence of two-element tuplesConvert a mapping object 
or a sequence of two-element tuples, which may contain str or bytes objects, to a percent-encoded ASCII text string. 
If the resultant string is to be used as a data for POST operation with the urlopen() function, then it should be encoded to 
bytes, otherwise it would result in a TypeError.
The resulting string is a series of key=value pairs separated by '&' characters, 
where both key and value are quoted using quote_plus() above
To reverse this encoding process, Use parse_qs() and parse_qsl() 

urllib.parse.parse_qs(qs, keep_blank_values=False, strict_parsing=False, encoding='utf-8', errors='replace')
Parse a query string given as a string argument (data of type application/x-www-form-urlencoded). 
Data are returned as a dictionary. The dictionary keys are the unique query variable names and the values are lists of values for each name.

urllib.parse.parse_qsl(qs, keep_blank_values=False, strict_parsing=False, encoding='utf-8', errors='replace')
Parse a query string given as a string argument (data of type application/x-www-form-urlencoded). 
Data are returned as a list of name, value pairs

urllib.parse.urlsplit(urlstring, scheme='', allow_fragments=True)
This is similar to urlparse(), but does not split the params from the URL

urllib.parse.urlparse(urlstring, scheme='', allow_fragments=True)
Parse a URL into six components, returning a 6-tuple. This corresponds to the general structure of a URL: scheme://netloc/path;parameters?query#fragment

urllib.parse.urljoin(base, url, allow_fragments=True)
Construct a full (“absolute”) URL by combining a “base URL” (base) with another URL (url). 


#fetching file from localhost only

import urllib.request
f = urllib.request.urlopen('file://localhost/d:/Desktop/PPT/var.txt') 
f.read(100).decode('ascii')  #or iterate or call readlines() which returns list of b'str'
f.close()

#ftp interface ftp://user:password@host:port/path 
import urllib.request
f = urllib.request.urlopen('ftp://ftpuser:ftpuser@localhost/' )
f.read(100).decode('ascii')  #or iterate or call readlines() which returns list of b'str'
f.close()


#Py2.x httplib module has been renamed to http.client in Python 3
#don't use directly, urllib.request uses httplib to handle URLs that use HTTP and HTTPS
#use 'pip3 install requests' and use requests package in real code



###Apache requests 

#basic usgae 

p1 = 'http://proxy_username:proxy_password@proxy_server.com:port'
p2 = 'https://proxy_username:proxy_password@proxy_server.com:port'
proxy = {'http': p1, 'https':p2}
r = requests.get(site, proxies=proxy, auth=('site_username', 'site_password'))

#example
import requests
r = requests.get("http://www.yahoo.com")
r.text
r.status_code   # status code
r.headers  		# dict object

#RESTful API
r = requests.post(site)
r = requests.put("site/put")
r = requests.delete("site/delete")
r = requests.head("site/get")
r = requests.options("site/get")


#Get

payload1 = {'key1': 'value1', 'key2': 'value2'}
r = requests.get("http://httpbin.org/get", params=payload1)
print(r.url)  #http://httpbin.org/get?key2=value2&key1=value1
r.headers
r.text
r.json()  # it's a python dict

#For Request debugging,
>>> r.request.url
'http://httpbin.org/forms/post?delivery=12&topping=onion&custtel=123&comments=ok&custname=das&custemail=ok%40com&size=small'
>>> r.request.headers
{'Content-Length': '0', 'User-Agent': 'Mozilla/5.0', 'Connection': 'keep-alive', 'Accept': '*/*', 'Accept-Encoding': 'gzip, deflate'}
>>> r.request.body


#POST
headers = {'User-Agent': 'Mozilla/5.0'}
payload = {'custname':'das', 'custtel': '123', 'custemail' : 'ok@com', 'size':'small',  'topping':'bacon',  'topping': 'onion',  'delivery':'12', 'comments': 'ok'}
r = requests.post("http://httpbin.org/post", data=payload, headers=headers)
r.text
r.headers
r.json()

r.request.headers
r.request.body #custname=das&custtel=123&custemail=ok@com&size=small&topping=bacon&topping=onion&delivery=12&comments=ok

#Content
r.text
r.content  # as bytes
r.json()  # json content

#Example to handle image Py3- http://www.lfd.uci.edu/~gohlke/pythonlibs/#pillow , for py2, http://www.pythonware.com/products/pil/

from PIL import Image
from io import StringIO   
i = Image.open(StringIO(r.content))




#Custom Headers

import json
payload = {'some': 'data'}
headers = {'content-type': 'application/json'}
r = requests.post(url, data=json.dumps(payload), headers=headers)

#POST a Multipart-Encoded File

files = {'file': open('report.xls', 'rb')}
r = requests.post(url, files=files)

#Cookies
r.cookies['example_cookie_name']

#or sending
cookies = dict(cookies_are='working')
r = requests.get(url, cookies=cookies)


#Or persisting across session
s = requests.Session()

s.get('http://httpbin.org/cookies/set/sessioncookie/123456789')
r = s.get("http://httpbin.org/cookies")

r.text # contains cookies from first access

#Example:
import requests
headers = {'User-Agent': 'Mozilla/5.0'}
payload = {'username':'niceusername','pass':'123456'}

session = requests.Session()
session.post('https://admin.example.com/login.php',headers=headers,data=payload)
# the session instance holds the cookie. So use it to get/post later.
# e.g. session.get('https://example.com/profile')

#Authentication with form
response = requests.get(url, auth = ('username', 'password')) 


#Example - Form conatins

<textarea id="text" class="wikitext" name="text" cols="80" rows="20">
This is where your edited text will go
</textarea>
<input type="submit" id="save" name="save" value="Submit changes">

#Code:

import requests
from bs4 import BeautifulSoup

url = "http://www.someurl.com"

username = "your_username"
password = "your_password"

session = requests.Session()
session.auth = (username, password)
response = session.get(url, verify=False)

# Getting the text of the page from the response data       
page = BeautifulSoup(response.text)

# Finding the text contained in a specific element, for instance, the 
# textarea element that contains the area where you would write a forum post
txt = page.find('textarea', id="text").string

# Finding the value of a specific attribute with name = "version" and 
# extracting the contents of the value attribute
tag = page.find('input', attrs = {'name':'version'})
ver = tag['value']

# Changing the text to whatever you want
txt = "Your text here, this will be what is written to the textarea for the post"

# construct the POST request
form_data = {
    'save' : 'Submit changes'
    'text' : txt
} 

post = session.post(url,data=form_data,verify=False)









###Mail handling using smtplib and poplib
#to use google access, make it less secure https://www.google.com/settings/security/lesssecureapps

Sending a mail
1. Create a message using email.mime.text.MIMEText
2. Create smptp object via SMTP and sendmail()

#example:
import smtplib
import email.utils
from email.mime.text import MIMEText
import getpass

# Create the message
msg = MIMEText('This is the body of the message.')
msg['To'] = email.utils.formataddr(('Recipient', 'ndas1971@gmail.com'))
msg['From'] = email.utils.formataddr(('Author', 'ndas1971@gmail.com'))
msg['Subject'] = 'Simple test message'


server = smtplib.SMTP('smtp.gmail.com', 587)  
server.set_debuglevel(True) # show communication with the server
try:
    # identify ourselves, prompting server for supported features
	server.ehlo()
	# If we can encrypt this session, do it
	if server.has_extn('STARTTLS'):
		server.starttls()
		server.ehlo() # re-identify ourselves over TLS connection
	
	p = getpass.getpass()
	server.login('ndas1971@gmail.com', p)
	server.sendmail('ndas1971@gmail.com', ['ndas1971@gmail.com',], msg.as_string())
finally:
    server.quit()

	
	
	
#Download message using poplib - server mmust support POP3



import getpass, poplib
import email
from io import StringIO   # for py2.x cStringIO
from email.generator import Generator


user = 'ndas1971@gmail.com' 

Mailbox = poplib.POP3_SSL('pop.gmail.com', '995') 
Mailbox.user(user) 

p = getpass.getpass()

Mailbox.pass_(p) 
numMessages = len(Mailbox.list()[1])

#(numMsgs, totalSize) = Mailbox.stat()


print(" No of messages=%s" % numMessages )

for i in range(5):
	msg = Mailbox.retr(i+1)[1]
	f_msg = email.message_from_bytes(b"\n".join(msg))   #in py2.x, use _string 
	for header in [ 'subject', 'to', 'from' ]:
		print('%-8s: %s' % (header.upper(), f_msg[header]))
	print("\n")

# complete message

num = int(input("Input msg # for seeing complete message>"))
msg = Mailbox.retr(num+1)[1]
c_msg = email.message_from_bytes(b"\n".join(msg))  #in py2.x, use _string 
fp = StringIO()
g = Generator(fp, mangle_from_=False, maxheaderlen=60)
g.flatten(c_msg)
text = fp.getvalue()
print(text)

Mailbox.quit()





#Fabric only for SSH - Py2.7 (install: paramiko at first)(cygwin or windows)
Fabric is a Python (2.5-2.7) library and command-line tool for streamlining 
the use of SSH for application deployment or systems administration tasks
#pip2 install fabric

#Execute by 'fab' command

env is a global dictionary-like object , has below important members
env.user
env.password
env.warn_only

#Note 
1) Use "run()" when you don't expect a prompt. eg run("sudo blah") instead of sudo("blah"), 
2) use "sudo()" when you do expect a prompt.

#Use context of env
from __future__ import with_statement
from fabric.api import local, settings, abort, run, cd
from fabric.contrib.console import confirm


def deploy():
	code_dir = '/srv/django/myproject'
	with settings(warn_only=True):
		if run("test -d %s" % code_dir).failed:
			run("git clone user@vcshost:/path/to/repo/.git %s" % code_dir)
	with cd(code_dir):
		run("git pull")
		run("touch app.wsgi")


		

#file name must be:  fabfile.py 
from fabric.api import env, run, local, cd

env.hosts = ['ftpuser@localhost', 'linuxbox']  		# format username@hostname:port
env.password = 'ftpuser'
def host_type(name="OK"):  	# just to showcase arg
	print("Executing on %s as %s" % (env.host, env.user))   #print() even in Py2.7
    run('uname -s')

def test():
    local("./manage.py test my_app")

def commit():
    local("git add -p && git commit")

def push():
    local("git push")

def prepare_deploy():
    test()
    commit()
    push()

#Execute by 
$ fab  host_type:name="NOK"

# command line Host option is given by -H localhost,linuxbox after removing env.hosts


#Running local command -can be done in interpretor
from fabric.api import local
local("ls-l")

#Complex example

from fabric.api import *

@hosts('username@host1:port')						# can specify host like this as well
def clean_and_upload():
    local('find assets/ -name "*.DS_Store" -exec rm '{}' \;')
    local('tar czf /tmp/assets.tgz assets/')
    put('/tmp/assets.tgz', '/tmp/assets.tgz')   # putting it in host1
    with cd('/var/www/myapp/'):					# remote
        run('tar xzf /tmp/assets.tgz')			# remote

		
		
#in py script
from fabric.api import run
from fabric.tasks import execute
from fabric.state import connections


def do_something():
    run("echo $RANDOM")


execute(do_something, hosts=["ftpuser@localhost"])
for key in connections.keys():
	connections[key].close()
	del connections[key]



	
	

#Exscript Module - Py2.7 (install Paramiko)(linux/cygwin only)
#Very good for ios and other defined drivers. Has template language 
#pip2 install exscript

#execute by
$ exscript --help
exscript [options] filename [hostname [hostname ...]]

#main class
SSH2_or_Telnet.__init__(self, driver=None, stdout=None, stderr=None, debug=0, connect_timeout=30, timeout=30, logfile=None, termtype='dumb', verify_fingerprint=True, account_factory=None)

#for logging
conn = SSH2(stdout=sys.stdout, debug =5, stderr=sys.stdout, logfile="exscript.log")  

#for ios 
conn = SSH2(driver="ios")
conn.autoinit() #for ios device

#drivers list
http://knipknap.github.io/exscript/api/Exscript.protocols-module.html


#sending
send(data) #Sends the given data to the remote host. Returns without waiting for a response.
execute(command) # send(command + '\r'); expect_prompt()

#for setting and getting various prompt
set_error_prompt(error=None)  #The pattern that, when matched, causes an error
set_login_error_prompt(error=None) #The pattern that, when matched, causes an error.
set_password_prompt(regex=None) 
set_prompt(self, prompt=None)  
 

#waiting for prompt
waitfor(prompt)  #buffer is not modified to remove prompt
expect(prompt) 	#buffer is  modified to remove prompt
expect_prompt()  # buffer is modified to remove prompt

#response 
conn.response   # Note it contains echo of command sent as well as prompt
__response__    #in template file, note there is no way to print , but can do echo $var



#General purpose SSH device 
from Exscript.util.interact import read_login
from Exscript.protocols import *
from Exscript import Account
import sys

account = Account('ftpuser','ftpuser')     #  use read_login()
conn = SSH2()               
conn.connect('127.0.0.1')  # Open the SSH connection
conn.set_prompt('\$')       #must to set prompt , Note this is RE, hence escape any RE character
conn.login(account)        # Authenticate on the remote host or 
conn.execute('dir')        # Execute the "uname -a" command
print(conn.response)              #Note many $PS1 prompt byte may come, set ftpuser's .profile file with PS1='$ ';export PS1
#or
conn.send('ls -l' + '\r\n')        # Execute the "uname -a" command
conn.waitfor('\$')
print(conn.response) 
conn.send('exit\r\n')       # Must have \r\n at the end 
conn.close()               # Wait for the connection to close

#for Telnet - genreal purpose Telnet device 

account = Account('ftpuser','ftpuser')     
conn = Telnet()             
conn.connect('127.0.0.1')  # Open the SSH connection
conn.set_prompt('>')
conn.login(account)
conn.execute('dir')  # conn.expect_prompt() automatically 
conn.response.split('\r\n') #works
conn.send('exit\r\n')
conn.close()   

#Using quickstart(asks prompts for user/pass) or start(no prompt for user/pass)
#Only for all given drivers - http://knipknap.github.io/exscript/api/Exscript.protocols-module.html

from Exscript.util.start import quickstart  
def do_something(job, host, conn):
    conn.execute('uname -a')

quickstart('ssh://localhost', do_something) # it propmts for username/password


#Running a script on multiple hosts

from Exscript.util.start import quickstart

def do_something(job, host, conn):
    conn.execute('uname -a')

hosts = ['ssh://localhost', 'telnet://myuser:mypassword@anotherhost']  # can give username password for telnet
# or load from file by hosts = get_hosts_from_file('myhosts.txt')
quickstart(hosts, do_something, max_threads = 2, logdir = './logs/')

#setting username, password programmatically
from Exscript import Host, Account
from Exscript import Host, Account

account1 = Account('myuser', 'mypassword')
host1    = Host('ssh://localhost')
host1.set_account(account1)

account2 = Account('myuser2', 'mypassword2')
host2    = Host('ssh://otherhost')
host2.set_account(account2)

quickstart([host1 , host2], do_something, max_threads = 2)
#OR
accounts = []  #accounts = get_accounts_from_file('accounts.cfg') , defined syntax
start(accounts, [host1 , host2], do_something, max_threads = 
#with loging
log=open('log.txt','w+')
start(accounts, hosts, do_something, stdout=log)

#Interacting with a device
from Exscript.util.start import quickstart
from Exscript.util.match import first_match

def do_something(job, host, conn):
    conn.execute('uname -a')
    print "The response was", repr(conn.response)
    os, hostname = first_match(conn, r'^(\S+)\s+(\S+)')
    print "The hostname is:", hostname
    print "Operating system:", os
	

quickstart('ssh://localhost', do_something)

#Using a file 
#file: test.exscript 
uname -a

#execute , -c 2 means two threads 
exscript -c 2 test.exscript ssh://localhost  telnet://localhost
exscript -c 2 test.exscript ssh://localhost ssh://user:password@otherhost
exscript -c 2 --hosts myhosts.txt test.exscript  #hostname from hostfile
exscript -l /tmp/logs -c 2 --hosts myhosts.txt --account-pool accounts.cfg test.exscript #logging

#Interacting with a device
uname -a{extract /^(\S+)\s+(\S+)/ as os, hostname}


#The Exscript Template Language - study more from https://github.com/knipknap/exscript/downloads
#exscript -d var=value -d var=value my_script.exscript host1 host2
#exscript --account-pool account.txt test.exscript ssh://localhost

#there is no need of print as each command output is displayed on screen
#extract must match the group of regex with number of variables
#each line can be extracted into a list

#set_prompt is problematic for authentication, check test.exscript file


{connection.set_prompt(/#/)}

uname -a {extract /(.+)\s+(.+)(.+)\s+(.+)(.+)\s+(.+)(.+)/ as item1, item2, item3, item4, item5, item6, item7}
echo "$item1,$item2,$item3,$item4"


ls -l{
  extract /^(.*)/ into lines
  extract /^(d.*)/ into directories from lines
}
{loop directories as dir}
  echo "$dir"
{end}


ls -l .profile {extract /(\.profile)$/ as found}
{if found is not ".profile"}
ls
{end}

#various files 
#--hosts file 
#contains
host1
host2
...
host20

#csv file
#--csv-hosts
#contains - first line is key line
#separated by tab 
address	var1	var2
telnet://myhost	value	another_value
ssh://yourhost	hello	world

#accounts file 
#--account-pool
#contains user=base64 encoding password
#if many lines are there, then round robin way choosen
[account-pool]
ftpuser=ZnRwdXNlcg==
   
   
   
   
   
  
  
#Web programming and automation using Python

#Common Gateway Interface - A CGI script is invoked by an HTTP server, eg to process <FORM> or <ISINDEX> element
#cgi script - print text to display in browser

#file :  cgi-bin/cgiEx.py

import cgi
import cgitb
cgitb.enable(display=0, logdir="logs") #display detailed reports in the Web browser if any errors occur, display=0 means don't display, but write to log 
form = cgi.FieldStorage()


#headers section
print("Content-Type: text/html")    # HTML is following
print()                             # blank line, end of headers

#content section
print("<HTML>")
print("<TITLE>CGI script output</TITLE>")
print("<BODY>")
print("<H1>This is my first CGI script</H1>")
print("Hello, world!")
print(form.getfirst("name","").upper(), ",".join(form.getlist("addr")) )
print("</BODY></HTML>")

#run http server
#python3 -m http.server --bind 127.0.0.1 --cgi 8080
#serves scripts from ['/cgi-bin', '/htbin']
#in browser http://localhost:8080/cgiEx.py
#or http://localhost:8080/cgi-bin/cgiEx.py?name=Joe+Blow&addr=At+Home&addr=At+Office



#To get at submitted form data, use FieldStorage class

form = cgi.FieldStorage()   #can be indexed like dict with key as form field 
if "name" not in form or "addr" not in form:
    print("<H1>Error</H1>")
    print("Please fill in the name and addr fields.")
    return
print("<p>name:", form["name"].value)   #form["name"] is instance of FieldStorage, access value attribute 
print("<p>addr:", form.getvalue("addr", " ")) #alternate way, can pass default value if not set 
...further form processing here...

#Form fields containing empty strings are ignored and do not appear in the dictionary; 
#to keep such values, provide keep_blank_values=true while creating FieldStorage 

#to get only one value for a form field -REcommended
user = form.getfirst("user", "").upper()    

#to get list of values for  a form field

value = form.getlist("username")
usernames = ",".join(value)


#for  an uploaded file field 

fileitem = form["userfile"]
if fileitem.file:
    # It's an uploaded file; count lines
    linecount = 0
    while True:
        line = fileitem.file.readline()
        if not line: break
        linecount = linecount + 1

#Other important methods
cgi.print_environ()
Format the shell environment in HTML.

cgi.print_form(form)  #form is dict of { key:value} which will be formatted
Format a form in HTML.

cgi.print_directory()
Format the current directory in HTML.

cgi.print_environ_usage()
Print a list of useful (used by CGI) environment variables in HTML.

cgi.escape(s, quote=False)
Convert the characters '&', '<' and '>'  and quotes (if quote=True) in string s to HTML-safe sequences
Deprecated, Must use quote=True or  Use html.escape(s, quote=True) for escaping and html.unescape(s) for unescaping 

cgi.test()
Robust test CGI script, usable as main program. Writes minimal HTTP headers and formats all information provided to the script in HTML form


cgi.parse_qs(qs, keep_blank_values=False, strict_parsing=False)
Use urllib.parse.parse_qs(qs, keep_blank_values=False, strict_parsing=False, encoding='utf-8', errors='replace')
Parse a query string given as a string argument (data of type application/x-www-form-urlencoded). 
Data are returned as a dictionary

cgi.parse_qsl(qs, keep_blank_values=False, strict_parsing=False)
Use urllib.parse.parse_qsl(qs, keep_blank_values=False, strict_parsing=False, encoding='utf-8', errors='replace')
Parse a query string given as a string argument (data of type application/x-www-form-urlencoded). Data are returned as a list of name, value pairs

		
		
#Web Server Gateway Interface (WSGI)
#Used for writing a server, or a py file used by  web server eg http
#Ref implementation is wsgiref module
#details of protocol between webserver and py file is  http://www.wsgi.org

#protocol
# The application interface is a callable object
def application ( # It accepts two arguments:
    # environ points to a dictionary containing CGI like environment
    # variables which is populated by the server for each
    # received request from the client
    environ,
    # start_response is a callback function supplied by the server
    # which takes the HTTP status and headers as arguments
    start_response
):

    # Build the response body possibly
    # using the supplied environ dictionary
    response_body = 'Request method: %s' % environ['REQUEST_METHOD']

    # HTTP response code and message
    status = '200 OK'

    # HTTP headers expected by the client
    # They must be wrapped as a list of tupled pairs:
    # [(Header name, Header value)].
    response_headers = [
        ('Content-Type', 'text/plain'),
        ('Content-Length', str(len(response_body)))
    ]

    # Send them to the server using the supplied function
    start_response(status, response_headers)

    # Return the response body. Notice it is wrapped
    # in a list although it could be any iterable.
    return [response_body]   #must encode in Py3.x eg response_body.encode('utf-8')
	
	
#example:

# Python's bundled WSGI server
from wsgiref.simple_server import make_server

def application (environ, start_response):

    # Sorting and stringifying the environment key, value pairs
    response_body = [
        '%s: %s' % (key, value) for key, value in sorted(environ.items())
    ]
    response_body = '\n'.join(response_body)

    status = '200 OK'
    response_headers = [
        ('Content-Type', 'text/plain'),
        ('Content-Length', str(len(response_body)))
    ]
    start_response(status, response_headers)

    return [response_body]   #must encode in Py3.x eg response_body.encode('utf-8')

# Instantiate the server
httpd = make_server (
    'localhost', # The host name
    8051, # A port number where to wait for the request
    application # The application object name, in this case a function
)

# Wait for a single request, serve it and quit
httpd.handle_request()


#wsgiref.handlers - provides base handler classes for implementing WSGI servers and gateways. 

class wsgiref.handlers.CGIHandler
CGI-based invocation via sys.stdin, sys.stdout, sys.stderr and os.environ. 
Convert a  WSGI application to  CGI script.

#example - cgi-bin/cgihandler.py
#http://localhost:8080/cgi-bin/cgihandler.py
#python3 -m http.server --bind 127.0.0.1 --cgi 8080


from cgi import parse_qs, escape

def application(environ, start_response):
	start_response('200 OK', [('Content-Type', 'text/plain')]) 
	return [b'Hello World!\n']   # can use all functionality of WSGI

#convert above application to CGI script
import wsgiref.handlers
wsgiref.handlers.CGIHandler().run(application)


#WSGI  can be executed inside http apache using mod_wsgi
#or using uwsgi under nginx (only for unix or cygwin)

#mod_wsgi - https://modwsgi.readthedocs.org/
#The default behaviour of mod_wsgi is to create a distinct Python sub interpreter for each WSGI application

#wsgi-bin/info.wsgi


#!C:/Python34/python3.exe

def application(environ, start_response):
	response_body = ['%s: %s' % (key, value)  for key, value in sorted(environ.items())]
	response_body = '<br>\n'.join(response_body)
	status = '200 OK'
	response_headers = [('Content-Type', 'text/html'), ('Content-Length', str(len(response_body)))]
	start_response(status, response_headers)
	return [response_body.encode('utf-8')]

if __name__ == '__main__':
	import wsgiref.handlers
	wsgiref.handlers.CGIHandler().run(application)
	
	
#Hosting wsgi script in mod_wsgi	 - conf/httpd.conf
	
LoadModule wsgi_module modules/mod_wsgi.so

#can be accessed like /info
WSGIScriptAlias /info "C:/indigoampp/apache-2.2.15/wsgi-bin/info.wsgi"

<Directory "C:/indigoampp/apache-2.2.15/wsgi-bin">
   Order allow,deny
   Allow from all
</Directory>


#With uwsgi (only for unix, cygwin)

uwsgi --http :9090 --wsgi-file foobar.py

#foobar.py:
def application(env, start_response):
    start_response('200 OK', [('Content-Type','text/html')])
    return [b"Hello World"]

#Stopping the server
kill -INT `cat /tmp/project-master.pid`
# or for convenience...
uwsgi --stop /tmp/project-master.pid


#Configuring Nginx
#uWSGI natively speaks HTTP, FastCGI, SCGI and its specific protocol named “uwsgi” 
Already supported by nginx and Cherokee (and  Apache modules )

#nginx config - This means “pass every request to the server bound to port 3031 speaking the uwsgi protocol”.
location / {
    include uwsgi_params;
    uwsgi_pass 127.0.0.1:3031;
}


# spawn uWSGI to natively speak the uwsgi protocol:
uwsgi --socket 127.0.0.1:3031 --wsgi-file foobar.py --master --processes 4 --threads 2 --stats 127.0.0.1:9191

#For  proxy/webserver/router speaks HTTP, 
uwsgi --http-socket 127.0.0.1:3031 --wsgi-file foobar.py --master --processes 4 --threads 2 --stats 127.0.0.1:9191

#Configuring Apache2, use mod_proxy_uwsgi, 
#It is a “proxy” module, so you will get all of the features exported by mod_proxy
#To load mod_proxy and mod_proxy_uwsgi modules in your apache config.

ProxyPass /foo uwsgi://127.0.0.1:3032/
ProxyPass /bar uwsgi://127.0.0.1:3033/
ProxyPass / uwsgi://127.0.0.1:3031/


#The first two forms set SCRIPT_NAME respectively to /foo and /bar 
#while the last one use an empty SCRIPT_NAME. 

#You can set additional uwsgi vars using the SetEnv directive and load balance requests using mod_proxy_balancer.


<Proxy balancer://mycluster>
   BalancerMember uwsgi://192.168.1.50:3031/
   BalancerMember uwsgi://192.168.1.51:3031/
</Proxy>
ProxyPass / balancer://mycluster

#Using Configuration file
#each command line option maps 1:1 with entries in the config files.
uwsgi --http-socket :9090 --psgi myapp.pl
#ini file: 
[uwsgi]
http-socket = :9090
psgi = myapp.pl

#There are many command line options, check by
uwsgi -h 

#example 
uwsgi --socket /tmp/uwsgi.sock --socket 127.0.0.1:8000 --master --workers 3

#To add concurrency (by default uWSGI starts with a single process and a single thread).
uwsgi --http :9090 --wsgi-file foobar.py --master --processes 4 --threads 2
#This will spawn 4 processes (each with 2 threads), a master process (will respawn your processes when they die) 

#For  monitoring. The stats subsystem allows you to export uWSGI’s internal statistics as JSON:
uwsgi --http :9090 --wsgi-file foobar.py --master --processes 4 --threads 2 --stats 127.0.0.1:9191
#Make some request to your app and then telnet to the port 9191, you’ll get lots of fun information. 


#Loading  Configuration file 
uwsgi --ini http://uwsgi.it/configs/myapp.ini # HTTP
uwsgi --xml - # standard input
uwsgi --yaml fd://0 # file descriptor
uwsgi --json 'exec://nc 192.168.11.2:33000' # arbitrary executable

#Deploying Django- http://uwsgi-docs.readthedocs.io/en/latest/tutorials/Django_and_nginx.html
#project at /home/foobar/myproject:
uwsgi --socket 127.0.0.1:3031 --chdir /home/foobar/myproject/ --wsgi-file myproject/wsgi.py --master --processes 
#Or ini
[uwsgi]
socket = 127.0.0.1:3031
chdir = /home/foobar/myproject/
wsgi-file = myproject/wsgi.py
processes = 4
threads = 2
stats = 127.0.0.1:9191

#Websocket (directly connecting to Server bypassing HTTP)- (see tests directory)

#file : websockets_echo.py 

#!./uwsgi --https :8443,foobar.crt,foobar.key --http-raw-body --gevent 100 --wsgi-file websockets_echo.py 

import uwsgi
import time


def application(env, sr):

    ws_scheme = 'ws'
    if 'HTTPS' in env or env['wsgi.url_scheme'] == 'https':
        ws_scheme = 'wss'

    if env['PATH_INFO'] == '/':
        sr('200 OK', [('Content-Type', 'text/html')])
        return """
    <html>
      <head>
          <script language="Javascript">
            var s = new WebSocket("%s://%s/foobar/");
            s.onopen = function() {
              alert("connected !!!");
              s.send("ciao");
            };
            s.onmessage = function(e) {
        var bb = document.getElementById('blackboard')
        var html = bb.innerHTML;
        bb.innerHTML = html + '<br/>' + e.data;
            };

        s.onerror = function(e) {
            alert(e);
        }

    s.onclose = function(e) {
        alert("connection closed");
    }

            function invia() {
              var value = document.getElementById('testo').value;
              s.send(value);
            }
          </script>
     </head>
    <body>
        <h1>WebSocket</h1>
        <input type="text" id="testo"/>
        <input type="button" value="invia" onClick="invia();"/>
    <div id="blackboard" style="width:640px;height:480px;background-color:black;color:white;border: solid 2px red;overflow:auto">
    </div>
    </body>
    </html>
        """ % (ws_scheme, env['HTTP_HOST'])
    elif env['PATH_INFO'] == '/foobar/':
        uwsgi.websocket_handshake(env['HTTP_SEC_WEBSOCKET_KEY'], env.get('HTTP_ORIGIN', ''))
        print "websockets..."
        while True:
            msg = uwsgi.websocket_recv()
            uwsgi.websocket_send("[%s] %s" % (time.time(), msg
			

#GET Request Handling - http://localhost:8051/?age=10&hobbies=software&hobbies=tunning

#!/usr/bin/env python

from wsgiref.simple_server import make_server
from cgi import parse_qs, escape

#Note the way html string is written , html is used for format string 
#Use %(dict_key)format_specifier   eg   %(age)s  ,means replace this by dict-key as age with %s 


html = """
<html>
<body>
   <form method="get" action="">
        <p>
           Age: <input type="text" name="age" value="%(age)s">
        </p>
        <p>
            Hobbies:
            <input   name="hobbies" type="checkbox" value="software" %(checked-software)s
            > Software
            <input
                name="hobbies" type="checkbox" value="tunning"  %(checked-tunning)s
            > Auto Tunning
        </p>
        <p>
            <input type="submit" value="Submit">
        </p>
    </form>
    <p>
        Age: %(age)s<br>
        Hobbies: %(hobbies)s
    </p>
</body>
</html>
"""

def application (environ, start_response):

    # Returns a dictionary in which the values are lists
    d = parse_qs(environ['QUERY_STRING'])  #returns a dict 

    # As there can be more than one value for a variable then
    # a list is provided as a default value.
    age = d.get('age', [''])[0]       # Returns the first age value or default 
    hobbies = d.get('hobbies', [])    # Returns a list of hobbies or default 

    # Always escape user input to avoid script injection
    age = escape(age)
    hobbies = [escape(hobby) for hobby in hobbies]

    response_body = html % { # Fill the above html template in
        'checked-software': ('', 'checked')['software' in hobbies],
        'checked-tunning': ('', 'checked')['tunning' in hobbies],
        'age': age or 'Empty',
        'hobbies': ', '.join(hobbies or ['No Hobbies?'])
    }

    status = '200 OK'

    # Now content type is text/html
    response_headers = [
        ('Content-Type', 'text/html'),
        ('Content-Length', str(len(response_body)))
    ]

    start_response(status, response_headers)
    return [response_body]        #must encode in Py3.x eg response_body.encode('utf-8')

httpd = make_server('localhost', 8051, application)

# Now it is serve_forever() in instead of handle_request()
httpd.serve_forever() #it would server forever 





#POST  Request handling - Send via firefox's httprequester add_on 

from wsgiref.simple_server import make_server
from cgi import parse_qs, escape

#Note the way html string is written , html is used for format string 
#Use %(dict_key)format_specifier   eg   %(age)s  ,means replace this by dict-key as age with %s 


html = """
<html>
<body>
   <form method="post" action="">
        <p>
           Age: <input type="text" name="age" value="%(age)s">
        </p>
        <p>
            Hobbies:
            <input
                name="hobbies" type="checkbox" value="software"
                %(checked-software)s
            > Software
            <input
                name="hobbies" type="checkbox" value="tunning"
                %(checked-tunning)s
            > Auto Tunning
        </p>
        <p>
            <input type="submit" value="Submit">
        </p>
    </form>
    <p>
        Age: %(age)s<br>
        Hobbies: %(hobbies)s
    </p>
</body>
</html>
"""

def application(environ, start_response):

    if environ['REQUEST_METHOD'] != 'POST':     #Only handle POST 
        start_response('403 Forbidden', [])
        return []

    # the environment variable CONTENT_LENGTH may be empty or missing
    try:
        request_body_size = int(environ.get('CONTENT_LENGTH', 0))
    except (ValueError):
        request_body_size = 0

    # When the method is POST the variable will be sent
    # in the HTTP request body which is passed by the WSGI server
    # in the file like wsgi.input environment variable.
    request_body = environ['wsgi.input'].read(request_body_size)  #must decode in Py3.x eg request_body.decode('utf-8')
    d = parse_qs(request_body)

    age = d.get('age', [''])[0] # Returns the first age value.
    hobbies = d.get('hobbies', []) # Returns a list of hobbies.

    # Always escape user input to avoid script injection
    age = escape(age)
    hobbies = [escape(hobby) for hobby in hobbies]

    response_body = html % { # Fill the above html template in
        'checked-software': ('', 'checked')['software' in hobbies],
        'checked-tunning': ('', 'checked')['tunning' in hobbies],
        'age': age or 'Empty',
        'hobbies': ', '.join(hobbies or ['No Hobbies?'])
    }

    status = '200 OK'

    response_headers = [
        ('Content-Type', 'text/html'),
        ('Content-Length', str(len(response_body)))
    ]

    start_response(status, response_headers)
    return [response_body]    #must encode in Py3.x eg response_body.encode('utf-8')

httpd = make_server('localhost', 8051, application)
httpd.serve_forever()





		
# CherryPy- webframework - Py2.7 and Py3.x  (conda install -c trentonoliphant cherrypy=3.5.0  for win32 chek in https://anaconda.org/)
#can be used as standalone for small application
#Can use uwsgi, twisted, tornado http server 
#can have virtual hosting
#can be hosted under apache with mod_wsgi and Nginx

#close shell and start again, if wrong behaviour

#example
import cherrypy

class Root(object):
    @cherrypy.expose    #user can access
    def index(self):		#user can access /index or by /
        return "Hello World!"  #this would be printed 

if __name__ == '__main__':
   cherrypy.quickstart(Root(), '/')  #mounted on / #can take last arg as 'conf'

   
#conf format - conf is dict
conf = {
         '/': {
             'tools.sessions.on': True,                #session on
             'tools.staticdir.root': os.path.abspath(os.getcwd()) # give abs dir as root
         },
         '/static': {
             'tools.staticdir.on': True,
             'tools.staticdir.dir': './public'   # map to cwd/Public directory
         }
     }

	 



#if file contains Multiple applications

import cherrypy

class Root:
    def index(self):
        return "Hello, world!"
    index.exposed = True

class Admin:
    def user(self, name=""):  #../user
        return "You asked for user '%s'" % name
    user.exposed = True

class Search:
    def index(self):
        return "some page"
    index.exposed = True

root = Root()              #/
root.admin = Admin()       #/admin
root.admin.search = Search() #/admin/search
cherrypy.quickstart(root, '/')

#OR
#use cherrypy.mount inplace of quickstart

import cherrypy
import web

class WebService(object):

    def __init__(self):
        app_config = {
            '/static': {
                # enable serving up static resource files
                'tools.staticdir.root': '/static',
                'tools.staticdir.on': True,
                'tools.staticdir.dir': "static",
            },
        }

        api_config = {
            '/': {
                # the api uses restful method dispatching
                'request.dispatch': cherrypy.dispatch.MethodDispatcher(),

                # all api calls require that the client passes HTTP basic authentication
                'tools.authorize.on': True,
            }
        }

        cherrypy.tree.mount(web.Application(), '/', config=app_config)
        cherrypy.tree.mount(web.API(), '/api', config=api_config)

    # a blocking call that starts the web application listening for requests
    def start(self, port=8080):
        cherrypy.config.update({'server.socket_host': '0.0.0.0', })
        cherrypy.config.update({'server.socket_port': port, })
        cherrypy.engine.start()
        cherrypy.engine.block()

    # stops the web application
    def stop(self):
        cherrypy.engine.stop()


#Routing



#index method 
The index method is special- it handles intermediate URI’s that end in a slash
for example,  /orders/items/ might maps to root.orders.items.index. 
The index method can take additional keyword arguments, 
it cannot take positional arguments 
#Eg Given
class StoryHandler(object):	
	@cherrypy.expose
	def index(self, name =""):  
		return "this is story %s" % (name,)
root.story = StoryHandler()  
#you can give below , not http://localhost:8080/story/1
http://localhost:8080/story/?name=1
OR
http://localhost:8080/story/index/1

#to have http://localhost:8080/story/1
Merge the above with Root handler and make method name as 'story'

#Keyword Arguments
#Unified calling of POST or GET method, same method is invoked
Any page handler that is called by CherryPy (index, or any other suitable method) 
can receive additional data from HTML or other forms using keyword arguments.

<form action="doLogin" method="post">
    <p>Username</p>
    <input type="text" name="username" value=""
        size="15" maxlength="40"/>
    <p>Password</p>
    <input type="password" name="password" value=""
        size="10" maxlength="40"/>
    <p><input type="submit" value="Login"/></p>
    <p><input type="reset" value="Clear"/></p>
</form>


#handlers


class Root:
	@cherrypy.expose
	def doLogin(self, username=None, password=None):
        # check the username & password
        ...




#Positional Arguments - index method can not take positional.
Any trailing components are passed as positional arguments. 
For example, the URI "/branch/leaf/4" might result in the call: 
app.root.branch.leaf(4), or app.root.index(branch, leaf, 4) 


#For example, http://localhost/blog/2005/01/17. 

class Root:
    def blog(self, year, month, day):
        ...
    blog.exposed = True

root = Root()


#So the URL above will be mapped as a call to:
root.blog('2005', '1', '17')



#Default methods- if no handler is found, 'default' method is invoked at that tree level

class Blog:
    def default(self, year, month, day):
        ...
    default.exposed = True

class Root: pass

root = Root()
root.blog = Blog()


http://localhost/blog/2005/01/17 will be mapped as a call to:
root.blog.default('2005', '1', '17')



#Special characters
#default dispatcher converts all dots in the URI to underscores
#hence /path/to/my.html at URL would search  def my_html

#Get param or Post param - unified way of handling

import cherrypy

class Root(object):
	@cherrypy.expose
	def mytest(self, param_1=None, param_2=None, *args, **kw):   #args collect positional 
		return repr(dict(param_1=param_1,
                         param_2=param_2,
                         args=args,
                         kw=kw))
	@cherrypy.expose
	def mytestex(self, field1, field2):   #must be sent by field1=b&field2=c
		return repr(dict(field1=field1,
                         field2=field2))


cherrypy.quickstart(Root(), '/')


#Use like 
http://localhost:8080/mytest?foo=ok&bar=nok
{'kw': {'foo': u'ok', 'bar': u'nok'}, 'args': (), 'param_1': None, 'param_2': None}

http://localhost:8080/mytest/1/2/foo=ok&bar=nok
{'kw': {}, 'args': ('foo=ok&bar=nok',), 'param_1': '1', 'param_2': '2'}

http://localhost:8080/mytest/foo/bar/baz=testing
{'kw': {}, 'args': ('baz=testing',), 'param_1': 'foo', 'param_2': 'bar'}

http://localhost:8080/mytest/foo/bar?baz=testing 
{'kw': {'baz': u'testing'}, 'args': (), 'param_1': 'foo', 'param_2': 'bar'}

http://localhost:8080/mytest/foo/bar/baz
{'kw': {}, 'args': ('baz',), 'param_1': 'foo', 'param_2': 'bar'}

http://localhost:8080/mytest/foo/bar/baz/ae
{'kw': {}, 'args': ('baz', 'ae'), 'param_1': 'foo', 'param_2': 'bar'}

http://localhost:8080/mytest/foo/bar?param_1=baz 
Error

#conventional use 

http://localhost:8080/mytestex?field1=b&field2=c
{'field2': u'c', 'field1': u'b'}

http://localhost:8080/mytestex/1/2
{'field1': '1', 'field2': '2'}

#for POST
#py2.7
import urllib
print urllib.urlopen('http://localhost:8080', 'field1=b&field2=c').read()

#py3.x*x
import urllib.request
>>> print(urllib.request.urlopen('http://localhost:8080/mytestex', b'field1=b&field2=c').read())
b"{'field2': u'c', 'field1': u'b'}"


#Host a WSGI application in CherryPy


def raw_wsgi_app(environ, start_response):
    status = '200 OK'
    response_headers = [('Content-type','text/plain')]
    start_response(status, response_headers)
    return ['Hello world!']

cherrypy.tree.graft(raw_wsgi_app, '/')


#Logging
#conf fields 
log.access_file for incoming requests using the common log format
log.error_file for the other log, use cherrypy.log("kaboom!", traceback=True) or cherrypy.log("hello there")

#To disable file logging
To disable, console logging, set log.screen to False.

cherrypy.config.update({'log.screen': False,
                        'log.access_file': '',
                        'log.error_file': ''})


#Global server configuration
use the cherrypy.config.update() method.

cherrypy.config.update({'server.socket_port': 9090})

#OR
#file: server.conf
[global]
server.socket_port: 9090

cherrypy.config.update("server.conf")

#Cookies
Uses standard lib Cookie.SimpleCookie object type to handle cookies.
•To send a cookie to a browser, set cherrypy.response.cookie[key] = value.
•To retrieve a cookie sent by a browser, use cherrypy.request.cookie[key].
•To delete a cookie (on the client side), you must send the cookie with its expiration time set to 0:
cherrypy.response.cookie[key] = value
cherrypy.response.cookie[key]['expires'] = 0


Note  cherrypy.request.cookie should be populated each time you send data to browser
But the server doesn’t need to send the same cookies with every response; 
therefore, cherrypy.response.cookie will usually be empty. 
 

#Using sessions
Sessions are used to identify users and synchronize their activity. 
By default, CherryPy does not activate sessions because it is not a mandatory feature to have, 
to enable it simply add the following settings in your configuration:

[/]
tools.sessions.on: True


#Filesystem backend
Using a filesystem is a simple to not lose your sessions between reboots. 
Each session is saved in its own file within the given directory.


[/]
tools.sessions.on: True
tools.sessions.storage_type = "file"
tools.sessions.storage_path = "/some/directorys"



#Memcached backend
Memcached is a popular key-store on top of your RAM, it is distributed and a good choice 
if you want to share sessions outside of the process running CherryPy.


[/]
tools.sessions.on: True
tools.sessions.storage_type = "memcached"




#BenchMark -checking webserver performance - Use apache/ab

# 100 concurrent - -n requests  -c concurrency -k   Use HTTP KeepAlive feature
C:\indigoampp\apache-2.2.15\bin\ab -c 100 -n 1000 -k localhost:8888/ | grep "Time taken for tests:" 

# 5 concurrent
C:\indigoampp\apache-2.2.15\bin\ab -c 5 -n 10000 -k localhost:8888/ 


#check http://klen.github.io/py-frameworks-bench/#results




#Complete Example

#cherrypy + elixir(ORM) + jinja2(template)
#for Py2.7 only
#for windows 
#mysql-python from   http://www.lfd.uci.edu/~gohlke/pythonlibs/
#pip2 install MySQL_python-1.2.5-cp27-none-win32.whl
#for linux use pip2 install mysql-python
#for all
#pip2 install sqlalchemy==0.6.0
#pip2 install elixir
#pip2 install Jinja2
#pip2 install cherrypy

#start sqlserver  eg at cygwin
#mysql -running deamon in cygwin
/usr/bin/mysqld_safe &
#shutting down
mysqladmin.exe -h 127.0.0.1 -u root   --connect-timeout=5 shutdown
#mysql admin #  default port 3306, 
mysql -u root    -h 127.0.0.1 

#Template file -  vars substitution- {{ }}  - code {%  %}
#loop, if exists 
#must be under templates 

#templates/first.tmpl

<!doctype html>
<html lang="en">
<head>
  <meta charset="UTF-8" />

  <title>{{ title }}</title>
  <meta name="description" content="{{ description }}" />
</head>

<body>

<h1>Contents of table ::: <i> {{ name | upper }} </i> </h1><br />
<table align="left" border="1" width="50%">
   <tr>
        <th bgcolor="#cccccc" align="center">ID</th> 
		<th bgcolor="#cccccc" align="center">NAME</th> 
		<th bgcolor="#cccccc" align="center">EMAIL</th>
   </tr>
   {% for author in authors %}
     <tr align="center">
        <td>{{author.id}}</td> 
		<td>{{author.name}}</td> 
		<td>{{author.email}}</td>
     </tr>
   {% endfor %}
</table>

</body>
</html>

#Model file 
#must be under elixirex
#elixirex/model.py

from elixir import *


#must be ip address, not localhost
metadata.bind = 'mysql://root:@127.0.0.1/books'  
#metadata.bind.echo = True




class Author(Entity):
	
	using_options(tablename = 'authors')
	id = Field(Integer, primary_key=True, colname= 'id')
	name = Field(String(20), colname= 'name')
	email = Field(String(20), colname= 'email')

	def __repr__(self):
		return '<Authors "%s" (%s)>' % (self.name, self.email)


#Cherry py file

#!/usr/bin/python2.7
# http://localhost:8080/

import jinja2
from elixirex.model import *
import string
import cherrypy

TEMPLATE_FILE = "first.tmpl"

setup_all()     



index = """
<!DOCTYPE html>
<html>
<head>
    <title>Input</title>
</head>
<body>
   <form action="/getTable" method="POST">
     ID:  <input type="text" name="id" />
    <br>
	<br>
	<br>
    Name:  <input type="text" name="name" />
<input type="submit" name="submit" value="Search" />
<input type="reset" name="reset"  value="Reset" />
</form>
</body>
</html>

"""



class StringGenerator(object):

		
	@cherrypy.expose
	def index(self):
		return index

	@cherrypy.expose
	def getTable(self, id, name, submit='Submit'):
		templateLoader = jinja2.FileSystemLoader( searchpath="./templates" )
		templateEnv = jinja2.Environment( loader=templateLoader )
		
		authors = Author.query.all()
		filter_str = ''
		if id :
			authors = Author.query.filter(Author.id == id).all()
		elif name :
			authors = Author.query.filter(Author.name == name).all()
		else:
			pass
		
			
		template = templateEnv.get_template( TEMPLATE_FILE )
		templateVars = { "title" : "Display tables",
                 "description" : " simple example.", "name" : "Authors",
                 "authors" : authors
               }
		outputText = template.render( templateVars )
		return outputText

if __name__ == '__main__':
    cherrypy.quickstart(StringGenerator())

	
	
	
	
	
	








#Bottle - Python Web Framework


#example:
from bottle import route, run, template

@route('/hello/<name>')   #dynamic routing
def index(name):
    return template('<b>Hello {{name}}</b>!', name=name)

run(host='localhost', port=8080)

#check at http://localhost:8080/hello/world

#other decorators : get(), post(), put(), delete() or patch().

from bottle import get, post, request # or route

@get('/login') # or @route('/login')
def login():
    return '''
        <form action="/login" method="post">
            Username: <input name="username" type="text" />
            Password: <input name="password" type="password" />
            <input value="Login" type="submit" />
        </form>
    '''

@post('/login') # or @route('/login', method='POST')
def do_login():
    username = request.forms.get('username')
    password = request.forms.get('password')
    if check_login(username, password):
        return "<p>Your login information was correct.</p>"
    else:
        return "<p>Login failed.</p>"

		
		

#Multiple Routing


from bottle import template


@route('/')
@route('/hello/<name>')
def greet(name='Stranger'):
    return template('Hello {{name}}, how are you?', name=name)


	
@route('/wiki/<pagename>')            # matches /wiki/Learning_Python
def show_wiki_page(pagename):
    ...

@route('/<action>/<user>')            # matches /follow/defnull
def user_api(action, user):
    ...


#The rule /<action>/<item> matches as follows:

Path			Result
/save/123 		{'action': 'save', 'item': '123'} 
/save/123/ 		No Match 
/save/ 			No Match 
//123 			No Match 




#A filtered wildcard is declared as <name:filter> or <name:filter:config>. 
•:int 		matches (signed) digits and converts the value to integer.
•:float 	similar to :int but for decimal numbers.
•:path 		matches all characters including the slash character in a non-greedy way and may be used to match more than one path segment.
•:re[:exp] 	allows you to specify a custom regular expression in the config field. The matched value is not modified.

route('/object/<id:int>')
def callback(id):
    assert isinstance(id, int)

@route('/show/<name:re:[a-z]+>')
def callback(name):
    assert name.isalpha()

@route('/static/<path:path>')
def callback(path):
    return static_file(path, ...)

	
#You can add your own filters to the router. 
app = Bottle()

def list_filter(config):
    ''' Matches a comma separated list of numbers. '''
    delimiter = config or ','
    regexp = r'\d+(%s\d)*' % re.escape(delimiter)

    def to_python(match):
        return map(int, match.split(delimiter))

    def to_url(numbers):
        return delimiter.join(map(str, numbers))

    return regexp, to_python, to_url

app.router.add_filter('list', list_filter)

@app.route('/follow/<ids:list>')
def follow_users(ids):
    for id in ids:



#Old Syntax			New Syntax
:name 				<name> 
:name#regexp#		<name:re:regexp> 
:#regexp# 			<:re:regexp> 
:## <:re> 

#Explicit routing configuration

#default route
def setup_routing():
    bottle.route('/', 'GET', index)
    bottle.route('/edit', ['GET', 'POST'], edit)


#any other route 

def setup_routing(app):
    app.route('/new', ['GET', 'POST'], form_new)
    app.route('/edit', ['GET', 'POST'], form_edit)

app = Bottle()
setup_routing(app)



#Routing Static Files


from bottle import static_file
@route('/static/<filename>')
def server_static(filename):
    return static_file(filename, root='/path/to/your/static/files')

	
#To serve files in subdirectories, change the wildcard to use the path filter:


@route('/static/<filepath:path>')
def server_static(filepath):
    return static_file(filepath, root='/path/to/your/static/files')

from bottle import static_file
@route('/images/<filename:re:.*\.png>')
def send_image(filename):
    return static_file(filename, root='/path/to/image/files', mimetype='image/png')

@route('/static/<filename:path>')
def send_static(filename):
    return static_file(filename, root='/path/to/static/files')


#Error Pages
override the default for a specific HTTP status code with the error() decorator:


from bottle import error
@error(404)
def error404(error):
    return 'Nothing here, sorry'


#Forced Download

Most browsers try to open downloaded files if the MIME type is known and assigned to an application (e.g. PDF files). 
If this is not what you want, you can force a download dialog and even suggest a filename to the user:
If the download parameter is just True, the original filename is used.

@route('/download/<filename:path>')
def download(filename):
    return static_file(filename, root='/path/to/static/files', download=filename)



#HTTP Errors and Redirects
The abort() function is a shortcut for generating HTTP error pages.


from bottle import route, abort
@route('/restricted')
def restricted():
    abort(401, "Sorry, access denied.")


#To redirect a client to a different URL, you can send a 303 

from bottle import redirect
@route('/wrong/url')
def wrong():
    redirect("/right/url")


#Response Header

@route('/wiki/<page>')
def wiki(page):
    response.set_header('Content-Language', 'en')

	

#Cookies
A cookie is a named piece of text stored in the user’s browser profile. 

@route('/hello')
def hello_again():
    if request.get_cookie("visited"):
        return "Welcome back! Nice to see you again"
    else:
        response.set_cookie("visited", "yes")
        return "Hello there! Nice to meet you"


#Request Data
Cookies(request.cookies.<<name>>, HTTP header, HTML <form> (request.forms.<<fieldName>>) and other request data is available through the global request object.
All values are also accessible as attributes. Error case: the string is empty, but still present:

from bottle import request, route, template

@route('/hello')
def hello():
    name = request.cookies.username or 'Guest'
    return template('Hello {{name}}', name=name)


#Example
name = request.cookies.name
# is a shortcut for:

name = request.cookies.getunicode('name') # encoding='utf-8' (default)

# which basically does this:

try:
    name = request.cookies.get('name', '').decode('utf-8')
except UnicodeError:
    name = u''

#Multiple values per key: 

for choice in request.forms.getall('multiple_choice'):
    do_something(choice)

#HTTP Headers

from bottle import route, request
@route('/is_ajax')
def is_ajax():
    if request.headers.get('X-Requested-With') == 'XMLHttpRequest':
        return 'This is an AJAX request'
    else:
        return 'This is a normal request'
		
		
#Query Variables

Attribute			GET Form fields		POST Form fields		File Uploads
BaseRequest.query 	yes 				no 						no 
BaseRequest.forms 	no 					yes 					no 
BaseRequest.files 	no 					no 						yes 
BaseRequest.params 	yes 				yes 					no 
BaseRequest.GET 	yes 				no 						no 
BaseRequest.POST 	no 					yes 					yes 


# Example: /forum?id=1&page=5

from bottle import route, request, response, template
@route('/forum')
def display_forum():
    forum_id = request.query.id
    page = request.query.page or '1'
    return template('Forum ID: {{id}} (page {{page}})', id=forum_id, page=page)

#form fields

from bottle import route, request

@route('/login')
def login():
    return '''
        <form action="/login" method="post">
            Username: <input name="username" type="text" />
            Password: <input name="password" type="password" />
            <input value="Login" type="submit" />
        </form>
    '''

@route('/login', method='POST')
def do_login():
    username = request.forms.get('username')
    password = request.forms.get('password')
    if check_login(username, password):
        return "<p>Your login information was correct.</p>"
    else:
        return "<p>Login failed.</p>"




#File uploaded



<form action="/upload" method="post" enctype="multipart/form-data">
  Category:      <input type="text" name="category" />
  Select a file: <input type="file" name="upload" />
  <input type="submit" value="Start upload" />
</form>


@route('/upload', method='POST')
def do_upload():
    category   = request.forms.get('category')
    upload     = request.files.get('upload')
    name, ext = os.path.splitext(upload.filename)
    if ext not in ('.png','.jpg','.jpeg'):
        return 'File extension not allowed.'

    save_path = get_save_path_for_category(category)
    upload.save(save_path) # appends upload.filename automatically
    return 'OK'

#WSGI Environment
Each BaseRequest instance wraps a WSGI environment dictionary. 
The original is stored in BaseRequest.environ, but the request object itself behaves like a dictionary, too.
 

@route('/my_ip')
def show_ip():
    ip = request.environ.get('REMOTE_ADDR')
    # or ip = request.get('REMOTE_ADDR')
    # or ip = request['REMOTE_ADDR']
    return template("Your IP is: {{ip}}", ip=ip)



#SimpleTemplate Engine

 

>>> from bottle import SimpleTemplate
>>> tpl = SimpleTemplate('Hello {{name}}!')
>>> tpl.render(name='World')
u'Hello World!'


#or 


>>> from bottle import template
>>> template('Hello {{name}}!', name='World')
u'Hello World!'


# pass a dictionary into the template using keyword arguments:


>>> from bottle import template
>>> my_dict={'number': '123', 'street': 'Fake St.', 'city': 'Fakeville'}
>>> template('I live at {{number}} {{street}}, {{city}}', **my_dict)
u'I live at 123 Fake St., Fakeville'

#any python expression is allowed within the curly brackets as long 
#as it evaluates to a string or something that has a string representation:


>>> template('Hello {{name}}!', name='World')
u'Hello World!'
>>> template('Hello {{name.title() if name else "stranger"}}!', name=None)
u'Hello stranger!'
>>> template('Hello {{name.title() if name else "stranger"}}!', name='mArC')
u'Hello Marc!'

#HTML special characters are escaped automatically to prevent XSS attacks. 
#You can start the expression with an exclamation mark to disable escaping for that expression:


>>> template('Hello {{name}}!', name='<b>World</b>')
u'Hello &lt;b&gt;World&lt;/b&gt;!'
>>> template('Hello {{!name}}!', name='<b>World</b>')
u'Hello <b>World</b>!'

#Embedded python code
Both the % and the <% tokens are only recognized if they are the first non-whitespace characters in a line


% name = "Bob"  # a line of python code
<p>Some plain text in between</p>
<%
  # A block of python code
  name = name.title().strip()
%>
<p>More plain text</p>

#Close explicitly with an end keyword.


<ul>
  % for item in basket:
    <li>{{item}}</li>
  % end
</ul>



#other functions
defined(name)
Return True if the variable is defined in the current template namespace, False otherwise.

get(name, default=None)
Return the variable, or a default value.

setdefault(name, default)
If the variable is not defined, create it with the given default value. Return the variable.

#Example:

% setdefault('text', 'No Text')
<h1>{{get('title', 'No Title')}}</h1>
<p> {{ text }} </p>
% if defined('author'):
  <p>By {{ author }}</p>
% end













# Parsing HTML using HTMLParser
#The HTMLParser module has been renamed to html.parser in Python 3.

#Example HTML Parser Application - Implement all callback routines

from html.parser import HTMLParser
from html.entities import name2codepoint

class MyHTMLParser(HTMLParser):
    def handle_starttag(self, tag, attrs):
        print("Start tag:", tag)
        for attr in attrs:
            print("     attr:", attr)
    def handle_endtag(self, tag):
        print("End tag  :", tag)
    def handle_data(self, data):
        print("Data     :", data)
    def handle_comment(self, data):
        print("Comment  :", data)
    def handle_entityref(self, name):
        c = chr(name2codepoint[name])
        print("Named ent:", c)
    def handle_charref(self, name):
        if name.startswith('x'):
            c = chr(int(name[1:], 16))
        else:
            c = chr(int(name))
        print("Num ent  :", c)
    def handle_decl(self, data):
        print("Decl     :", data)

parser = MyHTMLParser()


#Parsing a doctype:


>>> parser.feed('<!DOCTYPE HTML PUBLIC "-//W3C//DTD HTML 4.01//EN" '
...             '"http://www.w3.org/TR/html4/strict.dtd">')
Decl     : DOCTYPE HTML PUBLIC "-//W3C//DTD HTML 4.01//EN" "http://www.w3.org/TR/html4/strict.dtd"


#Parsing an element with a few attributes and a title:


>>> parser.feed('<img src="python-logo.png" alt="The Python logo">')
Start tag: img
     attr: ('src', 'python-logo.png')
     attr: ('alt', 'The Python logo')
>>>
>>> parser.feed('<h1>Python</h1>')
Start tag: h1
Data     : Python
End tag  : h1

parser.feed('<style type="text/css">#python { color: green }</style>')
Start tag: style
     attr: ('type', 'text/css')
Data     : #python { color: green }
End tag  : style
>>>
>>> parser.feed('<script type="text/javascript">'
...             'alert("<strong>hello!</strong>");</script>')




#other handler
HTMLParser.handle_startendtag(tag, attrs)
Similar to handle_starttag(), but called when the parser encounters an XHTML-style empty tag (<img ... />). 

HTMLParser.handle_charref(name)
This method is called to process decimal and hexadecimal numeric character references of the form &#NNN; 
and &#xNNN

HTMLParser.handle_entityref(name)
This method is called to process a named character reference of the form &name; (e.g. &gt;), 
where name is a general entity reference (e.g. 'gt'

HTMLParser.handle_comment(data)
This method is called when a comment is encountered (e.g. <!--comment-->).





###xml.etree.ElementTree — The ElementTree XML API (same in Py3.x and Py2.7)
#std module 
 
#Example: country_data.xml file 
  
<?xml version="1.0"?>
<data>
    <country name="Liechtenstein">
        <rank>1</rank>
        <year>2008</year>
        <gdppc>141100</gdppc>
        <neighbor name="Austria" direction="E"/>
        <neighbor name="Switzerland" direction="W"/>
    </country>
    <country name="Singapore">
        <rank>4</rank>
        <year>2011</year>
        <gdppc>59900</gdppc>
        <neighbor name="Malaysia" direction="N"/>
    </country>
    <country name="Panama">
        <rank>68</rank>
        <year>2011</year>
        <gdppc>13600</gdppc>
        <neighbor name="Costa Rica" direction="W"/>
        <neighbor name="Colombia" direction="E"/>
    </country>
</data>


#code 

import xml.etree.ElementTree as ET
tree = ET.parse('country_data.xml')
root = tree.getroot()


#Or directly from a string:
root = ET.fromstring(country_data_as_string)

#Every element has a tag and a dictionary of attributes:

>>> root.tag
'data'
>>> root.attrib
{}


#It also has children nodes over which we can iterate

>>> for child in root:
...   print(child.tag, child.attrib)
...
country {'name': 'Liechtenstein'}
country {'name': 'Singapore'}
country {'name': 'Panama'}


#Children are nested, and we can access specific child nodes by index:

>>> root[0][1].text


#Pull API for non-blocking parsing- XMLPullParser with events 
#To get the parsed XML elements, call XMLPullParser.read_events()


>>> parser = ET.XMLPullParser(['start', 'end'])

>>> parser.feed('<mytag>sometext')

>>> list(parser.read_events())

[('start', <Element 'mytag' at 0x7fa66db2be58>)]

>>> parser.feed(' more text</mytag>')

>>> for event, elem in parser.read_events():
		print(event)
		print(elem.tag, 'text=', elem.text)



#Finding interesting elements

#use Element.iter():


>>> for neighbor in root.iter('neighbor'):
		print(neighbor.attrib)

#Use Element.findall() finds only elements with a tag which are direct children of the current element. 
#Element.find() finds the first child with a particular tag, 
#Element.text accesses the element’s text content. 
#Element.get() accesses the element’s attributes:


>>> for country in root.findall('country'):
		rank = country.find('rank').text
		name = country.get('name')
		print(name, rank)


#Modifying an XML File
#To add new element, use Element.append()
#to update use Element.set()

>>> for rank in root.iter('rank'):
		new_rank = int(rank.text) + 1
		rank.text = str(new_rank)
		rank.set('updated', 'yes')
...
>>> tree.write('output.xml')


#remove elements using Element.remove(). 

>>> for country in root.findall('country'):
		rank = int(country.find('rank').text)
		if rank > 50:
			root.remove(country)

>>> tree.write('output.xml')

#Building XML documents

>>> a = ET.Element('a')
>>> b = ET.SubElement(a, 'b')
>>> c = ET.SubElement(a, 'c')
>>> d = ET.SubElement(c, 'd')
>>> ET.dump(a)
<a><b /><c><d /></c></a>


#Parsing XML with Namespaces

If the XML input has namespaces, tags and attributes with prefixes in the form prefix:sometag 
get expanded to {uri}sometag where the prefix is replaced by the full URI.
Also, if there is a default namespace, that full URI gets prepended to all of the non-prefixed tags

<?xml version="1.0"?>
<actors xmlns:fictional="http://characters.example.com"
        xmlns="http://people.example.com">
    <actor>
        <name>John Cleese</name>
        <fictional:character>Lancelot</fictional:character>
        <fictional:character>Archie Leach</fictional:character>
    </actor>
    <actor>
        <name>Eric Idle</name>
        <fictional:character>Sir Robin</fictional:character>
        <fictional:character>Gunther</fictional:character>
        <fictional:character>Commander Clement</fictional:character>
    </actor>
</actors>


#Option-1 

root = fromstring(xml_text)
for actor in root.findall('{http://people.example.com}actor'):
    name = actor.find('{http://people.example.com}name')
    print(name.text)
    for char in actor.findall('{http://characters.example.com}character'):
        print(' |-->', char.text)


#Option-2

ns = {'real_person': 'http://people.example.com',
      'role': 'http://characters.example.com'}

for actor in root.findall('real_person:actor', ns):
    name = actor.find('real_person:name', ns)
    print(name.text)
    for char in actor.findall('role:character', ns):
        print(' |-->', char.text)


		
#XPath support - limited support 


import xml.etree.ElementTree as ET

root = ET.fromstring(countrydata)

# Top-level elements
root.findall(".")

# All 'neighbor' grand-children of 'country' children of the top-level
# elements
root.findall("./country/neighbor")

# Nodes with name='Singapore' that have a 'year' child
root.findall(".//year/..[@name='Singapore']")

# 'year' nodes that are children of nodes with name='Singapore'
root.findall(".//*[@name='Singapore']/year")

# All 'neighbor' nodes that are the second child of their parent
root.findall(".//neighbor[2]")



#Element Objects methods 
tag
A string identifying what kind of data this element represents (the element type, in other words).

text
tail
attrib
get(key, default=None)
Gets the element attribute named key.
Returns the attribute value, or default if the attribute was not found.

items()
Returns the element attributes as a sequence of (name, value) pairs.
The attributes are returned in an arbitrary order.

keys()
Returns the elements attribute names as a list. The names are returned in an arbitrary order.

set(key, value)
findtext(match, default=None, namespaces=None)
Finds text for the first subelement matching match. match may be a tag name or a path. 
Returns the text content of the first matching element, or default if no element was found. 
Note that if the matching element has no text content an empty string is returned.

getchildren()
Deprecated since version 3.2: Use list(elem) or iteration.

iter(tag=None)
Creates a tree iterator with the current element as the root. 
The iterator iterates over this element and all elements below it, in document (depth first) order. 
If tag is not None or '*', only elements whose tag equals tag are returned from the iterator. 
If the tree structure is modified during iteration, the result is undefined.


iterfind(match, namespaces=None)
Finds all matching subelements, by tag name or path. 
Returns an iterable yielding all matching elements in document order. 

itertext()
Creates a text iterator. The iterator loops over this element and all subelements, in document order, 
and returns all inner text




####lxml - XML and HTML with Python - mimics and superior to ElementTree API
#lxml is a Pythonic, mature binding for the libxml2 and libxslt libraries. 
#It extends the ElementTree API significantly to offer support for XPath, RelaxNG, XML Schema, XSLT, C14N 


#installation - pip install lxml

from lxml import etree

#An Element is the main container object for the ElementTree API

from lxml import etree
page = etree.Element('html')  #root element 
doc = etree.ElementTree(page)  #create document tree 
headElt = etree.SubElement(page, 'head')
bodyElt = etree.SubElement(page, 'body')
title = etree.SubElement(headElt, 'title')
title.text = 'Your page title here'
#<link rel='stylesheet' href='mystyle.css' type='text/css'>
linkElt = etree.SubElement(headElt, 'link', rel='stylesheet',  href='mystyle.css', type='text/css')
outFile = open('homemade.xml', 'wb')  #binary because doc generates b''
doc.write(outFile)  #save the doc 
print(etree.tostring(page, pretty_print=True).decode('utf-8')) #decode for py3

#adding a attribute 
linkNode.attrib['href'] = 'http://www.nmt.edu/'


##namespace map 
nsm = {"xsl":  "http://www.w3.org/1999/XSL/Transform",
       "fo":   "http://www.w3.org/1999/XSL/Format",
       "date": "http://exslt.org/dates-and-times"}
           
#To define the NSURI of the blank namespace, use an entry whose key is None
#this namespace map would define elements without a namespace as belonging to XHTML, 
#and elements with namespace prefix “xl:” belong to the XLink namespace: 
    nsm = {None: "http://www.w3.org/1999/xhtml",
           "xl": "http://www.w3.org/1999/xlink"}

#Example 
import sys
from lxml import etree as et

HTML_NS  =  "http://www.w3.org/TR/xhtml1/DTD/xhtml1-strict.dtd"
XSL_NS   =  "http://www.w3.org/1999/XSL/Transform"
NS_MAP = {None:  HTML_NS,
          "xsl": XSL_NS}

rootName = et.QName(XSL_NS, 'stylesheet')
root = et.Element(rootName, nsmap=NS_MAP)
sheet = et.ElementTree(root)

top = et.SubElement(root, et.QName(XSL_NS, "template"), match='/')
html = et.SubElement(top, et.QName(HTML_NS, "html"))
head = et.SubElement(html, "head")
title = et.SubElement(head, "title")
title.text = "Heading title"
body = et.SubElement(html, "body")
h1 = et.SubElement(body, "h1")
h1.text = "Body heading"
p = et.SubElement(body, "p")
p.text = "Paragraph text"
sheet.write(sys.stdout, pretty_print=True)


#output is 
<xsl:stylesheet xmlns:xsl="http://www.w3.org/1999/XSL/Transform"
    xmlns="http://www.w3.org/TR/xhtml1/DTD/xhtml1-strict.dtd">
  <xsl:template match="/">
    <html>
      <head>
        <title>Heading title</title>
      </head>
      <body>
        <h1>Body heading</h1>
        <p>Paragraph text</p>
      </body>
    </html>
  </xsl:template>
</xsl:stylesheet>



           
           
           
##Features of the etree module
#The Comment() constructor
newComment = etree.Comment(s)
bodyElt.append(newComment)

#The Element() constructor - etree.Element(tag, attrib={}, nsmap=None, **extras)
# extras  - Any keyword arguments of the form name=value that you supply to the constructor are added to the element's attributes

newReed = etree.Element('reed', pitch='440', id='a4') #<reed pitch='440' id='a4'/>


#The ElementTree() constructor - etree.ElementTree(element=None, file=None)
#to transform a file named balrog.xml into an ElementTree
balrogTree = etree.ElementTree(file='balrog.xml')


#The fromstring() function: Create an element from a string
#it returns a new Element instance representing all that XML. 
etree.fromstring(s)



#The parse() function: build an ElementTree from a file
#to convert an XML file into an ElementTree is to use this function: 
etree.parse(source)



#The ProcessingInstruction() constructor - etree.ProcessingInstruction(target, text=None):

#The QName() constructor - etree.QName(text, tag=none)
#for combining the “namespace URI” part with the “local name” part.

qn = etree.QName("{http://www.w3.org/1999/XSL/Transform}template")  #clark's notation - {nsURI}local
qn = etree.QName("http://www.w3.org/1999/XSL/Transform", "template")





#The SubElement() constructor - SubElement(parent, tag, attrib={}, nsmap=None, **extras):
>>> st=etree.Element('state', name='New Mexico')
>>> co=etree.SubElement(st, 'county', name='Socorro')
>>> ppl=etree.SubElement(co, 'ppl', name='Luis Lopez')
>>> print etree.tostring(st)
<state name="New Mexico"><county name="Socorro"><ppl name="Luis Lopez"/>
</county></state>
>>>


#The tostring() function: Serialize as XML - etree.tostring(elt, pretty_print=False, encoding=None)
#To output Unicode, use the keyword argument encoding=unicode. 
print(etree.tostring(page, pretty_print=True, encoding='unicode'))


#The XMLID() function: Convert text to XML with a dictionary of id values - etree.XMLID(text)
#The return value is a tuple (E, D), where: 
#E is the converted XML as an Element instance rooting the converted tree,
#D is a dictionary whose keys are the values of id attributes in the converted tree, and each corresponding value is the Element instance that carried that id value

from lxml import etree

SOURCE = '''<dog id="Fido">
Woof!
<cat id="Fluff">Mao?</cat>
<rhino id="ZR"/>
</dog>'''
tree, idMap = etree.XMLID(SOURCE)

for id in sorted(idMap.keys()):
    elt = idMap[id].text or "(none)"
    print "Tag {0}, text is '{1}'".format(id, elt.strip())

#And its output: 
Tag Fido, text is 'Woof!'
Tag Fluff, text is 'Mao?'
Tag ZR, text is '(none)'




## class ElementTree: A complete XML document
#ElementTree.find() - ET.find(path[, namespaces=D]) - find single element as Element instance 
doc.find('h1')

#ElementTree.findall(): Find matching elements - ET.findall(path[, namespaces=N]) - find sequence of Element 


#ElementTree.findtext(): Retrieve the text content from any element - ET.findtext(path[, default=None][, namespaces=N])


#ElementTree.getiterator(): Make an iterator - ET.getiterator(tag=None)
for  elt in page.getiterator('div'):
    if  elt.attrib.has_key('class'):
        print elt.get('class')

#ElementTree.getroot(): Find the root element - ET.getroot()

#ElementTree.xpath(): Evaluate an XPath expression - ET.xpath(s)

#ElementTree.write(): Translate back to XML - ET.write(file, pretty_print=False)


##class Element: One element in the tree
#Attributes of an Element instance
.attrib         A dictionary containing the element attributes. 
                The keys are the attribute names, and each corresponding value is the attribute's value. 
.base           The base URI from an xml:base attribute that this element contains or inherits, if any; None otherwise. 
.prefix         The namespace prefix of this element, if any, otherwise None. 
.sourceline     The line number of this element when parsed, if known, otherwise None. 
.tag            The element's name. 
.tail           The text following this element's closing tag, up to the start tag of the next sibling element. If there was no text there, this attribute will have the value None. 
.text           The text inside the element, up to the start tag of the first child element. 
                If there was no text there, this attribute will have the value None. 


#Accessing the list of child elements
#an Element instance acts like a Python list, with its XML child elements acting as the members of that list
len(node)           to determine how many children an element has
E[i]                returns the child element of E at position i else IndexError 
E[i:j]              returns a list of the child elements between positions i and j. 
E[i] = c            replace one child of an element E with a new element c 
E[i:j] = seq        replace many 
del E[i]            delete 
del E[i:j]
for kid in node:    iteration 
    print kid.tag
    
#Not all children of an element are themselves elements. 
#• Processing instructions are instances of class etree._ProcessingInstruction. 
#• Comments are instances of class etree._Comment. 
issubclass(node, etree._Comment)

    
#Element.append(): Add a new element child - E.append(c)

#Element.clear(): Make an element empty

#Element.find(): Find a matching sub-element - E.find(path[, namespaces=D])
#path can be 
"tag"       Find the first child element whose name is "tag". 
"tag1/tag2/.../tagn" Find the first child element whose name is tag1; then, under that child element, find its first child named tag2; and so forth. 

node.find("county/seat")
#with nsmap 
nsd = {'mp': 'http://example.com/mphg/',
           'k':  'http://example.org/sirs/ns/'}
someNode.find('mp:roundtable/k:knight', namespaces=nsd}


#Element.findall(): Find all matching sub-elements - E.findall(path[, namespaces=N])


#Element.findtext(): Extract text content - E.findtext(path, default=None, namespaces=N)
>>> from lxml import etree
>>> node=etree.fromstring('<a><b>bum</b><b>ear</b><c/></a>')
>>> node.findtext('b')
'bum'


#Element.get(): Retrieve an attribute value with defaulting - E.get(key, default=None)
>>> from lxml import etree
>>> node = etree.fromstring('<mount species="Jackalope"/>')
>>> print node.get('species')
Jackalope


#Element.getchildren(): Get element children
>>> xml = '''<corral><horse n="2"/><cow n="17"/>
...  <cowboy n="2"/></corral>'''
>>> pen = etree.fromstring(xml)
>>> penContents = pen.getchildren()
>>> for  content in penContents:
...     print "%-10s %3s" % (content.tag, content.get("n", "0"))
... 
horse        2
cow         17
cowboy       2
>>> 

#Element.getiterator(): Make an iterator to walk a subtree - E.getiterator(tag=None)
        a
     b      e
  c     d
#A preorder traversal of this tree goes in this order: a, b, c, d, e. 
>>> xml = '''<a><b><c/><d/></b><e/></a>'''
>>> tree = etree.fromstring(xml)
>>> walkAll = tree.getiterator()
>>> for  elt in walkAll:
...     print elt.tag,
... 
a b c d e
>>> 



#Element.getroottree(): Find the ElementTree containing this element


#Element.insert(): Insert a new child element - E.insert(index, elt)
>>> node = etree.fromstring('<a><c0/><c1/><c2/></a>')
>>> newKid = etree.Element('c-1', laugh="Hi!")
>>> node.insert(0, newKid)
>>> etree.tostring(node)
'<a><c-1 laugh="Hi!"/><c0/><c1/><c2/></a>'


#Element.items(): Produce attribute names and values
>>> node = etree.fromstring("<event time='1830' cost='3.50' rating='nc-03'/>")
>>> node.items()
[('cost', '3.50'), ('time', '1830'), ('rating', 'nc-03')]
>>> 


#Element.iterancestors(): Find an element's ancestors - E.iterancestors(tag=None)
>>> xml = '''<class sci='Aves' eng='Birds'>
...   <order sci='Strigiformes' eng='Owls'>
...     <family sci='Tytonidae' eng='Barn-Owls'>
...       <genus sci='Tyto'>
...         <species sci='Tyto alba' eng='Barn Owl'/>
...       </genus>
...     </family>
...   </order>
... </class>'''
>>> root = etree.fromstring(xml)
>>> barney = root.xpath('//species') [0]
>>> print "%s: %s" % (barney.get('sci'), barney.get('eng'))
Tyto alba: Barn Owl
>>> for  ancestor in barney.iterancestors():
...     print ancestor.tag,
genus family order class
>>> for  fam in barney.iterancestors('family'):
...    print "%s: %s" % (fam.get('sci'), fam.get('eng'))
Tytonidae: Barn-Owls

#Element.iterchildren(): Find all children - E.iterchildren(reversed=False, tag=None)
>>> root=et.fromstring("<mom><aaron/><betty/><clarence/><dana/></mom>")
>>> for kid in root.getchildren():
...     print kid.tag
aaron
betty
clarence
dana
>>> for kid in root.iterchildren(reversed=True):
...     print kid.tag
... 
dana
clarence
betty
aaron
>>>



#Element.iterdescendants(): Find all descendants - E.iterdescendants(tag=None)
>>> xml = '''<root>
...   <grandpa>
...     <dad>
...       <yuo/>
...     </dad>
...   </grandpa>
... </root>'''
>>> root = etree.fromstring(xml)
>>> you = root.xpath('.//yuo')[0]
>>> for  anc in you.iterancestors():
...    print anc.tag,
dad grandpa root
>>> 



#Element.itersiblings(): Find other children of the same parent - E.itersiblings(preceding=False)
#preceding=True, the iterator will visit the siblings that precede E in document order. 
>>> root=etree.fromstring(
...   "<mom><aaron/><betty/><clarence/><dana/></mom>")
>>> betty=root.find('betty')
>>> for sib in betty.itersiblings(preceding=True):
...     print sib.tag
... 
aaron
>>> for sib in betty.itersiblings():
...     print sib.tag
... 
clarence
dana
>>>


#Element.keys(): Find all attribute names
>>> node = etree.fromstring("<event time='1830' cost='3.50' rating='nc-03'/>")
>>> node.keys()
['time', 'rating', 'cost']
>>> 


#Element.remove(): Remove a child element - E.remove(C)


#Element.set(): Set an attribute value - E.set(A, V)
#The other method is to store values into the .attrib dictionary of the Element instance. 
>>> node = etree.Element('div', id='u401')
>>> etree.tostring(node)
'<div id="u401"/>'
>>> node.set('class', 'flyer')
>>> etree.tostring(node)
'<div id="u401" class="flyer"/>'


#Element.xpath(): Evaluate an XPath expression - E.xpath(s[, namespaces=N][, var=value][, ...])
var=value  to define the values of XPath variables to be used in the evaluation of s. 
           For example, if you pass an argument count=17, the value of variable $count in the XPath expression will be 17. 

#Return value can be 
• A list of zero or more selected Element instances.
• A Python bool value for true/false tests. 
• A Python float value for numeric results. 
• A string for string results. 

>>> node=etree.fromstring('''<a>
...   a-text <b>b-text</b> b-tail <c>c-text</c> c-tail
... </a>''')
>>> alltext = node.xpath('descendant-or-self::text()')
>>> alltext
['\n  a-text ', 'b-text', ' b-tail ', 'c-text', ' c-tail\n']
>>> clump = "".join(alltext)
>>> clump
'\n  a-text b-text b-tail c-text c-tail\n'
>>> 

## lxml interfaces with BeautifulSoup through the lxml.html.soupparser module. 
soupparser.parse(input) 
The input argument specifies a Web page's HTML source as either a file name or a file-like object. 
The return value is an ElementTree instance whose root element is an html element as an Element instance. 

soupparser.fromstring(s) 
The s argument is a string containing some tag soup. 
The return value is a tree of nodes representing s. 
The root node of this tree will always be an html element as an Element instance. 

convert_tree() 
convert an existing BeautifulSoup tree into a list of top-level Element


from lxml.html import soupparser
tag_soup = '''<meta/><head><title>Hello</head><body onload=crash()>Hi all<p>'''

>>> root = soupparser.fromstring(tag_soup)

>>> from lxml.etree import tostring
>>> print(tostring(root, pretty_print=True).strip())
<html>
  <meta/>
  <head>
    <title>Hello</title>
  </head>
  <body onload="crash()">Hi all<p/></body>
</html>

##Automated validation of input files
#“Validation with a Relax NG schema”. 
#“Validation with an XSchema (XSD) schema”. 

etree.RelaxNG(S) or etree.XMLSchema(S) to convert that tree into a “schema instance,” 
where S is the ElementTree instance, containing the schema

Use the .validate(ET) method of the schema instance to validate ET. 
This method returns 1 if ET validates against the schema, or 0 if it does not. 


##lxml.etree.XSLT. 
#The class can be given an ElementTree or Element object to construct an XSLT transformer:
from lxml import etree
xslt_root = etree.XML('''\
<xsl:stylesheet version="1.0"
 xmlns:xsl="http://www.w3.org/1999/XSL/Transform">
 <xsl:template match="/">
 <foo><xsl:value-of select="/a/b/text()" /></foo>
 </xsl:template>
</xsl:stylesheet>''')
transform = etree.XSLT(xslt_root)

#usage 
from io import StringIO   
f = StringIO('<a><b>Text</b></a>')
doc = etree.parse(f)
result = transform(doc)

#The result of an XSL transformation can be accessed like a normal ElementTree document
>>> result.getroot().text
'Text'

#as opposed to normal ElementTree objects, can also be turned into an (XML or text) 
#string by applying the str() function:

>>> str(result)
'<?xml version="1.0"?>\n<foo>Text</foo>\n'




######################### Additional Topics #########################
#tornado - based on asynchronous networking library -Py2.7 and Py3.x
#download from http://www.lfd.uci.edu/~gohlke/pythonlibs/


#example - http://localhost:8888/
import tornado.ioloop
import tornado.web

class HelloHandler(tornado.web.RequestHandler):
    def get(self):                           #get method
        self.write("Hello, world")
		
class MainHandler(tornado.web.RequestHandler):
	def get(self):
		self.write('<a href="%s">link to story 1</a>' %  self.reverse_url("story", "1"))

class StoryHandler(tornado.web.RequestHandler):
	def initialize(self, db):    #Entry points
		self.db = db
	def get(self, story_id):
		self.write("this is story %s" % story_id)


if __name__ == "__main__":
	app = tornado.web.Application([
		(r"/", MainHandler),
		tornado.web.url(r"/story/([0-9]+)", StoryHandler, dict(db=2), name="story"),
		(r"/hello", HelloHandler),
    ])
	app.listen(8888)
	tornado.ioloop.IOLoop.current().start()
	
#The routing table(inside web.Application) is a list of URLSpec objects (or tuples), 
#each of which contains (at least) a regular expression and a handler class
#the URLSpec may have a name, which will allow it to be used with RequestHandler.reverse_url.



#Subclassing RequestHandler
entry point for a handler subclass is a method named after the HTTP method being handled: get(), post(), 
handler may define one or more of these methods(but one for each ) to handle different HTTP action

Within a handler, call methods such as RequestHandler.render or RequestHandler.write to produce a response.
render() loads a Template by name and renders it with the given arguments. 
write() is used for non-template-based output; it accepts strings, bytes, and dictionaries 
(dicts will be encoded as  JSON).


#Handling request input
Current request: self.request. 
Use self.get_query_argument(name, default=[], strip=True) for get and self.get_body_argument(name, default=[], strip=True) for post
For lists, use get_query_arguments(name, strip=True) and get_body_arguments(name, strip=True) instead of their singular counterparts

class MyFormHandler(tornado.web.RequestHandler):
    def get(self):
        self.write('<html><body><form action="/myform" method="POST">'
                   '<input type="text" name="message">'
                   '<input type="submit" value="Submit">'
                   '</form></body></html>')

    def post(self):
        self.set_header("Content-Type", "text/plain")
        self.write("You wrote " + self.get_body_argument("message"))

	
#File uploads
Files uploaded via a form are available in self.request.files, 
which maps names (the name of the HTML <input type="file"> element) to a list of files.
Each file is a dictionary of the form {"filename":..., "content_type":..., "body":...}. 
The files object is only present if the files were uploaded  with a form wrapper (i.e. a multipart/form-data Content-Type)
	
#Json handling	
Applications that wish to use JSON instead of form-encoding may override prepare to parse their requests

def prepare(self):
    if self.request.headers["Content-Type"].startswith("application/json"):
        self.json_args = json.loads(self.request.body)
    else:
        self.json_args = None
		

#Error Handling
If a handler raises an exception, Tornado will call RequestHandler.write_error to generate an error page. 
tornado.web.HTTPError can be used to generate a specified status code; 
all other exceptions return a 500 status

#Redirection
Use RequestHandler.redirect and with the RedirectHandler.
use self.redirect() within a RequestHandler method to redirect users elsewhere.
Unlike RequestHandler.redirect, RedirectHandler uses permanent redirects by default
#example

app = tornado.web.Application([
    url(r"/app", tornado.web.RedirectHandler,
        dict(url="http://itunes.apple.com/my-app-id")),
    ])


#RedirectHandler also supports regular expression substitutions. 

app = tornado.web.Application([
    url(r"/photos/(.*)", MyPhotoHandler),
    url(r"/pictures/(.*)", tornado.web.RedirectHandler,
        dict(url=r"/photos/\1")),
    ])




#Asynchronous handlers
Tornado handlers are synchronous by default

To make asynchronous, 
Use coroutines
or 
use callback style by using tornado.web.asynchronous 

class MainHandler(tornado.web.RequestHandler):
    @tornado.web.asynchronous
    def get(self):
        http = tornado.httpclient.AsyncHTTPClient()
        http.fetch("http://friendfeed-api.com/v2/feed/bret",  callback=self.on_response)

    def on_response(self, response):
        if response.error: raise tornado.web.HTTPError(500)
        json = tornado.escape.json_decode(response.body)
        self.write("Fetched " + str(len(json["entries"])) + " entries "
                   "from the FriendFeed API")
        self.finish()

#using a coroutine:


class MainHandler(tornado.web.RequestHandler):
    @tornado.gen.coroutine
    def get(self):
        http = tornado.httpclient.AsyncHTTPClient()
        response = yield http.fetch("http://friendfeed-api.com/v2/feed/bret")
        json = tornado.escape.json_decode(response.body)
        self.write("Fetched " + str(len(json["entries"])) + " entries "
                   "from the FriendFeed API")


#Configuring templates - in same directory of .py file


#Template syntax : file : template.html

<html>
   <head>
      <title>{{ title }}</title>
   </head>
   <body>
     <ul>
       {% for item in items %}
         <li>{{ escape(item) }}</li>
       {% end %}
     </ul>
   </body>
 </html>


#handler

class MainHandler(tornado.web.RequestHandler):
    def get(self):
        items = ["Item 1", "Item 2", "Item 3"]
        self.render("template.html", title="My title", items=items)


#Control statements are surrounded by {% and %}, e.g., {% if len(items) > 2 %}. 
#Expressions are surrounded by {{ and }}, e.g., {{ items[0] }}.

#builtin methods 
•escape: alias for tornado.escape.xhtml_escape
•xhtml_escape: alias for tornado.escape.xhtml_escape
•url_escape: alias for tornado.escape.url_escape
•json_encode: alias for tornado.escape.json_encode
•squeeze: alias for tornado.escape.squeeze
•linkify: alias for tornado.escape.linkify
•datetime: the Python datetime module
•handler: the current RequestHandler object
•request: alias for handler.request
•current_user: alias for handler.current_user
•locale: alias for handler.locale
•_: alias for handler.locale.translate
•static_url: alias for handler.static_url
•xsrf_form_html: alias for handler.xsrf_form_html
reverse_url: alias for Application.reverse_url
•All entries from the ui_methods and ui_modules Application settings
•Any keyword arguments passed to render or render_string

#UI modules

Tornado supports reusable UI widgets across your application. 
UI modules are like special function calls to render components of your page, 
and they can come packaged with their own CSS and JavaScript



#How make async in Tornado by using below
•Callback argument
•Return a placeholder (Future, Promise, Deferred)
•Deliver to a queue
•Callback registry (eg POSIX signals)

#Example-synchronous function:


from tornado.httpclient import HTTPClient

def synchronous_fetch(url):
    http_client = HTTPClient()
    response = http_client.fetch(url)
    return response.body


#async with a callback 


from tornado.httpclient import AsyncHTTPClient

def asynchronous_fetch(url, callback):
    http_client = AsyncHTTPClient()
    def handle_response(response):
        callback(response.body)
    http_client.fetch(url, callback=handle_response)

#async with a Future 

from tornado.concurrent import Future

def async_fetch_future(url):
    http_client = AsyncHTTPClient()
    my_future = Future()
    fetch_future = http_client.fetch(url)
    fetch_future.add_done_callback(
        lambda f: my_future.set_result(f.result()))
    return my_future

#using coroutines - Coroutines are the recommended way to write asynchronous code in Tornado

from tornado import gen

@gen.coroutine
def fetch_coroutine(url):
    http_client = AsyncHTTPClient()
    response = yield http_client.fetch(url)
    return response.body              #for <Py3.3 , use raise gen.Return(response.body)

 
#How to call a coroutine - must be in another coroutine, do 'yield other_coroutine'
#to return from co-routine, use raise gen.Return(o) for <p3.3 or return o for >=Py3.3 
@gen.coroutine
def divide(x, y):
    return x / y

def bad_call():
    # This should raise a ZeroDivisionError, but it won't because
    # the coroutine is called incorrectly.
    divide(1, 0)

#goo way

@gen.coroutine
def good_call():
    # yield will unwrap the Future returned by divide() and raise
    # the exception.
    yield divide(1, 0)

#to do "fire and forget" a coroutine 

# The IOLoop will catch the exception and print a stack trace in
# the logs. Note that this doesn't look like a normal call, since
# we pass the function object to be called by the IOLoop.
IOLoop.current().spawn_callback(divide, 1, 0)

#At top level or the first co-routine, execute as below

# run_sync() doesn't take arguments, so we must wrap the
# call in a lambda.
IOLoop.current().run_sync(lambda: divide(1, 0))



#Interaction with callbacks
To interact with asynchronous code that uses callbacks instead of Future, 
wrap the call in a Task. returns a Future which you can yield:


@gen.coroutine
def call_task():
    # Note that there are no parens on some_function.
    # This will be translated by Task into
    #   some_function(other_args, callback=callback)
    yield gen.Task(some_function, other_args)

#Calling blocking functions
use a ThreadPoolExecutor, which returns Futures that are compatible with coroutines:


thread_pool = ThreadPoolExecutor(4)

@gen.coroutine
def call_blocking():
    yield thread_pool.submit(blocking_func, args)

#Parallelism
The coroutine decorator recognizes lists and dicts whose values are Futures, 
and waits for all of those Futures in parallel:


@gen.coroutine
def parallel_fetch(url1, url2):
    resp1, resp2 = yield [http_client.fetch(url1),
                          http_client.fetch(url2)]

@gen.coroutine
def parallel_fetch_many(urls):
    responses = yield [http_client.fetch(url) for url in urls]
    # responses is a list of HTTPResponses in the same order

@gen.coroutine
def parallel_fetch_dict(urls):
    responses = yield {url: http_client.fetch(url)
                        for url in urls}
    # responses is a dict {url: HTTPResponse}

#Interleaving
to save a Future instead of yielding it immediately, 
so you can start another operation before waiting:


@gen.coroutine
def get(self):
    fetch_future = self.fetch_next_chunk()
    while True:
        chunk = yield fetch_future
        if chunk is None: break
        self.write(chunk)
        fetch_future = self.fetch_next_chunk()
        yield self.flush()


#Looping
Must separate the loop condition from accessing the results,


import motor
db = motor.MotorClient().test

@gen.coroutine
def loop_example(collection):
    cursor = db.collection.find()
    while (yield cursor.fetch_next):
        doc = cursor.next_object()

		

#periodically Running in the background

@gen.coroutine
def minute_loop():
    while True:
        yield do_something()
        yield gen.sleep(60)

# Coroutines that loop forever are generally started with
# spawn_callback().
IOLoop.current().spawn_callback(minute_loop)

#the previous loop runs every 60+N seconds, where N is the running time of do_something(). 
#To run exactly every 60 seconds, use the interleaving pattern from above:


@gen.coroutine
def minute_loop2():
    while True:
        nxt = gen.sleep(60)   # Start the clock.
        yield do_something()  # Run while the clock is ticking.
        yield nxt             # Wait for the timer to run out.


#Queue - Tornado’s tornado.queues module implements an asynchronous producer / consumer pattern for coroutines
#Because , all must be co-routine and inside coroutine, you do yield other coroutine
Producer does yield q.put(o)
Consumer does yield q.get() and then after work q.task_done()
#check webcrawler







#########################
###Use BeautifulSoup for HTML/XML processing
#no xpath, but css select 
(pip install beautifulsoup4)

from bs4 import BeautifulSoup
import requests

r  = requests.get("http://www.yahoo.com")
data = r.text

soup = BeautifulSoup(data)

for link in soup.find_all('a'):  # tag <a href=".."
		print(link.get('href'))	 # Attribute href

#pip install requests
#pip install BeautifulSoup4

#Signature: find_all(name, attrs, recursive, text, limit, **kwargs)
#	tag becomes attributes of soup object


soup = BeautifulSoup('<html><body><p class="title">data</p></body></html>')
soup.html
soup.html.body.p
soup.html.body.text
soup.html.body.attrs
soup.html.body.name
soup.html.body.p['class']
soup.body
soup.body.attrs
soup.p.text  		# can call nested .p directly as well 
soup.get_text()  
soup.html.name
soup.p.parent.name

#Example - multivalued attribute

# class defines multivalued attribute

css_soup = BeautifulSoup('<p class="body strikeout"></p>')
css_soup.p['class']
# ["body", "strikeout"]

css_soup = BeautifulSoup('<p class="body"></p>')
css_soup.p['class']
# ["body"]


# but id does not

id_soup = BeautifulSoup('<p id="my id"></p>')
id_soup.p['id']
# 'my id'


# returning back


rel_soup = BeautifulSoup('<p>Back to the <a rel="index">homepage</a></p>')
rel_soup.a['rel']
# ['index']
rel_soup.a['rel'] = ['index', 'contents']
print(rel_soup.p)
# <p>Back to the <a rel="index contents">homepage</a></p>

#	Pretty-printing

print(soup.prettify())

#	Non-pretty printing

print(str(soup))

#	Special tags

.contents and .children
A tag’s children are available in a list called .contents
The .contents and .children attributes only consider a tag’s direct children

.descendants
The .descendants attribute iterates over all of a tag’s children, recursively: its direct children, the children of its direct children, and so on

.parent
 an element’s parent 

.parents
iterate over all of an element’s parents 

#Example
soup = BeautifulSoup('<html><body><p class="title">data</p></body></html>')
>>> soup.html.contents
[<body><p class="title">data</p></body>]
>>> soup.html.children
<list_iterator object at 0x6fffeae4c50>
>>> list(soup.html.children)
[<body><p class="title">data</p></body>]
>>> list(soup.html.descendants)
[<body><p class="title">data</p></body>, <p class="title">data</p>, 'data']
>>>
>>> list(soup.p.parent)
[<p class="title">data</p>]
>>> list(soup.p.parents)
[<body><p class="title">data</p></body>, <html><body><p class="title">data</p></body></html>, <html><body><p class="title">data</p></body></html>]



.next_sibling and .previous_sibling
navigate to sibling

.next_siblings and .previous_siblings
iterate all siblings

.next_element and .previous_element
next or previous element 

.next_elements and .previous_elements
iterate all elements

#	Searching the string
find_all(tag_name, attrs_value, recursive, text, limit, **kwargs)  
The find_all() method looks through a tag’s descendants and retrieves all descendants that match your filters

soup.find_all('p')
soup.find_all("p", "title")

#	With RE
import re
for tag in soup.find_all(re.compile("^b")):
		print(tag.name)

#	With List
soup.find_all(["a", "b"])

#	True
The value True matches everything it can. This code finds all the tags in the document, but none of the text strings


for tag in soup.find_all(True):
    print(tag.name)


#	With function

def has_class_but_no_id(tag):
		return tag.has_attr('class') and not tag.has_attr('id')

soup.find_all(has_class_but_no_id)

#	Other find methods


find_parents(name, attrs, text, limit, **kwargs)
find_parent(name, attrs, text, **kwargs)
find_next_siblings(name, attrs, text, limit, **kwargs)
find_next_sibling(name, attrs, text, **kwargs)
find_previous_siblings(name, attrs, text, limit, **kwargs)
find_previous_sibling(name, attrs, text, **kwargs)
find_all_next(name, attrs, text, limit, **kwargs)
find_next(name, attrs, text, **kwargs)
find_all_previous(name, attrs, text, limit, **kwargs)
find_previous(name, attrs, text, **kwargs)

#	CSS Selector (supports a subset of CSS3)

soup.select("title")
soup.select("p nth-of-type(3)")
soup.select("body a")
soup.select("html head title")
soup.select("head > title")
soup.select("p > a")
soup.select("p > a:nth-of-type(2)")
soup.select("p > #link1")
soup.select("body > a")


#CSS Selector(contains all CSS3) Reference

Selector 			Example 		Description
.class 				.intro 			Selects all elements with class="intro" 
#id 				#firstname 		Selects the element with id="firstname" 
* 					* 				Selects all elements 
element 			p 				Selects all <p> elements 
element,element 	div, p 			Selects all <div> elements and all <p> elements 
element element 	div p 			Selects all <p> elements inside <div> elements 
element>element 	div > p 		Selects all <p> elements where the parent is a <div> element 
element+element 	div + p 		Selects all <p> elements that are placed immediately after <div> elements 
element1~element2 	p ~ ul 			Selects every <ul> element that are preceded by a <p> element 
[attribute] 		[target] 		Selects all elements with a target attribute 
[attribute=value] 	[target=_blank] Selects all elements with target="_blank" 
[attribute~=value] 	[title~=flower] Selects all elements with a title attribute containing the word "flower" 
[attribute|=value] 	[lang|=en] 		Selects all elements with a lang attribute value starting with "en" 
[attribute^=value] 	a[href^="https"] Selects every <a> element whose href attribute value begins with "https" 
[attribute$=value] 	a[href$=".pdf"] Selects every <a> element whose href attribute value ends with ".pdf" 
[attribute*=value] 	a[href*="w3schools"] Selects every <a> element whose href attribute value contains the substring "w3schools" 
:active 			a:active 		Selects the active link 
::after 			p::after 		Insert something after the content of each <p> element 
::before 			p::before 		Insert something before the content of each <p> element 
:checked 			input:checked 	Selects every checked <input> element 
:disabled 			input:disabled 	Selects every disabled <input> element 
:empty 				p:empty 		Selects every <p> element that has no children (including text nodes) 
:enabled 			input:enabled 	Selects every enabled <input> element 
:first-child 		p:first-child 	Selects every <p> element that is the first child of its parent 
::first-letter 		p::first-letter Selects the first letter of every <p> element 
::first-line 		p::first-line 	Selects the first line of every <p> element 
:first-of-type 		p:first-of-type Selects every <p> element that is the first <p> element of its parent 
:focus 				input:focus 	Selects the input element which has focus 
:in-range 			input:in-range 	Selects input elements with a value within a specified range 
:invalid 			input:invalid 	Selects all input elements with an invalid value 
:last-child 		p:last-child 	Selects every <p> element that is the last child of its parent 
:last-of-type 		p:last-of-type 	Selects every <p> element that is the last <p> element of its parent 
:link 				a:link 			Selects all unvisited links 
:not(selector) 		:not(p) 		Selects every element that is not a <p> element  
:nth-child(n) 		p:nth-child(2) 	Selects every <p> element that is the second child of its parent 
:nth-last-child(n) 	p:nth-last-child(2) Selects every <p> element that is the second child of its parent, counting from the last child 
:nth-last-of-type(n) p:nth-last-of-type(2) Selects every <p> element that is the second <p> element of its parent, counting from the last child 
:nth-of-type(n) 	p:nth-of-type(2) Selects every <p> element that is the second <p> element of its parent 
:only-of-type 		p:only-of-type 	Selects every <p> element that is the only <p> element of its parent 
:only-child 		p:only-child 	Selects every <p> element that is the only child of its parent 
:optional 			input:optional 	Selects input elements with no "required" attribute 
:out-of-range 		input:out-of-range Selects input elements with a value outside a specified range 
:required 			input:required Selects input elements with the "required" attribute specified 
:root 				:root 			Selects the document's root element 
::selection 		::selection 	Selects the portion of an element that is selected by a user   
:target 			#news:target  	Selects the current active #news element (clicked on a URL containing that anchor name) 3 
:valid 				input:valid 	Selects all input elements with a valid value 
:visited 			a:visited 		Selects all visited links 


#Modifying the tree

soup = BeautifulSoup('<b class="boldest">Extremely bold</b>')
tag = soup.b

tag.name = "blockquote"
tag['class'] = 'verybold'
tag['id'] = 1
tag
# <blockquote class="verybold" id="1">Extremely bold</blockquote>

del tag['class']
del tag['id']
tag.string = "New link text"
tag
# <blockquote> New link text </blockquote>







#########################
#SQLAlchemy (Py3.x and Py2.7)
#pip3 install sqlalchemy
#or windows from http://www.lfd.uci.edu/~gohlke/pythonlibs/#sqlalchemy (does not get installed)

#start cygwin sqld
#mysql -running daemon in cygwin
/usr/bin/mysqld_safe &
#shutting down
mysqladmin.exe -h 127.0.0.1 -u root   --connect-timeout=5 shutdown
#mysql admin #  default port 3306, 
mysql -u root    -h 127.0.0.1 
#few commands
show databases;
create database python;
use python;
show tables;
create table employes ( id INT, first_name VARCHAR(20), last_name VARCHAR(20), hire_date  DATE);
desc employes;
insert into employes values (3, "das", "das", '1999-03-30');
select * from employes; 
drop table employes;

##Quick Tutorial 
https://www.bytefish.de/blog/first_steps_with_sqlalchemy/

@@@
Prerequisites

virtualenv

I like working with virtualenv, because it keeps your Python installation clean and prevents you from messing around with Package dependencies. You can install virtualenv with pip (administrative privileges may be required):

pip install virtualenv


Once you have installed virtualenv, decide where you want to create the virtual environments at. virtualenv normally creates environments in the current working directory. I have created a separate folder virtualenv, where my Python environments go to.

To create the sqlalchemy virtual environment for this tutorial simply type:

PS D:\virtualenv> virtualenv sqlalchemy
New python executable in sqlalchemy\Scripts\python.exe
Installing setuptools, pip...done.


virtualenv created a new directory sqlalchemy and the scripts to activate the virtual environment:

PS D:\virtualenv> .\sqlalchemy\Scripts\activate


The virtual environment is activated and the name of the virtual environment is prepended to the command line prompt:

(sqlalchemy) PS D:\virtualenv>


ipython

ipython is an amazing command shell for Python. I use it for all my Python development and you can also use it for this tutorial.

Install ipython with:

pip install ipython


Then simply start it. And whenever you see a code snippet in this tutorial, you can copy it to the clipboard and then use ipythons paste magic:

>>> %paste


sqlalchemy

Finally install SQLAlchemy with:

pip install SQLAlchemy


SQLAlchemy

Every sufficiently complex application needs to persist data, and there are a million ways to persist data. You could use a flat file, a [document-oriented approach][3] or a [Relational database][4]. SQLAlchemy, as the name implies, is a SQL Toolkit and Object-Relational Mapper. 

Scope of this Article

The scope of this article is an introduction to SQLAlchemy. It's not a complete guide and it does not cover any essential parts for real database access, like sessions, caching, database migrations or any other advanced topics. It is only meant to be a quick introduction and show how to build mapping tables and query data with SQLAlchemy. 

The Example

The database application we are going to build should persist images with associated likes, tags and comments. In a later article I want to use it to build a small Web service around it.

So let's take look at the entities first. An image consist of a UUID and its associated number of likes. Each image can be associated with many tags, a tag can be associated with many images. That's a many-to-many relationship, so we need a mapping table. Finally each image can have multiple comments, a one-to-many relation with a foreign key on the comments side.

Model

This tutorial uses the declarative extensions of SQLAlchemy. declarative_base is a factory function, that returns a base class (actually a metaclass), and the entities are going to inherit from it. Once the definition of the class is done, the Table and mapper will be generated automatically. There is some magic involved, but on the other hand SQLAlchemy forces you to explicitly define things like the table name, primary keys and relationships.

First create the Base class:

from sqlalchemy.ext.declarative import declarative_base

Base = declarative_base()


The entities are classes, which derive from the Base class. We are also using the relationship function to define the relationships between the entities. The many-to-many relationship between tags and images requires us to define an association table, which we'll be joining over. 

When defining the images relationships, we are also using the backref parameter, which adds the image properties to the tags and comments entities. We want those references to be dynamically loaded, because we probably don't want to load all images, when accessing these entities.

The code is relatively straightforward to read:

from datetime import datetime, timedelta
from sqlalchemy import Table, Column, Integer, String, DateTime, ForeignKey
from sqlalchemy.orm import relationship, backref

tags = Table('tag_image', Base.metadata,
    Column('tag_id', Integer, ForeignKey('tags.id')),
    Column('image_id', Integer, ForeignKey('images.id'))
)

class Image(Base):

    __tablename__ = 'images'

    id          =   Column(Integer, primary_key=True)
    uuid        =   Column(String(36), unique=True, nullable=False)
    likes       =   Column(Integer, default=0)
    created_at  =   Column(DateTime, default=datetime.utcnow)
    tags        =   relationship('Tag', secondary=tags, 
                        backref = backref('images', lazy='dynamic'))
    comments    =   relationship('Comment', backref='image', lazy='dynamic')

    def __repr__(self):
        str_created_at = self.created_at.strftime("%Y-%m-%d %H:%M:%S")
        return "<Image (uuid='%s', likes='%d', created_at=%s)>" % (self.uuid, self.likes, str_created_at)

class Tag(Base):

    __tablename__ = 'tags'

    id      =   Column(Integer, primary_key=True)
    name    =   Column(String(255), unique=True, nullable=False)

    def __repr__(self):
        return "<Tag (name='%s')>" % (self.name)

class Comment(Base):

    __tablename__ = 'comments'

    id          =   Column(Integer, primary_key=True)
    text        =   Column(String(2000))
    image_id    =   Column(Integer, ForeignKey('images.id'))

    def __repr__(self):
        return "<Comment (text='%s')>" % (self.text)


Connecting and Creating the Schema

First of all we need to create the engine, which is used to connect to the database. This example uses SQLite3, which should already be included in your Python installation.

from sqlalchemy import create_engine

engine = create_engine('sqlite:///:memory:', echo=True)


A call to the metadata of the Base class then generates the Schema:

Base.metadata.create_all(engine)


Since we have set echo=True for the engine, we can see the generated SQL:

CREATE TABLE tags (
        id INTEGER NOT NULL,
        name VARCHAR(255) NOT NULL,
        PRIMARY KEY (id),
        UNIQUE (name)
)

COMMIT

CREATE TABLE images (
        id INTEGER NOT NULL,
        uuid VARCHAR(36) NOT NULL,
        likes INTEGER,
        created_at DATETIME,
        PRIMARY KEY (id),
        UNIQUE (uuid)
)

COMMIT

CREATE TABLE tag_image (
        tag_id INTEGER,
        image_id INTEGER,
        FOREIGN KEY(tag_id) REFERENCES tags (id),
        FOREIGN KEY(image_id) REFERENCES images (id)
)

COMMIT

CREATE TABLE comments (
        id INTEGER NOT NULL,
        text VARCHAR(2000),
        image_id INTEGER,
        PRIMARY KEY (id),
        FOREIGN KEY(image_id) REFERENCES images (id)
)
COMMIT


Sessions

A Session is a Python class, which handles the conversation with the database for us. It implements the Unit of Work pattern for synchronizing changes to the database. Basically it tracks all records you add or modify. We can acquire a Session class with the sessionmaker, which simplifies the configuration, since we only bind our database engine to it:

from sqlalchemy.orm import sessionmaker

Session = sessionmaker(bind=engine)


And now whenever we want to talk to the database, we can create a new Session:

session = Session()


Insert

First of all, we'll create some data:

tag_cool = Tag(name='cool')
tag_car = Tag(name='car')
tag_animal = Tag(name='animal')

comment_rhino = Comment(text='Rhinoceros, often abbreviated as rhino, is a group of five extant species of odd-toed ungulates in the family Rhinocerotidae.')

image_car = Image(uuid='uuid_car', \
    tags=[tag_car, tag_cool], \
    created_at=(datetime.utcnow() - timedelta(days=1)))

image_another_car = Image(uuid='uuid_anothercar', \
    tags=[tag_car])

image_rhino = Image(uuid='uuid_rhino', \
    tags=[tag_animal], \
    comments=[comment_rhino])


And then we can get a new session object, add the records and commit the work:

session = Session()

session.add(tag_cool)
session.add(tag_car)
session.add(tag_animal)

session.add(comment_rhino)

session.add(image_car)
session.add(image_another_car)
session.add(image_rhino)

session.commit()


The generated SQL appears in the command prompt:

BEGIN

INSERT INTO tags (name) VALUES (?)
('cool',)

INSERT INTO tags (name) VALUES (?)
('car',)

INSERT INTO tags (name) VALUES (?)
('animal',)

INSERT INTO images (uuid, likes, created_at) VALUES (?, ?, ?)
('uuid_car', 0, '2014-12-20 19:16:19.822000')

INSERT INTO images (uuid, likes, created_at) VALUES (?, ?, ?)
('uuid_anothercar', 0, '2014-12-21 19:16:19.828000')

INSERT INTO images (uuid, likes, created_at) VALUES (?, ?, ?)
('uuid_rhino', 0, '2014-12-21 19:16:19.829000')

INSERT INTO tag_image (tag_id, image_id) VALUES (?, ?)
((2, 1), (1, 1), (3, 3), (2, 2))

INSERT INTO comments (text, image_id) VALUES (?, ?)
('Rhinoceros, often abbreviated as rhino, is a group of five extant species of odd-toed ungulates in the family Rhinocerotidae.', 3)

COMMIT


Update

Updating a record is easy... Imagine someone upvoted an image, and we get UUID:

# Find the image with the given uuid:
image_to_update = session.query(Image).filter(Image.uuid == 'uuid_rhino').first()
# Increase the number of upvotes:
image_to_update.likes = image_to_update.likes + 1
# And commit the work:
session.commit()


SQLAlchemy translates this to SQL as:

SELECT images.id AS images_id, images.uuid AS images_uuid, images.likes AS images_likes, images.created_at AS images_created_at
FROM images
WHERE images.uuid = ?
LIMIT ? OFFSET ?
('uuid_rhino', 1, 0)

UPDATE images SET likes=? WHERE images.id = ?
(1, 3)

COMMIT


Delete

Deleting an entity is as easy as calling delete on the session object:

session.delete(image_rhino)


In the following SQL we can see that the comments of an image do not get deleted once the image is deleted. We can also see, that orphaned tags are not deleted, see this stackoverflow post for the very detailed reason.

SELECT images.id AS images_id, images.uuid AS images_uuid, images.likes AS images_likes, images.created_at AS images_created_at
FROM images
WHERE images.id = ?
(3,)

SELECT comments.id AS comments_id, comments.text AS comments_text, comments.image_id AS comments_image_id
FROM comments
WHERE ? = comments.image_id
(3,)

sqlalchemy.engine.base.Engine SELECT tags.id AS tags_id, tags.name AS tags_name 
FROM tags, tag_image
WHERE ? = tag_image.image_id AND tags.id = tag_image.tag_id
(3,)

DELETE FROM tag_image WHERE tag_image.tag_id = ? AND tag_image.image_id = ?
(3, 3)

UPDATE comments SET image_id=? WHERE comments.id = ?
(None, 1)

DELETE FROM images WHERE images.id = ?
(3,)

COMMIT
(1,)


It is up to you to decide wether this is an acceptable situation for your application or not. If you want to prevent the foreign key from being set to null, then declare the column as not nullable (Column(Integer, ForeignKey('images.id'), nullable=False). A delete to an image will then fail with an IntegrityError.

If you want the comments to be deleted, when the parent image is deleted, then add a cascade = "all,delete" to the relationship declaration:

comments = relationship('Comment', cascade = "all,delete", backref='image', lazy='dynamic')


I know there are cascaded deletes and updates in a lot of databases, but my SQLite version doesn't seem to respect them.

Queries

If you are familiar with SQL, then writing queries with SQLAlchemy is easy for you. Here are some queries you can fire against the database of this article.

# Get a list of tags:
for name in session.query(Tag.name).order_by(Tag.name):
    print name

# How many tags do we have?
session.query(Tag).count()

# Get all images created yesterday:
session.query(Image) \
    .filter(Image.created_at < datetime.utcnow().date()) \
    .all()

# Get all images, that belong to the tag 'car' or 'animal', using a subselect:
session.query(Image) \
    .filter(Image.tags.any(Tag.name.in_(['car', 'animal']))) \
    .all()

# This can also be expressed with a join:
session.query(Image) \
    .join(Tag, Image.tags) \
    .filter(Tag.name.in_(['car', 'animal'])) \
    .all()

# Play around with functions:
from sqlalchemy.sql import func, desc

max_date = session.query(func.max(Image.created_at))
session.query(Image).filter(Image.created_at == max_date).first()

# Get a list of tags with the number of images:
q = session.query(Tag, func.count(Tag.name)) \
    .outerjoin(Image, Tag.images) \
    .group_by(Tag.name) \
    .order_by(desc(func.count(Tag.name))) \
    .all()

for tag, count in q:
    print 'Tag "%s" has %d images.' % (tag.name, count) 

# Get images created in the last two hours and zero likes so far:
session.query(Image) \
    .join(Tag, Image.tags) \
    .filter(Image.created_at > (datetime.utcnow() - timedelta(hours=2))) \
    .filter(Image.likes == 0) \
    .all()


Conclusion

SQLAlchemy is fun to work with. Defining a schema is dead simple and the query language feels like writing native SQL. In the next article we are going to see how Flask uses SQLAlchemy and write a small Web service with it.

Appendix

Model

from sqlalchemy.ext.declarative import declarative_base

Base = declarative_base()

from datetime import datetime, timedelta
from sqlalchemy import Table, Column, Integer, String, DateTime, ForeignKey
from sqlalchemy.orm import relationship, backref

tags = Table('tag_image', Base.metadata,
    Column('tag_id', Integer, ForeignKey('tags.id')),
    Column('image_id', Integer, ForeignKey('images.id'))
)

class Image(Base):

    __tablename__ = 'images'

    id          =   Column(Integer, primary_key=True)
    uuid        =   Column(String(36), unique=True, nullable=False)
    likes       =   Column(Integer, default=0)
    created_at  =   Column(DateTime, default=datetime.utcnow)
    tags        =   relationship('Tag', secondary=tags, 
                        backref = backref('images', lazy='dynamic'))
    comments    =   relationship('Comment', backref='image', lazy='dynamic')

    def __repr__(self):
        str_created_at = self.created_at.strftime("%Y-%m-%d %H:%M:%S")
        return "<Image (uuid='%s', likes='%d', created_at=%s)>" % (self.uuid, self.likes, str_created_at)

class Tag(Base):

    __tablename__ = 'tags'

    id      =   Column(Integer, primary_key=True)
    name    =   Column(String(255), unique=True, nullable=False)

    def __repr__(self):
        return "<Tag (name='%s')>" % (self.name)

class Comment(Base):

    __tablename__ = 'comments'

    id          =   Column(Integer, primary_key=True)
    text        =   Column(String(2000))
    image_id    =   Column(Integer, ForeignKey('images.id'))

    def __repr__(self):
        return "<Comment (text='%s')>" % (self.text)


Data

#----------------------------
# Turn Foreign Key Constraints ON for
# each connection.
#----------------------------

from sqlalchemy.engine import Engine
from sqlalchemy import event

@event.listens_for(Engine, "connect")
def set_sqlite_pragma(dbapi_connection, connection_record):
    cursor = dbapi_connection.cursor()
    cursor.execute("PRAGMA foreign_keys=ON")
    cursor.close()

#----------------------------
# Create the engine
#----------------------------

from sqlalchemy import create_engine
engine = create_engine('sqlite:///:memory:', echo=True)

#----------------------------
# Create the Schema
#----------------------------

Base.metadata.create_all(engine)

#----------------------------
# Create the Session class 
#----------------------------

from sqlalchemy.orm import sessionmaker
Session = sessionmaker(bind=engine)

#----------------------------
# Populate the database 
#----------------------------

tag_cool = Tag(name='cool')
tag_car = Tag(name='car')
tag_animal = Tag(name='animal')

comment_rhino = Comment(text='Rhinoceros, often abbreviated as rhino, is a group of five extant species of odd-toed ungulates in the family Rhinocerotidae.')

image_car = Image(uuid='uuid_car', \
    tags=[tag_car, tag_cool], \
    created_at=(datetime.utcnow() - timedelta(days=1)))

image_another_car = Image(uuid='uuid_anothercar', \
    tags=[tag_car])

image_rhino = Image(uuid='uuid_rhino', \
    tags=[tag_animal], \
    comments=[comment_rhino])

# Create a new Session and add the images:
session = Session()

session.add(tag_cool)
session.add(tag_car)
session.add(tag_animal)

session.add(comment_rhino)

session.add(image_car)
session.add(image_another_car)
session.add(image_rhino)

# Commit the changes:
session.commit()

#----------------------------
# Update a Record
#----------------------------

image_to_update = session.query(Image).filter(Image.uuid == 'uuid_rhino').first()
image_to_update.likes = image_to_update.likes + 1
session.commit()

#----------------------------
# Query the database
#
# List of common filter: 
#
#   *http://docs.sqlalchemy.org/en/rel_0_9/orm/tutorial.html#common-filter-operators
#
#----------------------------

# Get a list of tags:
for name in session.query(Tag.name).order_by(Tag.name):
    print name

# How many tags do we have?
session.query(Tag).count()

# Get all images created yesterday:
session.query(Image) \
    .filter(Image.created_at < datetime.utcnow().date()) \
    .all()

# Get all images, that belong to the tag 'car' or 'animal', using a subselect:
session.query(Image) \
    .filter(Image.tags.any(Tag.name.in_(['car', 'animal']))) \
    .all()

# This can also be expressed with a join:
session.query(Image) \
    .join(Tag, Image.tags) \
    .filter(Tag.name.in_(['car', 'animal'])) \
    .all()

# Play around with functions:
from sqlalchemy.sql import func, desc

max_date = session.query(func.max(Image.created_at))
session.query(Image).filter(Image.created_at == max_date).first()

# Get a list of tags with the number of images:
q = session.query(Tag, func.count(Tag.name)) \
    .outerjoin(Image, Tag.images) \
    .group_by(Tag.name) \
    .order_by(desc(func.count(Tag.name))) \
    .all()

for tag, count in q:
    print 'Tag "%s" has %d images.' % (tag.name, count) 

# Get images created in the last two hours and zero likes so far:
session.query(Image) \
    .join(Tag, Image.tags) \
    .filter(Image.created_at > (datetime.utcnow() - timedelta(hours=2))) \
    .filter(Image.likes == 0) \
    .all()

    

###SQLAlchemy - Details 
#Creating table

import os
import sys
from sqlalchemy import Column, ForeignKey, Integer, String
from sqlalchemy.ext.declarative import declarative_base
from sqlalchemy.orm import relationship
from sqlalchemy import create_engine

 
#Derive from Base
Base = declarative_base()
 
class Person(Base):
	__tablename__ = 'person'
	# Here we define columns for the table person
	# Notice that each column is also a normal Python instance attribute.
	id = Column(Integer, primary_key=True)
	name = Column(String(250), nullable=False)
 


class Address(Base):
	__tablename__ = 'address'
	# Here we define columns for the table address.
	# Notice that each column is also a normal Python instance attribute.
	id = Column(Integer, primary_key=True)
	street_name = Column(String(250))
	street_number = Column(String(250))
	post_code = Column(String(250), nullable=False)
	#create relation ship
	person_id = Column(Integer, ForeignKey('person.id'))
	person = relationship(Person)   #this is not stored in DB, but SQLAlchemy keeps internally for association
	

 
# Create an engine that stores 
engine = create_engine('mysql://root:@127.0.0.1/python', echo=True)		# for sqllite :sqlite:///sqlalchemy_example.db
 
# Create all tables in the engine. This is equivalent to "Create Table"
# statements in raw SQL.

#You should put your model class before create_all() call
#If your models are declared in a separate module, import them 
#drop_all() method does the exact opposite of create_all() 
#you don't need to call below if no need to create in DB
Base.metadata.create_all(engine)


#.create_all(tables=tableList, checkfirst=True) 
#Create the tables described by this metadata. 
#If the checkfirst argument is True (the default value), no existing tables will be destroyed and recreated empty.
#To create only a specific set of tables, provide a list of the Table instances 


#insertions
from sqlalchemy import create_engine
from sqlalchemy.orm import sessionmaker
 
from sqlalchemy_declarative import Address, Base, Person
 
engine = create_engine('mysql://root:@127.0.0.1/python', echo=True)
# Bind the engine to the metadata of the Base class so that the
# declaratives can be accessed through a DBSession instance
Base.metadata.bind = engine
 
DBSession = sessionmaker(bind=engine)
# A DBSession() instance establishes all conversations with the database
# and represents a "staging zone" for all the objects loaded into the
# database session object. Any change made against the objects in the
# session won't be persisted into the database until you call
# session.commit(). If you're not happy about the changes, you can
# revert all of them back to the last commit by calling
# session.rollback()
session = DBSession()
 
# Insert a Person in the person table
new_person = Person(name='new person')
session.add(new_person)
session.commit()
 
# Insert an Address in the address table
new_address = Address(post_code='00000', person=new_person)
session.add(new_address)
session.commit()



#Query

from sqlalchemy_declarative import Person, Base, Address
from sqlalchemy import create_engine

engine = create_engine('mysql://root:@127.0.0.1/python')
Base.metadata.bind = engine

from sqlalchemy.orm import sessionmaker
DBSession = sessionmaker()
DBSession.bind = engine
session = DBSession()


# Make a query to find all Persons in the database
session.query(Person).all()  #List if Persons, can be iterated
for p in session.query(Person).all():
	print(p.name)
	
	
# Return the first Person from all Persons in the database
person = session.query(Person).first()
person.name #u'new person'

# Find all Address whose person field is pointing to the person object
#Address.person comes from association , not stored in DB
session.query(Address).filter(Address.person == person).all()  #[<sqlalchemy_declarative.Address object at 0x2ee3cd0>]

# Retrieve one Address whose person field is point to the person object
session.query(Address).filter(Address.person == person).one()  #<sqlalchemy_declarative.Address object at 0x2ee3cd0>
address = session.query(Address).filter(Address.person == person).one()
address.post_code #u'00000'

 
 
 

#Creation of Table without declarative class definition - individual creation

from sqlalchemy import create_engine, MetaData, Table
from sqlalchemy import Table, Column, Integer, String, MetaData, ForeignKey

#Scalar Defaults
Table("mytable", meta,
    Column("somecolumn", Integer, default=12)
)
#OR 
Table("mytable", meta,
    Column("somecolumn", Integer, onupdate=25)
)
#Or Python executed method
# a function which counts upwards
i = 0
def mydefault():
    global i
    i += 1
    return i

t = Table("mytable", meta,
    Column('id', Integer, primary_key=True, default=mydefault),
)
#or with onupdate
import datetime

t = Table("mytable", meta,
    Column('id', Integer, primary_key=True),
    # define 'last_updated' to be populated with datetime.now()
    Column('last_updated', DateTime, onupdate=datetime.datetime.now),
)

#Example 

from sqlalchemy import create_engine, MetaData, Table
from sqlalchemy import Table, Column, Integer, String, MetaData, ForeignKey

engine = create_engine('mysql://root:@127.0.0.1/books')
metadata = MetaData(bind=engine)



users = Table('users', metadata,
     Column('id', Integer, primary_key=True),
     Column('name', String),
     Column('fullname', String),
 )

addresses = Table('addresses', metadata,
	Column('id', Integer, primary_key=True),
	Column('user_id', None, ForeignKey('users.id')),
	Column('email_address', String, nullable=False)
)

users.create(engine)  #creates in DB if not existing
addresses.create(engine)  
#users.drop(engine) #drops in DB
#or 
metadata.create_all(engine)


#Iterating metadata to get all tables 
for t in metadata.sorted_tables:
	print t.name

#Then querying 
# access the column "ID":
users.columns.id

# or just
users.c.id

# via string
users.c['id']

# iterate through all columns
for c in users.c:
    print c

# get the table's primary key columns
for primary_key in addresses.primary_key:
    print primary_key

# get the table's foreign key objects:
for fkey in addresses.foreign_keys:
    print fkey

# access the table's MetaData:
addresses.metadata

# access the table's bound Engine or Connection, if its MetaData is bound:
addresses.bind

# access a column's name, type, nullable, primary key, foreign key
addresses.c.id.name
addresses.c.id.type
addresses.c.id.nullable
addresses.c.id.primary_key
addresses.c.user_id.foreign_keys

#insert

ins = users.insert()
str(ins)   #'INSERT INTO users (id, name, fullname) VALUES (:id, :name, :fullname)'

ins = users.insert().values(name='jack', fullname='Jack Jones')
str(ins)  #'INSERT INTO users (name, fullname) VALUES (:name, :fullname)'

#Executing
conn = engine.connect()
result = conn.execute(ins)
result.inserted_primary_key  # auto generated, 1

#Executing Multiple Statements
ins = users.insert()
conn.execute(ins, id=2, name='wendy', fullname='Wendy Williams')
conn.execute(addresses.insert(), [
 {'user_id': 1, 'email_address' : 'jack@yahoo.com'},
 {'user_id': 1, 'email_address' : 'jack@msn.com'},
 {'user_id': 2, 'email_address' : 'www@www.org'},
 {'user_id': 2, 'email_address' : 'wendy@aol.com'},
 ])
 
#update
 stmt = users.update().values(fullname="Fullname: " + users.c.name)
conn.execute(stmt)

#means 
UPDATE users SET fullname=(? || users.name)
('Fullname: ',)
COMMIT
#example
stmt = users.update().where(users.c.name == 'jack').values(name='ed')

conn.execute(stmt)
#multiple
stmt = users.update().where(users.c.name == bindparam('oldname')).values(name=bindparam('newname'))
conn.execute(stmt, [
	{'oldname':'jack', 'newname':'ed'},
	{'oldname':'wendy', 'newname':'mary'},
	{'oldname':'jim', 'newname':'jake'},
	])

#Correlated Updates
A correlated update lets you update a table using selection from another table, or the same table:


stmt = select([addresses.c.email_address]).\
	where(addresses.c.user_id == users.c.id).limit(1)
	
conn.execute(users.update().values(fullname=stmt))
#means 
UPDATE users SET fullname=(SELECT addresses.email_address
    FROM addresses
    WHERE addresses.user_id = users.id
    LIMIT ? OFFSET ?)
(1, 0)
COMMIT

#Deletes

result = conn.execute(addresses.delete())
#means 
DELETE FROM addresses
()
COMMIT

#Matched Row Counts
result.rowcount  #1


#with where clause e
conn.execute(users.delete().where(users.c.name > 'm'))
#means
DELETE FROM users WHERE users.name > ?
('m',)
COMMIT




#Selecting
from sqlalchemy.sql import select
s = select([users])  #SELECT users.id, users.name, users.fullname FROM users
result = conn.execute(s)
for row in result:
	print(row)

#or fetch one
result = conn.execute(s)
row = result.fetchone()
#or
for row in conn.execute(s):
	print("name:", row[users.c.name], "; fullname:", row[users.c.fullname])

#or Selecting few columns
s = select([users.c.name, users.c.fullname])
result = conn.execute(s)
for row in result:
	print(row)

#Using two tables join
for row in conn.execute(select([users, addresses])):
	print(row)

#Including where clause for join
s = select([users, addresses]).where(users.c.id == addresses.c.user_id)
for row in conn.execute(s):
	print(row)

#Where clause operators
print(users.c.id == addresses.c.user_id)  #users.id = addresses.user_id

print(users.c.id == 7)  #users.id = :id_1
 
print(users.c.id != 7)  #users.id != :id_1

# None converts to IS NULL
print(users.c.name == None) #users.name IS NULL

# reverse works too   
print('fred' > users.c.name) #users.name < :name_1

#addition  , number - add, for string, concate ie || 
print(users.c.id + addresses.c.id) #users.id + addresses.id
print(users.c.name + users.c.fullname)  #users.name || users.fullname


#for unknown operator use .op()
print(users.c.name.op('tiddlywinks')('foo'))   # users.name tiddlywinks :name_1
somecolumn.op('&')(0xff)  #bitwise operators

#And, or, not

from sqlalchemy.sql import and_, or_, not_
print(and_(
	users.c.name.like('j%'),
	users.c.id == addresses.c.user_id,
	or_(
		addresses.c.email_address == 'wendy@aol.com',
		addresses.c.email_address == 'jack@yahoo.com'
		),
		not_(users.c.id > 5)
		)
	)

#outputs 
users.name LIKE :name_1 AND users.id = addresses.user_id AND
(addresses.email_address = :email_address_1
   OR addresses.email_address = :email_address_2)
AND users.id <= :id_1

#with  bitwise AND, OR and NOT operators, 

print(users.c.name.like('j%') & (users.c.id == addresses.c.user_id) &
	(
		(addresses.c.email_address == 'wendy@aol.com') | \
		(addresses.c.email_address == 'jack@yahoo.com')
		) \
		& ~(users.c.id>5)
	)
#outputs 
users.name LIKE :name_1 AND users.id = addresses.user_id AND
(addresses.email_address = :email_address_1
    OR addresses.email_address = :email_address_2)
AND users.id <= :id_1


#example - select all users who have an email address at AOL or MSN, 
#whose name starts with a letter between “m” and “z”, a
#use between() and label(). 
#between() produces a BETWEEN clause, 
#label() is used in a column expression to produce labels using the AS keyword; 
it’s recommended when selecting from expressions that otherwise would not have a name:

 s = select([(users.c.fullname + ", " + addresses.c.email_address).label('title')]).\
	where(
	and_(
		users.c.id == addresses.c.user_id,
		users.c.name.between('m', 'z'),
		or_(
			addresses.c.email_address.like('%@aol.com'),
			addresses.c.email_address.like('%@msn.com')
			)
		)
	)
	
#means
SELECT users.fullname || ? || addresses.email_address AS title
FROM users, addresses
WHERE users.id = addresses.user_id AND users.name BETWEEN ? AND ? AND
(addresses.email_address LIKE ? OR addresses.email_address LIKE ?)
(', ', 'm', 'z', '%@aol.com', '%@msn.com')
[(u'Wendy Williams, wendy@aol.com',)]

#execute 
conn.execute(s).fetchall()


#or using multiple where() clauses


s = select([(users.c.fullname + ", " + addresses.c.email_address).label('title')]).\
	where(users.c.id == addresses.c.user_id).\
	where(users.c.name.between('m', 'z')).\
	where(
		or_(
			addresses.c.email_address.like('%@aol.com'),
			addresses.c.email_address.like('%@msn.com')
			)
	)
#means 
SELECT users.fullname || ? || addresses.email_address AS title
FROM users, addresses
WHERE users.id = addresses.user_id AND users.name BETWEEN ? AND ? AND
(addresses.email_address LIKE ? OR addresses.email_address LIKE ?)
(', ', 'm', 'z', '%@aol.com', '%@msn.com')
[(u'Wendy Williams, wendy@aol.com',)]


#Or Using Textual SQL

>>> from sqlalchemy.sql import text
s = text(
	"SELECT users.fullname || ', ' || addresses.email_address AS title "
	"FROM users, addresses "
	"WHERE users.id = addresses.user_id "
	"AND users.name BETWEEN :x AND :y "
	"AND (addresses.email_address LIKE :e1 "
	"OR addresses.email_address LIKE :e2)")
conn.execute(s, x='m', y='z', e1='%@aol.com', e2='%@msn.com').fetchall()

#Specifying Bound Parameter Behaviors

stmt = text("SELECT * FROM users WHERE users.name BETWEEN :x AND :y")
stmt = stmt.bindparams(x="m", y="z")

#or 
stmt = stmt.bindparams(bindparam("x", String), bindparam("y", String))
result = conn.execute(stmt, {"x": "m", "y": "z"})

#Specifying Result-Column Behaviors

stmt = text("SELECT id, name FROM users")
stmt = stmt.columns(users.c.id, users.c.name)

j = stmt.join(addresses, stmt.c.id == addresses.c.user_id)

new_stmt = select([stmt.c.id, addresses.c.id]). select_from(j).where(stmt.c.name == 'x')

#Example
stmt = text("SELECT users.id, addresses.id, users.id, "
	"users.name, addresses.email_address AS email "
	"FROM users JOIN addresses ON users.id=addresses.user_id "
	"WHERE users.id = 1").columns(
	users.c.id,
	addresses.c.id,
	addresses.c.user_id,
	users.c.name,
	addresses.c.email_address
	)
result = conn.execute(stmt)
row = result.fetchone()
row[addresses.c.email_address]  #'jack@yahoo.com'

#Ordering or Grouping by a Label
from sqlalchemy import func
from sqlalchemy import func, desc
>>> stmt = select([
	addresses.c.user_id,
	func.count(addresses.c.id).label('num_addresses')]).\
	order_by(desc("num_addresses"))  #or use asc()
#outputs 
SELECT addresses.user_id, count(addresses.id) AS num_addresses
FROM addresses ORDER BY num_addresses

conn.execute(stmt).fetchall()
#group by
stmt = select([users.c.name, func.count(addresses.c.id)]).\
	select_from(users.join(addresses)).group_by(users.c.name)
conn.execute(stmt).fetchall()
#means
SELECT users.name, count(addresses.id) AS count_1
FROM users JOIN addresses
    ON users.id = addresses.user_id
GROUP BY users.name

#Having clause
stmt = select([users.c.name, func.count(addresses.c.id)]).\
	select_from(users.join(addresses)).group_by(users.c.name).having(func.length(users.c.name) > 4)
conn.execute(stmt).fetchall()
#means 
SELECT users.name, count(addresses.id) AS count_1
FROM users JOIN addresses
    ON users.id = addresses.user_id
GROUP BY users.name
HAVING length(users.name) > ?

#distinct
stmt = select([users.c.name]).\
	where(addresses.c.email_address.contains(users.c.name)).distinct()

conn.execute(stmt).fetchall()
#means 
SELECT DISTINCT users.name
FROM users, addresses
WHERE (addresses.email_address LIKE '%%' || users.name || '%%')
()

#limit, offset

stmt = select([users.c.name, addresses.c.email_address]).\
	select_from(users.join(addresses)).limit(1).offset(1)

conn.execute(stmt).fetchall()

#means
SELECT users.name, addresses.email_address
FROM users JOIN addresses ON users.id = addresses.user_id
 LIMIT ? OFFSET ?
(1, 1)


#using alias 

u1a, u1b = users.alias(), users.alias()
tmt = select([u1a, u1b]).\
	where(u1a.c.name > u1b.c.name).\
	order_by(u1a.c.name)  # using "name" here would be ambiguous

conn.execute(stmt).fetchall()

#Using Joins
print(users.join(addresses)) #users JOIN addresses ON users.id = addresses.user_id

s = select([users.c.fullname]).select_from(users.outerjoin(addresses))
print(s)
SELECT users.fullname
    FROM users
    LEFT OUTER JOIN addresses ON users.id = addresses.user_id
#Example
s = select([users.c.fullname]).select_from(
		users.join(addresses,
		addresses.c.email_address.like(users.c.name + '%'))
	)
	
#outputs 
SELECT users.fullname
FROM users JOIN addresses ON addresses.email_address LIKE (users.name || ?)
('%',)


conn.execute(s).fetchall()

#Functions
SQL functions are created using the func keyword, which generates functions using attribute access:

from sqlalchemy.sql import func
print(func.now()) #now()

print(func.concat('x', 'y')) #concat(:concat_1, :concat_2)

#By “generates”, we mean that any SQL function is created based on the word you choose:
print(func.xyz_my_goofy_function()) #xyz_my_goofy_function()

#Example:
conn.execute(
	select([
		func.max(addresses.c.email_address, type_=String).label('maxemail')
		])
	).scalar()

#Unions and Other Set Operations
from sqlalchemy.sql import union
u = union(
	addresses.select().where(addresses.c.email_address == 'foo@bar.com'),
	addresses.select().where(addresses.c.email_address.like('%@yahoo.com')),
	).order_by(addresses.c.email_address)

#outputs 
SELECT addresses.id, addresses.user_id, addresses.email_address
FROM addresses
WHERE addresses.email_address = ?
UNION
SELECT addresses.id, addresses.user_id, addresses.email_address
FROM addresses
WHERE addresses.email_address LIKE ? ORDER BY addresses.email_address
('foo@bar.com', '%@yahoo.com')

conn.execute(u).fetchall()

#Scalar Selects -A scalar select is a SELECT that returns exactly one row and one column

stmt = select([func.count(addresses.c.id)]).\
	where(users.c.id == addresses.c.user_id).as_scalar()
	
#outputs 
SELECT users.name, (SELECT count(addresses.id) AS count_1
FROM addresses
WHERE users.id = addresses.user_id) AS anon_1
FROM users


conn.execute(select([users.c.name, stmt])).fetchall()

#Correlated Subqueries

stmt = select([addresses.c.user_id]).\
		where(addresses.c.user_id == users.c.id).where(addresses.c.email_address == 'jack@yahoo.com')

enclosing_stmt = select([users.c.name]).where(users.c.id == stmt)
conn.execute(enclosing_stmt).fetchall()

#means 
SELECT users.name
FROM users
WHERE users.id = (SELECT addresses.user_id
    FROM addresses
    WHERE addresses.user_id = users.id
    AND addresses.email_address = ?)
('jack@yahoo.com',)







#Complex example using declarative

#A department has many employees while an employee belongs to at most one department.
#A department has many employees - association is created in Employee!!
#creations

from sqlalchemy import Column, String, Integer, ForeignKey
from sqlalchemy.orm import relationship, backref
from sqlalchemy.ext.declarative import declarative_base


Base = declarative_base()


class Department(Base):
	__tablename__ = 'department'
	id = Column(Integer, primary_key=True)
	name = Column(String)


class Employee(Base):
	__tablename__ = 'employee'
	id = Column(Integer, primary_key=True)
	name = Column(String)
	department_id = Column(Integer, ForeignKey('department.id'))
	department = relationship(Department, backref=backref('employees', uselist=True))
	
	
#The column 'employee.department_id' is a foreign key to the column 'department.id' 
#the relationship 'department.employees' include all the employees in that department(uselist = True)


from sqlalchemy import create_engine
engine = create_engine('mysql://root:@127.0.0.1/python')

#create all metadata, tables would be created only once 
Base.metadata.create_all(engine)

#create session
from sqlalchemy.orm import sessionmaker
Base.metadata.bind = engine
session = sessionmaker(bind=engine)





#Insertions

john = Employee(name='john')
it_department = Department(name='IT')
john.department = it_department     #adding at only Employee site. Only one department, but department.employees would iterate all employees

s = session()
s.add(john)
s.add(it_department)
s.commit()

#Query

it = s.query(Department).filter(Department.name == 'IT').one()
it.employees  #[...]
it.employees[0].name #u'john'

#Using select 
from sqlalchemy import select
find_it = select([Department.id]).where(Department.name == 'IT')
rs = s.execute(find_it)
rs
 
rs.fetchone()  #(1,)
rs.fetchone()  # Only one result is returned from the query, so getting one more returns None.
rs.fetchone()  # Since the previous fetchone() returned None, fetching more would lead to a result-closed exception


find_john = select([Employee.id]).where(Employee.department_id == 1)
rs = s.execute(find_john)
 
rs.fetchone()  # Employee John's ID (1,)
rs.fetchone()


#Many-to-many between the departments and the employees
#create association in both class using third class

#creations
from sqlalchemy import Column, String, Integer, ForeignKey
from sqlalchemy.orm import relationship, backref
from sqlalchemy.ext.declarative import declarative_base


Base = declarative_base()


class Department(Base):
	__tablename__ = 'department'
	id = Column(Integer, primary_key=True)
	name = Column(String)
	employees = relationship('Employee', secondary='department_employee')
	
	
class Employee(Base):
	__tablename__ = 'employee'
	id = Column(Integer, primary_key=True)
	name = Column(String)
	departments = relationship('Department', secondary='department_employee')


class DepartmentEmployee(Base):
	__tablename__ = 'department_employee'
	department_id = Column(Integer, ForeignKey('department.id'), primary_key=True)
	employee_id = Column(Integer, ForeignKey('employee.id'), primary_key=True)

	
from sqlalchemy import create_engine
engine = create_engine('mysql://root:@127.0.0.1/python')

from sqlalchemy.orm import sessionmaker
session = sessionmaker()
session.configure(bind=engine)
Base.metadata.create_all(engine)


#Insert
s = session()
john = Employee(name='john')
s.add(john)

it_department = Department(name='IT')
it_department.employees.append(john)   #adding at only Department site 

s.add(it_department)
s.commit()

 

#Query


john = s.query(Employee).filter(Employee.name == 'john').one()
john.departments  #[..]
john.departments[0].name #u'IT'

it = s.query(Department).filter(Department.name == 'IT').one()
it.employees #[]
it.employees[0].name #u'john' 

 
#New insert 
marry = Employee(name='marry')
financial_department = Department(name='financial')
financial_department.employees.append(marry)
s.add(marry)
s.add(financial_department)
s.commit()

 


#To find all the employees in the IT department,

s.query(Employee).filter(Employee.departments.any(Department.name == 'IT')).one().name  #u'john'

 

#or Using the Expression Language
find_employees = select([DepartmentEmployee.employee_id]).select_from(Department.__table__.join(DepartmentEmployee)).where(Department.name == 'IT')
rs = s.execute(find_employees)
rs.fetchone() #(1,)
rs.fetchone() #
 


#assign employee marry into the IT department so that she will be part of two departments

s.refresh(marry)
s.refresh(it)
it.employees #[..]

it.employees.append(marry) 
s.commit()
it.employees #[.. ]

 


#to find  all the employees who belong to at least two departments, we use group_by 

from sqlalchemy import func
s.query(Employee).join(Employee.departments).group_by(Employee.id).having(func.count(Department.id) > 1).one().name

 
#using expression
find_marry = select([Employee.id]).select_from(Employee.__table__.join(DepartmentEmployee)).group_by(Employee.id).having(func.count(DepartmentEmployee.department_id) > 1)

rs = s.execute(find_marry)
rs.fetchall() #[(2,)]

 
#and then close the db
s.close()


#SQLAlchemy Schema Reflection / Introspection - binding to exiting table 

#Option-1 : Without using mapper

#suppose you have a table 

c.execute('''
          CREATE TABLE authors
          (id, name, email)
          ''')
c.execute("INSERT INTO authors VALUES (2, 'john', 'john@example.com')")
c.close()

#Get the table using below
from sqlalchemy import create_engine, MetaData, Table


engine = create_engine('mysql://root:@127.0.0.1/books')
metadata = MetaData(bind=engine)

#from DB
authors = Table("authors", metadata, autoload=True, autoload_with=engine)
authors   #Table('person', MetaData(bind=), Column(u'name', TEXT(), table=), Column(u'email', TEXT(), table=), schema=None)
[c.name for c in authors.columns]  #['id', 'name', 'email']


#insert
ins = authors.insert().values(id = 100, name='jack', email='Jack@Jones')

#Executing
conn = engine.connect()
result = conn.execute(ins)


#Executing Multiple Statements
conn.execute(authors.insert(), [
 {'id': 101, 'name' : " das1", 'email' : 'jack@yahoo.com'},
 {'id': 102, 'name' : " das1", 'email' : 'jack@yahoo.com'},
 {'id': 103, 'name' : " das1", 'email' : 'jack@yahoo.com'},
 {'id': 104,'name' : " das1", 'email' : 'jack@yahoo.com'},
 ])
 

#update
stmt = authors.update().where(authors.c.name == 'jack').values(name='jacks')
conn.execute(stmt)

#Selecting
from sqlalchemy.sql import select

s = select([authors])  #SELECT users.id, users.name, users.fullname FROM users
result = conn.execute(s)
for row in result:
	print(row)

#or fetch one
result = conn.execute(s)
row = result.fetchone()
#or
for row in conn.execute(s):
	print("name:", row[authors.c.name], "; fullname:", row[authors.c.email])

#or Selecting few columns
s = select([authors.c.name, authors.c.email])
result = conn.execute(s)
for row in result:
	print(row)

#Using two tables join
for row in conn.execute(select([authors]).where(authors.c.id > 100)):
	print(row)

	
	
#reflect all tables in the database using the MetaData.reflect method.
meta = MetaData()
meta.reflect(bind=engine)
person = meta.tables['person']
person  #Table(u'person', MetaData(bind=None), Column(u'name', TEXT 





#Option-2 - Using mapper 

from sqlalchemy import *
from sqlalchemy.orm import mapper, sessionmaker
 
#in SQL
class Authors(object):
	def __init__(self, id, name, email):
		self.id = id
		self.name = name
		self.email = email
	


def loadSession(url = 'mysql://root:@127.0.0.1/books' ):   
    engine = create_engine(url, echo=True) 
    metadata = MetaData(engine)
    authors = Table('authors', metadata, Column("id", Integer, primary_key=True), autoload=True)
    mapper(Authors, authors) 
    Session = sessionmaker(bind=engine)
    session = Session()
    return session
 

session = loadSession()
res = session.query(Authors).all()
res
res[1].name
session.close()

#insert - Note id existing in table, hence no auto increment and you must provide this 
#also underlying table does not have any constraint on id, hence same row can be created multiple times
n = Authors(id = 25, name="Das", email = "das@das.com")
session.add(n)
session.commit()
session.refresh(n)


Note:
1. if errors comes as 'already has a primary mapper define', exit python interpretor and do it again
2. For mapper to work, primary key is must, 
If any existing column, col_name  is primary key, add 'Column("col_name", Integer, primary_key=True)' in Table
If there is no primary key, add 'Column("id", Integer, primary_key=True)', SQLAlchemy would id col and would increment
If exiting multiple cols, col1, col2 are primary keys , add 'Column("col1", Integer, primary_key=True), Column("col2", Integer, primary_key=True)'
3. You can add any other constraint as well
eg Column("foo_id", Integer, ForeignKey('foo.id'))  where 'foo' is another table in DB

#Option-3 : Using declarative way 

from sqlalchemy import *
from sqlalchemy.ext.declarative import declarative_base
from sqlalchemy.orm import sessionmaker
 
engine = create_engine('mysql://root:@127.0.0.1/books', echo=True)
Base = declarative_base(engine)

#Can add any constraint, existing key would not be recreated 
#might provide a  __init__ function, but declarative would provide one
#Note declarive provide a class var, but no issue accessing with self 
#could provide __str__


class Authors(Base):
	__tablename__ = 'authors'
	__table_args__ = {'autoload':True}
	id = Column("id", Integer, primary_key=True)
	def __str__(self):
		return "Authors(id=%d,name=%s,email=%s)" % (self.id, self.name, self.email)
	def __repr__(self): return self.__str__()
 

def loadSession():    
    metadata = Base.metadata
    Session = sessionmaker(bind=engine)
    session = Session()
    return session
 

session = loadSession()
res = session.query(Authors).all()
res[1].name

session.query(Authors).filter(Authors.id == None).all() #check NULL, but no output here because id is taken as PK

session.query(Authors).filter(Authors.id > 100 ).all()



#Now SQlAlchemy is more strict on pk, it fails on below as already existing
n = Authors(id = 25, name="Das", email = "das@das.com")
session.add(n)
session.commit()

#must give new id which is not existing
n = Authors(id = 107, name="Das", email = "das@das.com")
session.add(n)
session.commit()


#Updating
session.query(Authors).filter_by(id=107).update({"name": "NDas"})
session.commit()
#OR
u = session.query(Authors).get(107)
u.name = "New Das"
session.commit()
#mass updating
session.query(Authors).filter(Authors.id > 101 ).update({"name": "NDas"})
session.commit()


#deleting
n = session.query(Authors).get(107)
session.delete(n)
session.commit()
#or
session.query(Authors).filter_by(id=104).delete()
session.commit()
#mass deleting
session.query(Authors).filter(Authors.id > 101 ).delete()
session.commit()


#Note any constraint can be added in class, existing key would not be recreated 

class Foo(Base):
	__tablename__ = 'foo'
	__table_args__ = {'autoload':True}
	bars = relationship("Bar")

class Bar(Base):
	__tablename__ = 'bar'
	__table_args__ = {'autoload':True}
	# "bar.foo_id" to have a foreign key constraint otherwise not existing in DB	
	foo_id = Column(Integer, ForeignKey('foo.id'))

	
#Note all the existing table columns can be added with any constraints
#in above forms

moz_bookmarks = Table('moz_bookmarks', metadata, 
                          Column('id', Integer, primary_key=True),
                          Column('type', Integer),
                          Column('fk', Integer),
                          Column('parent', Integer),
                          Column('position', Integer),
                          Column('title', String),
                          Column('keyword_id', Integer),
                          Column('folder_type', Text),
                          Column('dateAdded', Integer),
                          Column('lastModified', Integer)
                          )

		

#Important Methods on Session

add(instance, _warn=True)

add_all(instances)

expunge(instance)
	Remove the instance from this Session.
	
delete(instance)

execute(clause, params=None, mapper=None, bind=None, **kw)
	Execute a SQL expression construct or string statement within the current transaction.

result = session.execute(
            user_table.select().where(user_table.c.id == 5)
        )

accepts any executable clause construct, such as select(), insert(), update(), delete(), and text(). 
Plain SQL strings can be passed as well, 
result = session.execute(
            "SELECT * FROM user WHERE id=:param",
            {"param":5}
        )
is equivalent to:
from sqlalchemy import text
result = session.execute(
            text("SELECT * FROM user WHERE id=:param"),
            {"param":5}
        )
For inserting
result = session.execute( users.insert(), {"id": 7, "name": "somename"})
or for multiple rows:
result = session.execute(users.insert(), [
                        {"id": 7, "name": "somename7"},
                        {"id": 8, "name": "somename8"},
                        {"id": 9, "name": "somename9"}
                    ])



is_modified(instance, include_collections=True, passive=True)
Return True if the given instance has locally modified attributes.

query(*entities, **kwargs)
Return a new Query object corresponding to this Session

refresh(instance, attribute_names=None, lockmode=None)
Expire and refresh the attributes on the given instance

scalar(clause, params=None, mapper=None, bind=None, **kw)
Like execute() but return a scalar result.

#Important methods of Query -class sqlalchemy.orm.query.Query(entities, session=None)
Can have multiple entities

add_columns(*column)
Add one or more column expressions to the list of result columns to be returned.

add_entity(entity, alias=None)
add a mapped entity to the list of result columns to be returned.

all()
Return the results represented by this Query as a list.


as_scalar()
Return the full SELECT statement represented by this Query, converted to a scalar subquery

column_descriptions
Return metadata about the columns which would be returned by this Query
user_alias = aliased(User, name='user2')
q = sess.query(User, User.id, user_alias)
q.column_descriptions

count()
Return a count of rows this Query would return
For fine grained control over specific columns to count, use below
from sqlalchemy import func
# count User records, without
# using a subquery.
session.query(func.count(User.id))
# return count of user "id" grouped
# by "name"
session.query(func.count(User.id)). group_by(User.name)

from sqlalchemy import distinct
# count distinct "name" values
session.query(func.count(distinct(User.name)))


delete(synchronize_session='evaluate')  
Perform a bulk delete query. False means don't synchronize with session
Deletes rows matched by this query from the database.

sess.query(User).filter(User.age == 25).delete(synchronize_session=False)


distinct(*criterion)
Apply a DISTINCT to the query and return the newly resulting Query

exists()
EXISTS subquery of the form EXISTS (SELECT 1 FROM ... WHERE ...).
q = session.query(User).filter(User.name == 'fred')
session.query(q.exists())

Producing SQL similar to:
SELECT EXISTS (
    SELECT 1 FROM users WHERE users.name = :name_1
) AS anon_1

The EXISTS construct is usually used in the WHERE clause:
session.query(User.id).filter(q.exists()).scalar()


filter(*criterion)
apply the given filtering criterion to a copy of this Query, using SQL expressions.

filter_by(**kwargs)
apply the given filtering criterion to a copy of this Query, using keyword expressions.
session.query(MyClass).filter_by(name = 'some name')
session.query(MyClass).filter_by(name = 'some name', id = 5)

first()
Return the first result of this Query or None if the result doesn’t contain any row.

get(ident)
Return an instance based on the given primary key identifier, or None if not found.
my_user = session.query(User).get(5)
some_object = session.query(VersionedFoo).get((5, 10))

group_by(*criterion)
apply one or more GROUP BY criterion to the query and return the newly resulting Query

having(criterion)
apply a HAVING criterion to the query and return the newly resulting Query.
HAVING criterion makes it possible to use filters on aggregate functions like COUNT, SUM, AVG, MAX, and MIN, 

q = session.query(User.id).\
            join(User.addresses).\
            group_by(User.id).\
            having(func.count(Address.id) > 2)
			
			

join(*props, **kwargs)
Simple join
q = session.query(User).join(User.addresses)
#SELECT user.* FROM user JOIN address ON user.id = address.user_id
#or
q = session.query(User).join("addresses")
#multiple 
q = session.query(User).join("orders", "items", "keywords")
#or
q = session.query(User).\
        join(User.orders).\
        join(Order.items).\
        join(Item.keywords)

#Joins to a Target Entity or Selectable
join() will attempt to create a JOIN along the natural foreign key relationship between two entities else error
q = session.query(User).join(Address)
q = session.query(User).join(addresses_table) #using address table

#Joins to a Target with an ON Clause explicitely

a_alias = aliased(Address)

q = session.query(User).\
        join(User.addresses).\
        join(a_alias, User.addresses).\
        filter(Address.email_address=='ed@foo.com').\
        filter(a_alias.email_address=='ed@bar.com')

#means 
SELECT user.* FROM user
    JOIN address ON user.id = address.user_id
    JOIN address AS address_1 ON user.id=address_1.user_id
    WHERE address.email_address = :email_address_1
    AND address_1.email_address = :email_address_2
	
	
limit(limit)
Apply a LIMIT to the query and return the newly resulting Query.

offset(offset)
Apply an OFFSET to the query and return the newly resulting Query.

one()
Return exactly one result or raise an exception.

one_or_none()
Return at most one result or raise an exception.

order_by(*criterion)
apply one or more ORDER BY criterion to the query and return the newly resulting Query

scalar()
Return the first element of the first result or None if no rows present. 
If multiple rows are returned, raises MultipleResultsFound.

session.query(Item).scalar()
session.query(Item.id).scalar()
session.query(Item.id).filter(Item.id < 0).scalar()
session.query(Item.id, Item.name).scalar()
session.query(func.count(Parent.id)).scalar()


select_from(*from_obj)
Set the FROM clause of this Query explicitly.
q = session.query(Address).select_from(User).\
    join(User.addresses).\
    filter(User.name == 'ed')

Which produces SQL equivalent to:
SELECT address.* FROM user
JOIN address ON user.id=address.user_id
WHERE user.name = :name_1


slice(start, stop)
Computes the “slice” of the Query represented by the given indices and returns the resulting Query.
session.query(User).order_by(User.id).slice(1, 3)
renders as
SELECT users.id AS users_id,
       users.name AS users_name
FROM users ORDER BY users.id
LIMIT ? OFFSET ?
(2, 1)


update(values, synchronize_session='evaluate', update_args=None)
Perform a bulk update query.
sess.query(User).filter(User.age == 25).update({User.age: User.age - 10}, synchronize_session=False)


value(column)
Return a scalar result corresponding to the given column expression.

values(*columns)
Return an iterator yielding result tuples corresponding to the given list of columns

		
#########################
#Web Testing by using selenium
#pip3 install -U selenium


from selenium import webdriver
from selenium.common.exceptions import NoSuchElementException
from selenium.webdriver.common.keys import Keys
import time

browser = webdriver.Firefox()

browser.get("http://www.yahoo.com") # Load page

assert "Yahoo" in browser.title

elem = browser.find_element_by_name("p") # Find the query box
elem.send_keys("seleniumhq" + Keys.RETURN)

time.sleep(0.2) # Let the page load, or use explicit or Implicit wait

try:
    elem = browser.find_element_by_xpath("//a[text()='Selenium - Web Browser Automation']")
    print(elem.get_attribute("href"))
    for i in browser.find_elements_by_xpath("//a"):
        print(i.get_attribute("href"))
except NoSuchElementException:
    assert 0, "can't find seleniumhq"
browser.close()


#Webelement - important methods other than searching
clear() 
Clears the text if it’s a text entry element.

click() 
Clicks the element.

get_attribute(name) 
Gets the given attribute or property of the element.

is_displayed() 
Whether the element is visible to a user.

is_enabled() 
Returns whether the element is enabled.

is_selected() 
Returns whether the element is selected.

send_keys(*value) 
Use this to send simple key events or to fill out form fields:

form_textfield = driver.find_element_by_name('username')
form_textfield.send_keys("admin")

#or with file input
file_input = driver.find_element_by_name('profilePic')
file_input.send_keys("path/to/profilepic.gif")
# Generally it's better to wrap the file path in one of the methods
# in os.path to return the actual path to support cross OS testing.
# file_input.send_keys(os.path.abspath("path/to/profilepic.gif"))

submit() 
Submits a form.

value_of_css_property(property_name) 
The value of a CSS property.

id 
Internal ID used by selenium.

size 
The size of the element.

tag_name 
This element’s tagName property.

text 
The text of the element.

#UI Support for Select 
class selenium.webdriver.support.select.Select(webelement)
deselect_all() 
deselect_by_index(index) 
deselect_by_value(value) 
select_by_index(index) 
select_by_value(value) 
all_selected_options 
Returns a list of all selected options belonging to this select tag

first_selected_option 
The first selected option in this select tag (or the currently selected option in a normal select)

options 
Returns a list of all options belonging to this select tag

#Filling in forms


element = driver.find_element_by_xpath("//select[@name='name']")
all_options = element.find_elements_by_tag_name("option")
for option in all_options:
    print("Value is: %s" % option.get_attribute("value"))
    option.click()


#or using these methods 

from selenium.webdriver.support.ui import Select
select = Select(driver.find_element_by_name('name'))
select.select_by_index(index)
select.select_by_visible_text("text")
select.select_by_value(value)


# for deselecting all the selected options:
select = Select(driver.find_element_by_id('id'))
select.deselect_all()


# returns a list:

select = Select(driver.find_element_by_xpath("xpath"))
all_selected_options = select.all_selected_options


#To get all available options:

options = select.options
# Assume the button has the ID "submit" :)
driver.find_element_by_id("submit").click()

#or  WebDriver has the convenience method “submit” on every element. 


#Note for check-box etc, standard selenium code can be written
#selecting a checkbox with text 'One CheckBox Value'
browser.find_element_by_xpath(
    ".//*[contains(text(), 'One CheckBox Value')]"
).click()

#for text input element, get value by 'value'
<input type="text" name="inputbox" value="name" class="box">
input = driver.find_element_by_name("inputbox").get_attribute('value').encode('utf-8')

#OR setting 
inputElement.send_keys('admin ')
inputElement.send_keys(Keys.ENTER)
inputElement.submit() 



#xpath Searching

<html> <body>  <form id="loginForm">
   <input name="username" type="text" />
   <input name="password" type="password" />
   <input name="continue" type="submit" value="Login" />
   <input name="continue" type="button" value="Clear" />
  </form></body><html>

#The form elements can be located like this:
login_form = driver.find_element_by_xpath("/html/body/form[1]")
login_form = driver.find_element_by_xpath("//form[1]")
login_form = driver.find_element_by_xpath("//form[@id='loginForm']")

#The username element can be located like this:
username = driver.find_element_by_xpath("//form[input/@name='username']")
username = driver.find_element_by_xpath("//form[@id='loginForm']/input[1]")
username = driver.find_element_by_xpath("//input[@name='username']")

#The “Clear” button element can be located like this:
clear_button = driver.find_element_by_xpath("//input[@name='continue'][@type='button']")
clear_button = driver.find_element_by_xpath("//form[@id='loginForm']/input[4]")

#Note
If you need to write an XPath like //div[text()="hello world"] 
but the HTML of the link is really "hello&nbsp;world", 
you'll need to insert a real "&nbsp;" , use  like this: //div[text()="hello${nbsp}world"].

#Locating element via various methods 

find_element_by_id
find_element_by_name
find_element_by_xpath
find_element_by_link_text
find_element_by_partial_link_text
find_element_by_tag_name
find_element_by_class_name
find_element_by_css_selector

#To find multiple elements (these methods will return a list):
find_elements_by_name
find_elements_by_xpath
find_elements_by_link_text
find_elements_by_partial_link_text
find_elements_by_tag_name
find_elements_by_class_name
find_elements_by_css_selector

#Example: Locating Hyperlinks by Link Text

<html>
 <body>
  <p>Are you sure you want to do this?</p>
  <a href="continue.html">Continue</a>
  <a href="cancel.html">Cancel</a>
</body>
<html>

#code 
continue_link = driver.find_element_by_link_text('Continue')
continue_link = driver.find_element_by_partial_link_text('Conti')


#Example: Locating Elements by Tag Name

<html>
 <body>
  <h1>Welcome</h1>
  <p>Site content goes here.</p>
</body>
<html>


#code 
heading1 = driver.find_element_by_tag_name('h1')

#Locating Elements by Class Name

<html>
 <body>
  <p class="content">Site content goes here.</p>
</body>
<html>


content = driver.find_element_by_class_name('content')

#Locating Elements by CSS Selectors

<html>
 <body>
  <p class="content">Site content goes here.</p>
</body>
<html>

content = driver.find_element_by_css_selector('p.content')

#Drag and drop
either moving an element by a certain amount, or on to another element:


element = driver.find_element_by_name("source")
target = driver.find_element_by_name("target")

from selenium.webdriver import ActionChains
action_chains = ActionChains(driver)
action_chains.drag_and_drop(element, target).perform()



#Implicit Waits (default=0) for loading the page and waiting fto locate a element

from selenium import webdriver
driver = webdriver.Firefox()
driver.implicitly_wait(10) # seconds
driver.get("http://somedomain/url_that_delays_loading")
myDynamicElement = driver.find_element_by_id("myDynamicElement")

#Explicit Wait- using expected_conditions

from selenium import webdriver
from selenium.webdriver.common.by import By
from selenium.webdriver.support.ui import WebDriverWait
from selenium.webdriver.support import expected_conditions as EC

driver = webdriver.Firefox()
driver.get("http://somedomain/url_that_delays_loading")
try:
    element = WebDriverWait(driver, 10).until(
        EC.presence_of_element_located((By.ID, "myDynamicElement"))
    )
finally:
    driver.quit()

#Example
from selenium.webdriver.support import expected_conditions as EC

wait = WebDriverWait(driver, 10)
element = wait.until(EC.element_to_be_clickable((By.ID,'someid')))

#expected_conditions methods
•title_is
•title_contains
•presence_of_element_located
•visibility_of_element_located
•visibility_of
•presence_of_all_elements_located
•text_to_be_present_in_element
•text_to_be_present_in_element_value
•frame_to_be_available_and_switch_to_it
•invisibility_of_element_located
•element_to_be_clickable - it is Displayed and Enabled.
•staleness_of
•element_to_be_selected
•element_located_to_be_selected
•element_selection_state_to_be
•element_located_selection_state_to_be
•alert_is_present


#Cookies
# Go to the correct domain
driver.get("http://www.example.com")
# Now set the cookie. This one's valid for the entire domain
cookie = {'name' : 'foo', 'value' : 'bar'}
driver.add_cookie(cookie)
# And now output all the available cookies for the current URL
driver.get_cookies()

#To move backwards and forwards in your browser’s history:
driver.forward()
driver.back()

#Moving between windows and frames
driver.switch_to_window("windowName")

Get windowname from the javascript or link that opened it, eg
<a href="somewhere.html" target="windowName">Click here to open a new window</a>

You can also swing from frame to frame (or into iframes):
driver.switch_to_frame("frameName")

#Popup dialogs
interface works equally well on alerts, confirms, prompts. 
alert = driver.switch_to_alert()
return the currently open alert object. 

#Accepting / Dismissing alert prompts:
Alert(driver).accept()
Alert(driver).dismiss()


#Inputting a value into an alert prompt:
name_prompt = Alert(driver) 
name_prompt.send_keys('Willian Shakesphere') 
name_prompt.accept()

#Reading a the text of a prompt for verification:


alert_text = Alert(driver).text 
self.assertEqual('Do you wish to quit?', alert_text)

#Authenticating a alert by authenticate(username, password) 
Implicitly ‘clicks ok’

driver.switch_to_alert().authenticate('cheese', 'secretGouda')


#Special Keys - eg  Keys.ADD
ADD= u'\ue025' ALT= u'\ue00a' ARROW_DOWN= u'\ue015' ARROW_LEFT= u'\ue012' ARROW_RIGHT= u'\ue014' ARROW_UP= u'\ue013'
BACKSPACE= u'\ue003' BACK_SPACE= u'\ue003' CANCEL= u'\ue001' CLEAR= u'\ue005' COMMAND= u'\ue03d' CONTROL= u'\ue009'
DECIMAL= u'\ue028' DELETE= u'\ue017' DIVIDE= u'\ue029' DOWN= u'\ue015' END= u'\ue010' ENTER= u'\ue007' EQUALS= u'\ue019'
ESCAPE= u'\ue00c' F1= u'\ue031' F10= u'\ue03a' F11= u'\ue03b' F12= u'\ue03c' F2= u'\ue032' F3= u'\ue033'
F4= u'\ue034' F5= u'\ue035' F6= u'\ue036' F7= u'\ue037' F8= u'\ue038' F9= u'\ue039' HELP= u'\ue002'
HOME= u'\ue011' INSERT= u'\ue016' LEFT= u'\ue012' LEFT_ALT= u'\ue00a' LEFT_CONTROL= u'\ue009' LEFT_SHIFT= u'\ue008' META= u'\ue03d' 


#Locate elements by selenium.webdriver.common.by.By,  By.xxx, xxx can be 
CLASS_NAME= 'class name' 
CSS_SELECTOR= 'css selector' 
ID= 'id' 
LINK_TEXT= 'link text' 
NAME= 'name' 
PARTIAL_LINK_TEXT= 'partial link text' 
TAG_NAME= 'tag name'
XPATH= 'xpath'
#usage 
from selenium.webdriver.common.by import By
driver.find_element(By.XPATH, '//button[text()="Some text"]')
driver.find_elements(By.XPATH, '//button')






###XPath Reference
starts with the slash / ,  represents an absolute path to the required element. 
/AAA
/AAA/CCC
/AAA/DDD/BBB

starts with // , all elements in the document which fulfill following criteria are selected.
//BBB
//DDD/BBB

The star * selects all elements located by preceeding path
/AAA/CCC/DDD/*
/*/*/*/BBB
//*

The .. moves backward the element just like file system
/AAA/BBB/CCC/../../B1B1  (B1B1 is sibling of BBB)



#Expression in square brackets can further specify an element.

The function last() selects the last element

function text() selects text child element. It can be used for matching and selecting
/AAA/BBB[1]      -> 1 selects first element
/AAA/BBB[last()]
/AAA/BBB[text() = 'Tom']

Attributes are accessed by @ prefix.
//@id
//BBB[@id]
//BBB[@name = 'Tom']   -> BBB with name is Tom
//BBB[@*]      		-> all attributes

Function normalize-space removes leading and trailing spaces and replaces sequences of whitespace characters by a single space. 
//BBB[normalize-space(@name)='bbb']

Function count() counts the number of selected elements 
//*[count(BBB)=2]  -> returns all Nodes with BBB is 2
//*[count(*)=3]    -> returns all nodes with nodes is 3


Function name() returns name of the element,
//*[name()='BBB']
//*[starts-with(name(),'B')]
//*[contains(name(),'C')]

The string-length function returns the number of characters in the string. 
//*[string-length(name()) = 3]
//*[string-length(name()) < 3]
//*[string-length(name()) > 3]

 Several paths can be combined with | separator.
//CCC | //BBB
/AAA/EEE | //BBB
/AAA/EEE | //DDD/CCC | /AAA | //BBB

The child axis contains the children of the context node. 
The child axis is the default axis and it can be omitted. 
/AAA same as /child::AAA
/AAA/BBB same as /child::AAA/child::BBB same as /child::AAA/BBB

The descendant axis contains the descendants of the context node; 
a descendant is a child or a child of a child and so on; thus the descendant axis never contains attribute or namespace nodes 
/descendant::*
/AAA/BBB/descendant::*
//CCC/descendant::*
//CCC/descendant::DDD

The parent axis contains the parent of the context node, if there is one.
//DDD/parent::*

The ancestor axis contains the ancestors of the context node; 
the ancestors of the context node consist of the parent of context node and the parent's parent and so on; 
thus, the ancestor axis will always include the root node, unless the context node is the root node. 
/AAA/BBB/DDD/CCC/EEE/ancestor::*
//FFF/ancestor::*


The following-sibling axis contains all the following siblings of the context node.
/AAA/BBB/following-sibling::*
//CCC/following-sibling::*

The preceding-sibling axis contains all the preceding siblings of the context node
/AAA/XXX/preceding-sibling::*
//CCC/preceding-sibling::*

The following axis contains all nodes in the same document as the context node that are after the context node in document order, 
excluding any descendants and excluding attribute nodes and namespace nodes. 
/AAA/XXX/following::*
//ZZZ/following::*

The preceding axis contains all nodes in the same document as the context node that are before the context node in document order, 
excluding any ancestors and excluding attribute nodes and namespace nodes 
/AAA/XXX/preceding::*
//GGG/preceding::*

  
The descendant-or-self axis contains the context node and the descendants of the context node
/AAA/XXX/descendant-or-self::*
//CCC/descendant-or-self::*

The ancestor-or-self axis contains the context node and the ancestors of the context node; thus, the ancestor-or-self axis will always include the root node. 
/AAA/XXX/DDD/EEE/ancestor-or-self::*
//GGG/ancestor-or-self::*

 Orthogonal axes
The ancestor, descendant, following, preceding and self axes partition a document they do not overlap and together they contain all the nodes in the document. 
//GGG/ancestor::*
//GGG/descendant::*
//GGG/following::*
//GGG/preceding::*
//GGG/self::*
//GGG/ancestor::* | //GGG/descendant::* | //GGG/following::* | //GGG/preceding::* | //GGG/self::*


########################
To make tab work as indention
Create file .pythonrc

import readline  
readline.parse_and_bind("set disable-completion on")  
readline.parse_and_bind("tab: self-insert")

Edit your ~/.bashrc file and add the following line
export PYTHONSTARTUP=$HOME/.pythonrc


#### PEP-8, automatic in pycharms ###

1. Use 4 spaces per indentation level instead of TAB 

2. Hanging Indent 
	#Yes
	# Aligned with opening delimiter.
	foo = long_function_name(var_one, var_two,
							 var_three, var_four)

	#No
	# Arguments on first line forbidden when not using vertical alignment.
	foo = long_function_name(var_one, var_two,
		var_three, var_four)
	#Yes
	my_list = [
		1, 2, 3,
		4, 5, 6,
		]
	

3. Limit all lines to a maximum of 79 characters. 

4. Break Line before binary operator
	# Yes: easy to match operators with operands
	income = (gross_wages
			  + taxable_interest
			  + (dividends - qualified_dividends)
			  - ira_deduction
			  - student_loan_interest)
		  
5. Surround top-level function and class definitions with two blank lines. 
   Method definitions inside a class are surrounded by a single blank line. 

6. Imports should usually be on separate lines
	#Yes: 
		 import os
		 import sys

	#No:  
		 import sys, os

	#Yes 
		from subprocess import Popen, PIPE
		from foo.bar.yourclass import YourClass

7. Wildcard imports (  from <module> import *  ) should be avoided

8. Imports are always put at the top of the file, just after any module comments and docstrings, 
   and before module globals and constants.  

9. Folllow below placements for Module 
   Modules that are designed for use via  from M import *  should use the  __all__  mechanism to prevent exporting globals

"""
This is the example module.
This module does stuff.
"""

from __future__ import barry_as_FLUFL

__all__ = ['a', 'b', 'c']
__version__ = '0.1'
__author__ = 'Cardinal Biggles'

import os
import sys

10. White space - follow below style 

	#Yes
		spam(ham[1], {eggs: 2})
	#No
		spam( ham[ 1 ], { eggs: 2 } )     			#extra space around 1

	#Yes
		if x == 4: print x, y; x, y = y, x
	#No
		if x == 4 : print x , y ; x , y = y , x  	#extra space before , 

	#Yes 
		ham[lower+offset : upper+offset]
	#No
		ham[lower+offset:upper+offset]        		#No space around :

	#Yes
		spam(1)
	#No:
		spam (1)						#extra space after name 

	#Yes
		dct['key'] = lst[index]
	#No
		dct ['key'] = lst [index]		#extra space after dct,lst 

	#Yes: 
		x = 1
	#No: 
		x             = 1               #Many spaces

	#Yes: 
		def complex(real, imag=0.0):
			return magic(r=real, i=imag)
	#No: 
		def complex(real, imag = 0.0):           #space 
			return magic(r = real, i = imag)     #Space around keyword arg


11. It is okay to put an if/for/while with a small body on the same line, never do this for multi-clause statements
#OK
	if foo == 'blah': do_blah_thing()
#No
	if foo == 'blah': do_blah_thing()
	else: do_non_blah_thing()

12. Inline comments are unnecessary and in fact distracting if they state the obvious. 
	Atleast two spaces must be between statement and #
	#No 
		x = x + 1                 # Increment x

13. The  """  that ends a multiline docstring should be on a line by itself
    For one liner docstrings, keep the closing  """  on the same line. 
	
"""Return a foobang
Optional plotz says to frobnicate the bizbaz first.
"""


14. Nameing style of variable, function must be good, No specific recomendation to use only one style(eg CamelCase)
    Package and Module Names - the use of underscores is discouraged
	Class names should normally use the CapWords convention
	Function names should be lowercase, with words separated by underscores as necessary to improve readability
    Note below special rule 	
	_single_leading_underscore  : "internal use" indicator. E.g.  from M import *  does not import these elements 
	single_trailing_underscore_  : used by convention to avoid conflicts with Python keyword, e.g. 
		Tkinter.Toplevel(master, class_='ClassName')
	__double_leading_underscore  : when naming a class attribute, invokes name mangling (inside class FooBar,  __boo  becomes  _FooBar__boo )
	__double_leading_and_trailing_underscore__  : "magic" objects or attributes E.g.  __init__  , etc . Never invent such names; only use them as documented. 


15. Always use  self  for the first argument to instance methods. 
	Always use  cls  for the first argument to class methods. 
	Method Names and Instance Variables should be lowercase, with words separated by underscores
	Use one leading underscore only for non-public methods and instance variables. 
	To avoid name clashes with subclasses, use two leading underscores to invoke Python's name mangling rules
	f class Foo has an attribute named  __a  , it cannot be accessed by  Foo.__a  . (but by   Foo._Foo__a  .)

16. Programming Recommendations
	Comparisons to singletons like None should always be done with  is  or  is not  , never the equality operators. 
	Use  is not  operator rather than  not ... is
	When implementing ordering operations with rich comparisons, use the  functools.total_ordering()
	Always use a def statement instead of an assignment statement that binds a lambda expression directly to an identifier. 
	#Yes: 
		def f(x): return 2*x
	#No: 
		f = lambda x: 2*x
	Derive exceptions from  Exception  rather than  BaseException
	When raising an exception in Python 2, use  raise  ValueError('message')   instead of the older form  raise ValueError, 'message'
	When catching exceptions, mention specific exceptions whenever possible instead of using a bare  except:  clause
	When binding caught exceptions to a name, prefer the explicit name binding syntax(except Exception as exc) added in Python 2.6:
	When a resource is local to a particular section of code, use a  with  statement to ensure it is cleaned up promptly 
	Be consistent in return statements. Either all return statements in a function should return an expression, or none of them should
	Use  ''.startswith()   and  ''.endswith()   instead of string slicing to check for prefixes or suffixes. 	
	Object type comparisons should always use isinstance() instead of comparing types directly. 
	#Yes: 
		if isinstance(obj, int):
	#No:  
		if type(obj) is type(1):
	Py2.x:  When checking if an object is a string, check for unicode string as well  
			str and unicode have a common base class, basestring,  
				if isinstance(obj, basestring):
			Note in Py3,x,  unicode  and  basestring  no longer exist, there is only  str  and a bytes object 
	For sequences, (strings, lists, tuples), use the fact that empty sequences are false. 
	#Yes: 
		if not seq:
		if seq:
	#No: 
		if len(seq):
		if not len(seq):
	Don't compare boolean values to True or False using  ==  . 
	#Yes:   
		if greeting:
	#No:    
		if greeting == True:
	#Worse: 
		if greeting is True:



	
	
	
	
	
	
	
	
#########################
#Module -logging 
#default level is WARNING. hence WARNING and above are printed by default
#level DEBUG, INFO, WARNING, ERROR, CRITICAL

import logging
logging.warning('this is warning') 
logging.info('this would not be printed') 
logging.warning('%s before you %s', 'Look', 'leap!')

#With config 
import logging
logging.basicConfig(filename='l.log', filemode='w', level=logging.DEBUG)  # default filemode='a'
logging.debug('This message should go to the log file')
logging.info('So should this')
logging.warning('And this, too')

#Changing format
import logging
logging.basicConfig(format='%(levelname)s:%(message)s', level=logging.DEBUG)

#Format

asctime 	%(asctime)s  
created 	%(created)f 
filename 	%(filename)s 
funcName 	%(funcName)s 
levelname 	%(levelname)s 
levelno 	%(levelno)s 
lineno 		%(lineno)d 
module 		%(module)s 
msecs 		%(msecs)d 
message 	%(message)s 
name 		%(name)s  
pathname 	%(pathname)s 
process 	%(process)d 
processName %(processName)s 
thread 		%(thread)d 
threadName 	%(threadName)s 

#A good convention to use when naming loggers is to use a module-level logger, 
#in each module which uses logging, named as follows:
logger = logging.getLogger("filename")
logger.setLevel(logging.DEBUG) 
#call logger.debug, warning, info etc


#The default format set by basicConfig() for messages is:
severity:logger name:message


#Logging levels
DEBUG 	    Detailed information, typically of interest only when diagnosing problems. 
INFO 		Confirmation that things are working as expected. 
WARNING 	An indication that something unexpected happened, or indicative of some problem in the near future (e.g. ‘disk space low’). The software is still working as expected. 
ERROR 	    Due to a more serious problem, the software has not been able to perform some function. 
CRITICAL 	A serious error, indicating that the program itself may be unable to continue running.


#Command line

#If you want to control from command line, eg 
--log INFO

#code 
from optparse import OptionParser
p = OptionParser()
p.add_option("--log", action="store", type="string", dest="loglevel")

(o,r) = p.parse_args()  #pass commandline

numeric_level = getattr(logging, o.loglevel.upper(), None)
if not isinstance(numeric_level, int):
	raise ValueError('Invalid log level: %s' % loglevel)
logging.basicConfig(level=numeric_level, ...other_args...)



#########################
#Command line options - Module -optparse 

#To parse command line options 

(options, args) = parser.parse_args()  #defaults args taken from  sys.argv[1:]

#options, if --file takes a single string argument, then options.file will be the filename supplied by the user, 
#or None if the user did not supply that option
#args, the list of positional arguments leftover after parsing options

#Use below to add options 
parser.add_option(short_form, long_form, type="type_string", action="action_string", , default=value, help=help_string, 
                  dest= "variable_name", metavar="metavartype")
#	type can be
"string", "int", "long", "choice", "float" and "complex". 

#	Action can be

"store"  							store the value
"store_true", "store_false"			store boolean true or false if arg exists or not
"store_const"						store a constant value
"append"							append this option’s argument to a list
"count"								increment a counter by one
"callback"							call a specified function def my_callback(option, opt_str, value, parser, *args, **kwargs)


#Option style
-f foo
--file foo

#or included in the same argument:

-ffoo
--file=foo

#Example 


from optparse import OptionParser
...
parser = OptionParser()
parser.add_option("-f", "--file", dest="filename",
                  help="write report to FILE", metavar="FILE")
parser.add_option("-q", "--quiet",
                  action="store_false", dest="verbose", default=True,
                  help="don't print status messages to stdout")

(options, args) = parser.parse_args()


With these few lines of code, users of your script can now do the “usual thing” on the command-line, for example:


<yourscript> --file=outfile -q


As it parses the command line, optparse sets attributes of the options object returned by parse_args() based on user-supplied command-line values. When parse_args() returns from parsing this command line, options.filename will be "outfile" and options.verbose will be False. optparse supports both long and short options, allows short options to be merged together, and allows options to be associated with their arguments in a variety of ways. Thus, the following command lines are all equivalent to the above example:


<yourscript> -f outfile --quiet
<yourscript> --quiet --file outfile
<yourscript> -q -foutfile
<yourscript> -qfoutfile


Additionally, users can run one of


<yourscript> -h
<yourscript> --help


and optparse will print out a brief summary of your script’s options:


Usage: <yourscript> [options]

Options:
  -h, --help            show this help message and exit
  -f FILE, --file=FILE  write report to FILE
  -q, --quiet           don't print status messages to stdout


where the value of yourscript is determined at runtime (normally from sys.argv[0]).




#Example 

from optparse import OptionParser
p = OptionParser()
p.add_option("-f", "--file", action="store", type="string", dest="filename")
args = ["-f", "foo.txt"]
(o,r) = p.parse_args(args)
>>> o.filename
'foo.txt'



def record_foo_seen(option, opt_str, value, parser):
    parser.values.saw_foo = True

p.add_option("--foo", action="callback", callback=record_foo_seen)

p.add_option("-n", type="int", dest="num")
(o, args) = parser.parse_args(["-n42"])
print(o.num)

p.add_option("-v", action="store_true", dest="verbose", default=True)
(o, args) = parser.parse_args(["-v"])
print(o.verbose)

### optparse Hands-On
import os
import re
from optparse import OptionParser

def trim (str):
   tStr1 = re.sub(r'^\s+','',str)
   tStr2 = re.sub(r'\s+$','',tStr1)
   return tStr2;

def tuplize(str):
		return tuple(str.replace(',','').split("  "))




def get_exe(name):
	command = "tasklist /FI 'MODULES eq " + trim(name) + "' | tail -n +4 | awk '{print $1 \"  \"$5}' | tr '\n' ':'"
	result = os.popen (command).read()
	return ( dict(map(tuplize,(filter( bool, map(trim, result.split(':')))) )))


p = OptionParser()
p.add_option("-d", "--dll", action="store", type="string", dest="dll", help="give dll name")
args = ["-d", "ntdll.dll"]
(o,r) = p.parse_args(args)
get_exe(o.dll)

>>> (o,r) = p.parse_args(["-h"])

Usage:  [options]

Options:
-h, --help         show this help message and exit
-d DLL, --dll=DLL  give dll name

###Module Json 
#JSON (JavaScript Object Notation) specified by RFC 7159 

# to and from file 
json.dump(obj, fp, skipkeys=False, ensure_ascii=True, check_circular=True, 
   allow_nan=True, cls=None, indent=None, separators=None, default=None, sort_keys=False, **kw)
#fp is opened as fp = open(filename, "w")
#indent is string which is used for indentation

json.load(fp, cls=None, object_hook=None, parse_float=None, parse_int=None, 
    parse_constant=None, object_pairs_hook=None, **kw)
#fp is opened as fp = open(filename, "r")
#parse_type, if specified, will be called with the string of every JSON 'type' to be decoded


#To and from string  
json.dumps(obj, skipkeys=False, ensure_ascii=True, check_circular=True, 
    allow_nan=True, cls=None, indent=None, separators=None, default=None, sort_keys=False, **kw)

json.loads(s, encoding=None, cls=None, object_hook=None, parse_float=None, 
    parse_int=None, parse_constant=None, object_pairs_hook=None, **kw)
    
#Jason syntax 
JSON data is written as - "name":value pairs, eg -  "firstName":"John"
name is in double quote
JSON values can be:
    A number (integer or floating point)
    A string (in double quotes)
    A Boolean (true or false)
    An array (in square brackets with , as separator )
    An object (in curly braces with "name":value)
    null


#The file type for JSON files is ".json"
#The MIME type for JSON text is "application/json"
#No comment is allowed  even with /* */, // or #
#Only one root element or object is allowed, no multiple root elements 

#conversion table
#Python			JSON
dict 			object 
list, tuple 	array 
str 			string 
int, float,     int,float
True 			true 
False 			false 
None 			null 

#Note, separators=(',', ': ') as default if indent is not None.
#The json module always produces str objects, not bytes objects
#Keys in key/value pairs of JSON are always of the type str. 
#When a dictionary is converted into JSON, all the keys of the dictionary are coerced to strings
#That is, loads(dumps(x)) != x if x has non-string keys.
 

#Example file: example.json  
[
 { "empId: 1, "details": 
                        {                       
                          "firstName": "John",
                          "lastName": "Smith",
                          "isAlive": true,
                          "age": 25,
                          "salary": 123.5,
                          "address": {
                            "streetAddress": "21 2nd Street",
                            "city": "New York",
                            "state": "NY",
                            "postalCode": "10021-3100"
                          },
                          "phoneNumbers": [
                            {
                              "type": "home",
                              "number": "212 555-1234"
                            },
                            {
                              "type": "office",
                              "number": "646 555-4567"
                            },
                            {
                              "type": "mobile",
                              "number": "123 456-7890"
                            }
                          ],
                          "children": [],
                          "spouse": null
                        }
  } , { "empId: 20, "details": 
                            {                       
                              "firstName": "Johns",
                              "lastName": "Smith",
                              "isAlive": true,
                              "age": 25,
                              "salary": 123.5,
                              "address": {
                                "streetAddress": "21 2nd Street",
                                "city": "New York",
                                "state": "NY",
                                "postalCode": "10021-3100"
                              },
                              "phoneNumbers": [
                                {
                                  "type": "home",
                                  "number": "212 555-1234"
                                },
                                {
                                  "type": "office",
                                  "number": "646 555-4567"
                                },
                                {
                                  "type": "mobile",
                                  "number": "123 456-7890"
                                }
                              ],
                              "children": [],
                              "spouse": null
                            }
    }
]

#Example reading file 
import json 
import pprint 
fp = open("data/example.json", "r")
obj = json.load(fp)
fp.close()
pprint.pprint(obj)  #check size 
[{'details': {'address': {'city': 'New York',
                          'postalCode': '10021-3100',
                          'state': 'NY',
                          'streetAddress': '21 2nd Street'},
              'age': 25,
              'children': [],
              'firstName': 'John',
              'isAlive': True,
              'lastName': 'Smith',
              'phoneNumbers': [{'number': '212 555-1234', 'type': 'home'},
                               {'number': '646 555-4567', 'type': 'office'},
                               {'number': '123 456-7890', 'type': 'mobile'}],
              'salary': 123.5,
              'spouse': None},
  'empId': 1},
 {'details': {'address': {'city': 'New York',
                          'postalCode': '10021-3100',
                          'state': 'NY',
                          'streetAddress': '21 2nd Street'},
              'age': 25,
              'children': [],
              'firstName': 'Johns',
              'isAlive': True,
              'lastName': 'Smith',
              'phoneNumbers': [{'number': '212 555-1234', 'type': 'home'},
                               {'number': '646 555-4567', 'type': 'office'},
                               {'number': '123 456-7890', 'type': 'mobile'}],
              'salary': 123.5,
              'spouse': None},
  'empId': 20}]

#manipulations 
len(obj)        #2
type(obj)       #<class 'list'>
type(obj[0])    #<class 'dict'>
with open("data/example1.json", "w") as fp1:
    json.dump(obj, fp1, indent='\t')
  
#Obj is array , all array manipulations can be used 
[emp['details']['address']['state']   for emp in obj if emp['empId'] > 10] 






### Few file manipulations methods 

shutil.copyfileobj(fsrc, fdst) #with file object 
shutil.copyfile(src, dst)      #with string 
shutil.copystat(src, dst)
shutil.copy(src, dst)          #dst can be destination directory 
shutil.copy2(src, dst)         #same as copy(), but metadata preserving 
shutil.copytree(src, dst, ignore=None)      #recursive copy, dst must not exist and it is created with any missing parents 

shutil.move(src, dst)          #recursive move from src dir to dst dir 
shutil.rmtree(path, ignore_errors=False, onerror=None) #recursive delete, path is dir 
shutil.which(cmd) 

os.remove(path)             #remove a file 
os.rename(src, dst)         #rename a file 
os.stat(path)               #returns os.stat_result which is result of stat
os.scandir(path='.')        #returns os.DirEntry after opening a dir (except . and ..), DirEntry has name, path, is_dir(), is_file(), stat()(returns os.stat_result) attributes 
os.truncate(path, length)   #Truncate the file  with at most length bytes in size.

#Example 
from shutil import copytree, ignore_patterns
copytree(source, destination, ignore=ignore_patterns('*.pyc', 'tmp*'))

#Removing readonly files in windows 

import os, stat
import shutil

def remove_readonly(func, path, _):
    "Clear the readonly bit and reattempt the removal"
    os.chmod(path, stat.S_IWRITE)
    func(path)

shutil.rmtree(directory, onerror=remove_readonly)

#Example 
>>> import os
>>> statinfo = os.stat('somefile.txt')
>>> statinfo
os.stat_result(st_mode=33188, st_ino=7876932, st_dev=234881026,
st_nlink=1, st_uid=501, st_gid=501, st_size=264, st_atime=1297230295,
st_mtime=1297230027, st_ctime=1297230027)
>>> statinfo.st_size
264

#Example 
for entry in os.scandir(path):
   if not entry.name.startswith('.') and entry.is_file():
       print(entry.name)
       
       
       
### High performance collections - import collections 
namedtuple()        factory function for creating tuple subclasses with named fields 
deque               list-like container with fast appends and pops on either end 
ChainMap            dict-like class for creating a single view of multiple mappings 
Counter             dict subclass for counting hashable objects 
OrderedDict         dict subclass that remembers the order entries were added 
defaultdict         dict subclass that calls a factory function to supply missing values 

UserDict            wrapper around dictionary objects for easier dict subclassing 
UserList            wrapper around list objects for easier list subclassing 
UserString          wrapper around string objects for easier string subclassing 

##A ChainMap groups multiple dicts or other mappings together to create a single, updateable view.
#Lookups search the underlying mappings successively until a key is found. 
#In contrast, writes, updates, and deletions only operate on the first mapping.

#Example of simulating Python’s internal lookup chain:

import builtins
pylookup = ChainMap(locals(), globals(), vars(builtins))


##A counter tool is provided to support convenient and rapid tallies. 
#For example:
#A Counter is a dict subclass for counting hashable objects. 
#It is an unordered collection where elements are stored as dictionary keys 
#and their counts are stored as dictionary values.

>>>>>> # Tally occurrences of words in a list
>>> cnt = Counter()
>>> for word in ['red', 'blue', 'red', 'green', 'blue', 'blue']:
...     cnt[word] += 1
>>> cnt
Counter({'blue': 3, 'red': 2, 'green': 1})

>>> # Find the ten most common words in Hamlet
>>> import re
>>> words = re.findall(r'\w+', open('hamlet.txt').read().lower())
>>> Counter(words).most_common(10)
[('the', 1143), ('and', 966), ('to', 762), ('of', 669), ('i', 631),
 ('you', 554),  ('a', 546), ('my', 514), ('hamlet', 471), ('in', 451)]

#Counter objects support three methods beyond those available for all dictionaries:
elements()
most_common([n])
subtract([iterable-or-mapping])


##deque: Deques are a generalization of stacks and queues 
#Deques support thread-safe, memory efficient appends and pops from either side of the deque with approximately the same O(1) performance in either direction.

>>> from collections import deque
>>> d = deque('ghi')                 # make a new deque with three items
>>> for elem in d:                   # iterate over the deque's elements
...     print(elem.upper())
G
H
I

>>> d.append('j')                    # add a new entry to the right side
>>> d.appendleft('f')                # add a new entry to the left side
>>> d                                # show the representation of the deque
deque(['f', 'g', 'h', 'i', 'j'])

>>> d.pop()                          # return and remove the rightmost item
'j'
>>> d.popleft()                      # return and remove the leftmost item
'f'
>>> list(d)                          # list the contents of the deque
['g', 'h', 'i']
>>> d[0]                             # peek at leftmost item
'g'
>>> d[-1]                            # peek at rightmost item
'i'

>>> list(reversed(d))                # list the contents of a deque in reverse
['i', 'h', 'g']
>>> 'h' in d                         # search the deque
True
>>> d.extend('jkl')                  # add multiple elements at once
>>> d
deque(['g', 'h', 'i', 'j', 'k', 'l'])
>>> d.rotate(1)                      # right rotation
>>> d
deque(['l', 'g', 'h', 'i', 'j', 'k'])
>>> d.rotate(-1)                     # left rotation
>>> d
deque(['g', 'h', 'i', 'j', 'k', 'l'])

>>> deque(reversed(d))               # make a new deque in reverse order
deque(['l', 'k', 'j', 'i', 'h', 'g'])
>>> d.clear()                        # empty the deque
>>> d.pop()                          # cannot pop from an empty deque
Traceback (most recent call last):
    File "<pyshell#6>", line 1, in -toplevel-
        d.pop()
IndexError: pop from an empty deque

>>> d.extendleft('abc')              # extendleft() reverses the input order
>>> d
deque(['c', 'b', 'a'])



##defaultdict objects - defaultdict(arg)
#The first argument provides class which would be used for default value 
# default value for non existent key = arg()

#for example, in below code 

>>> s = [('yellow', 1), ('blue', 2), ('yellow', 3), ('blue', 4), ('red', 1)]
>>> d = defaultdict(list)  #default value is empty list , =list()
>>> for k, v in s: 
        d[k].append(v) #hence any list method can be used

>>> sorted(d.items())
[('blue', [2, 4]), ('red', [1]), ('yellow', [1, 3])]

>>> s = 'mississippi'
>>> d = defaultdict(int) #default value is zero, =int()
>>> for k in s:
        d[k] += 1 #hence any int method can be used
...
>>> d.items()
[('i', 4), ('p', 2), ('s', 4), ('m', 1)]


>>> s = [('red', 1), ('blue', 2), ('red', 3), ('blue', 4), ('red', 1), ('blue', 4)]
>>> d = defaultdict(set) #default value is is empty set, = set()
>>> for k, v in s:
        d[k].add(v) ##hence any set method can be used
...
>>> d.items()
[('blue', set([2, 4])), ('red', set([1, 3]))]


##namedtuple(): Factory Function for Tuples with Named Fields

>>> # Basic example
>>> Point = namedtuple('Point', ['x', 'y'])
>>> p = Point(11, y=22)     # instantiate with positional or keyword arguments
>>> p[0] + p[1]             # indexable like the plain tuple (11, 22)
33
>>> x, y = p                # unpack like a regular tuple
>>> x, y
(11, 22)
>>> p.x + p.y               # fields also accessible by name
33
>>> p                       # readable __repr__ with a name=value style
Point(x=11, y=22)

#example 
EmployeeRecord = namedtuple('EmployeeRecord', 'name, age, title, department, paygrade')

import csv
for emp in map(EmployeeRecord._make, csv.reader(open("employees.csv", "rb"))):
    print(emp.name, emp.title)

import sqlite3
conn = sqlite3.connect('/companydata')
cursor = conn.cursor()
cursor.execute('SELECT name, age, title, department, paygrade FROM employees')
for emp in map(EmployeeRecord._make, cursor.fetchall()):
    print(emp.name, emp.title)

    
##OrderedDict : Ordered dictionaries are just like regular dictionaries 
#but they remember the order that items were inserted. 

#When iterating over an ordered dictionary, 
#the items are returned in the order their keys were first added.
>>> # regular unsorted dictionary
>>> d = {'banana': 3, 'apple': 4, 'pear': 1, 'orange': 2}

>>> # dictionary sorted by key
>>> OrderedDict(sorted(d.items(), key=lambda t: t[0]))
OrderedDict([('apple', 4), ('banana', 3), ('orange', 2), ('pear', 1)])

>>> # dictionary sorted by value
>>> OrderedDict(sorted(d.items(), key=lambda t: t[1]))
OrderedDict([('pear', 1), ('orange', 2), ('banana', 3), ('apple', 4)])

>>> # dictionary sorted by length of the key string
>>> OrderedDict(sorted(d.items(), key=lambda t: len(t[0])))
OrderedDict([('pear', 1), ('apple', 4), ('orange', 2), ('banana', 3)])



##Classes that help to in builtin class inheritance 
#Note , directly dict, list, str can be inherited with self as underlying data 

class collections.UserDict([initialdata])
    Class that simulates a dictionary
    data
        A real dictionary used to store the contents of the UserDict class.

class collections.UserList([list])
    Class that simulates a list
    Subclasses of UserList are expected to offer a constructor 
    which can be called with either no arguments 
    or one argument(which is a sequence object used as a data source)
    data
        A real list object used to store the contents of the UserList class
    

class collections.UserString([sequence])
    Class that simulates a string or a Unicode string object
    data
        A real object used to store the contents 



###itertools — Functions creating iterators for efficient looping

#Infinite Iterators:
#Iterator        Arguments       Results
count()          start, [step]  start, start+step, start+2*step, 
    count(10) --> 10 11 12 13 14 ... 
cycle()          p              p0, p1, ... plast, p0, p1, 
    cycle('ABCD') --> A B C D A B C D ... 
repeat()        elem [,n]       elem, elem, elem, ... endlessly or up to n times 
    repeat(10, 3) --> 10 10 10 

##Iterators terminating on the shortest input sequence:
#Iterator        Arguments       Results
accumulate()     p [,func]      p0, p0+p1, p0+p1+p2, 
    accumulate([1,2,3,4,5]) --> 1 3 6 10 15 
chain() p, q, ... p0, p1, ... plast, q0, q1, ... chain('ABC', 'DEF') --> A B C D E F 
chain.from_iterable() iterable p0, p1, ... plast, q0, q1, ... chain.from_iterable(['ABC', 'DEF']) --> A B C D E F 
compress() data, selectors (d[0] if s[0]), (d[1] if s[1]), ... compress('ABCDEF', [1,0,1,0,1,1]) --> A C E F 
dropwhile() pred, seq seq[n], seq[n+1], starting when pred fails dropwhile(lambda x: x<5, [1,4,6,4,1]) --> 6 4 1 
filterfalse() pred, seq elements of seq where pred(elem) is false filterfalse(lambda x: x%2, range(10)) --> 0 2 4 6 8 
groupby() iterable[, keyfunc] sub-iterators grouped by value of keyfunc(v)   
islice() seq, [start,] stop [, step] elements from seq[start:stop:step] islice('ABCDEFG', 2, None) --> C D E F G 
starmap() func, seq func(*seq[0]), func(*seq[1]), ... starmap(pow, [(2,5), (3,2), (10,3)]) --> 32 9 1000 
takewhile() pred, seq seq[0], seq[1], until pred fails takewhile(lambda x: x<5, [1,4,6,4,1]) --> 1 4 
tee() it, n it1, it2, ... itn splits one iterator into n   
zip_longest() p, q, ... (p[0], q[0]), (p[1], q[1]), ... zip_longest('ABCD', 'xy', fillvalue='-') --> Ax By C- D- 

Combinatoric generators:






Iterator

Arguments

Results


product() p, q, ... [repeat=1] cartesian product, equivalent to a nested for-loop 
permutations() p[, r] r-length tuples, all possible orderings, no repeated elements 
combinations() p, r r-length tuples, in sorted order, no repeated elements 
combinations_with_replacement() p, r r-length tuples, in sorted order, with repeated elements 
product('ABCD', repeat=2)   AA AB AC AD BA BB BC BD CA CB CC CD DA DB DC DD 
permutations('ABCD', 2)   AB AC AD BA BC BD CA CB CD DA DB DC 
combinations('ABCD', 2)   AB AC AD BC BD CD 
combinations_with_replacement('ABCD', 2)   AA AB AC AD BB BC BD CC CD DD 


###################
#ipaddress — IPv4/IPv6 manipulation library

import ipaddress
ipaddress.ip_address('192.0.2.1') #IPv4Address('192.0.2.1')
ipaddress.ip_address(3221225985) #IPv4Address('192.0.2.1')
ipaddress.IPv4Address(1)  #IPv4Address('0.0.0.1')

#V6
>>> ipaddress.ip_address('2001:DB8::1')
IPv6Address('2001:db8::1')


#“network address/network prefix” 
#Network objects cannot have any host bits set
'''
A prefix /<nbits> is a notation that denotes 
how many high-order bits are set in the network mask. 
A net mask is an IP address with some number of high-order bits set.
Thus the prefix /24 is equivalent to the net mask 255.255.255.0 in IPv4, or ffff:ff00:: in IPv6. 
a host mask is the logical inverse of a net mask, to denote a network mask. 
The host mask equivalent to /24 in IPv4 is 0.0.0.255.
'''

ipaddress.ip_network('192.0.2.0/24')  #IPv4Network('192.0.2.0/24')
#By default, attempting to create a network object with host bits set will result in ValueError being raised
>>> ipaddress.ip_network('192.0.2.1/24')
Traceback (most recent call last):
   ...
ValueError: 192.0.2.1/24 has host bits set
>>> ipaddress.ip_network('192.0.2.1/24', strict=False)
IPv4Network('192.0.2.0/24')
# networks can also be defined with integers
>>> ipaddress.ip_network(3221225984)
IPv4Network('192.0.2.0/32')



#to describe an address on a particular network
#Notation like 192.0.2.1/24 is commonly used
#means “the host 192.0.2.1 on the network 192.0.2.0/24”, 

>>> ipaddress.ip_interface('192.0.2.1/24')
IPv4Interface('192.0.2.1/24')



##Inspecting Address/Network/Interface Objects

#Extracting the IP version:


>>> addr4 = ipaddress.ip_address('192.0.2.1')
>>> addr6 = ipaddress.ip_address('2001:db8::1')
>>> addr6.version
6
>>> addr4.version
4


#Obtaining the network from an interface:


>>> host4 = ipaddress.ip_interface('192.0.2.1/24')
>>> host4.network
IPv4Network('192.0.2.0/24')
>>> host6 = ipaddress.ip_interface('2001:db8::1/96')
>>> host6.network
IPv6Network('2001:db8::/96')


#Finding out how many individual addresses are in a network:

>>> net4 = ipaddress.ip_network('192.0.2.0/24')
>>> net4.num_addresses
256
>>> net6 = ipaddress.ip_network('2001:db8::0/96')
>>> net6.num_addresses
4294967296

#Iterating through the “usable” addresses on a network:
>>> net4 = ipaddress.ip_network('192.0.2.0/24')
>>> for x in net4.hosts():
...     print(x)  
192.0.2.1
192.0.2.2
192.0.2.3
192.0.2.4
...
192.0.2.252
192.0.2.253
192.0.2.254

#Obtaining the netmask (i.e. set bits corresponding to the network prefix) 
#or the hostmask (any bits that are not part of the netmask):

>>> net4 = ipaddress.ip_network('192.0.2.0/24')
>>> net4.netmask
IPv4Address('255.255.255.0')
>>> net4.hostmask
IPv4Address('0.0.0.255')
>>> net6 = ipaddress.ip_network('2001:db8::0/96')
>>> net6.netmask
IPv6Address('ffff:ffff:ffff:ffff:ffff:ffff::')
>>> net6.hostmask
IPv6Address('::ffff:ffff')


#Exploding or compressing the address:
#IPv4 doesn’t support explosion or compression
>>> addr6.exploded
'2001:0db8:0000:0000:0000:0000:0000:0001'
>>> addr6.compressed
'2001:db8::1'
>>> net6.exploded
'2001:0db8:0000:0000:0000:0000:0000:0000/96'
>>> net6.compressed
'2001:db8::/96'


#Networks as lists of Addresses
>>> net4[1]
IPv4Address('192.0.2.1')
>>> net4[-1]
IPv4Address('192.0.2.255')
>>> net6[1]
IPv6Address('2001:db8::1')
>>> net6[-1]
IPv6Address('2001:db8::ffff:ffff')

#membership checking 

>>> addr4 = ipaddress.ip_address('192.0.2.1')
>>> addr4 in ipaddress.ip_network('192.0.2.0/24')
True
>>> addr4 in ipaddress.ip_network('192.0.3.0/24')
False


#Comparisons
>>> ipaddress.ip_address('192.0.2.1') < ipaddress.ip_address('192.0.2.2')
True

>>> IPv4Address('127.0.0.2') > IPv4Address('127.0.0.1')
True
>>> IPv4Address('127.0.0.2') == IPv4Address('127.0.0.1')
False
>>> IPv4Address('127.0.0.2') != IPv4Address('127.0.0.1')
True

#Arithmetic operators
#Integers can be added to or subtracted from address objects
>>> IPv4Address('127.0.0.2') + 3
IPv4Address('127.0.0.5')
>>> IPv4Address('127.0.0.2') - 3
IPv4Address('126.255.255.255')
>>> IPv4Address('255.255.255.255') + 1
Traceback (most recent call last):
  File "<stdin>", line 1, in <module>

#Other modules that use IP addresses (such as socket) usually won’t accept objects from this module directly. Instead, they must be coerced to an integer or string that the other module will accept:

>>> addr4 = ipaddress.ip_address('192.0.2.1')
>>> str(addr4)
'192.0.2.1'
>>> int(addr4)
3221225985

#Other Important attributes of class ipaddress.IPv4Address or V6

exploded
The string representation in dotted decimal notation. Leading zeroes are never included in the representation

packed
The binary representation of this address - a bytes object of the appropriate length (most significant octet first). This is 4 bytes for IPv4 and 16 bytes for IPv6.

reverse_pointer (3.5)
The name of the reverse DNS PTR record for the IP address, e.g.:
>>> ipaddress.ip_address("127.0.0.1").reverse_pointer
'1.0.0.127.in-addr.arpa'

is_multicast
True if the address is reserved for multicast use. See RFC 3171 (for IPv4) or RFC 2373 (for IPv6).
is_private
True if the address is allocated for private networks. See iana-ipv4-special-registry (for IPv4) or iana-ipv6-special-registry (for IPv6).
is_global(3.4)
True if the address is allocated for public networks. See iana-ipv4-special-registry (for IPv4) or iana-ipv6-special-registry (for IPv6).
is_unspecified
True if the address is unspecified. See RFC 5735 (for IPv4) or RFC 2373 (for IPv6).
is_reserved
True if the address is otherwise IETF reserved.
is_loopback
True if this is a loopback address. See RFC 3330 (for IPv4) or RFC 2373 (for IPv6).
is_link_local
True if the address is reserved for link-local usage. See RFC 3927.


#Important methods of class ipaddress.IPv4Network(address, strict=True)
network_address
The network address for the network. The network address and the prefix length together uniquely define a network.
broadcast_address
The broadcast address for the network. Packets sent to the broadcast address should be received by every host on the network.
hostmask
The host mask, as a string.
with_netmask
A string representation of the network, with the mask in net mask notation.
with_hostmask
A string representation of the network, with the mask in host mask notation.
num_addresses
The total number of addresses in the network.
prefixlen
Length of the network prefix, in bits.
hosts()
Returns an iterator over the usable hosts in the network. The usable hosts are all the IP addresses that belong to the network, except the network address itself and the network broadcast address.
>>> list(ip_network('192.0.2.0/29').hosts())  
[IPv4Address('192.0.2.1'), IPv4Address('192.0.2.2'),
 IPv4Address('192.0.2.3'), IPv4Address('192.0.2.4'),
 IPv4Address('192.0.2.5'), IPv4Address('192.0.2.6')]

overlaps(other)
True if this network is partly or wholly contained in other or other is wholly contained in this network.

address_exclude(network)
Computes the network definitions resulting from removing the given network from this one. Returns an iterator of network objects. Raises ValueError if network is not completely contained in this network.
>>> n1 = ip_network('192.0.2.0/28')
>>> n2 = ip_network('192.0.2.1/32')
>>> list(n1.address_exclude(n2))  
[IPv4Network('192.0.2.8/29'), IPv4Network('192.0.2.4/30'),
 IPv4Network('192.0.2.2/31'), IPv4Network('192.0.2.0/32')]

subnets(prefixlen_diff=1, new_prefix=None)
The subnets that join to make the current network definition
prefixlen_diff is the amount our prefix length should be increased by. 
new_prefix is the desired new prefix of the subnets; 
it must be larger than our prefix. 
One and only one of prefixlen_diff and new_prefix must be set. 
Returns an iterator of network objects.

>>> list(ip_network('192.0.2.0/24').subnets())
[IPv4Network('192.0.2.0/25'), IPv4Network('192.0.2.128/25')]
>>> list(ip_network('192.0.2.0/24').subnets(prefixlen_diff=2))  
[IPv4Network('192.0.2.0/26'), IPv4Network('192.0.2.64/26'),
 IPv4Network('192.0.2.128/26'), IPv4Network('192.0.2.192/26')]
>>> list(ip_network('192.0.2.0/24').subnets(new_prefix=26))  
[IPv4Network('192.0.2.0/26'), IPv4Network('192.0.2.64/26'),
 IPv4Network('192.0.2.128/26'), IPv4Network('192.0.2.192/26')]
>>> list(ip_network('192.0.2.0/24').subnets(new_prefix=23))
Traceback (most recent call last):
  File "<stdin>", line 1, in <module>
    raise ValueError('new prefix must be longer')
ValueError: new prefix must be longer
>>> list(ip_network('192.0.2.0/24').subnets(new_prefix=25))
[IPv4Network('192.0.2.0/25'), IPv4Network('192.0.2.128/25')]

supernet(prefixlen_diff=1, new_prefix=None)
The supernet containing this network definition
prefixlen_diff is the amount our prefix length should be decreased by. 
 
>>> ip_network('192.0.2.0/24').supernet()
IPv4Network('192.0.2.0/23')
>>> ip_network('192.0.2.0/24').supernet(prefixlen_diff=2)
IPv4Network('192.0.0.0/22')
>>> ip_network('192.0.2.0/24').supernet(new_prefix=20)
IPv4Network('192.0.0.0/20')

compare_networks(other)
Compare this network to other. 
In this comparison only the network addresses are considered; 
host bits aren’t. Returns either -1, 0 or 1.

>>> ip_network('192.0.2.1/32').compare_networks(ip_network('192.0.2.2/32'))
-1
>>> ip_network('192.0.2.1/32').compare_networks(ip_network('192.0.2.0/32'))
1
>>> ip_network('192.0.2.1/32').compare_networks(ip_network('192.0.2.1/32'))
0

#Logical operators ipaddress.IPv4Network(address, strict=True)

#Network objects can be compared with the usual set of logical operators, similarly to address objects.


#Iteration ipaddress.IPv4Network(address, strict=True)
#Network objects can be iterated to list all the addresses belonging to the network. 
#For iteration, all hosts are returned, including unusable hosts 
#(for usable hosts, use the hosts() method)


>>>>>> for addr in IPv4Network('192.0.2.0/28'):
...     addr
...
IPv4Address('192.0.2.0')
IPv4Address('192.0.2.1')
IPv4Address('192.0.2.2')
IPv4Address('192.0.2.3')
IPv4Address('192.0.2.4')
IPv4Address('192.0.2.5')
IPv4Address('192.0.2.6')
IPv4Address('192.0.2.7')
IPv4Address('192.0.2.8')
IPv4Address('192.0.2.9')
IPv4Address('192.0.2.10')
IPv4Address('192.0.2.11')
IPv4Address('192.0.2.12')
IPv4Address('192.0.2.13')
IPv4Address('192.0.2.14')
IPv4Address('192.0.2.15')


##class ipaddress.IPv4Interface(address)
#IPv4Interface is a subclass of IPv4Address
#Additional attributes 
ip
The address (IPv4Address) without network information.

>>> interface = IPv4Interface('192.0.2.5/24')
>>> interface.ip
IPv4Address('192.0.2.5')

network
The network (IPv4Network) this interface belongs to.

>>> interface = IPv4Interface('192.0.2.5/24')
>>> interface.network
IPv4Network('192.0.2.0/24')

with_prefixlen
A string representation of the interface with the mask in prefix notation.

>>> interface = IPv4Interface('192.0.2.5/24')
>>> interface.with_prefixlen
'192.0.2.5/24'

with_netmask
A string representation of the interface with the network as a net mask.

>>> interface = IPv4Interface('192.0.2.5/24')
>>> interface.with_netmask
'192.0.2.5/255.255.255.0'

with_hostmask
A string representation of the interface with the network as a host mask.
>>> interface = IPv4Interface('192.0.2.5/24')
>>> interface.with_hostmask
'192.0.2.5/0.0.0.255'

##Other Module Level Functions
ipaddress.v4_int_to_packed(address)
Represent an address as 4 packed bytes in network (big-endian) order. address is an integer representation of an IPv4 IP address. A ValueError is raised if the integer is negative or too large to be an IPv4 IP address.

>>> ipaddress.ip_address(3221225985)
IPv4Address('192.0.2.1')
>>> ipaddress.v4_int_to_packed(3221225985)
b'\xc0\x00\x02\x01'

ipaddress.v6_int_to_packed(address)
Represent an address as 16 packed bytes in network (big-endian) order. address is an integer representation of an IPv6 IP address. A ValueError is raised if the integer is negative or too large to be an IPv6 IP address.

ipaddress.summarize_address_range(first, last)
Return an iterator of the summarized network range given the first and last IP addresses. first is the first IPv4Address or IPv6Address in the range and last is the last IPv4Address or IPv6Address in the range. A TypeError is raised if first or last are not IP addresses or are not of the same version. A ValueError is raised if last is not greater than first or if first address version is not 4 or 6.

>>> [ipaddr for ipaddr in ipaddress.summarize_address_range(
...    ipaddress.IPv4Address('192.0.2.0'),
...    ipaddress.IPv4Address('192.0.2.130'))]
[IPv4Network('192.0.2.0/25'), IPv4Network('192.0.2.128/31'), IPv4Network('192.0.2.130/32')]

ipaddress.collapse_addresses(addresses)
Return an iterator of the collapsed IPv4Network or IPv6Network objects. addresses is an iterator of IPv4Network or IPv6Network objects. A TypeError is raised if addresses contains mixed version objects.

>>> [ipaddr for ipaddr in
... ipaddress.collapse_addresses([ipaddress.IPv4Network('192.0.2.0/25'),
... ipaddress.IPv4Network('192.0.2.128/25')])]
[IPv4Network('192.0.2.0/24')]

ipaddress.get_mixed_type_key(obj)
Return a key suitable for sorting between networks and addresses. Address and Network objects are not sortable by default; they’re fundamentally different, so the expression:


##Private IP network
•192.168.0.0 - 192.168.255.255 (65,536 IP addresses)
•172.16.0.0 - 172.31.255.255 (1,048,576 IP addresses)
•10.0.0.0 - 10.255.255.255 (16,777,216 IP addresses)

#Subnet Masks and Subnets:  
For a "Class C" or "8-bit" subnet (32-24=8), 
So we use 255.255.255.0, or its shorthand equivalent, /24. 
For a Class B or "16-bit" subnet (32-16=16), 
So we use 255.255.0.0, or /16. 

#Classless Inter-Domain Routing (CIDR) Chart
#notation  resulting subnet  
#netmask       shorthand      number of addresses  
255.255.255.0    /24 [8-bit]  2**8 =  256  = 254 hosts + 1 bcast + 1 net base  
255.255.255.128  /25 [7-bit]  2**7 =  128  = 126 hosts + 1 bcast + 1 net base  
255.255.255.192  /26 [6-bit]  2**6 =  64  = 62 hosts + 1 bcast + 1 net base  
255.255.255.224  /27 [5-bit]  2**5 =  32  = 30 hosts + 1 bcast + 1 net base  
255.255.255.240  /28 [4-bit]  2**4 =  16  = 14 hosts + 1 bcast + 1 net base  
255.255.255.248  /29 [3-bit]  2**3 =  8  = 6 hosts + 1 bcast + 1 net base  
255.255.255.252  /30 [2-bit]  2**2 =  4  = 2 hosts + 1 bcast + 1 net base  
255.255.255.254  /31 [1-bit]  2**1 =  -  invalid (no possible hosts)  
255.255.255.255  /32 [0-bit]  2**0 =  1  a host route (odd duck case)  


#How many 9-bit subnets can fit into a 13-bit subnet?  
2**13
-- 
2**9  
= 2**(13 - 9) =  2**4 = 16  

#Example 
 207.199.153.192/27 is a "5-bit" subnet (32-27=5). 
There are 32 IP's in the subnet. 
The "base" address or first IP of the range is simply 207.199.153.192, and is unusable as a host address. 
The 30 Usable IPs are 207.199.153.193..207.199.153.223. 
The last one, 207.199.153.224, is the broadcast address for the subnet. 
the broadcast address is not usable as a host address. 

#General Network Architecture:  
The internal subnetting uses the private "Class B" network, 
172.16.0.0/16, divided up as follows
#network/mask usable IP address range  bcast address
172.16.1.0/24 172.16.1.1..172.16.1.254 172.16.1.255 
172.16.2.0/25 172.16.2.1..172.16.2.126 172.16.2.127 
172.16.3.0/25 172.16.3.1..172.16.3.126 172.16.3.127 
172.16.4.0/24 172.16.4.1..172.16.4.254 172.16.4.255 
172.16.5.0/24 172.16.5.1..172.16.5.254 172.16.5.255 
172.16.6.0/24 172.16.6.1..172.16.6.254 172.16.6.255 


##Network Base Address and Broadcast Address: 
The network base address is the first IP address in a given subnet; 
the broadcast address is the last. 
All NICs have to listen for traffic directed at their specific IP address(es) 
and the broadcast address for their subnet. 
The base network address is all 0 for the hostid and refers to the subnet itself; 
the broadcast address is all 1 and refers to all hosts on the subnet. 


#Routing table 
When an IP packet is to be forwarded, the routing table is used to determine:
1.The forwarding or next-hop IP address: 
For a direct delivery, the forwarding IP address is the destination IP address in the IP packet.
For an indirect delivery, the forwarding IP address is the IP address of a router. 

2.The interface to be used for the forwarding: 
The interface identifies the physical or logical interface such as a network adapter 
that is used to forward the packet to either its destination or the next router. 


•For each entry in a routing table, 
perform a bit-wise logical AND between the destination IP address and the network mask. 
Compare the result with the network ID of the entry for a match.

•The list of matching routes is compiled. 
The route that has the longest match (the route that matched the most amount of bits with the destination IP address) is chosen. 
The longest matching route is the most specific route to the destination IP address. 
If multiple entries with the longest match are found (multiple routes to the same network ID, for example), 
the router uses the lowest metric to select the best route



Network Destination Netmask         Gateway         Interface           Metric      Purpose
0.0.0.0             0.0.0.0         157.55.16.1     157.55.27.90        1           Default Route
127.0.0.0           255.0.0.0       127.0.0.1       127.0.0.1           1           Loopback Network
157.55.16.0         255.255.240.0   157.55.27.90    157.55.27.90        1           Directly Attached Network
157.55.27.90        255.255.255.255 127.0.0.1       127.0.0.1           1           Local Host
157.55.255.255      255.255.255.255 157.55.27.90    157.55.27.90        1           Network Broadcast
224.0.0.0           224.0.0.0       157.55.27.90    157.55.27.90        1           Multicast Address
255.255.255.255     255.255.255.255 157.55.27.90    157.55.27.90        1           Limited Broadcast
 

Direct delivery occurs when the IP node (either the sending node or an IP router)
forwards a packet to the final destination on a directly attached network. 
The IP node encapsulates the IP datagram in a frame format for the Network Interface layer 
(such as Ethernet or Token Ring) addressed to the destination's physical address.

Indirect delivery occurs when the IP node (either the sending node or an IP router) 
forwards a packet to an intermediate node (an IP router) 
because the final destination is not on a directly attached network. 
The IP node encapsulates the IP datagram in a frame format, 
addressed to the IP router's physical address, for the Network Interface layer (such as Ethernet or Token Ring).

#IP on the Sending Host
When a packet is sent by a sending host, 
the packet is handed from an upper layer protocol (TCP, UDP, or ICMP) to IP. 
IP on the sending host does the following:

1.Sets the Time-to-Live (TTL) value to either a default or application-specified value.

2.IP checks its routing table for the best route to the destination IP address. 
If no route is found, IP indicates a routing error to the upper layer protocol (TCP, UDP, or ICMP). 

3.Based on the most specific route, IP determines the forwarding IP address 
and the interface to be used for forwarding the packet.

4.IP hands the packet, the forwarding IP address, and the interface 
to Address Resolution Protocol (ARP), 
and then ARP resolves the forwarding IP address to its media access control (MAC) address and forwards the packet.

#IP on the Router

When a packet is received at a router, the packet is passed to IP. 
IP on the router does the following:

1.IP verifies the IP header checksum. 
If the IP header checksum fails, the IP packet is discarded without notification to the user. This is known as a silent discard . 

2.IP verifies whether the destination IP address in the IP datagram corresponds to an IP address assigned to a router interface. 
If so, the router processes the IP datagram as the destination host 

3.If the destination IP address is not the router, IP decreases the time-to-live (TTL) by 1. 
If the TTL is 0, the router discards the packet and sends an ICMP Time Expired-TTL Expired message to the sender. 

4.If the TTL is 1 or greater, IP updates the TTL field and calculates a new IP header checksum.

5.IP checks its routing table for the best route to the destination IP address in the IP datagram. 
If no route is found, the router discards the packet and sends an ICMP Destination Unreachable-Network Unreachable message to the sender. 

6.Based on the best route found, IP determines the forwarding IP address 
and the interface to be used for forwarding the packet.

7.IP hands the packet, the forwarding IP address, and the interface to ARP, 
and then ARP forwards the packet to the appropriate MAC address.

This entire process is repeated at each router in the path between the source 
and destination host.

#IP on the Destination Host
When a packet is received at the destination host, it is passed up to IP. 
IP on the destination host does the following:

1.IP verifies the IP header checksum. 
If the IP header checksum fails, the IP packet is silently discarded. 

2.IP verifies that the destination IP address in the IP datagram corresponds to an IP address assigned to the host. 
If the destination IP address is not assigned to the host, the IP packet is silently discarded. 

3.Based on the IP protocol field, IP passes the IP datagram without the IP header to the appropriate upper-level protocol. 
If the protocol does not exist, ICMP sends a Destination Unreachable-Protocol Unreachable message back to the sender. 

4.For TCP and UDP packets, the destination port is checked and the TCP segment or UDP header is processed. 
If no application exists for the UDP port number, ICMP sends a Destination Unreachable-Port Unreachable message back to the sender. 
If no application exists for the TCP port number, TCP sends a Connection Reset segment back to the sender. 

#There are two ways of maintaining routing table entries on IP routers:

•Manually—Static IP routers have routing tables that do not change unless manually changed by a network administrator. 

•Automatically—Dynamic IP routers have routing tables that change automatically based on the communication of routing information with other routers. 
Dynamic routing employs the use of routing protocols, such as Routing Information Protocol (RIP) and Open Shortest Path First (OSPF), 
to dynamically update the routing table through the exchange of routing information between routers. Remote network IDs are discovered by dynamic routers and automatically entered into the routing table


###CSV handling - std module 
#csv.reader(csvfile, dialect='excel', **fmtparams)
#Each row read from the csv file is returned as a list of strings. 
#No automatic data type conversion is performed unless the QUOTE_NONNUMERIC format option is specified 
#(in which case unquoted fields are transformed into floats).

>>> import csv
>>> with open('eggs.csv', newline='') as csvfile:
...     spamreader = csv.reader(csvfile, delimiter=' ', quotechar='|')
...     for row in spamreader:
...         print(', '.join(row))
Spam, Spam, Spam, Spam, Spam, Baked Beans
Spam, Lovely Spam, Wonderful Spam

#reders attributes 
csvreader.line_num
csvreader.fieldnames

#csv.writer(csvfile, dialect='excel', **fmtparams)

import csv
with open('eggs.csv', 'w', newline='') as csvfile:
    spamwriter = csv.writer(csvfile, delimiter=' ',
                            quotechar='|', quoting=csv.QUOTE_MINIMAL)
    spamwriter.writerow(['Spam'] * 5 + ['Baked Beans'])
    spamwriter.writerow(['Spam', 'Lovely Spam', 'Wonderful Spam'])

#writer attributes 
csvwriter.writerow(row)
csvwriter.writerows(rows)
DictWriter.writeheader()
    
    
    
#To know dialect 
csv.get_dialect(name)
csv.list_dialects()


#csv.DictReader(f, fieldnames=None, restkey=None, restval=None, dialect='excel', *args, **kwds)
#maps the information in each row to an OrderedDict whose keys are given by the optional fieldnames parameter.
#fieldnames is sequence of fieldName, for more columns than in fieldnames are put into a list 
#If fieldnames is omitted, the values in the first row of file f will be used as the fieldnames


>>> import csv
>>> with open('names.csv') as csvfile:
...     reader = csv.DictReader(csvfile)
...     for row in reader:
...         print(row['first_name'], row['last_name'])
...
Eric Idle
John Cleese

>>> print(row)
OrderedDict([('first_name', 'John'), ('last_name', 'Cleese')])



#class csv.DictWriter(f, fieldnames, restval='', extrasaction='raise', dialect='excel', *args, **kwds)
#If extrasaction is set to 'raise', the default value, a ValueError is raised. 
#If extrasaction is set to 'ignore', extra values in the dictionary are ignored.

import csv

with open('names.csv', 'w') as csvfile:
    fieldnames = ['first_name', 'last_name']
    writer = csv.DictWriter(csvfile, fieldnames=fieldnames)

    writer.writeheader()
    writer.writerow({'first_name': 'Baked', 'last_name': 'Beans'})
    writer.writerow({'first_name': 'Lovely', 'last_name': 'Spam'})
    writer.writerow({'first_name': 'Wonderful', 'last_name': 'Spam'})



##Quoting fields 
csv.QUOTE_ALL           Instructs writer objects to quote all fields.
csv.QUOTE_MINIMAL       Instructs writer objects to only quote those fields which contain special characters such as delimiter, quotechar or any of the characters in lineterminator.
csv.QUOTE_NONNUMERIC    Instructs writer objects to quote all non-numeric fields.
                        Instructs the reader to convert all non-quoted fields to type float.
csv.QUOTE_NONE




# specific formatting parameters are grouped together into dialects
#A dialect is a subclass of the Dialect class having a set of specific methods and a single validate() method. 
#When creating reader or writer objects, specify a string or a subclass of the Dialect class 
#OR specify individual formatting parameters, which have the same names as the attributes defined below for the Dialect class.

Dialect.delimiter
Dialect.doublequote
Dialect.escapechar
Dialect.lineterminator
Dialect.quotechar
Dialect.quoting
Dialect.skipinitialspace
Dialect.strict

#Registering a new dialect:


import csv
csv.register_dialect('unixpwd', delimiter=':', quoting=csv.QUOTE_NONE)
with open('passwd', newline='') as f:
    reader = csv.reader(f, 'unixpwd')


#catching and reporting errors:


import csv, sys
filename = 'some.csv'
with open(filename, newline='') as f:
    reader = csv.reader(f)
    try:
        for row in reader:
            print(row)
    except csv.Error as e:
        sys.exit('file {}, line {}: {}'.format(filename, reader.line_num, e))





###pywin32
# download MSIinstaller from https://sourceforge.net/projects/pywin32/files/pywin32/
#Check what win32API is available - http://timgolden.me.uk/pywin32-docs/contents.html

#check Python35\Lib\site-packages\win32\test  for usage pattern 
# or check Python35\Lib\site-packages\win32\Demos

#example 
import win32file, win32api
import win32api, win32process, win32security


#for COM related , check Python35\Lib\site-packages\win32com\demos
#For shell related check Python35\Lib\site-packages\win32comext\shell\demos
#for many other COM extention, check demos under Python35\Lib\site-packages\win32comext

#Example file Demo 
import win32file, win32api, win32con
import os

testName = os.path.join( win32api.GetTempPath(), "win32file_demo_test_file")
if os.path.exists(testName): os.unlink(testName)
# Open the file for writing.
handle = win32file.CreateFile(testName, win32file.GENERIC_WRITE, 0, None, win32con.CREATE_NEW, 0, None)
test_data = "Hello\0there".encode("ascii")
win32file.WriteFile(handle, test_data)
handle.Close()
# Open it for reading.
handle = win32file.CreateFile(testName, win32file.GENERIC_READ, 0, None, win32con.OPEN_EXISTING, 0, None)
rc, data = win32file.ReadFile(handle, 1024)
handle.Close()
if data == test_data:
    print("Successfully wrote and read a file")
else:
    raise Exception("Got different data back???")
os.unlink(testName)



###openpyxl - A Python library to read/write Excel 2010 xlsx/xlsm files

#install 
pip install openpyxl pillow
pip install pillow  #for image manipulation 

#few Helper methods 

openpyxl.utils.cell.absolute_coordinate(coord_string)
Convert a coordinate to an absolute coordinate string (B12 -> $B$12)

openpyxl.utils.cell.cols_from_range(range_string)
Get individual addresses for every cell in a range. Yields one row at a time.

openpyxl.utils.cell.column_index_from_string(str_col)
Convert a column name into a numerical index (‘A’ -> 1)

openpyxl.utils.cell.coordinate_from_string(coord_string)
Convert a coordinate string like ‘B12’ to a tuple (‘B’, 12)

openpyxl.utils.cell.coordinate_to_tuple(coordinate)
Convert an Excel style coordinate to (row, colum) tuple

openpyxl.utils.cell.get_column_interval(start, end)
Given the start and end colums, return all the columns in the series.
The start and end columns can be either column letters or 1-based indexes.

openpyxl.utils.cell.get_column_letter(idx)
Convert a column index into a column letter (3 -> ‘C’)

openpyxl.utils.cell.quote_sheetname(sheetname)
Add quotes around sheetnames if they contain spaces.

openpyxl.utils.cell.range_boundaries(range_string)
Convert a range string into a tuple of boundaries: (min_col, min_row, max_col, max_row) Cell coordinates will be converted into a range with the cell at both end

openpyxl.utils.cell.range_to_tuple(range_string)
Convert a worksheet range to the sheetname and maximum and minimum coordinate indices

openpyxl.utils.cell.rows_from_range(range_string)
Get individual addresses for every cell in a range. Yields one row at a time.


openpyxl.utils.datetime.W3CDTF_to_datetime(formatted_string)[source]
Convert from a timestamp string to a datetime object.

openpyxl.utils.datetime.datetime_to_W3CDTF(dt)[source]
Convert from a datetime to a timestamp string.

openpyxl.utils.datetime.days_to_time(value)[source]

openpyxl.utils.datetime.from_excel(value, offset=2415018.5)[source]

openpyxl.utils.datetime.time_to_days(value)[source]
Convert a time value to fractions of day

openpyxl.utils.datetime.timedelta_to_days(value)[source]
Convert a timedelta value to fractions of a day

openpyxl.utils.datetime.to_excel(dt, offset=2415018.5)[source]


openpyxl.utils.escape.escape(value)[source]
Convert ASCII < 31 to OOXML: n == _x + hex(ord(n)) +_

openpyxl.utils.escape.unescape(value)[source]
Convert escaped strings to ASCIII: _x000a_ == n

##Write a workbook

>>> from openpyxl import Workbook
>>> from openpyxl.compat import range
>>> from openpyxl.utils import get_column_letter
>>>
>>> wb = Workbook()
>>>
>>> dest_filename = 'empty_book.xlsx'
>>>
>>> ws1 = wb.active                #Get default active sheet 
>>> ws1.title = "range names"      #name the sheet 
>>>
>>> for row in range(1, 40):
        ws1.append(range(600))
>>>
>>> ws2 = wb.create_sheet(title="Pi")
>>>
>>> ws2['F5'] = 3.14
>>>
>>> ws3 = wb.create_sheet(title="Data")
>>> for row in range(10, 20):
        for col in range(27, 54):
            _ = ws3.cell(column=col, row=row, value="{0}".format(get_column_letter(col)))
>>> print(ws3['AA10'].value)
AA
>>> wb.save(filename = dest_filename)


#few more exmaples 
>>> ws1 = wb.create_sheet("Mysheet") # insert at the end (default)
# or
>>> ws2 = wb.create_sheet("Mysheet", 0) # insert at first position

ws.title = "New Title"  #give ws a name 

#update tabColor properties 
ws.sheet_properties.tabColor = "1072BA"

#Once you gave a worksheet a name, you can get it as a key of the workbook:
>>> ws3 = wb["New Title"]

#iterate 
>>> for sheet in wb:
...     print(sheet.title)

#create copies of worksheets within a single workbook:
>>> source = wb.active
>>> target = wb.copy_worksheet(source)

#Cells can be accessed directly as keys of the worksheet
>>> c = ws['A4']
>>> ws['A4'] = 4

#access to cells using row and column notation:
>>> d = ws.cell(row=4, column=2, value=10)

#Ranges of cells can be accessed using slicing
>>> cell_range = ws['A1':'C2']

#Ranges of rows or columns can be obtained similarly:
>>> colC = ws['C']
>>> col_range = ws['C:D']
>>> row10 = ws[10]
>>> row_range = ws[5:10]

>>> for row in ws.iter_rows(min_row=1, max_col=3, max_row=2):
        for cell in row:
            print(cell)
<Cell Sheet1.A1>
<Cell Sheet1.B1>
<Cell Sheet1.C1>
<Cell Sheet1.A2>
<Cell Sheet1.B2>
<Cell Sheet1.C2>

>>> for col in ws.iter_cols(min_row=1, max_col=3, max_row=2):
...     for cell in col:
...         print(cell)
<Cell Sheet1.A1>
<Cell Sheet1.A2>
<Cell Sheet1.B1>
<Cell Sheet1.B2>
<Cell Sheet1.C1>
<Cell Sheet1.C2>

#to iterate through all the rows or columns of a file
>>> ws = wb.active
>>> ws['C9'] = 'hello world'
>>> tuple(ws.rows)
((<Cell Sheet.A1>, <Cell Sheet.B1>, <Cell Sheet.C1>),
(<Cell Sheet.A2>, <Cell Sheet.B2>, <Cell Sheet.C2>),
(<Cell Sheet.A3>, <Cell Sheet.B3>, <Cell Sheet.C3>),
(<Cell Sheet.A4>, <Cell Sheet.B4>, <Cell Sheet.C4>),
(<Cell Sheet.A5>, <Cell Sheet.B5>, <Cell Sheet.C5>),
(<Cell Sheet.A6>, <Cell Sheet.B6>, <Cell Sheet.C6>),
(<Cell Sheet.A7>, <Cell Sheet.B7>, <Cell Sheet.C7>),
(<Cell Sheet.A8>, <Cell Sheet.B8>, <Cell Sheet.C8>),
(<Cell Sheet.A9>, <Cell Sheet.B9>, <Cell Sheet.C9>))

#or the openpyxl.worksheet.Worksheet.columns property:

>>> tuple(ws.columns)
((<Cell Sheet.A1>,
<Cell Sheet.A2>,
<Cell Sheet.A3>,
<Cell Sheet.A4>,
<Cell Sheet.A5>,
<Cell Sheet.A6>,
...
<Cell Sheet.B7>,
<Cell Sheet.B8>,
<Cell Sheet.B9>),
(<Cell Sheet.C1>,
<Cell Sheet.C2>,
<Cell Sheet.C3>,
<Cell Sheet.C4>,
<Cell Sheet.C5>,
<Cell Sheet.C6>,
<Cell Sheet.C7>,
<Cell Sheet.C8>,
<Cell Sheet.C9>))

#Once we have a openpyxl.cell.Cell, we can assign it a value:

>>> c.value = 'hello, world'
>>> print(c.value)
'hello, world'

>>> d.value = 3.14
>>> print(d.value)
3.14

#enable type and format inference:
>>> wb = Workbook(guess_types=True)
>>> c.value = '12%'
>>> print(c.value)
0.12

>>> import datetime
>>> d.value = datetime.datetime.now()
>>> print d.value
datetime.datetime(2010, 9, 10, 22, 25, 18)

>>> c.value = '31.50'
>>> print(c.value)
31.5



##Read an existing workbook
#guess_types        enable or disable (default) type inference when reading cells.
data_only           controls whether cells with formulae have either the formula (default) or the value stored the last time Excel read the sheet.
keep_vba            controls whether any Visual Basic elements are preserved or not (default). If they are preserved they are still not editable.
#openpyxl does currently not read all possible items in an Excel file 
#so images and charts will be lost from existing files if they are opened and saved with the same name.


>>> from openpyxl import load_workbook
>>> wb = load_workbook(filename = 'empty_book.xlsx')
>>> sheet_ranges = wb['range names']
>>> print(sheet_ranges['D18'].value)
3
>>> print wb.get_sheet_names()

#Using number formats

>>> import datetime
>>> from openpyxl import Workbook
>>> wb = Workbook()
>>> ws = wb.active
>>> # set date using a Python datetime
>>> ws['A1'] = datetime.datetime(2010, 7, 21)
>>>
>>> ws['A1'].number_format
'yyyy-mm-dd h:mm:ss'
>>> # You can enable type inference on a case-by-case basis
>>> wb.guess_types = True
>>> # set percentage using a string followed by the percent sign
>>> ws['B1'] = '3.14%'
>>> wb.guess_types = False
>>> ws['B1'].value
0.031400000000000004
>>>
>>> ws['B1'].number_format
'0%'

#Using formulae
#If you’re trying to use a formula that isn’t known 
#this could be because you’re using a formula that was not included in the initial specification. 
#Such formulae must be prefixed with xlfn. to work.

>>> from openpyxl import Workbook
>>> wb = Workbook()
>>> ws = wb.active
>>> # add a simple formula
>>> ws["A1"] = "=SUM(1, 1)"
>>> wb.save("formula.xlsx")

#openpyxl never evaluates formula but it is possible to check the name of a formula:

>>> from openpyxl.utils import FORMULAE
>>> "HEX2DEC" in FORMULAE
True



#Merge / Unmerge cells
#When you merge cells all cells but the top-left one are removed from the worksheet.


>>> from openpyxl.workbook import Workbook
>>>
>>> wb = Workbook()
>>> ws = wb.active
>>>
>>> ws.merge_cells('A1:B1')
>>> ws.unmerge_cells('A1:B1')
>>>
>>> # or
>>> ws.merge_cells(start_row=2,start_column=1,end_row=2,end_column=4)
>>> ws.unmerge_cells(start_row=2,start_column=1,end_row=2,end_column=4)

#Inserting an image
>>> from openpyxl import Workbook
>>> from openpyxl.drawing.image import Image
>>>
>>> wb = Workbook()
>>> ws = wb.active
>>> ws['A1'] = 'You should see three logos below'

>>> # create an image
>>> img = Image('logo.png')

>>> # add to worksheet and anchor next to cells
>>> ws.add_image(img, 'A1')
>>> wb.save('logo.xlsx')

#Fold columns (outline)
>>> import openpyxl
>>> wb = openpyxl.Workbook()
>>> ws = wb.create_sheet()
>>> ws.column_dimensions.group('A','D', hidden=True)
>>> wb.save('group.xlsx')


#With Pandas and NumPy
#openpyxl has builtin support for the NumPy types float, integer and boolean. 
#DateTimes are supported using the Pandas’ Timestamp type.
#Note pandas support excel reading

from openpyxl.utils.dataframe import dataframe_to_rows
wb = Workbook()
ws = wb.active

for r in dataframe_to_rows(df, index=True, header=True):
    ws.append(r)

#To convert a dataframe into a worksheet highlighting the header and index:

wb = Workbook()
ws = wb.active

for r in dataframe_to_rows(df, index=True, header=True):
    ws.append(r)

for cell in ws['A'] + ws[1]:
    cell.style = 'Pandas'

wb.save("pandas_openpyxl.xlsx")

#OR you just want to convert the data you can use write-only mode:

from openpyxl.cell.cell import WriteOnlyCell
wb = Workbook(write_only=True)
ws = wb.create_sheet()

cell = WriteOnlyCell(ws)
cell.style = 'Pandas'

def format_first_row(row, cell):
    for c in row:
        cell.value = c
        yield cell

rows = dataframe_to_rows(df)
first_row = format_first_row(next(rows), cell)
ws.append(first_row)

for row in rows:
    row = list(row)
    cell.value = row[0]
    row[0] = cell
    ws.append(row)

wb.save("openpyxl_stream.xlsx")


##Converting a worksheet to a Dataframe
# if the worksheet has no headers or indices:

df = pandas.DataFrame(ws.values)

#with headers or indices
data = ws.values
cols = next(data)[1:]
data = list(data)
idx = [r[0] for r in data]
data = (islice(r, 1, None) for r in data)
df = DataFrame(data, index=idx, columns=cols)

##Creating a chart - https://openpyxl.readthedocs.io/en/default/charts/introduction.html
#Charts are composed of at least one series of one or more data points. 
#Series themselves are comprised of references to cell ranges.

from openpyxl import Workbook
wb = Workbook()
ws = wb.active
for i in range(10):
    ws.append([i])    #append each number to new row 

from openpyxl.chart import BarChart, Reference, Series
values = Reference(ws, min_col=1, min_row=1, max_col=1, max_row=10) #give range of data 
chart = BarChart()
chart.add_data(values)
ws.add_chart(chart, "E15")
wb.save("SampleChart.xlsx")


##Line Charts
#Similar to bar charts there are three kinds of line charts: standard, stacked and percentStacked.

from datetime import date

from openpyxl import Workbook
from openpyxl.chart import   LineChart,  Reference
from openpyxl.chart.axis import DateAxis

wb = Workbook()
ws = wb.active

rows = [
    ['Date', 'Batch 1', 'Batch 2', 'Batch 3'],
    [date(2015,9, 1), 40, 30, 25],
    [date(2015,9, 2), 40, 25, 30],
    [date(2015,9, 3), 50, 30, 45],
    [date(2015,9, 4), 30, 25, 40],
    [date(2015,9, 5), 25, 35, 30],
    [date(2015,9, 6), 20, 40, 35],
]

for row in rows:
    ws.append(row)

c1 = LineChart()
c1.title = "Line Chart"
c1.style = 13
c1.y_axis.title = 'Size'
c1.x_axis.title = 'Test Number'

data = Reference(ws, min_col=2, min_row=1, max_col=4, max_row=7)
c1.add_data(data, titles_from_data=True)

# Style the lines
s1 = c1.series[0]
s1.marker.symbol = "triangle"
s1.marker.graphicalProperties.solidFill = "FF0000" # Marker filling
s1.marker.graphicalProperties.line.solidFill = "FF0000" # Marker outline

s1.graphicalProperties.line.noFill = True

s2 = c1.series[1]
s2.graphicalProperties.line.solidFill = "00AAAA"
s2.graphicalProperties.line.dashStyle = "sysDot"
s2.graphicalProperties.line.width = 100050 # width in EMUs

s2 = c1.series[2]
s2.smooth = True # Make the line smooth

ws.add_chart(c1, "A10")
#Other type of Line charts 
from copy import deepcopy
stacked = deepcopy(c1)
stacked.grouping = "stacked"
stacked.title = "Stacked Line Chart"
ws.add_chart(stacked, "A27")

percent_stacked = deepcopy(c1)
percent_stacked.grouping = "percentStacked"
percent_stacked.title = "Percent Stacked Line Chart"
ws.add_chart(percent_stacked, "A44")


# Chart with date axis
c2 = LineChart()
c2.title = "Date Axis"
c2.style = 12
c2.y_axis.title = "Size"
c2.y_axis.crossAx = 500
c2.x_axis = DateAxis(crossAx=100)
c2.x_axis.number_format = 'd-mmm'
c2.x_axis.majorTimeUnit = "days"
c2.x_axis.title = "Date"

c2.add_data(data, titles_from_data=True)
dates = Reference(ws, min_col=1, min_row=2, max_row=7)
c2.set_categories(dates)

ws.add_chart(c2, "A61")

wb.save("line.xlsx")

##Pie Charts

from openpyxl import Workbook

from openpyxl.chart import (
    PieChart,
    ProjectedPieChart,
    Reference
)
from openpyxl.chart.series import DataPoint

data = [
    ['Pie', 'Sold'],
    ['Apple', 50],
    ['Cherry', 30],
    ['Pumpkin', 10],
    ['Chocolate', 40],
]

wb = Workbook()
ws = wb.active

for row in data:
    ws.append(row)

pie = PieChart()
labels = Reference(ws, min_col=1, min_row=2, max_row=5)
data = Reference(ws, min_col=2, min_row=1, max_row=5)
pie.add_data(data, titles_from_data=True)
pie.set_categories(labels)
pie.title = "Pies sold by category"

# Cut the first slice out of the pie
slice = DataPoint(idx=0, explosion=20)
pie.series[0].data_points = [slice]

ws.add_chart(pie, "D1")


ws = wb.create_sheet(title="Projection")

data = [
    ['Page', 'Views'],
    ['Search', 95],
    ['Products', 4],
    ['Offers', 0.5],
    ['Sales', 0.5],
]

for row in data:
    ws.append(row)

projected_pie = ProjectedPieChart()
projected_pie.type = "pie"
projected_pie.splitType = "val" # split by value
labels = Reference(ws, min_col=1, min_row=2, max_row=5)
data = Reference(ws, min_col=2, min_row=1, max_row=5)
projected_pie.add_data(data, titles_from_data=True)
projected_pie.set_categories(labels)

ws.add_chart(projected_pie, "A10")

from copy import deepcopy
projected_bar = deepcopy(projected_pie)
projected_bar.type = "bar"
projected_bar.splitType = 'pos' # split by position

ws.add_chart(projected_bar, "A27")

wb.save("pie.xlsx")


##Vertical, Horizontal and Stacked Bar Charts
#Switch between vertical and horizontal bar charts by setting type to col or bar respectively.
#When using stacked charts the overlap needs to be set to 100.
#If bars are horizontal, x and y axes are revesed.


from openpyxl import Workbook
from openpyxl.chart import BarChart, Series, Reference

wb = Workbook(write_only=True)
ws = wb.create_sheet()

rows = [
    ('Number', 'Batch 1', 'Batch 2'),
    (2, 10, 30),
    (3, 40, 60),
    (4, 50, 70),
    (5, 20, 10),
    (6, 10, 40),
    (7, 50, 30),
]


for row in rows:
    ws.append(row)


chart1 = BarChart()
chart1.type = "col"
chart1.style = 10
chart1.title = "Bar Chart"
chart1.y_axis.title = 'Test number'
chart1.x_axis.title = 'Sample length (mm)'

data = Reference(ws, min_col=2, min_row=1, max_row=7, max_col=3)
cats = Reference(ws, min_col=1, min_row=2, max_row=7)
chart1.add_data(data, titles_from_data=True)
chart1.set_categories(cats)
chart1.shape = 4
ws.add_chart(chart1, "A10")

from copy import deepcopy

#Other bar chart 
chart2 = deepcopy(chart1)
chart2.style = 11
chart2.type = "bar"
chart2.title = "Horizontal Bar Chart"

ws.add_chart(chart2, "G10")


chart3 = deepcopy(chart1)
chart3.type = "col"
chart3.style = 12
chart3.grouping = "stacked"
chart3.overlap = 100
chart3.title = 'Stacked Chart'

ws.add_chart(chart3, "A27")


chart4 = deepcopy(chart1)
chart4.type = "bar"
chart4.style = 13
chart4.grouping = "percentStacked"
chart4.overlap = 100
chart4.title = 'Percent Stacked Chart'

ws.add_chart(chart4, "G27")

wb.save("bar.xlsx")

##3D Bar Charts- Many charts can be made 3D, check Reference 

from openpyxl import Workbook
from openpyxl.chart import (
    Reference,
    Series,
    BarChart3D,
)

wb = Workbook()
ws = wb.active

rows = [
    (None, 2013, 2014),
    ("Apples", 5, 4),
    ("Oranges", 6, 2),
    ("Pears", 8, 3)
]

for row in rows:
    ws.append(row)

data = Reference(ws, min_col=2, min_row=1, max_col=3, max_row=4)
titles = Reference(ws, min_col=1, min_row=2, max_row=4)
chart = BarChart3D()
chart.title = "3D Bar Chart"
chart.add_data(data=data, titles_from_data=True)
chart.set_categories(titles)

ws.add_chart(chart, "E5")
wb.save("bar3d.xlsx")


##Parsing Formulas
#openpyxl supports limited parsing of formulas embedded in cells. 
#The openpyxl.formula package contains a Tokenizer class 
#to break formulas into their consitutuent tokens. 

>>> from openpyxl.formula import Tokenizer
>>> tok = Tokenizer("""=IF($A$1,"then True",MAX(DEFAULT_VAL,'Sheet 2'!B1))""")
>>> print("\n".join("%12s%11s%9s" % (t.value, t.type, t.subtype) for t in tok.items))
         IF(       FUNC     OPEN
        $A$1    OPERAND    RANGE
           ,        SEP      ARG
 "then True"    OPERAND     TEXT
           ,        SEP      ARG
        MAX(       FUNC     OPEN
 DEFAULT_VAL    OPERAND    RANGE
           ,        SEP      ARG
'Sheet 2'!B1    OPERAND    RANGE
           )       FUNC    CLOSE
           )       FUNC    CLOSE

 
##Using filters and sorts
#To add a filter you define a range and then add columns and sort conditions:

from openpyxl import Workbook

wb = Workbook()
ws = wb.active

data = [
    ["Fruit", "Quantity"],
    ["Kiwi", 3],
    ["Grape", 15],
    ["Apple", 3],
    ["Peach", 3],
    ["Pomegranate", 3],
    ["Pear", 3],
    ["Tangerine", 3],
    ["Blueberry", 3],
    ["Mango", 3],
    ["Watermelon", 3],
    ["Blackberry", 3],
    ["Orange", 3],
    ["Raspberry", 3],
    ["Banana", 3]
]

for r in data:
    ws.append(r)

ws.auto_filter.ref = "A1:B15"
ws.auto_filter.add_filter_column(0, ["Kiwi", "Apple", "Mango"])
ws.auto_filter.add_sort_condition("B2:B15")

wb.save("filtered.xlsx")

##Worksheet Tables
#Worksheet tables are references to groups of cells.
#By default tables are created with a header from the first row and filters for all the columns.

from openpyxl import Workbook
from openpyxl.worksheet.table import Table, TableStyleInfo

wb = Workbook()
ws = wb.active

data = [
    ['Apples', 10000, 5000, 8000, 6000],
    ['Pears',   2000, 3000, 4000, 5000],
    ['Bananas', 6000, 6000, 6500, 6000],
    ['Oranges',  500,  300,  200,  700],
]

# add column headings. NB. these must be strings
ws.append(["Fruit", "2011", "2012", "2013", "2014"])
for row in data:
    ws.append(row)

tab = Table(displayName="Table1", ref="A1:E5")

# Add a default style with striped rows and banded columns
style = TableStyleInfo(name="TableStyleMedium9", showFirstColumn=False,
                       showLastColumn=False, showRowStripes=True, showColumnStripes=True)
tab.tableStyleInfo = style
ws.add_table(tab)
wb.save("table.xlsx")

##Cell Styles
#Cell styles are shared between objects and once they have been assigned they cannot be changed. 
#This stops unwanted side-effects such as changing the style for lots of cells when instead of only one.

>>> from openpyxl.styles import colors
>>> from openpyxl.styles import Font, Color
>>> from openpyxl import Workbook
>>> wb = Workbook()
>>> ws = wb.active
>>>
>>> a1 = ws['A1']
>>> d4 = ws['D4']
>>> ft = Font(color=colors.RED)
>>> a1.font = ft
>>> d4.font = ft
>>>
>>> a1.font.italic = True # is not allowed 
>>>
>>> # If you want to change the color of a Font, you need to reassign it::
>>>
>>> a1.font = Font(color=colors.RED, italic=True) # the change only affects A1

#Copying styles

>>> from openpyxl.styles import Font
>>> from copy import copy
>>>
>>> ft1 = Font(name='Arial', size=14)
>>> ft2 = copy(ft1)
>>> ft2.name = "Tahoma"
>>> ft1.name
'Arial'
>>> ft2.name
'Tahoma'
>>> ft2.size # copied from the
14.0

#Basic Font Colors
#Colors are usually RGB or aRGB hexvalues. 
#The colors module contains some handy constants

>>> from openpyxl.styles import Font
>>> from openpyxl.styles.colors import RED
>>> font = Font(color=RED)
>>> font = Font(color="FFBB00")

#Applying Styles

>>> from openpyxl.workbook import Workbook
>>> from openpyxl.styles import Font, Fill
>>> wb = Workbook()
>>> ws = wb.active
>>> c = ws['A1']
>>> c.font = Font(size=12)

#Styles can also applied to columns and rows 
#but note that this applies only to cells created (in Excel) after the file is closed. 

#OR If you want to apply styles to entire rows and columns 
#then you must apply the style to each cell yourself

>>> col = ws.column_dimensions['A']
>>> col.font = Font(bold=True)
>>> row = ws.row_dimensions[1]
>>> row.font = Font(underline="single")

##Conditional Formatting

>>> from openpyxl.formatting import Rule
>>> from openpyxl.styles import Font, PatternFill, Border
>>> from openpyxl.styles.differential import DifferentialStyle
>>> dxf = DifferentialStyle(font=Font(bold=True), fill=PatternFill(start_color='EE1111', end_color='EE1111'))
>>> rule = Rule(type='cellIs', dxf=dxf, formula=["10"])

#The builtins conditional formats are:
        ColorScale
        IconSet
        DataBar
#Builtin formats contain a sequence of formatting settings which combine a type with an integer for comparison. 
#Possible types are: ‘num’, ‘percent’, ‘max’, ‘min’, ‘formula’, ‘percentile’.
#Example - ColorScale
#2 color scales produce a gradient from one color to another; 
#3 color scales use an additional color for 2 gradients.


>>> from openpyxl.formatting.rule import ColorScale, FormatObject
>>> from openpyxl.styles import Color
>>> first = FormatObject(type='min')
>>> last = FormatObject(type='max')
>>> # colors match the format objects:
>>> colors = [Color('AA0000'), Color('00AA00')]
>>> cs2 = ColorScale(cfvo=[first, last], color=colors)
>>> # a three color scale would extend the sequences
>>> mid = FormatObject(type='num', val=40)
>>> colors.insert(1, Color('00AA00'))
>>> cs3 = ColorScale(cfvo=[first, mid, last], color=colors)
>>> # create a rule with the color scale
>>> from openpyxl.formatting.rule import Rule
>>> rule = Rule(type='colorScale', colorScale=cs3)

#There is a convenience function for creating ColorScale rules

>>> from openpyxl.formatting.rule import ColorScaleRule
>>> rule = ColorScaleRule(start_type='percentile', start_value=10, start_color='FFAA0000',
...                       mid_type='percentile', mid_value=50, mid_color='FF0000AA',
...                       end_type='percentile', end_value=90, end_color='FF00AA00')


##Standard conditional formats
        Average
        Percent
        Unique or duplicate
        Value
        Rank

>>> from openpyxl import Workbook
>>> from openpyxl.styles import Color, PatternFill, Font, Border
>>> from openpyxl.styles.differential import DifferentialStyle
>>> from openpyxl.formatting.rule import ColorScaleRule, CellIsRule, FormulaRule
>>>
>>> wb = Workbook()
>>> ws = wb.active
>>>
>>> # Create fill
>>> redFill = PatternFill(start_color='EE1111',
...                end_color='EE1111',
...                fill_type='solid')
>>>
>>> # Add a two-color scale
>>> # Takes colors in excel 'RRGGBB' style.
>>> ws.conditional_formatting.add('A1:A10',
...             ColorScaleRule(start_type='min', start_color='AA0000',
...                           end_type='max', end_color='00AA00')
...                           )
>>>
>>> # Add a three-color scale
>>> ws.conditional_formatting.add('B1:B10',
...                ColorScaleRule(start_type='percentile', start_value=10, start_color='AA0000',
...                            mid_type='percentile', mid_value=50, mid_color='0000AA',
...                            end_type='percentile', end_value=90, end_color='00AA00')
...                              )
>>>
>>> # Add a conditional formatting based on a cell comparison
>>> # addCellIs(range_string, operator, formula, stopIfTrue, wb, font, border, fill)
>>> # Format if cell is less than 'formula'
>>> ws.conditional_formatting.add('C2:C10',
...             CellIsRule(operator='lessThan', formula=['C$1'], stopIfTrue=True, fill=redFill))
>>>
>>> # Format if cell is between 'formula'
>>> ws.conditional_formatting.add('D2:D10',
...             CellIsRule(operator='between', formula=['1','5'], stopIfTrue=True, fill=redFill))
>>>
>>> # Format using a formula
>>> ws.conditional_formatting.add('E1:E10',
...             FormulaRule(formula=['ISBLANK(E1)'], stopIfTrue=True, fill=redFill))
>>>
>>> # Aside from the 2-color and 3-color scales, format rules take fonts, borders and fills for styling:
>>> myFont = Font()
>>> myBorder = Border()
>>> ws.conditional_formatting.add('E1:E10',
...             FormulaRule(formula=['E1=0'], font=myFont, border=myBorder, fill=redFill))
>>>
>>> # Highlight cells that contain particular text by using a special formula
>>> red_text = Font(color="9C0006")
>>> red_fill = PatternFill(bgColor="FFC7CE")
>>> dxf = DifferentialStyle(font=red_text, fill=red_fill)
>>> rule = Rule(type="containsText", operator="containsText", text="highlight", dxf=dxf)
>>> rule.formula = ['NOT(ISERROR(SEARCH("highlight",A1)))']
>>> ws.conditional_formatting.add('A1:F40', rule)
>>> wb.save("test.xlsx")








###python-pptx
#pip install python-pptx


##Hello World! example


from pptx import Presentation

prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Hello, World!"
subtitle.text = "python-pptx was here!"

prs.save('test.pptx')

##Bullet slide example
#Not all shapes can contain text, but those that do always have at least one paragraph,

from pptx import Presentation

prs = Presentation()
bullet_slide_layout = prs.slide_layouts[1]

slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes

title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = 'Adding a Bullet Slide'

tf = body_shape.text_frame
tf.text = 'Find the bullet slide layout'

p = tf.add_paragraph()
p.text = 'Use _TextFrame.text for first bullet'
p.level = 1

p = tf.add_paragraph()
p.text = 'Use _TextFrame.add_paragraph() for subsequent bullets'
p.level = 2

prs.save('test.pptx')

 even if that paragraph is empty and no text is visible within the shape. _BaseShape.has_text_frame can be used to determine whether a shape can contain text. (All shapes subclass _BaseShape.) When _BaseShape.has_text_frame is True, _BaseShape.text_frame.paragraphs[0] returns the first paragraph. The text of the first paragraph can be set using text_frame.paragraphs[0].text. As a shortcut, the writable properties _BaseShape.text and _TextFrame.text are provided to accomplish the same thing. Note that these last two calls delete all the shape’s paragraphs except the first one before setting the text it contains.



## add_textbox() example

from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)

left = top = width = height = Inches(1)
txBox = slide.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame

tf.text = "This is text inside a textbox"

p = tf.add_paragraph()
p.text = "This is a second paragraph that's bold"
p.font.bold = True

p = tf.add_paragraph()
p.text = "This is a third paragraph that's big"
p.font.size = Pt(40)

prs.save('test.pptx')

##add_picture() example
from pptx import Presentation
from pptx.util import Inches

img_path = 'monty-truth.png'

prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)

left = top = Inches(1)
pic = slide.shapes.add_picture(img_path, left, top)

left = Inches(5)
height = Inches(5.5)
pic = slide.shapes.add_picture(img_path, left, top, height=height)

prs.save('test.pptx')

##add_shape() example
#http://python-pptx.readthedocs.io/en/latest/api/enum/MsoAutoShapeType.html#msoautoshapetype

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

prs = Presentation()
title_only_slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(title_only_slide_layout)
shapes = slide.shapes

shapes.title.text = 'Adding an AutoShape'

left = Inches(0.93)  # 0.93" centers this overall set of shapes
top = Inches(3.0)
width = Inches(1.75)
height = Inches(1.0)

shape = shapes.add_shape(MSO_SHAPE.PENTAGON, left, top, width, height)
shape.text = 'Step 1'

left = left + width - Inches(0.4)
width = Inches(2.0)  # chevrons need more width for visual balance

for n in range(2, 6):
    shape = shapes.add_shape(MSO_SHAPE.CHEVRON, left, top, width, height)
    shape.text = 'Step %d' % n
    left = left + width - Inches(0.4)

prs.save('test.pptx')

##add_table() example
from pptx import Presentation
from pptx.util import Inches

prs = Presentation()
title_only_slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(title_only_slide_layout)
shapes = slide.shapes

shapes.title.text = 'Adding a Table'

rows = cols = 2
left = top = Inches(2.0)
width = Inches(6.0)
height = Inches(0.8)

table = shapes.add_table(rows, cols, left, top, width, height).table

# set column widths
table.columns[0].width = Inches(2.0)
table.columns[1].width = Inches(4.0)

# write column headings
table.cell(0, 0).text = 'Foo'
table.cell(0, 1).text = 'Bar'

# write body cells
table.cell(1, 0).text = 'Baz'
table.cell(1, 1).text = 'Qux'

prs.save('test.pptx')

##Extract all text from slides in presentation

from pptx import Presentation

prs = Presentation(path_to_presentation)

# text_runs will be populated with a list of strings,
# one for each text run in presentation
text_runs = []

for slide in prs.slides:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                text_runs.append(run.text)


##Opening and saving a presentation
prs = Presentation('existing-prs-file.pptx')
prs.save('new-file-name.pptx')

##Slide layout basics
#A slide layout is like a template for a slide
#The presentation themes that come with PowerPoint have about nine slide layouts, 
#The slide layouts in a standard PowerPoint theme always occur in the same sequence. 
#In python-pptx, these are prs.slide_layouts[0] through prs.slide_layouts[8]. 
    Title (presentation title slide)
    Title and Content
    Section Header (sometimes called Segue)
    Two Content (side by side bullet textboxes)
    Comparison (same but additional title for each side by side content box)
    Title Only
    Blank
    Content with Caption
    Picture with Caption

#Each has zero or more placeholders (mostly not zero)
# - preformatted areas into which you can place a title, multi-level bullets, an image, etc. 


#Example                 
                
SLD_LAYOUT_TITLE_AND_CONTENT = 1

prs = Presentation()
slide_layout = prs.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT]
slide = prs.slides.add_slide(slide_layout)


##six different types of shapes that can be placed on a slide:
auto shape
    This is a regular shape, like a rectangle, an ellipse, or a block arrow. 
    An auto shape can have a fill and an outline, and can contain text
    A text box is also an autoshape, a rectangular one, just by default without a fill and without an outline.
picture
    A raster image, like a photograph or clip art is referred to as a picture in PowerPoint
graphic frame
    This is the technical name for the container that holds a table, a chart, a smart art diagram, or media clip. 
group shape
    In PowerPoint, a set of shapes can be grouped, allowing them to be selected, moved, resized, and even filled as a unit. 
line/connector
    Lines are different from auto shapes because, well, they’re linear.    
content part   
   It has something to do with embedding “foreign” XML like SVG in with the presentation. 
    
#In reality, use below (each belongs to one of above)
shape shapes – auto shapes with fill and an outline
text boxes – auto shapes with no fill and no outline
placeholders – auto shapes that can appear on a slide layout or master and be inherited on slides that use that layout, allowing content to be added that takes on the formatting of the placeholder
line/connector – as described above
picture – as described above
table – that row and column thing
chart – pie chart, line chart, etc. python-pptx doesn’t support creating these yet.
smart art – not supported yet, although preserved if present
media clip – not supported yet, although preserved if present
    
    
#Accessing the shapes on a slide       
shapes = slide.shapes
    
#Adding an auto shape
#Check  https://python-pptx.readthedocs.io/en/latest/api/enum/MsoAutoShapeType.html#msoautoshapetype
#for a list of all 182 auto shape types.

#For example - To add  a rounded rectangle shape, one inch square, 
#and positioned one inch from the top-left corner of the slide:

from pptx.enum.shapes import MSO_SHAPE

shapes = slide.shapes
left = top = width = height = Inches(1.0)
shape = shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
)


##Understanding English Metric Unit
    
>>> from pptx.util import Inches, Pt
>>> length = Inches(1)
>>> length
914400
>>> length.inches
1.0
>>> length.cm
2.54
>>> length.pt
72.0
>>> length = Pt(72)
>>> length
914400   
    

##Shape - position and dimensions
#All shapes have a position on their slide and have a size. 

>>> from pptx.enum.shapes import MSO_SHAPE
>>> left = top = width = height = Inches(1.0)
>>> shape = shapes.add_shape(
>>>     MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
>>> )
>>> shape.left, shape.top, shape.width, shape.height
(914400, 914400, 914400, 914400)
>>> shape.left.inches
1.0
>>> shape.left = Inches(2.0)
>>> shape.left.inches
2.0

##Shape - Fill
#AutoShapes have an outline around their outside edge. 
#What appears within that outline is called the shape’s fill.
#A shape may also be filled with a gradient, a picture, a pattern (like cross-hatching for example), or may have no fill (transparent).

#When a color is used, it may be specified as a specific RGB value or a color from the theme palette.
>>> fill = shape.fill
>>> fill.solid()
>>> fill.fore_color.rgb = RGBColor(255, 0, 0)

#below sets it to the theme color that appears as ‘Accent 1 - 25% Darker’ in the toolbar palette:

>>> from pptx.enum.dml import MSO_THEME_COLOR
>>> fill = shape.fill
>>> fill.solid()
>>> fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_1
>>> fill.fore_color.brightness = -0.25

#Below sets the shape fill to transparent, or ‘No Fill’ as it’s called in the PowerPoint UI
>>> shape.fill.background()



##Shape - Line - for outline
#setting its color, width, dash (solid, dashed, dotted, etc.), line style (single, double, thick-thin, etc.), end cap, join type, and others. 
#As of now, only color and width can be set using python-pptx:

>>> line = shape.line
>>> line.color.rgb = RGBColor(255, 0, 0)
>>> line.color.brightness = 0.5  # 50% lighter
>>> line.width = Pt(2.5)

#Theme colors can be used on lines too:

>>> line.color.theme_color = MSO_THEME_COLOR.ACCENT_6

#Shape.line has the attribute .color - a shortcut for - default solid line 
>>> line.fill.solid()
>>> line.fill.fore_color

#Accessing the fill directly is required, for example, to set the line to transparent:
>>> line.fill.background()


##Shape - Line width - read/write width property
>>> line.width
9525
>>> line.width.pt
0.75
>>> line.width = Pt(2.0)
>>> line.width.pt
2.0

##Shape - Adjusting an autoshape
#Many auto shapes have adjustments. 
#In PowerPoint, these show up as little yellow diamonds you can drag to change the look of the shape. 

#Adjustment values are large integers, each based on a nominal value of 100,000. 
#The effective value of an adjustment is proportional to the width or height of the shape. 
#So a value of 50,000 for an x-coordinate adjustment corresponds to half the width of the shape; 
#a value of 75,000 for a y-coordinate adjustment corresponds to 3/4 of the shape height.

#Adjustment values can be negative, 
#generally indicating the coordinate is to the left or above the top left corner (origin) of the shape. 

#Values can also be subject to limits, meaning their effective value cannot be outside a prescribed range. 
#In practice this corresponds to a point not being able to extend beyond the left side of the shape, 


#The following code formats a callout shape using its adjustments:

callout_sp = shapes.add_shape(
    MSO_SHAPE.LINE_CALLOUT_2_ACCENT_BAR, left, top, width, height
)

# get the callout line coming out of the right place
adjs = callout_sp.adjustments
adjs[0] = 0.5   # vert pos of junction in margin line, 0 is top
adjs[1] = 0.0   # horz pos of margin ln wrt shape width, 0 is left side
adjs[2] = 0.5   # vert pos of elbow wrt margin line, 0 is top
adjs[3] = -0.1  # horz pos of elbow wrt shape width, 0 is margin line
adjs[4] = 3.0   # vert pos of line end wrt shape height, 0 is top
a5 = adjs[3] - (adjs[4] - adjs[0]) * height/width
adjs[5] = a5    # horz pos of elbow wrt shape width, 0 is margin line

# rotate 45 degrees counter-clockwise
callout_sp.rotation = -45.0


##Understanding placeholders
#a placeholder is a pre-formatted container into which content can be placed
#A placeholder is a shape

#the auto shape (p:sp element), picture (p:pic element), and graphic frame (p:graphicFrame) shape types can be a placeholder. 
#The group shape (p:grpSp), connector (p:cxnSp), and content part (p:contentPart) shapes cannot be a placeholder. 

#A graphic frame placeholder can contain a table, a chart, or SmartArt.

#Placeholder types
#There are 18 types of placeholder.
Title, Center Title, Subtitle, Body
Content
Picture, Clip Art
Chart, Table, Smart Art
Media Clip
Date, Footer, Slide Number
Header
Vertical Body, Vertical Object, Vertical Title
Placholders inherit

#Placeholder behavior requires three different categories of placeholder shape; 
#those that exist on a slide master, those on a slide layout, 
#and those that ultimately appear on a slide in a presentation.

#These three categories of placeholder participate in a property inheritance hierarchy, 


##Access a placeholder - placeholders property or since placeholder is a Shape, by .shapes property
#to access a known placeholder is by its idx value. 
#The idx value of a placeholder is the integer key of the slide layout placeholder it inherits properties from. 
#As such, it remains stable throughout the life of the slide 
#and will be the same for any slide created using that layout.

>>> prs = Presentation()
>>> slide = prs.slides.add_slide(prs.slide_layouts[8])
>>> for shape in slide.placeholders:
...     print('%d %s' % (shape.placeholder_format.idx, shape.name))
...
0  Title 1
1  Picture Placeholder 2
2  Text Placeholder 3

#In general, the idx value of a placeholder from a built-in slide layout  will be between 0 and 5. 
#The title placeholder will always have idx 0 if present 
#and any other placeholders will follow in sequence, top to bottom and left to right. 
#A placeholder added to a slide layout by a user in PowerPoint will receive an idx value starting at 10.


#to access it directly:

>>> slide.placeholders[1]
<pptx.parts.slide.PicturePlaceholder object at 0x10d094590>
>>> slide.placeholders[2].name
'Text Placeholder 3'


#Identify and Characterize a placeholder

>>> prs = Presentation()
>>> slide = prs.slides.add_slide(prs.slide_layouts[8])
>>> for shape in slide.shapes:
...     print('%s' % shape.shape_type)
...
PLACEHOLDER (14)
PLACEHOLDER (14)
PLACEHOLDER (14)

#inspect the contents of the placeholder’s placeholder_format attribute. 
#All shapes have this attribute, but accessing it on a non-placeholder shape raises ValueError. 
#The is_placeholder attribute can be used to determine whether a shape is a placeholder:

>>> for shape in slide.shapes:
...     if shape.is_placeholder:
...         phf = shape.placeholder_format
...         print('%d, %s' % (phf.idx, phf.type))
...
0, TITLE (1)
1, PICTURE (18)
2, BODY (2)

##Insert content into a placeholder
#Note Text can be inserted into title and body placeholders by .text attribute 

##PicturePlaceholder.insert_picture()
#A reference to a picture placeholder becomes invalid after its insert_picture() method is called
>>> prs = Presentation()
>>> slide = prs.slides.add_slide(prs.slide_layouts[8])
>>> placeholder = slide.placeholders[1]  # idx key, not position
>>> placeholder.name
'Picture Placeholder 2'
>>> placeholder.placeholder_format.type
PICTURE (18)
>>> picture = placeholder.insert_picture('my-image.png') #returns new placegolder

#A picture inserted in this way is stretched proportionately and cropped to fill the entire placeholder
#Cropping can be adjusted using the crop properties on the placeholder, such as crop_bottom.


##TablePlaceholder.insert_table()
#The table placeholder has an insert_table() method. 

#The built-in template has no layout containing a table placeholder, 
#so this example assumes a starting presentation named having-table-placeholder.pptx 
#having a table placeholder with idx 10 on its second slide layout:

>>> prs = Presentation('having-table-placeholder.pptx')
>>> slide = prs.slides.add_slide(prs.slide_layouts[1])
>>> placeholder = slide.placeholders[10]  # idx key, not position
>>> placeholder.name
'Table Placeholder 1'
>>> placeholder.placeholder_format.type
TABLE (12)
>>> graphic_frame = placeholder.insert_table(rows=2, cols=2) #return new updated platholders
>>> table = graphic_frame.table
>>> len(table.rows), len(table.columns)
(2, 2)

##ChartPlaceholder.insert_chart()
#The chart placeholder has an insert_chart() method. 

#The presentation template built into python-pptx has no layout containing a chart placeholder, 
#so as per example above 

>>> from pptx.chart.data import ChartData
>>> from pptx.enum.chart import XL_CHART_TYPE

>>> prs = Presentation('having-chart-placeholder.pptx')
>>> slide = prs.slides.add_slide(prs.slide_layouts[1])

>>> placeholder = slide.placeholders[10]  # idx key, not position
>>> placeholder.name
'Chart Placeholder 9'
>>> placeholder.placeholder_format.type
CHART (12)

>>> chart_data = ChartData()
>>> chart_data.categories = ['Yes', 'No']
>>> chart_data.add_series('Series 1', (42, 24))

>>> graphic_frame = placeholder.insert_chart(XL_CHART_TYPE.PIE, chart_data)  #returns new updated placeholder
>>> chart = graphic_frame.chart
>>> chart.chart_type
PIE (5)

##Setting the slide title
#Almost all slide layouts have a title placeholder, 


title_placeholder = slide.shapes.title
title_placeholder.text = 'Air-speed Velocity of Unladen Swallows'

##Working with text
#Auto shapes and table cells can contain text. 
#Other shapes can’t. 
#Text is always manipulated the same way, regardless of its container.

#Text exists in a hierarchy of three levels:
    Shape.text_frame
    TextFrame.paragraphs
    _Paragraph.runs

#All the text in a shape is contained in its text frame. 


for shape in slide.shapes:
    if not shape.has_text_frame:
        continue
    text_frame = shape.text_frame
    # do things with the text frame
    ...

#A text frame always contains at least one paragraph.

paragraph_strs = [
    'Egg, bacon, sausage and spam.',
    'Spam, bacon, sausage and spam.',
    'Spam, egg, spam, spam, bacon and spam.'
]

text_frame = shape.text_frame
text_frame.clear()  # remove any existing paragraphs, leaving one empty one

p = text_frame.paragraphs[0]
p.text = paragraph_strs[0]

for para_str in paragraph_strs[1:]:
    p = text_frame.add_paragraph()
    p.text = para_str

    
#Only runs can actually contain text. 
#Assigning a string to the .text attribute on a shape, text frame, or paragraph is a shortcut method 

shape.text = 'foobar'

# is equivalent to ...
text_frame = shape.text_frame
text_frame.clear()
p = text_frame.paragraphs[0]
run = p.add_run()
run.text = 'foobar'

##Applying text frame-level formatting
#The following produces a shape with a single paragraph, a slightly wider bottom than top margin 
#(these default to 0.05”), no left margin, text aligned top, and word wrapping turned off. 
#In addition, the auto-size behavior is set to adjust the width and height of the shape to fit its text. 

#Note that vertical alignment is set on the text frame. 
#Horizontal alignment is set on each paragraph

#The possible values for TextFrame.auto_size and TextFrame.vertical_anchor 
#are specified by the enumeration MSO_AUTO_SIZE and MSO_VERTICAL_ANCHOR respectively.


from pptx.util import Inches
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE

text_frame = shape.text_frame
text_frame.text = 'Spam, eggs, and spam'
text_frame.margin_bottom = Inches(0.08)
text_frame.margin_left = 0
text_frame.vertical_anchor = MSO_ANCHOR.TOP
text_frame.word_wrap = False
text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

##Applying paragraph formatting

#The following produces a shape containing three left-aligned paragraphs, 
#the second and third indented (like sub-bullets) under the first:

from pptx.enum.text import PP_ALIGN

paragraph_strs = [
    'Egg, bacon, sausage and spam.',
    'Spam, bacon, sausage and spam.',
    'Spam, egg, spam, spam, bacon and spam.'
]

text_frame = shape.text_frame
text_frame.clear()

p = text_frame.paragraphs[0]
p.text = paragraph_strs[0]
p.alignment = PP_ALIGN.LEFT

for para_str in paragraph_strs[1:]:
    p = text_frame.add_paragraph()
    p.text = para_str
    p.alignment = PP_ALIGN.LEFT
    p.level = 1

##Applying character formatting
#Character level formatting is applied at the run level, using the .font attribute. 

#The following formats a sentence in 18pt Calibri Bold and applies the theme color Accent 1.

from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.util import Pt

text_frame = shape.text_frame
text_frame.clear()  # not necessary for newly-created shape

p = text_frame.paragraphs[0]
run = p.add_run()
run.text = 'Spam, eggs, and spam'

font = run.font
font.name = 'Calibri'
font.size = Pt(18)
font.bold = True
font.italic = None  # cause value to be inherited from theme
font.color.theme_color = MSO_THEME_COLOR.ACCENT_1

#Or set the font color to an absolute RGB value. 
#Note that this will not change color when the theme is changed:

font.color.rgb = RGBColor(0xFF, 0x7F, 0x50)

#A run can also be made into a hyperlink by providing a target URL:
run.hyperlink.address = 'https://github.com/scanny/python-pptx'


##Working with charts
#python-pptx supports adding charts and modifying existing ones. 
#2D bar and column, line, and pie charts are supported.

#The following code adds a single-series column chart in a new presentation:

from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

# create presentation with 1 slide ------
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[5])

# define chart data ---------------------
chart_data = ChartData()
chart_data.categories = ['East', 'West', 'Midwest']
chart_data.add_series('Series 1', (19.2, 21.4, 16.7))

# add chart to slide --------------------
x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
)

prs.save('chart-01.pptx')


#To create a multi-series chart 

chart_data = ChartData()
chart_data.categories = ['East', 'West', 'Midwest']
chart_data.add_series('Q1 Sales', (19.2, 21.4, 16.7))
chart_data.add_series('Q2 Sales', (22.3, 28.6, 15.2))
chart_data.add_series('Q3 Sales', (20.4, 26.3, 14.2))

graphic_frame = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
)

chart = graphic_frame.chart  #Get the chart object for any formatting

##XY and Bubble charts
# XY (aka. scatter) charts. 

chart_data = XyChartData()

series_1 = chart_data.add_series('Model 1')
series_1.add_data_point(0.7, 2.7)
series_1.add_data_point(1.8, 3.2)
series_1.add_data_point(2.6, 0.8)

series_2 = chart_data.add_series('Model 2')
series_2.add_data_point(1.3, 3.7)
series_2.add_data_point(2.7, 2.3)
series_2.add_data_point(1.6, 1.8)

chart = slide.shapes.add_chart(
    XL_CHART_TYPE.XY_SCATTER, x, y, cx, cy, chart_data
).chart

#A bubble chart is essentially an XY chart where the marker size is used to reflect an additional value, effectively adding a third dimension to the chart.

chart_data = BubbleChartData()

series_1 = chart_data.add_series('Series 1')
series_1.add_data_point(0.7, 2.7, 10)
series_1.add_data_point(1.8, 3.2, 4)
series_1.add_data_point(2.6, 0.8, 8)

chart = slide.shapes.add_chart(
    XL_CHART_TYPE.BUBBLE, x, y, cx, cy, chart_data
).chart


##Chart - Axes

from pptx.enum.chart import XL_TICK_MARK
from pptx.util import Pt

category_axis = chart.category_axis
category_axis.has_major_gridlines = True
category_axis.minor_tick_mark = XL_TICK_MARK.OUTSIDE
category_axis.tick_labels.font.italic = True
category_axis.tick_labels.font.size = Pt(24)

value_axis = chart.value_axis
value_axis.maximum_scale = 50.0
value_axis.minor_tick_mark = XL_TICK_MARK.OUTSIDE
value_axis.has_minor_gridlines = True

tick_labels = value_axis.tick_labels
tick_labels.number_format = '0"%"'
tick_labels.font.bold = True
tick_labels.font.size = Pt(14)


##Chart - Data Labels
#access a Plot object to gain access to the data labels
#A plot is like a sub-chart, containing one or more series and drawn as a particular chart type, like column or line. 
#This distinction is needed for charts that combine more than one type, like a line chart appearing on top of a column chart. 
#A chart like this would have two plot objects, one for the series appearing as columns and the other for the lines. 

#Most charts only have a single plot and python-pptx doesn’t yet support creating multi-plot charts, 
#but you can access multiple plots on a chart that already has them.

from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_LABEL_POSITION

plot = chart.plots[0]
plot.has_data_labels = True
data_labels = plot.data_labels

data_labels.font.size = Pt(13)
data_labels.font.color.rgb = RGBColor(0x0A, 0x42, 0x80)
data_labels.position = XL_LABEL_POSITION.INSIDE_END

##Chart - Legend


from pptx.enum.chart import XL_LEGEND_POSITION

chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.RIGHT
chart.legend.include_in_layout = False



##Chart - Line Chart

chart_data = ChartData()
chart_data.categories = ['Q1 Sales', 'Q2 Sales', 'Q3 Sales']
chart_data.add_series('West',    (32.2, 28.4, 34.7))
chart_data.add_series('East',    (24.3, 30.6, 20.2))
chart_data.add_series('Midwest', (20.4, 18.3, 26.2))

x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
chart = slide.shapes.add_chart(
    XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
).chart

chart.has_legend = True
chart.legend.include_in_layout = False
chart.series[0].smooth = True

##Chart - Pie Chart
#it only ever has a single series and doesn’t have any axes:

chart_data = ChartData()
chart_data.categories = ['West', 'East', 'North', 'South', 'Other']
chart_data.add_series('Series 1', (0.135, 0.324, 0.180, 0.235, 0.126))

chart = slide.shapes.add_chart(
    XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
).chart

chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.BOTTOM
chart.legend.include_in_layout = False

chart.plots[0].has_data_labels = True
data_labels = chart.plots[0].data_labels
data_labels.number_format = '0%'
data_labels.position = XL_LABEL_POSITION.OUTSIDE_END



###Important Methods of .shapes ( class pptx.shapes.shapetree.SlideShapes)
add_chart(chart_type, x, y, cx, cy, chart_data)  #cx, cy are size , chart_type - http://python-pptx.readthedocs.io/en/latest/api/enum/XlChartType.html#xlcharttype
add_connector(connector_type, begin_x, begin_y, end_x, end_y) #connector_type http://python-pptx.readthedocs.io/en/latest/api/enum/MsoConnectorType.html#msoconnectortype
add_picture(image_file, left, top, width=None, height=None) #a path to a file (a string) or a file-like object.
add_shape(autoshape_type_id, left, top, width, height)  # autoshape_type_id from http://python-pptx.readthedocs.io/en/latest/api/enum/MsoAutoShapeType.html
add_table(rows, cols, left, top, width, height)
add_textbox(left, top, width, height)
index(shape)            Return the index of shape in this sequence, raising ValueError if shape is not in the collection.
placeholders            Instance of SlidePlaceholders containing sequence of placeholder shapes in this slide.
title                   The title placeholder shape on the slide or None if the slide has no title placeholder.

#Methods on any Shape object 
element                 Reference to the lxml element for this shape, e.g. a CT_Shape instance.
has_chart               True if this shape is a graphic frame containing a chart object. False otherwise. When True, the chart object can be accessed using the .chart property.
has_table               True if this shape is a graphic frame containing a table object. False otherwise. When True, the table object can be accessed using the .table property.
has_text_frame          True if this shape can contain text.
height                  Read/write. Integer distance between top and bottom extents of shape in EMUs
is_placeholder          True if this shape is a placeholder. A shape is a placeholder if it has a <p:ph> element.
left                    Read/write. Integer distance of the left edge of this shape from the left edge of the slide, in English Metric Units (EMU)
name                    Name of this shape, e.g. ‘Picture 7’
placeholder_format      A _PlaceholderFormat object providing access to placeholder-specific properties such as placeholder type. Raises ValueError on access if the shape is not a placeholder.
rotation                Read/write float. Degrees of clockwise rotation. Negative values can be assigned to indicate counter-clockwise rotation, e.g. assigning -45.0 will change setting to 315.0.
shape_id                Read-only positive integer identifying this shape.
                        The id of a shape is unique among all shapes on a slide.
shape_type              Unique integer identifying the type of this shape, like MSO_SHAPE_TYPE.CHART. Must be implemented by subclasses.
top                     Read/write. Integer distance of the top edge of this shape from the top edge of the slide, in English Metric Units (EMU)
width                   Read/write. Integer distance between left and right extents of shape in EMUs


##Shape - AutoShapes - class pptx.shapes.autoshape.Shape - Corresponds to the <p:sp>
adjustments         Read-only reference to AdjustmentCollection instance for this shape
auto_shape_type     Enumeration value identifying the type of this auto shape, like MSO_SHAPE.ROUNDED_RECTANGLE. Raises ValueError if this shape is not an auto shape.
fill                FillFormat instance for this shape, providing access to fill properties such as fill color.
has_text_frame      True if this shape can contain text. Always True for an AutoShape.
line                LineFormat instance for this shape, providing access to line properties such as line color.
shape_type          Unique integer identifying the type of this shape, like MSO_SHAPE_TYPE.TEXT_BOX.
text                Read/write. All the text in this shape as a single string. A line feed character (‘\n’) appears in the string for each paragraph and line break in the shape, except the last paragraph. A shape containing a single paragraph with no line breaks will produce a string having no line feed characters. Assigning a string to this property replaces all text in the shape with a single paragraph containing the assigned text. The assigned value can be a 7-bit ASCII string, a UTF-8 encoded 8-bit string, or unicode. String values are converted to unicode assuming UTF-8 encoding. Each line feed character in an assigned string is translated into a line break within the single resulting paragraph.
text_frame          TextFrame instance for this shape, containing the text of the shape and providing access to text formatting properties.


##Shape - class pptx.shapes.connector.Connector
# A connector is a linear shape having end-points that can be connected to other objects (but not to other connectors). 
#A line can be straight, have elbows, or can be curved.

begin_connect(shape, cxn_pt_idx)
Connect the beginning of this connector to shape at the connection point specified by cxn_pt_idx. 
Each shape has zero or more connection points and they are identified by index, starting with 0. 
Generally, the first connection point of a shape is at the top center of its bounding box 
and numbering proceeds counter-clockwise from there. 
However this is only a convention and may vary, especially with non built-in shapes.

begin_x
Return the X-position of the begin point of this connector, in English Metric Units (as a Length object).

begin_y
Return the Y-position of the begin point of this connector, in English Metric Units (as a Length object).

end_connect(shape, cxn_pt_idx)[source]
Connect the ending of this connector to shape at the connection point specified by cxn_pt_idx.

end_x
Return the X-position of the end point of this connector, in English Metric Units (as a Length object).

end_y
Return the Y-position of the end point of this connector, in English Metric Units (as a Length object).


##Shape - class pptx.shapes.picture.Picture

crop_bottom
A float representing the relative portion cropped from the bottom of this picture where 1.0 represents 100%. For example, 25% is represented by 0.25. Negative values are valid as are values greater than 1.0.

crop_left
A float representing the relative portion cropped from the left side of this picture where 1.0 represents 100%.

crop_right
A float representing the relative portion cropped from the right side of this picture where 1.0 represents 100%.

crop_top
A float representing the relative portion cropped from the top of this picture where 1.0 represents 100%.

image
An Image object providing access to the properties and bytes of the image in this picture shape.

line
An instance of LineFormat, providing access to the properties of the outline bordering this picture, such as its color and width.

shape_type
Unique integer identifying the type of this shape, unconditionally MSO_SHAPE_TYPE.PICTURE in this case.


##Shape - class pptx.shapes.graphfrm.GraphicFrame
#A graphic frame is the shape containing a table, chart, or smart art.

chart
The Chart object containing the chart in this graphic frame. Raises ValueError if this graphic frame does not contain a chart.

click_action
An ActionSetting instance providing access to the mouse click behaviors defined on this shape. An ActionSetting object is always returned, even when no click behavior is defined on the shape.

element
Reference to the lxml element for this shape, e.g. a CT_Shape instance.

has_chart
True if this graphic frame contains a chart object. False otherwise. When True, the chart object can be accessed using the .chart property.

has_table
True if this graphic frame contains a table object. False otherwise. When True, the table object can be accessed using the .table property.

height
Read/write. Integer distance between top and bottom extents of shape in EMUs

left
Read/write. Integer distance of the left edge of this shape from the left edge of the slide, in English Metric Units (EMU)

name
Name of this shape, e.g. ‘Picture 7’

rotation
Read/write float. Degrees of clockwise rotation. Negative values can be assigned to indicate counter-clockwise rotation, e.g. assigning -45.0 will change setting to 315.0.

shape_id
Read-only positive integer identifying this shape.
The id of a shape is unique among all shapes on a slide.

table
The Table object contained in this graphic frame. Raises ValueError if this graphic frame does not contain a table.

top
Read/write. Integer distance of the top edge of this shape from the top edge of the slide, in English Metric Units (EMU)

width
Read/write. Integer distance between left and right extents of shape in EMUs


##Table objects - class pptx.shapes.table.Table

cell(row_idx, col_idx)[source]
Return table cell at row_idx, col_idx location. 
Indexes are zero-based, e.g. cell(0, 0) is the top, left cell.

columns
Read-only reference to collection of _Column objects representing the table’s columns. 
_Column objects are accessed using list notation, e.g. col = tbl.columns[0].

first_col
Read/write boolean property which, when true, indicates the first column should be formatted differently, as for a side-heading column at the far left of the table.

first_row
Read/write boolean property which, when true, indicates the first row should be formatted differently, e.g. for column headings.

horz_banding
Read/write boolean property which, when true, indicates the rows of the table should appear with alternating shading.

last_col
Read/write boolean property which, when true, indicates the last column should be formatted differently, as for a row totals column at the far right of the table.

last_row
Read/write boolean property which, when true, indicates the last row should be formatted differently, as for a totals row at the bottom of the table.

rows
Read-only reference to collection of _Row objects representing the table’s rows. 
_Row objects are accessed using list notation, e.g. col = tbl.rows[0].

vert_banding
Read/write boolean property which, when true, indicates the columns of the table should appear with alternating shading.


##Table objects - _Column objects - class pptx.shapes.table._Column[source]
width
Width of column in EMU.


##Table objects -_Row objects - class pptx.shapes.table._Row[source]
cells
Read-only reference to collection of cells in row. 
An individual cell is referenced using list notation, e.g. cell = row.cells[0].

height
Height of row in EMU.


##Table objects -_Cell objects - class pptx.shapes.table._Cell
fill
FillFormat instance for this cell, providing access to fill properties such as foreground color.

margin_left
Read/write integer value of left margin of cell as a Length value object. If assigned None, the default value is used, 0.1 inches for left and right margins and 0.05 inches for top and bottom.

margin_right
Right margin of cell.

margin_top
Top margin of cell.

margin_bottom
Bottom margin of cell.

text
Write-only. Assignment to text replaces all text currently contained in the cell, 
resulting in a text frame containing exactly one paragraph, itself containing a single run. 
The assigned value can be a 7-bit ASCII string, a UTF-8 encoded 8-bit string, or unicode. 
String values are converted to unicode assuming UTF-8 encoding.

text_frame
TextFrame instance containing the text that appears in the cell.

vertical_anchor
Vertical anchor of this table cell, determines the vertical alignment of text in the cell. 
Value is like MSO_ANCHOR.MIDDLE. Can be None, meaning the cell has no vertical anchor setting and its effective value is inherited from a higher-level object.





###pyPdf2  

##The PdfFileReader Class
class PyPDF2.PdfFileReader(stream, strict=True, warndest=None, overwriteWarnings=True)
    stream – A File object or an object that supports the standard read and seek methods similar to a File object. Could also be a string representing a path to a PDF file.

    decrypt(password)
    
    getDestinationPageNumber(destination)


    getDocumentInfo()
    Retrieves the PDF file’s document information dictionary, if it exists
    Return type:DocumentInformation or None if none exists 

    getFields(tree=None, retval=None, fileobj=None)
    Extracts field data if this PDF contains interactive form fields. 
    The tree and retval parameters are for recursive use.

    getFormTextFields()
    Retrieves form fields from the document with textual data (inputs, dropdowns)

    getNamedDestinations(tree=None, retval=None)
    Retrieves the named destinations present in the document.
    Returns: a dictionary which maps names to Destinations. 

    getNumPages()
    Calculates the number of pages in this PDF file.

    getOutlines(node=None, outlines=None)
    Retrieves the document outline present in the document.
    Returns:a nested list of Destinations. 

    getPage(pageNumber)
    Retrieves a page by number from this PDF file.
    Returns:a PageObject instance. 

    getPageLayout()
    Get the page layout. See setPageLayout() for a description of valid layouts.
    Return type:str, None if not specified 

    getPageMode()
    Get the page mode. See setPageMode() for a description of valid modes.
    Return type:str, None if not specified 

    getPageNumber(page)
    Retrieve page number of a given PageObject

    getXmpMetadata()
    Retrieves XMP (Extensible Metadata Platform) data from the PDF document root.
    Returns:a XmpInformation instance that can be used to access XMP metadata from the document. 

    isEncrypted
    Read-only boolean property showing whether this PDF file is encrypted. 
    Note that this property, if true, will remain true even after the decrypt() method is called.





class PyPDF2.PdfFileMerger(strict=True)
    Initializes a PdfFileMerger object. 
    PdfFileMerger merges multiple PDFs into a single PDF. 
    It can concatenate, slice, insert, or any combination of the above.

    addBookmark(title, pagenum, parent=None)
    Add a bookmark to this PDF file.
    •title (str) – Title to use for this bookmark.
    •pagenum (int) – Page number this bookmark will point to.
    •parent – A reference to a parent bookmark to create nested bookmarks.
     
    addMetadata(infos)
    Add custom metadata to the output.
    infos (dict) – a Python dictionary where each key is a field 
    and each value is your new metadata. Example: {u'/Title': u'My title'} 

    addNamedDestination(title, pagenum)
    Add a destination to the output.
    •title (str) – Title to use
    •pagenum (int) – Page number this destination points at.
     
    append(fileobj, bookmark=None, pages=None, import_bookmarks=True)
    Identical to the merge() method, 
    but assumes you want to concatenate all pages onto the end of the file instead of specifying a position.
    •fileobj – A File Object or an object that supports the standard read and seek methods similar to a File Object. Could also be a string representing a path to a PDF file.
    •bookmark (str) – Optionally, you may specify a bookmark to be applied at the beginning of the included file by supplying the text of the bookmark.
    •pages – can be a Page Range or a (start, stop[, step]) tuple to merge only the specified range of pages from the source document into the output document.
    •import_bookmarks (bool) – You may prevent the source document’s bookmarks from being imported by specifying this as False.
     
    close()
    Shuts all file descriptors (input and output) and clears all memory usage.

    merge(position, fileobj, bookmark=None, pages=None, import_bookmarks=True)
    Merges the pages from the given file into the output file at the specified page number.
    •position (int) – The page number to insert this file. File will be inserted after the given number.
    •fileobj – A File Object or an object that supports the standard read and seek methods similar to a File Object. Could also be a string representing a path to a PDF file.
    •bookmark (str) – Optionally, you may specify a bookmark to be applied at the beginning of the included file by supplying the text of the bookmark.
    •pages – can be a Page Range or a (start, stop[, step]) tuple to merge only the specified range of pages from the source document into the output document.
    •import_bookmarks (bool) – You may prevent the source document’s bookmarks from being imported by specifying this as False.
     
    setPageLayout(layout)
    Set the page layout
    layout (str) – The page layout to be used 
    Valid layouts are:
        /NoLayout Layout explicitly not specified 
        /SinglePage Show one page at a time 
        /OneColumn Show one column at a time 
        /TwoColumnLeft Show pages in two columns, odd-numbered pages on the left 
        /TwoColumnRight       Show pages in two columns, odd-numbered pages on the right 
        /TwoPageLeft Show two pages at a time, odd-numbered pages on the left 
        /TwoPageRight Show two pages at a time, odd-numbered pages on the right 
        
    setPageMode(mode)
    Set the page mode.
    Valid modes are:
        /UseNone Do not show outlines or thumbnails panels 
        /UseOutlines Show outlines (aka bookmarks) panel 
        /UseThumbs Show page thumbnails panel 
        /FullScreen Fullscreen view 
        /UseOC Show Optional Content Group (OCG) panel 
        /UseAttachments   Show attachments panel 
      
    write(fileobj)
    Writes all data that has been merged to the given output file.
    fileobj – Output file. Can be a filename or any kind of file-like object. 




class PyPDF2.PdfFileWriter
    This class supports writing PDF files out, given pages produced by another class 
    (typically PdfFileReader).

    addAttachment(fname, fdata)
    Embed a file inside the PDF.
    •fname (str) – The filename to display.
    •fdata (str) – The data in the file.
     
    addBlankPage(width=None, height=None)
    Appends a blank page to this PDF file and returns it. 
    If no page size is specified, use the size of the last page.

    addBookmark(title, pagenum, parent=None, color=None, bold=False, italic=False, fit='/Fit', *args)
    Add a bookmark to this PDF file.

    addJS(javascript)
    Add Javascript which will launch upon opening this PDF.
    javascript (str) – Your Javascript. 

    addLink(pagenum, pagedest, rect, border=None, fit='/Fit', *args)
    Add an internal link from a rectangular area to the specified page.
    pagenum (int) – index of the page on which to place the link.
    •pagedest (int) – index of the page to which the link should go.
    •rect – RectangleObject or array of four integers specifying the clickable rectangular area [xLL, yLL, xUR, yUR], or string in the form "[ xLL yLL xUR yUR ]".

    addMetadata(infos)
    Add custom metadata to the output.
    infos (dict) – a Python dictionary where each key is a field and each value is your new metadata. 

    addPage(page)
    Adds a page to this PDF file. The page is usually acquired from a PdfFileReader instance.
    page (PageObject) – The page to add to the document. Should be an instance of PageObject 

    appendPagesFromReader(reader, after_page_append=None)
    Copy pages from reader to writer. 
    Includes an optional callback parameter which is invoked after pages are appended to the writer.
    reader – a PdfFileReader object from which to copy page annotations to this writer object. The writer’s annots 
    after_page_append (function): Callback function that is invoked after each page is appended to the writer. Callback signature:

    cloneDocumentFromReader(reader, after_page_append=None)
    Create a copy (clone) of a document from a PDF file reader

    cloneReaderDocumentRoot(reader)
    Copy the reader document root to the writer.

    encrypt(user_pwd, owner_pwd=None, use_128bit=True)
    Encrypt this PDF file with the PDF Standard encryption handler.
    •user_pwd (str) – The “user password”, which allows for opening and reading the PDF file with the restrictions provided.
    •owner_pwd (str) – The “owner password”, which allows for opening the PDF files without any restrictions. By default, the owner password is the same as the user password.
    •use_128bit (bool) – flag as to whether to use 128bit encryption. When false, 40bit encryption will be used. By default, this flag is on.
     
    getNumPages()
    Returns:the number of pages. 

    getPage(pageNumber)
    Retrieves a page by number from this PDF file.
    pageNumber (int) – The page number to retrieve (pages begin at zero) 
    Return type:PageObject 

    getPageLayout()
    Get the page layout. See setPageLayout() for a description of valid layouts.
    Return type:str, None if not specified 

    getPageMode()
    Get the page mode. See setPageMode() for a description of valid modes.
    Return type:str, None if not specified 

    insertBlankPage(width=None, height=None, index=0)
    Inserts a blank page to this PDF file and returns it. 
    If no page size is specified, use the size of the last page.
    •width (float) – The width of the new page expressed in default user space units.
    •height (float) – The height of the new page expressed in default user space units.
    •index (int) – Position to add the page.
    Return type:PageObject

    insertPage(page, index=0)
    Insert a page in this PDF file. The page is usually acquired from a PdfFileReader instance.
    •page (PageObject) – The page to add to the document. This argument should be an instance of PageObject.
    •index (int) – Position at which the page will be inserted.
     

    removeImages(ignoreByteStringObject=False)
    Removes images from this output.
    ignoreByteStringObject (bool) – optional parameter to ignore ByteString Objects. 

    removeLinks()
    Removes links and annotations from this output.

    removeText(ignoreByteStringObject=False)
    Removes text from this output.

    setPageLayout(layout)
    Set the page layout(str, check Merge reference to get the values).

    setPageMode(mode)
    Set the page mode(str, check Merge reference to get the values).

    updatePageFormFieldValues(page, fields)
    Update the form field values for a given page from a fields dictionary. 
    Copy field texts and values from fields to page.
    •page – Page reference from PDF writer where the annotations and field data will be updated.
    •fields – a Python dictionary of field names (/T) and text values (/V)
     
    write(stream)
    Writes the collection of pages added to this object out as a PDF file.
    stream – An object to write the file to. The object must support the write method and the tell method, similar to a file object. 




class PyPDF2.generic.RectangleObject(arr)
    This class is used to represent page boxes in PyPDF2. 
    These boxes include
        •artBox
        •bleedBox
        •cropBox
        •mediaBox
        •trimBox
        
    lowerLeft
    Property to read and modify the lower left coordinate of this box in (x,y) form.

    lowerRight
    Property to read and modify the lower right coordinate of this box in (x,y) form.

    upperLeft
    Property to read and modify the upper left coordinate of this box in (x,y) form.

    upperRight
    Property to read and modify the upper right coordinate of this box in (x,y) form.




class PyPDF2.generic.Field(data)
    A class representing a field dictionary. 
    This class is accessed through getFields()

    additionalActions
    Read-only property accessing the additional actions dictionary. 
    This dictionary defines the field’s behavior in response to trigger events. 

    altName
    Read-only property accessing the alternate name of this field.

    defaultValue
    Read-only property accessing the default value of this field.

    fieldType
    Read-only property accessing the type of this field.

    flags
    Read-only property accessing the field flags, specifying various characteristics of the field (see Table 8.70 of the PDF 1.7 reference).

    kids
    Read-only property accessing the kids of this field.

    mappingName
    Read-only property accessing the mapping name of this field. This name is used by PyPDF2 as a key in the dictionary returned by getFields()

    name
    Read-only property accessing the name of this field.

    parent
    Read-only property accessing the parent of this field.

    value
    Read-only property accessing the value of this field. Format varies based on field type.



class PyPDF2.pdf.PageObject(pdf=None, indirectRef=None)
    This class represents a single page within a PDF file. 
    Typically this object will be created by accessing the getPage() method of the PdfFileReader class, 
    or  create an empty page with the createBlankPage() static method.
     
    addTransformation(ctm)
    Applies a transformation matrix to the page.
    ctm (tuple) – A 6-element tuple containing the operands of the transformation matrix. 

    artBox
    A RectangleObject, expressed in default user space units, 
    defining the extent of the page’s meaningful content as intended by the page’s creator.

    bleedBox
    A RectangleObject, expressed in default user space units, 
    defining the region to which the contents of the page should be clipped when output in a production enviroment.

    compressContentStreams()
    Compresses the size of this page by joining all content streams 
    and applying a FlateDecode filter.

    static createBlankPage(pdf=None, width=None, height=None)
    Returns a new blank page. If width or height is None, 
    try to get the page size from the last page of pdf.
    •pdf – PDF file the page belongs to
    •width (float) – The width of the new page expressed in default user space units.
    •height (float) – The height of the new page expressed in default user space units.
    Return type:PageObject
     
    cropBox
    A RectangleObject, expressed in default user space units, 
    defining the visible region of default user space. 
    When the page is displayed or printed, 
    its contents are to be clipped (cropped) to this rectangle 
    and then imposed on the output medium in some implementation-defined manner. 
    Default value: same as mediaBox.

    extractText()
    Locate all text drawing commands, in the order they are provided in the content stream, 
    and extract the text. 
    This works well for some PDF files, but poorly for others, depending on the generator used. 
    Returns:a unicode string object. 

    getContents()
    Accesses the page contents.
    Returns: the /Contents object, or None if it doesn’t exist. 

    mediaBox
    A RectangleObject, expressed in default user space units, 
    defining the boundaries of the physical medium on 
    which the page is intended to be displayed or printed.

    mergePage(page2)
    Merges the content streams of two pages into one. 
    Resource references (i.e. fonts) are maintained from both pages. 
    The mediabox/cropbox/etc of this page are not altered. 

    mergeRotatedPage(page2, rotation, expand=False)
    This is similar to mergePage, 
    but the stream to be merged is rotated by appling a transformation matrix.
    •page2 (PageObject) – the page to be merged into this one. Should be an instance of PageObject.
    •rotation (float) – The angle of the rotation, in degrees
    •expand (bool) – Whether the page should be expanded to fit the dimensions of the page to be merged.
     
    mergeRotatedScaledPage(page2, rotation, scale, expand=False)
    This is similar to mergePage, but the stream to be merged is rotated 
    and scaled by appling a transformation matrix.
    •page2 (PageObject) – the page to be merged into this one. Should be an instance of PageObject.
    •rotation (float) – The angle of the rotation, in degrees
    •scale (float) – The scaling factor
    •expand (bool) – Whether the page should be expanded to fit the dimensions of the page to be merged.
     
    mergeRotatedScaledTranslatedPage(page2, rotation, scale, tx, ty, expand=False)
    This is similar to mergePage, but the stream to be merged is translated, rotated and scaled 
    •page2 (PageObject) – the page to be merged into this one. Should be an instance of PageObject.
    •tx (float) – The translation on X axis
    •ty (float) – The translation on Y axis
    •rotation (float) – The angle of the rotation, in degrees
    •scale (float) – The scaling factor
    •expand (bool) – Whether the page should be expanded to fit the dimensions of the page to be merged.
     
    mergeRotatedTranslatedPage(page2, rotation, tx, ty, expand=False)
    This is similar to mergePage, but the stream to be merged is rotated and translated 
    •page2 (PageObject) – the page to be merged into this one. Should be an instance of PageObject.
    •tx (float) – The translation on X axis
    •ty (float) – The translation on Y axis
    •rotation (float) – The angle of the rotation, in degrees
    •expand (bool) – Whether the page should be expanded to fit the dimensions of the page to be merged.
     
    mergeScaledPage(page2, scale, expand=False)
    This is similar to mergePage, but the stream to be merged is scaled 
    •page2 (PageObject) – The page to be merged into this one. Should be an instance of PageObject.
    •scale (float) – The scaling factor
    •expand (bool) – Whether the page should be expanded to fit the dimensions of the page to be merged.
     
    mergeScaledTranslatedPage(page2, scale, tx, ty, expand=False)
    This is similar to mergePage, but the stream to be merged is translated and scaled 
    •page2 (PageObject) – the page to be merged into this one. Should be an instance of PageObject.
    •scale (float) – The scaling factor
    •tx (float) – The translation on X axis
    •ty (float) – The translation on Y axis
    •expand (bool) – Whether the page should be expanded to fit the dimensions of the page to be merged.
     
    mergeTransformedPage(page2, ctm, expand=False)
    This is similar to mergePage, but a transformation matrix is applied to the merged stream.
    •page2 (PageObject) – The page to be merged into this one. Should be an instance of PageObject.
    •ctm (tuple) – a 6-element tuple containing the operands of the transformation matrix
    •expand (bool) – Whether the page should be expanded to fit the dimensions of the page to be merged.
     
    mergeTranslatedPage(page2, tx, ty, expand=False)
    This is similar to mergePage, but the stream to be merged is translated 
    •page2 (PageObject) – the page to be merged into this one. Should be an instance of PageObject.
    •tx (float) – The translation on X axis
    •ty (float) – The translation on Y axis
    •expand (bool) – Whether the page should be expanded to fit the dimensions of the page to be merged.
     
    rotateClockwise(angle)
    Rotates a page clockwise by increments of 90 degrees.
    angle (int) – Angle to rotate the page. Must be an increment of 90 deg. 

    rotateCounterClockwise(angle)
    Rotates a page counter-clockwise by increments of 90 degrees.
    angle (int) – Angle to rotate the page. Must be an increment of 90 deg. 

    scale(sx, sy)
    Scales a page by the given factors 
    •sx (float) – The scaling factor on horizontal axis.
    •sy (float) – The scaling factor on vertical axis.
     
    scaleBy(factor)
    Scales a page by the given factor 
    factor (float) – The scaling factor (for both X and Y axis)

    scaleTo(width, height)
    Scales a page to the specified dimentions 
    •width (float) – The new width.
    •height (float) – The new heigth.
     
    trimBox
    A RectangleObject, expressed in default user space units, 
    defining the intended dimensions of the finished page after trimming.




class PyPDF2.generic.Destination(title, page, typ, *args)
    A class representing a destination within a PDF file. 
    •title (str) – Title of this destination.
    •page (int) – Page number of this destination.
    •typ (str) – How the destination is displayed.
    •args – Additional arguments may be necessary depending on the type.
    Valid typ arguments 
        /Fit No additional arguments 
        /XYZ [left] [top] [zoomFactor] 
        /FitH [top] 
        /FitV [left] 
        /FitR [left] [bottom] [right] [top] 
        /FitB No additional arguments 
        /FitBH [top] 
        /FitBV [left] 
        
    bottom
    Read-only property accessing the bottom vertical coordinate.
    Return type:int, or None if not available. 

    left
    Read-only property accessing the left horizontal coordinate.

    page
    Read-only property accessing the destination page number.
    Return type:int 

    right
    Read-only property accessing the right horizontal coordinate.

    title
    Read-only property accessing the destination title.

    top
    Read-only property accessing the top vertical coordinate.

    typ
    Read-only property accessing the destination type.

    zoom
    Read-only property accessing the zoom factor
    Return type:int, or None if not available. 


class PyPDF2.pdf.DocumentInformation
    author
    Read-only property accessing the document’s author. 
    Returns a unicode string (TextStringObject) or None if the author is not specified.

    author_raw
    The “raw” version of author; can return a ByteStringObject.

    creator
    Read-only property accessing the document’s creator. If the document was converted to PDF from another format, this is the name of the application (e.g. OpenOffice) that created the original document from which it was converted. Returns a unicode string (TextStringObject) or None if the creator is not specified.

    creator_raw
    The “raw” version of creator; can return a ByteStringObject.

    producer
    Read-only property accessing the document’s producer. If the document was converted to PDF from another format, this is the name of the application (for example, OSX Quartz) that converted it to PDF. Returns a unicode string (TextStringObject) or None if the producer is not specified.

    producer_raw
    The “raw” version of producer; can return a ByteStringObject.

    subject
    Read-only property accessing the document’s subject. Returns a unicode string (TextStringObject) or None if the subject is not specified.

    subject_raw
    The “raw” version of subject; can return a ByteStringObject.

    title
    Read-only property accessing the document’s title. Returns a unicode string (TextStringObject) or None if the title is not specified.

    title_raw
    The “raw” version of title; can return a ByteStringObject.

#Example 

>>> from PyPDF2 import PdfFileReader
>>> inputPdf = PdfFileReader(open("test.pdf", "rb"))
>>> docInfo = inputPdf.getDocumentInfo()
>>> docInfo.author
Anonymous
>>> docInfo.creator
Hewlett Packard MFP
>>> docInfo.producer
Acrobat Distiller 10.0.0 (Windows)
>>> docInfo.title
A Test
>>> docInfo.subject
testing

###Easy Concatenation with pdfcat(under Scripts folder)
#we refer to the slices as page ranges.
#Page range expression examples:
:           all pages -1 last page 
22          just the 23rd page 
:-1         all but the last page 
0:3         the first three pages -2 second-to-last page 
:3          the first three pages -2: last two pages 
5:          from the sixth page onward -3:-1 third & second to last 

#step number is also recognized:
::2         0 2 4 ... to the end 
1:10:2      1 3 5 7 9 
::-1        all pages in reverse order 
3:0:-1      3 2 1 but not 0 
2::-1       2 1 0 

$ pdfcat [-h] [-o output.pdf] [-v] input.pdf [page_range...] ...

#Example 
#Concatenates all of head.pdf, all but page seven of content.pdf, and the last page of tail.pdf, producing output.pdf.
$ pdfcat -o output.pdf head.pdf content.pdf :6 7: tail.pdf -1


$ pdfcat chapter*.pdf >book.pdf
$ pdfcat chapter?.pdf chapter10.pdf >book.pdf




##Example to extract text 
#might not work for all pdfs 
import PyPDF2
pdf_file = open('sample.pdf', 'rb')
read_pdf = PyPDF2.PdfFileReader(pdf_file)
number_of_pages = read_pdf.getNumPages()
page = read_pdf.getPage(0)
page_content = page.extractText()
print page_content.encode('utf-8')


#Or use textract,  
$ pip install textract
#textract needs many external dependencies as it works on many formats 
#https://textract.readthedocs.io/en/stable/installation.html

#Needs ebooklib-0.15, but ebooklib has error for 0.15, 
#download latest, change version to 0.15 and then do python setup.py install 

# some python file
import textract
text = textract.process("path/to/file.extension")


#Currently supporting

.csv via python builtins
.doc via antiword
.docx via python-docx2txt
.eml via python builtins
.epub via ebooklib
.gif via tesseract-ocr
.jpg and .jpeg via tesseract-ocr
.json via python builtins
.html and .htm via beautifulsoup4
.mp3 via sox, SpeechRecognition, and pocketsphinx
.msg via msg-extractor
.odt via python builtins
.ogg via sox, SpeechRecognition, and pocketsphinx
.pdf via pdftotext (default) or pdfminer.six
.png via tesseract-ocr
.pptx via python-pptx
.ps via ps2text
.rtf via unrtf
.tiff and .tif via tesseract-ocr
.txt via python builtins
.wav via SpeechRecognition and pocketsphinx
.xlsx via xlrd
.xls via xlrd

#Note following is required 
libxml2 2.6.21 or later is required by the .docx parser which uses lxml via python-docx.
libxslt 1.1.15 or later is required by the .docx parser which users lxml via python-docx.
antiword is required by the .doc parser.
pdftotext is optionally required by the .pdf parser (there is a pure python fallback that works if pdftotext isn’t installed).
pstotext is required by the .ps parser.
tesseract-ocr is required by the .jpg, .png and .gif parser.
sox is required by the .mp3 and .ogg parser. You need to install ffmpeg, lame, libmad0 and libsox-fmt-mp3, before building sox, for these filetypes to work.

$ pip install  python-docx      
#install lxml from http://www.lfd.uci.edu/~gohlke/pythonlibs/ 

#install  tesseract-ocr
#https://github.com/tesseract-ocr/tesseract/wiki/Downloads
$ tesseract --help
$ tesseract imagename outputbase [-l lang] [-psm pagesegmode] [configfile...]


#So basic usage to do OCR on an image called 'myscan.png' 
#and save the result to 'out.txt' would be:
$ tesseract myscan.png out


#Using One Language
$ tesseract  --tessdata-dir ./     ./testing/eurotext.png  ./testing/eurotext-eng   -l eng
$ tesseract  --tessdata-dir ./ ./testing/eurotext.png ./testing/eurotext-engdeu -l eng+deu


##xPDF 
x86, Windows -- pdftops, pdftotext, pdftohtml, pdftoppm, pdftopng, pdfimages, pdfinfo, pdffonts, and pdfdetach only:
#http://mirror.unl.edu/ctan/support/xpdf/
xpdfbin-win-3.04.zip (includes 32-bit and 64-bit binaries) (10794993 bytes) 


##Example of pypdf 
from PyPDF2 import PdfFileWriter, PdfFileReader

output = PdfFileWriter()
input1 = PdfFileReader(open("document1.pdf", "rb"))

# print how many pages input1 has:
print "document1.pdf has %d pages." % input1.getNumPages()

# add page 1 from input1 to output document, unchanged
output.addPage(input1.getPage(0))

# add page 2 from input1, but rotated clockwise 90 degrees
output.addPage(input1.getPage(1).rotateClockwise(90))

# add page 3 from input1, rotated the other way:
output.addPage(input1.getPage(2).rotateCounterClockwise(90))
# alt: output.addPage(input1.getPage(2).rotateClockwise(270))

# add page 4 from input1, but first add a watermark from another PDF:
page4 = input1.getPage(3)
watermark = PdfFileReader(open("watermark.pdf", "rb"))
page4.mergePage(watermark.getPage(0))
output.addPage(page4)


# add page 5 from input1, but crop it to half size:
page5 = input1.getPage(4)
page5.mediaBox.upperRight = (
    page5.mediaBox.getUpperRight_x() / 2,
    page5.mediaBox.getUpperRight_y() / 2
)
output.addPage(page5)

# add some Javascript to launch the print window on opening this PDF.
# the password dialog may prevent the print dialog from being shown,
# comment the the encription lines, if that's the case, to try this out
output.addJS("this.print({bUI:true,bSilent:false,bShrinkToFit:true});")

# encrypt your new PDF and add a password
password = "secret"
output.encrypt(password)

# finally, write "output" to document-output.pdf
outputStream = file("PyPDF2-output.pdf", "wb")
output.write(outputStream)

##Example of Merging 

from PyPDF2 import PdfFileMerger

merger = PdfFileMerger()

input1 = open("document1.pdf", "rb")
input2 = open("document2.pdf", "rb")
input3 = open("document3.pdf", "rb")

# add the first 3 pages of input1 document to output
merger.append(fileobj = input1, pages = (0,3))

# insert the first page of input2 into the output beginning after the second page
merger.merge(position = 2, fileobj = input2, pages = (0,1))

# append entire input3 document to the end of the output document
merger.append(input3)

# Write to an output PDF document
output = open("document-output.pdf", "wb")
merger.write(output)

##Creation of Multipage pdf 
#pip install reportlab

from __future__ import print_function
from sys import argv

from reportlab.pdfgen import canvas

point = 1
inch = 72

TEXT = """%s    page %d of %d

a wonderful file
created with Sample_Code/makesimple.py"""


def make_pdf_file(output_filename, np):
    title = output_filename
    c = canvas.Canvas(output_filename, pagesize=(8.5 * inch, 11 * inch))
    c.setStrokeColorRGB(0,0,0)
    c.setFillColorRGB(0,0,0)
    c.setFont("Helvetica", 12 * point)
    for pn in range(1, np + 1):
        v = 10 * inch
        for subtline in (TEXT % (output_filename, pn, np)).split( '\n' ):
            c.drawString( 1 * inch, v, subtline )
            v -= 12 * point
        c.showPage()
    c.save()

if __name__ == "__main__":
    nps = [None, 5, 11, 17]
    for i, np in enumerate(nps):
        if np:
            filename = "simple%d.pdf" % i
            make_pdf_file(filename, np)
            print ("Wrote", filename)



##Examples of reportlab.pdfgen 
#https://www.reportlab.com/docs/reportlab-reference.pdf

##Example - Simple formatting is shown by the odyssey.py script
#Copyright ReportLab Europe Ltd. 2000-2012
#see license.txt for license details
__version__=''' $Id: odyssey.py 3959 2012-09-27 14:39:39Z robin $ '''
___doc__=''
#odyssey.py
#
#Demo/benchmark of PDFgen rendering Homer's Odyssey.



#results on my humble P266 with 64MB:
# Without page compression:
# 239 pages in 3.76 seconds = 77 pages per second

# With textOut rather than textLine, i.e. computing width
# of every word as we would for wrapping:
# 239 pages in 10.83 seconds = 22 pages per second

# With page compression and textLine():
# 239 pages in 39.39 seconds = 6 pages per second

from reportlab.pdfgen import canvas
import time, os, sys

#find out what platform we are on and whether accelerator is
#present, in order to print this as part of benchmark info.
try:
    import _rl_accel
    ACCEL = 1
except ImportError:
    ACCEL = 0




from reportlab.lib.units import inch, cm    #1inch = inch , 1cm= cm 
from reportlab.lib.pagesizes import A4      #A4 is tuple(hight, width), check other sizes, dir(reportlab.lib.pagesizes)

#precalculate some basics
top_margin = A4[1] - inch
bottom_margin = inch
left_margin = inch
right_margin = A4[0] - inch
frame_width = right_margin - left_margin


def drawPageFrame(canv):
    canv.line(left_margin, top_margin, right_margin, top_margin) #x1,y1, x2,y2
    canv.setFont('Times-Italic',12)
    canv.drawString(left_margin, top_margin + 2, "Homer's Odyssey")
    canv.line(left_margin, top_margin, right_margin, top_margin)
    canv.line(left_margin, bottom_margin, right_margin, bottom_margin)
    canv.drawCentredString(0.5*A4[0], 0.5 * inch,"Page %d" % canv.getPageNumber())



def run(verbose=1):
    if sys.platform[0:4] == 'java':
        impl = 'Jython'
    else:
        impl = 'Python'
    verStr = '%d.%d' % (sys.version_info[0:2])
    if ACCEL:
        accelStr = 'with _rl_accel'
    else:
        accelStr = 'without _rl_accel'
    print 'Benchmark of %s %s %s' % (impl, verStr, accelStr)

    started = time.time()
    canv = canvas.Canvas('odyssey.pdf', invariant=1)
    canv.setPageCompression(1)  #Possible values None, 1 or 0
    drawPageFrame(canv)

    #do some title page stuff
    canv.setFont("Times-Bold", 36)
    canv.drawCentredString(0.5 * A4[0], 7 * inch, "Homer's Odyssey")

    canv.setFont("Times-Bold", 18)
    canv.drawCentredString(0.5 * A4[0], 5 * inch, "Translated by Samuel Burton")

    canv.setFont("Times-Bold", 12)
    tx = canv.beginText(left_margin, 3 * inch)  #x,y 
    tx.textLine("This is a demo-cum-benchmark for PDFgen.  It renders the complete text of Homer's Odyssey")
    tx.textLine("from a text file.  On my humble P266, it does 77 pages per secondwhile creating a 238 page")
    tx.textLine("document.  If it is asked to computer text metrics, measuring the width of each word as ")
    tx.textLine("one would for paragraph wrapping, it still manages 22 pages per second.")
    tx.textLine("")
    tx.textLine("Andy Robinson, Robinson Analytics Ltd.")
    canv.drawText(tx)

    canv.showPage()  #Close the current page and possibly start on a new page
    #on with the text...
    drawPageFrame(canv)

    canv.setFont('Times-Roman', 12)
    tx = canv.beginText(left_margin, top_margin - 0.5*inch)

    for fn in ('odyssey.full.txt','odyssey.txt'):
        if os.path.isfile(fn):
            break

    data = open(fn,'r').readlines()
    for line in data:
        #this just does it the fast way...
        tx.textLine(line.rstrip())

        #page breaking
        y = tx.getY()   #get y coordinate
        if y < bottom_margin + 0.5*inch:
            canv.drawText(tx)
            canv.showPage()
            drawPageFrame(canv)
            canv.setFont('Times-Roman', 12)
            tx = canv.beginText(left_margin, top_margin - 0.5*inch)

            #page
            pg = canv.getPageNumber()
            if verbose and pg % 10 == 0:
                print 'formatted page %d' % canv.getPageNumber()

    if tx:
        canv.drawText(tx)
        canv.showPage()
        drawPageFrame(canv)

    if verbose:
        print 'about to write to disk...'

    canv.save()  #Saves and close the PDF document in the file

    finished = time.time()
    elapsed = finished - started
    pages = canv.getPageNumber()-1
    speed =  pages / elapsed
    fileSize = os.stat('odyssey.pdf')[6] / 1024
    print '%d pages in %0.2f seconds = %0.2f pages per second, file size %d kb' % (
                pages, elapsed, speed, fileSize)
    import md5
    print 'file digest: %s' % md5.md5(open('odyssey.pdf','rb').read()).hexdigest()

if __name__=='__main__':
    quiet = ('-q' in sys.argv)
    run(verbose = not quiet)

    
    
    
    
    
    
    
    
##FOrmatting, Template etc 
class SimpleDocTemplate(filename, **kw)
    Bases: BaseDocTemplate
    1)  The document has a list of data associated with it
        this data should derive from flowables. We'll have
        special classes like PageBreak, FrameBreak to do things
        like forcing a page end etc.

    2)  The document has one or more page templates.

    3)  Each page template has one or more frames.

    4)  The document class provides base methods for handling the
        story events and some reasonable methods for getting the
        story flowables into the frames.

    5)  The document instances can override the base handler routines.

    EXCEPTION: doctemplate.build(...) must be called for most reasonable uses
    since it builds a document using the page template.

    Each document template builds exactly one document into a file specified
    by the filename argument on initialization.

    Possible keyword arguments for the initialization:

    - pageTemplates: A list of templates.  Must be nonempty.  Names
      assigned to the templates are used for referring to them so no two used
      templates should have the same name.  For example you might want one template
      for a title page, one for a section first page, one for a first page of
      a chapter and two more for the interior of a chapter on odd and even pages.
      If this argument is omitted then at least one pageTemplate should be provided
      using the addPageTemplates method before the document is built.
    - pageSize: a 2-tuple or a size constant from reportlab/lib/pagesizes.pu.
      Used by the SimpleDocTemplate subclass which does NOT accept a list of
      pageTemplates but makes one for you; ignored when using pageTemplates.

    - showBoundary: if set draw a box around the frame boundaries.
    - leftMargin:
    - rightMargin:
    - topMargin:
    - bottomMargin:  Margin sizes in points (default 1 inch).  These margins may be
      overridden by the pageTemplates.  They are primarily of interest for the
      SimpleDocumentTemplate subclass.

    - allowSplitting:  If set flowables (eg, paragraphs) may be split across frames or pages
      (default: 1)
    - title: Internal title for document (does not automatically display on any page)
    - author: Internal author for document (does not automatically display on any page)

    def build(self,flowables,onFirstPage=_doNothing, onLaterPages=_doNothing, canvasmaker=canvas.Canvas):
        """build the document using the flowables.  
            Annotate the first page using the onFirstPage
               function and later pages using the onLaterPages function.  
               The onXXX pages should follow the signature

                  def myOnFirstPage(canvas, document):
                      # do annotations and modify the document
                      ...

               The functions can do things like draw logos, page numbers,
               footers, etcetera. They can use external variables to vary
               the look (for example providing page numbering or section names).
        """

class Frame(x1, y1, width,height, leftPadding=6, bottomPadding=6,
            rightPadding=6, topPadding=6, id=None, showBoundary=0,
            overlapAttachedSpace=None,_debug=None):
    '''
    A Frame is a piece of space in a document that is filled by the
    "flowables" in the story.  For example in a book like document most
    pages have the text paragraphs in one or two frames.  For generality
    a page might have several frames (for example for 3 column text or
    for text that wraps around a graphic).

    After creation a Frame is not usually manipulated directly by the
    applications program -- it is used internally by the platypus modules.

    Here is a diagramatid abstraction for the definitional part of a Frame::

                width                    x2,y2
        +---------------------------------+
        | l  top padding                r | h
        | e +-------------------------+ i | e
        | f |                         | g | i
        | t |                         | h | g
        |   |                         | t | h
        | p |                         |   | t
        | a |                         | p |
        | d |                         | a |
        |   |                         | d |
        |   +-------------------------+   |
        |    bottom padding               |
        +---------------------------------+
        (x1,y1) <-- lower left corner

    NOTE!! Frames are stateful objects.  No single frame should be used in
    two documents at the same time (especially in the presence of multithreading.
    '''
    

##Flowables
#flowables are list of class Flowable
Flowables are defined here as objects which know how to determine their size and which
can draw themselves onto a page with respect to a relative "origin" position determined
at a higher level.

#subclass of class Flowable
TraceInfo Flowable XBox Preformatted Image NullDraw Spacer UseUpSpace PageBreak SlowPageBreak
PageBreakIfNotEmpty CondPageBreak KeepTogether Macro CallerMacro ParagraphAndImage FailOnWrap
FailOnDraw HRFlowable PTOContainer KeepInFrame ImageAndFlowables AnchorFlowable FrameBG
FrameSplitter BulletDrawer DDIndenter LIIndenter ListItem ListFlowable TopPadder DocAssign
DocExec DocPara DocAssert DocIf DocWhile SetTopFlowables splitLines splitLine

#Other Flowable
class Paragraph(Flowable):
    Paragraph(text, style, bulletText=None, caseSensitive=1)
        text a string of stuff to go into the paragraph.
        style is a style definition as in reportlab.lib.styles.
        bulletText is an optional bullet defintion.
        caseSensitive set this to 0 if you want the markup tags and their attributes to be case-insensitive.

        This class is a flowable that can format a block of text
        into a paragraph with a given style.

        The paragraph Text can contain XML-like markup including the tags:
        <b> ... </b> - bold
        <i> ... </i> - italics
        <u> ... </u> - underline
        <strike> ... </strike> - strike through
        <super> ... </super> - superscript
        <sub> ... </sub> - subscript
        <font name=fontfamily/fontname color=colorname size=float>
        <span name=fontfamily/fontname color=colorname backcolor=colorname size=float style=stylename>
        <onDraw name=callable label="a label"/>
        <index [name="callablecanvasattribute"] label="a label"/>
        <link>link text</link>
        attributes of links
        size/fontSize=num
        name/face/fontName=name
        fg/textColor/color=color
        backcolor/backColor/bgcolor=color
        dest/destination/target/href/link=target
        <a>anchor text</a>
        attributes of anchors
        fontSize=num
        fontName=name
        fg/textColor/color=color
        backcolor/backColor/bgcolor=color
        href=href
        <a name="anchorpoint"/>
        <unichar name="unicode character name"/>
        <unichar value="unicode code point"/>
        <img src="path" width="1in" height="1in" valign="bottom"/>
                width="w%" --> fontSize*w/100   idea from Roberto Alsina
                height="h%" --> linewidth*h/100 <ralsina@netmanagers.com.ar>

        The whole may be surrounded by <para> </para> tags

        The <b> and <i> tags will work for the built-in fonts (Helvetica
        /Times / Courier).  For other fonts you need to register a family
        of 4 fonts using reportlab.pdfbase.pdfmetrics.registerFont; then
        use the addMapping function to tell the library that these 4 fonts
        form a family e.g.
        from reportlab.lib.fonts import addMapping
        addMapping('Vera', 0, 0, 'Vera')    #normal
        addMapping('Vera', 0, 1, 'Vera-Italic')    #italic
        addMapping('Vera', 1, 0, 'Vera-Bold')    #bold
        addMapping('Vera', 1, 1, 'Vera-BoldItalic')    #italic and bold

        It will also be able to handle any MathML specified Greek characters.
    
class Figure(Flowable):
    Figure(width, height, caption="",
                 captionFont=_baseFontNameI, captionSize=12,
                 background=None,
                 captionTextColor=toColor('black'),
                 captionBackColor=None,
                 border=None,
                 spaceBefore=12,
                 spaceAfter=12,
                 captionGap=None,
                 captionAlign='centre',
                 captionPosition='bottom',
                 hAlign='CENTER',
                 )
                 
    Subclasses of Figure 
        PageFigure(background=None):
            Shows a blank page in a frame, and draws on that.  

        ImageFigure(filename, caption, background=None,scaleFactor=None,hAlign='CENTER',border=None):
            Image with a caption below it


                
##StyleSheet1
styles = reportlab.lib.styles.getSampleStyleSheet()  
#returns StyleSheet1 instance , 
#has ParagraphStyle with name 'Normal','BodyText','Italic','Heading1','Title','Heading2','Heading3','Heading4','Heading5','Heading6',
#'Bullet','Definition','Code' 
#and ListStyle UnorderedList

#Access like dict 
paraStyle = copy.copy(styles["Heading1"])

##ParagraphStyle, has many attributes 
'fontName':_baseFontName,
'fontSize':10,
'leading':12,
'leftIndent':0,
'rightIndent':0,
'firstLineIndent':0,
'alignment':TA_LEFT,
'spaceBefore':0,
'spaceAfter':0,
'bulletFontName':_baseFontName,
'bulletFontSize':10,
'bulletIndent':0,
#'bulletColor':black,
'textColor': black,
'backColor':None,
'wordWrap':None,        #None means do nothing special
                        #CJK use Chinese Line breaking
                        #LTR RTL use left to right / right to left
                        #with support from pyfribi2 if available
'borderWidth': 0,
'borderPadding': 0,
'borderColor': None,
'borderRadius': None,
'allowWidows': 1,
'allowOrphans': 0,
'textTransform':None,   #uppercase lowercase (captitalize not yet) or None or absent
'endDots':None,         #dots on the last line of left/right justified paras
                        #string or object with text and optional fontName, fontSize, textColor & backColor
                        #dy
'splitLongWords':1,     #make best efforts to split long words
'underlineProportion': _baseUnderlineProportion,    #set to non-zero to get proportional
'bulletAnchor': 'start',    #where the bullet is anchored ie start, middle, end or numeric
'justifyLastLine': 0,   #n allow justification on the last line for more than n words 0 means don't bother
'justifyBreaks': 0,     #justify lines broken with <br/>
'spaceShrinkage': spaceShrinkage,   #allow shrinkage of percentage of space to fit on line

##Examples 

from reportlab.platypus import *
from reportlab.lib.units import inch
from reportlab.lib.styles import *
from reportlab.lib.enums import TA_LEFT, TA_RIGHT, TA_CENTER, TA_JUSTIFY

class PlatPropFigure1(PageFigure):
    """This shows a page with a frame on it"""
    def __init__(self):
        PageFigure.__init__(self)
        self.caption = "Figure 1 - a page with a simple frame"
    def drawVirtualPage(self):
        demo1(self.canv)

def demo1(canvas):
    frame = Frame(
                    2*inch,     # x
                    4*inch,     # y at bottom
                    4*inch,     # width
                    5*inch,     # height
                    showBoundary = 1  # helps us see what's going on
                    )
    bodyStyle = ParagraphStyle('Body', fontName=_baseFontName, fontSize=24, leading=28, spaceBefore=6)
    para1 = Paragraph('Spam spam spam spam. ' * 5, bodyStyle)
    para2 = Paragraph('Eggs eggs eggs. ' * 5, bodyStyle)
    mydata = [para1, para2]

    #this does the packing and drawing.  The frame will consume
    #items from the front of the list as it prints them
    frame.addFromList(mydata,canvas)


c  = Canvas('figures.pdf')
f = Frame(inch, inch, 6*inch, 9*inch, showBoundary=1)
v = PlatPropFigure1()
v.captionTextColor = toColor('blue')
v.captionBackColor = toColor('lightyellow')
f.addFromList([v],c)
c.save()

##Example -fodyssey.py handle paragraph formatting 
#REPORTLAB_TEST_SCRIPT
import sys, copy, os
from reportlab.platypus import *
from reportlab.lib.units import inch
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib.enums import TA_LEFT, TA_RIGHT, TA_CENTER, TA_JUSTIFY

import reportlab.rl_config
reportlab.rl_config.invariant = 1



def myFirstPage(canvas, doc):
    canvas.saveState()  #Save the current graphics state to be restored later by restoreState
    canvas.restoreState()

def myLaterPages(canvas, doc):
    canvas.saveState()
    canvas.setFont('Times-Roman',9)
    canvas.drawString(inch, 0.75 * inch, "Page %d" % doc.page)
    canvas.restoreState()

def go():
    doc = SimpleDocTemplate('fodyssey.pdf',showBoundary='showboundary' in sys.argv)
    doc.allowSplitting = not 'nosplitting' in sys.argv
    doc.build(Elements,myFirstPage,myLaterPages)

Elements = []


Title = "The Odyssey"
Author = "Homer"

styles = getSampleStyleSheet()  
#returns StyleSheet1 instance , 
#has ParagraphStyle with name 'Normal','BodyText','Italic','Heading1','Title','Heading2','Heading3','Heading4','Heading5','Heading6',
#'Bullet','Definition','Code' 
#and ListStyle UnorderedList

#Access like dict 
ChapterStyle = copy.copy(styles["Heading1"])


ChapterStyle.alignment = TA_CENTER
ChapterStyle.fontsize = 16
InitialStyle = copy.deepcopy(ChapterStyle)
InitialStyle.fontsize = 16
InitialStyle.leading = 20
PreStyle = styles["Code"]

def newPage():
    Elements.append(PageBreak())

def chapter(txt, style=ChapterStyle):
    newPage()
    Elements.append(Paragraph(txt, style))
    Elements.append(Spacer(0.2*inch, 0.3*inch))

def fTitle(txt,style=InitialStyle):
    Elements.append(Paragraph(txt, style))

ParaStyle = copy.deepcopy(styles["Normal"])
ParaStyle.spaceBefore = 0.1*inch
if 'right' in sys.argv:
    ParaStyle.alignment = TA_RIGHT
elif 'left' in sys.argv:
    ParaStyle.alignment = TA_LEFT
elif 'justify' in sys.argv:
    ParaStyle.alignment = TA_JUSTIFY
elif 'center' in sys.argv or 'centre' in sys.argv:
    ParaStyle.alignment = TA_CENTER
else:
    ParaStyle.alignment = TA_JUSTIFY

def spacer(inches):
    Elements.append(Spacer(0.1*inch, inches*inch))

def p(txt, style=ParaStyle):
    Elements.append(Paragraph(txt, style))

def pre(txt, style=PreStyle):
    spacer(0.1)
    p = Preformatted(txt, style)
    Elements.append(p)

def parseOdyssey(fn):
    from time import time
    E = []
    t0=time()
    text = open(fn,'r').read()
    i0 = text.index('Book I')
    endMarker = 'covenant of peace between the two contending parties.'
    i1 = text.index(endMarker)+len(endMarker)
    PREAMBLE=map(str.strip,text[0:i0].split('\n'))
    L=map(str.strip,text[i0:i1].split('\n'))
    POSTAMBLE=map(str.strip,text[i1:].split('\n'))

    def ambleText(L):
        while L and not L[0]: L.pop(0)
        while L:
            T=[]
            while L and L[0]:
                T.append(L.pop(0))
            yield T
            while L and not L[0]: L.pop(0)

    def mainText(L):
        while L:
            B = L.pop(0)
            while not L[0]: L.pop(0)
            T=[]
            while L and L[0]:
                T.append(L.pop(0))
            while not L[0]: L.pop(0)
            P = []
            while L and not (L[0].startswith('Book ') and len(L[0].split())==2):
                E=[]
                while L and L[0]:
                    E.append(L.pop(0))
                P.append(E)
                if L:
                    while not L[0]: L.pop(0)
            yield B,T,P

    t1 = time()
    print "open(%s,'r').read() took %.4f seconds" %(fn,t1-t0)

    E.append([spacer,2])
    E.append([fTitle,'<font color=red>%s</font>' % Title, InitialStyle])
    E.append([fTitle,'<font size=-4>by</font> <font color=green>%s</font>' % Author, InitialStyle])

    for T in ambleText(PREAMBLE):
        E.append([p,'\n'.join(T)])

    for (B,T,P) in mainText(L):
        E.append([chapter,B])
        E.append([p,'<font size="+1" color="Blue"><b>%s</b></font>' % '\n'.join(T),ParaStyle])
        for x in P:
            E.append([p,' '.join(x)])
    firstPre = 1
    for T in ambleText(POSTAMBLE):
        E.append([p,'\n'.join(T)])  #add p, calling is later 

    t3 = time()
    print "Parsing into memory took %.4f seconds" %(t3-t1)
    del L
    t4 = time()
    print "Deleting list of lines took %.4f seconds" %(t4-t3)
    for i in range(len(E)):
        E[i][0](*E[i][1:])  ##call all functions 
    t5 = time()
    print "Moving into platypus took %.4f seconds" %(t5-t4)
    del E
    t6 = time()
    print "Deleting list of actions took %.4f seconds" %(t6-t5)
    go()
    t7 = time()
    print "saving to PDF took %.4f seconds" %(t7-t6)
    print "Total run took %.4f seconds"%(t7-t0)

for fn in ('odyssey.full.txt','odyssey.txt'):
    if os.path.isfile(fn):
        break
if __name__=='__main__':
    parseOdyssey(fn)
    

###Examples of Table 
class Table(Flowable):
    Table(data, colWidths=None, rowHeights=None, style=None,
                repeatRows=0, repeatCols=0, splitByRow=1, emptyTableAction=None, ident=None,
                hAlign=None,vAlign=None, normalizedData=0, cellStyles=None, rowSplitRange=None,
                spaceBefore=None,spaceAfter=None, longTableOptimize=None)                 

All the cell values should be convertible to strings; 
embedded newline '\\n' characters cause the value to wrap (ie are like a traditional linefeed).

None values in the sequence of row heights or column widths, 
mean that the corresponding rows or columns should be automatically sized.      
#Subclass 
class LongTable(Table):
    '''Henning von Bargen's changes will be active'''
    _longTableOptimize = 1


#Various styles - reportlab.platypus.tables
GRID_STYLE = TableStyle(
    [('GRID', (0,0), (-1,-1), 0.25, colors.black),
     ('ALIGN', (1,1), (-1,-1), 'RIGHT')]
    )
BOX_STYLE = TableStyle(
    [('BOX', (0,0), (-1,-1), 0.50, colors.black),
     ('ALIGN', (1,1), (-1,-1), 'RIGHT')]
    )
LABELED_GRID_STYLE = TableStyle(
    [('INNERGRID', (0,0), (-1,-1), 0.25, colors.black),
     ('BOX', (0,0), (-1,-1), 2, colors.black),
     ('LINEBELOW', (0,0), (-1,0), 2, colors.black),
     ('LINEAFTER', (0,0), (0,-1), 2, colors.black),
     ('ALIGN', (1,1), (-1,-1), 'RIGHT')]
    )
COLORED_GRID_STYLE = TableStyle(
    [('INNERGRID', (0,0), (-1,-1), 0.25, colors.black),
     ('BOX', (0,0), (-1,-1), 2, colors.red),
     ('LINEBELOW', (0,0), (-1,0), 2, colors.black),
     ('LINEAFTER', (0,0), (0,-1), 2, colors.black),
     ('ALIGN', (1,1), (-1,-1), 'RIGHT')]
    )
LIST_STYLE = TableStyle(
    [('LINEABOVE', (0,0), (-1,0), 2, colors.green),
     ('LINEABOVE', (0,1), (-1,-1), 0.25, colors.black),
     ('LINEBELOW', (0,-1), (-1,-1), 2, colors.green),
     ('ALIGN', (1,1), (-1,-1), 'RIGHT')]
    )
    
##TableStyle command 
#command syntax 
(COMMAND, (start_column, start_row), (end_column, end_row), value1, value2,...)
#Note if start_*,end_* are negative, then number is taking from last 

#Line commands , values are 'weight, colour, cap, dashes, join'
'GRID','BOX','OUTLINE','INNERGRID','LINEBELOW','LINEABOVE','LINEBEFORE','LINEAFTER'

#Cell syle command - value is only one corresponding value, except for 'FONT'
'FONT'(fontname, fontsize,leading), 'FONTNAME'/'FACE', 'SIZE'/'FONTSIZE', 'LEADING', 'TEXTCOLOR'
'ALIGN'/'ALIGNMENT', 'VALIGN', 
'LEFTPADDING', 'RIGHTPADDING', 'TOPPADDING', 'BOTTOMPADDING'
'HREF', 'DESTINATION'

#Other commands 
'BACKGROUND' : value is single color or for gradient fill, values are 'HORIZONTAL'/'VERTICAL', color1, color2
'ROWBACKGROUNDS','COLBACKGROUNDS'  : values are color1, color2,...., these would be cycled 
'SPAN'  : no value, used for spanning number of rows or columns 
'NOSPLIT': no value, used for nosplit range of rows or columns 







#Example - basic 
from reportlab.lib import colors
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Table, TableStyle
 
doc = SimpleDocTemplate("simple_table.pdf", pagesize=letter)
# container for the 'Flowable' objects
elements = []
 
data= [['00', '01', '02', '03', '04'],
       ['10', '11', '12', '13', '14'],
       ['20', '21', '22', '23', '24'],
       ['30', '31', '32', '33', '34']]
t=Table(data)
t.setStyle(TableStyle([('BACKGROUND',(1,1),(-2,-2),colors.green),
                       ('TEXTCOLOR',(0,0),(1,-1),colors.red)]))
elements.append(t)
# write the document to disk
doc.build(elements)

#Example - adding grid 
from reportlab.lib import colors
from reportlab.lib.pagesizes import letter, inch
from reportlab.platypus import SimpleDocTemplate, Table, TableStyle
 
doc = SimpleDocTemplate("simple_table_grid.pdf", pagesize=letter)
# container for the 'Flowable' objects
elements = []
 
data= [['00', '01', '02', '03', '04'],
       ['10', '11', '12', '13', '14'],
       ['20', '21', '22', '23', '24'],
       ['30', '31', '32', '33', '34']]
t=Table(data,5*[0.4*inch], 4*[0.4*inch])
t.setStyle(TableStyle([('ALIGN',(1,1),(-2,-2),'RIGHT'),
                       ('TEXTCOLOR',(1,1),(-2,-2),colors.red),
                       ('VALIGN',(0,0),(0,-1),'TOP'),
                       ('TEXTCOLOR',(0,0),(0,-1),colors.blue),
                       ('ALIGN',(0,-1),(-1,-1),'CENTER'),
                       ('VALIGN',(0,-1),(-1,-1),'MIDDLE'),
                       ('TEXTCOLOR',(0,-1),(-1,-1),colors.green),
                       ('INNERGRID', (0,0), (-1,-1), 0.25, colors.black),
                       ('BOX', (0,0), (-1,-1), 0.25, colors.black),
                       ]))
 
elements.append(t)
# write the document to disk
doc.build(elements)
    
#Example - complex value 
from reportlab.lib import colors
from reportlab.lib.pagesizes import letter, inch
from reportlab.platypus import Image, Paragraph, SimpleDocTemplate, Table
from reportlab.lib.styles import getSampleStyleSheet
 
doc = SimpleDocTemplate("complex_cell_values.pdf", pagesize=letter)
# container for the 'Flowable' objects
elements = []
 
styleSheet = getSampleStyleSheet()
 
I = Image('replogo.gif')
I.drawHeight = 1.25*inch*I.drawHeight / I.drawWidth
I.drawWidth = 1.25*inch
P0 = Paragraph('''
               <b>A pa<font color=red>r</font>a<i>graph</i></b>
               <super><font color=yellow>1</font></super>''',
               styleSheet["BodyText"])
P = Paragraph('''
    <para align=center spaceb=3>The <b>ReportLab Left
    <font color=red>Logo</font></b>
    Image</para>''',
    styleSheet["BodyText"])
data= [['A', 'B', 'C', P0, 'D'],
       ['00', '01', '02', [I,P], '04'],
       ['10', '11', '12', [P,I], '14'],
       ['20', '21', '22', '23', '24'],
       ['30', '31', '32', '33', '34']]
 
t=Table(data,style=[('GRID',(1,1),(-2,-2),1,colors.green),
                    ('BOX',(0,0),(1,-1),2,colors.red),
                    ('LINEABOVE',(1,2),(-2,2),1,colors.blue),
                    ('LINEBEFORE',(2,1),(2,-2),1,colors.pink),
                    ('BACKGROUND', (0, 0), (0, 1), colors.pink),
                    ('BACKGROUND', (1, 1), (1, 2), colors.lavender),
                    ('BACKGROUND', (2, 2), (2, 3), colors.orange),
                    ('BOX',(0,0),(-1,-1),2,colors.black),
                    ('GRID',(0,0),(-1,-1),0.5,colors.black),
                    ('VALIGN',(3,0),(3,0),'BOTTOM'),
                    ('BACKGROUND',(3,0),(3,0),colors.limegreen),
                    ('BACKGROUND',(3,1),(3,1),colors.khaki),
                    ('ALIGN',(3,1),(3,1),'CENTER'),
                    ('BACKGROUND',(3,2),(3,2),colors.beige),
                    ('ALIGN',(3,2),(3,2),'LEFT'),
])
t._argW[3]=1.5*inch
 
elements.append(t)
# write the document to disk
doc.build(elements)



###Drawing Graphics/chart 
#Many exmaples - Lib\site-packages\reportlab\graphics\samples


##Graphics 
from reportlab.graphics.shapes import *
from reportlab.lib import colors
from reportlab.lib.units import cm
from reportlab.lib.utils import asNative
from reportlab.pdfgen.canvas import Canvas
from reportlab.pdfbase.pdfmetrics import stringWidth
from reportlab.platypus import Flowable
from reportlab.graphics.shapes import *
from reportlab.graphics.renderPDF import _PDFRenderer

def getDrawing1():
    """Hello World, on a rectangular background"""

    D = Drawing(400, 200)
    D.add(Rect(50, 50, 300, 100, fillColor=colors.yellow))  #round corners
    D.add(String(180,100, 'Hello World', fillColor=colors.red))


    return D


def getDrawing2():
    """This demonstrates the basic shapes.  There are
    no groups or references.  Each solid shape should have
    a purple fill."""
    D = Drawing(400, 200) #, fillColor=colors.purple)

    D.add(Line(10,10,390,190))
    D.add(Circle(100,100,20, fillColor=colors.purple))
    D.add(Circle(200,100,20, fillColor=colors.purple))
    D.add(Circle(300,100,20, fillColor=colors.purple))

    D.add(Wedge(330,100,40, -10,40, fillColor=colors.purple))

    D.add(PolyLine([120,10,130,20,140,10,150,20,160,10,
                    170,20,180,10,190,20,200,10]))

    D.add(Polygon([300,20,350,20,390,80,300,75, 330, 40]))

    D.add(Ellipse(50, 150, 40, 20))

    D.add(Rect(120, 150, 60, 30,
               strokeWidth=10,
               strokeColor=colors.red,
               fillColor=colors.yellow))  #square corners

    D.add(Rect(220, 150, 60, 30, 10, 10))  #round corners

    D.add(String(10,50, 'Basic Shapes', fillColor=colors.black))

    return D

    
def getDrawing07():
    """This tests the ability to translate and rotate groups.  The first set of axes should be
    near the bottom left of the drawing.  The second should be rotated counterclockwise
    by 15 degrees.  The third should be rotated by 30 degrees."""
    D = Drawing(400, 200)

    Axis = Group(
        Line(0,0,100,0), #x axis
        Line(0,0,0,50),   # y axis
        Line(0,10,10,10), #ticks on y axis
        Line(0,20,10,20),
        Line(0,30,10,30),
        Line(0,40,10,40),
        Line(10,0,10,10), #ticks on x axis
        Line(20,0,20,10),
        Line(30,0,30,10),
        Line(40,0,40,10),
        Line(50,0,50,10),
        Line(60,0,60,10),
        Line(70,0,70,10),
        Line(80,0,80,10),
        Line(90,0,90,10),
        String(20, 35, 'Axes', fill=colors.black)
        )

    firstAxisGroup = Group(Axis)
    firstAxisGroup.translate(10,10)
    D.add(firstAxisGroup)

    secondAxisGroup = Group(Axis)
    secondAxisGroup.translate(150,10)
    secondAxisGroup.rotate(15)

    D.add(secondAxisGroup)


    thirdAxisGroup = Group(Axis, transform=mmult(translate(300,10), rotate(30)))
    D.add(thirdAxisGroup)

    return D
    
    
def writePDF(drawings):
    "Create and save a PDF file containing some drawings."

    pdfPath = os.path.splitext(sys.argv[0])[0] + '.pdf'
    c = Canvas(pdfPath)
    c.setFont(_FONTS[0], 32)
    c.drawString(80, 750, 'ReportLab Graphics-Shapes Test')

    # Print drawings in a loop, with their doc strings.
    c.setFont(_FONTS[0], 12)
    y = 740
    i = 1
    for (drawing, docstring, funcname) in drawings:
        if y < 300:  # Allows 5-6 lines of text.
            c.showPage()
            y = 740
        # Draw a title.
        y = y - 30
        c.setFont(_FONTS[2],12)
        c.drawString(80, y, '%s (#%d)' % (funcname, i))
        c.setFont(_FONTS[0],12)
        y = y - 14
        textObj = c.beginText(80, y)
        textObj.textLines(docstring)
        c.drawText(textObj)
        y = textObj.getY()
        y = y - drawing.height
        drawing.drawOn(c, 80, y)
        i = i + 1

    c.save()
    print('wrote %s ' % pdfPath)

drawings = [(getDrawing1(), getDrawing1.__doc__, getDrawing1.__name__), ...]
writePDF(drawings)


##Example of Line charts 
#All the plot examples are similar
#use below class 
class LinePlot(AbstractLineChart)
class LinePlot3D(LinePlot)
class SimpleTimeSeriesPlot(LinePlot)
class GridLinePlot(SimpleTimeSeriesPlot)
class AreaLinePlot(LinePlot)
class SplitLinePlot(AreaLinePlot)
class ScatterPlot(LinePlot)

class BarChart(PlotArea)
class VerticalBarChart(BarChart)
class HorizontalBarChart(BarChart)
class BarChart3D(BarChart)
class VerticalBarChart3D(BarChart3D,VerticalBarChart)
class HorizontalBarChart3D(BarChart3D,HorizontalBarChart)
        
class Pie(AbstractPieChart)
class Pie3d(Pie)

class SpiderChart(PlotArea)

#LinePlot() demo
def demo(self):
    """Shows basic use of a line chart."""

    drawing = Drawing(400, 200)

    data = [
        ((1,1), (2,2), (2.5,1), (3,3), (4,5)),
        ((1,2), (2,3), (2.5,2), (3.5,5), (4,6))
        ]

    lp = LinePlot()

    lp.x = 50
    lp.y = 50
    lp.height = 125
    lp.width = 300
    lp.data = data
    lp.joinedLines = 1
    lp.lineLabelFormat = '%2.0f'
    lp.strokeColor = colors.black

    lp.lines[0].strokeColor = colors.red
    lp.lines[0].symbol = makeMarker('FilledCircle')
    lp.lines[1].strokeColor = colors.blue
    lp.lines[1].symbol = makeMarker('FilledDiamond')

    lp.xValueAxis.valueMin = 0
    lp.xValueAxis.valueMax = 5
    lp.xValueAxis.valueStep = 1

    lp.yValueAxis.valueMin = 0
    lp.yValueAxis.valueMax = 7
    lp.yValueAxis.valueStep = 1

    drawing.add(lp)

    return drawing

#Exmaple line_charts.py 
from reportlab.graphics.charts.legends import Legend
from reportlab.graphics.charts.lineplots import LinePlot
from reportlab.graphics.shapes import Drawing, _DrawingEditorMixin, String
from reportlab.graphics.charts.textlabels import Label
from reportlab.graphics.samples.excelcolors import *

class LineChart(_DrawingEditorMixin,Drawing):
    def __init__(self,width=200,height=150,*args,**kw):
        Drawing.__init__(self,width,height,*args,**kw)
        self._add(self,LinePlot(),name='chart',validate=None,desc="The main chart")
        self.chart.width      = 115
        self.chart.height     = 80
        self.chart.x          = 30
        self.chart.y          = 40
        self.chart.lines[0].strokeColor = color01
        self.chart.lines[1].strokeColor = color02
        #below are not used
        self.chart.lines[2].strokeColor = color03
        self.chart.lines[3].strokeColor = color04
        self.chart.lines[4].strokeColor = color05
        self.chart.lines[5].strokeColor = color06
        self.chart.lines[6].strokeColor = color07
        self.chart.lines[7].strokeColor = color08
        self.chart.lines[8].strokeColor = color09
        self.chart.lines[9].strokeColor = color10
        self.chart.fillColor         = backgroundGrey
        self.chart.lineLabels.fontName              = 'Helvetica'
        self.chart.xValueAxis.labels.fontName       = 'Helvetica'
        self.chart.xValueAxis.labels.fontSize       = 7
        self.chart.xValueAxis.forceZero             = 0
        #Data to be plotted, list of (lists of) x/y tuples.
        #Only two lines are drawn 
        self.chart.data             = [((0, 50), (100,100), (200,200), (250,210), (300,300), (400,500)), ((0, 150), (100,200), (200,300), (250,200), (300,400), (400, 600))]
        self.chart.xValueAxis.avoidBoundFrac           = 1
        self.chart.xValueAxis.gridEnd                  = 115
        self.chart.xValueAxis.tickDown                 = 3
        self.chart.xValueAxis.visibleGrid              = 1
        self.chart.yValueAxis.tickLeft              = 3
        self.chart.yValueAxis.labels.fontName       = 'Helvetica'
        self.chart.yValueAxis.labels.fontSize       = 7
        self._add(self,Label(),name='Title',validate=None,desc="The title at the top of the chart")
        self.Title.fontName   = 'Helvetica-Bold'
        self.Title.fontSize   = 7
        self.Title.x          = 100
        self.Title.y          = 135
        self.Title._text      = 'Chart Title'
        self.Title.maxWidth   = 180
        self.Title.height     = 20
        self.Title.textAnchor ='middle'
        self._add(self,Legend(),name='Legend',validate=None,desc="The legend or key for the chart")
        self.Legend.colorNamePairs = [(color01, 'Widgets'), (color02, 'Sprockets')]
        self.Legend.fontName       = 'Helvetica'
        self.Legend.fontSize       = 7
        self.Legend.x              = 153
        self.Legend.y              = 85
        self.Legend.dxTextSpace    = 5
        self.Legend.dy             = 5
        self.Legend.dx             = 5
        self.Legend.deltay         = 5
        self.Legend.alignment      ='right'
        self._add(self,Label(),name='XLabel',validate=None,desc="The label on the horizontal axis")
        self.XLabel.fontName       = 'Helvetica'
        self.XLabel.fontSize       = 7
        self.XLabel.x              = 85
        self.XLabel.y              = 10
        self.XLabel.textAnchor     ='middle'
        self.XLabel.maxWidth       = 100
        self.XLabel.height         = 20
        self.XLabel._text          = "X Axis"
        self._add(self,Label(),name='YLabel',validate=None,desc="The label on the vertical axis")
        self.YLabel.fontName       = 'Helvetica'
        self.YLabel.fontSize       = 7
        self.YLabel.x              = 12
        self.YLabel.y              = 80
        self.YLabel.angle          = 90
        self.YLabel.textAnchor     ='middle'
        self.YLabel.maxWidth       = 100
        self.YLabel.height         = 20
        self.YLabel._text          = "Y Axis"
        self.chart.yValueAxis.forceZero           = 1
        self.chart.xValueAxis.forceZero           = 1
        self._add(self,0,name='preview',validate=None,desc=None)

if __name__=="__main__": #NORUNTESTS
    LineChart().save(formats=['pdf'],outDir=None,fnRoot='line_chart')


    
    
    
###############Eventlet ##############
$ pip install eventlet

Eventlet is a concurrent networking library 
• It uses epoll or kqueue or libevent for highly scalable non-blocking I/O.
• uses Coroutines , called greenthreads , basically a python methods used inside spawn 
• Green threads cooperatively yield to each other instead of preemptively being scheduled, hence no locking is required
• The event dispatch is implicit


#Greening a application, Use only Green enabled standard  library
#eventlet.green has asynchat, asyncore, builtin(for file open), ftplib. httplib. os, 
#Queue, select, SimpleHTTPServer, socket, subprocess, ssl, threading, time, urllib2

#Example 
from eventlet.green import socket
from eventlet.green import threading
from eventlet.green import asyncore

#Or patch a module using 
eventlet.import_patched(module_name, *additional_modules, **kw_additional_modules)

#Example 
import eventlet
httplib2 = eventlet.import_patched('httplib2')

#or patch key modules using 
eventlet.patcher.monkey_patch(os=None, select=None, socket=None, thread=None, time=None, psycopg=None)

#Example 
import eventlet
eventlet.monkey_patch()
eventlet.monkey_patch(socket=True, select=True)  #selective 

#Basic API - greenthread – Green Thread Implementation
import eventlet 

eventlet.sleep(seconds=0)
    Yield control to another eligible coroutine until at least seconds have elapsed.
    Called inside a function which is used for spawn 
    Don't use time.sleep()

eventlet.spawn(func, *args, **kwargs)
    Create a greenthread to run func(*args, **kwargs). 
    Returns a GreenThread object which you can use to get the results of the call

eventlet.spawn_n(func, *args, **kwargs)
    Same as spawn(), but no return 

eventlet.spawn_after(seconds, func, *args, **kwargs)
    Spawns func after seconds have elapsed. It runs as scheduled even if the current greenthread has completed

eventlet.spawn_after_local(seconds, func, *args, **kwargs)
    Spawns func after seconds have elapsed. The function will NOT be called if the current greenthread has exited

#GreenThread contains following methods 
cancel(*throw_args)     							Kills the greenthread using kill(), but only if it hasn’t already started running.
kill(*throw_args)									Kills the greenthread using kill(). 
unlink(func, *curried_args, **curried_kwargs)  		Remove linked function set by link()
wait()												Returns the result of the main function of this GreenThread
link(func, *curried_args, **curried_kwargs) 		Set up a 'def func(gt, [curried args/kwargs])' to be called with the results of the GreenThread.
													Inside function to get result, use gt.wait()


# greenpool – Green Thread Pools - Create Pool and use them for greenthread creation 
pool = eventlet.GreenPool(size=1000)  #Returns GreenPool, is a pool of green threads
#pool Contains below methods 
free()								Returns the number of greenthreads available for use.
imap(function, *iterables)  		same as map(fun, *iterables) returns iterator of results 
resize(new_size)					Change the max number of greenthreads doing work at any given time.
running()							Returns the number of greenthreads that are currently executing functions in the GreenPool.
spawn(function, *args, **kwargs)	Run the function with its arguments in its own green thread. Returns the GreenThread object 
spawn_n(function, *args, **kwargs)  Same as spawn, but no return 
starmap(func, seq)					func(*seq[0]), func(*seq[1]), ... eg starmap(pow, [(2,5), (3,2), (10,3)]) --> 32 9 1000
waitall()							Waits until all greenthreads in the pool are finished working.
waiting()							Return the number of greenthreads waiting to spawn.

#GreenPile - pile acts as a iterator of return values from the functions
pile = eventlet.GreenPile(pool_or_size=1000) #Construct a GreenPile with an existing GreenPool object or standalone 
#pile conatins below methods
next()							Wait for the next result, suspending the current greenthread until it is available. 
spawn(func, *args, **kw)		Runs func in its own green thread, with the result available by iterating over the GreenPile object.

#Example of GreenPile 

import eventlet
from eventlet.green import socket


def geturl(url):
    c = socket.socket()
    ip = socket.gethostbyname(url)
    c.connect((ip, 80))
    print('%s connected' % url)
    c.sendall('GET /\r\n\r\n')
    return c.recv(1024)


urls = ['www.google.com', 'www.yandex.ru', 'www.python.org']
pile = eventlet.GreenPile()
for x in urls:
    pile.spawn(geturl, x)

# note that the pile acts as a collection of return values from the functions
# if any exceptions are raised by the function they'll get raised here
for url, result in zip(urls, pile):
    print('%s: %s' % (url, repr(result)[:50]))


#Example web crawler - Client Pattern

import eventlet
from eventlet.green import urllib2  #import a cooperatively-yielding version of urllib2, uses green sockets for its communication


urls = [
    "https://www.google.com/intl/en_ALL/images/logo.gif",
    "http://python.org/images/python-logo.gif",
    "http://us.i1.yimg.com/us.yimg.com/i/ww/beta/y3.gif",
]


def fetch(url):
    print("opening", url)
    body = urllib2.urlopen(url).read()
    print("done with", url)
    return url, body


pool = eventlet.GreenPool(200)  #constructs a GreenPool of a 200 green threads
for url, body in pool.imap(fetch, urls):
    print("got body from", url, "of length", len(body))

	

#Example Echo server - Server Pattern

from __future__ import print_function

import eventlet


def handle(fd):
    print("client connected")
    while True:
        # pass through every non-eof line
        x = fd.readline()
        if not x:
            break
        fd.write(x)
        fd.flush()
        print("echoed", x, end=' ')
    print("client disconnected")

print("server socket listening on port 6000")
server = eventlet.listen(('0.0.0.0', 6000))
pool = eventlet.GreenPool()
while True:
    try:
        new_sock, address = server.accept()
        print("accepted", address)
        pool.spawn_n(handle, new_sock.makefile('rw'))  #launches a green thread to handle the new client
    except (SystemExit, KeyboardInterrupt):
        break


##Socket Connect

"""Spawn multiple workers and collect their results.

Demonstrates how to use the eventlet.green.socket module.
"""
from __future__ import print_function

import eventlet
from eventlet.green import socket


def geturl(url):
    c = socket.socket()
    ip = socket.gethostbyname(url)
    c.connect((ip, 80))
    print('%s connected' % url)
    c.sendall('GET /\r\n\r\n')
    return c.recv(1024)


urls = ['www.google.com', 'www.yandex.ru', 'www.python.org']
pile = eventlet.GreenPile()
for x in urls:
    pile.spawn(geturl, x)

# note that the pile acts as a collection of return values from the functions
# if any exceptions are raised by the function they'll get raised here
for url, result in zip(urls, pile):
    print('%s: %s' % (url, repr(result)[:50]))



##Multi-User Chat Server
#This is a little different from the echo server, 
#in that it broadcasts the messages to all participants, not just the sender.


import eventlet
from eventlet.green import socket

PORT = 3001
participants = set()


def read_chat_forever(writer, reader):
    line = reader.readline()
    while line:
        print("Chat:", line.strip())
        for p in participants:
            try:
                if p is not writer:  # Don't echo
                    p.write(line)
                    p.flush()
            except socket.error as e:
                # ignore broken pipes, they just mean the participant
                # closed its connection already
                if e[0] != 32:
                    raise
        line = reader.readline()
    participants.remove(writer)
    print("Participant left chat.")

try:
    print("ChatServer starting up on port %s" % PORT)
    server = eventlet.listen(('0.0.0.0', PORT))
    while True:
        new_connection, address = server.accept()
        print("Participant joined chat.")
        new_writer = new_connection.makefile('w')
        participants.add(new_writer)
        eventlet.spawn_n(read_chat_forever,
                         new_writer,
                         new_connection.makefile('r'))
except (KeyboardInterrupt, SystemExit):
    print("ChatServer exiting.")



##Feed Scraper
#This example requires Feedparser to be installed or on the PYTHONPATH.


"""A simple web server that accepts POSTS containing a list of feed urls,
and returns the titles of those feeds.
"""
import eventlet
feedparser = eventlet.import_patched('feedparser')

# the pool provides a safety limit on our concurrency
pool = eventlet.GreenPool()


def fetch_title(url):
    d = feedparser.parse(url)
    return d.feed.get('title', '')


def app(environ, start_response):
    if environ['REQUEST_METHOD'] != 'POST':
        start_response('403 Forbidden', [])
        return []

    # the pile collects the result of a concurrent operation -- in this case,
    # the collection of feed titles
    pile = eventlet.GreenPile(pool)
    for line in environ['wsgi.input'].readlines():
        url = line.strip()
        if url:
            pile.spawn(fetch_title, url)
    # since the pile is an iterator over the results,
    # you can use it in all sorts of great Pythonic ways
    titles = '\n'.join(pile)
    start_response('200 OK', [('Content-type', 'text/plain')])
    return [titles]


if __name__ == '__main__':
    from eventlet import wsgi
    wsgi.server(eventlet.listen(('localhost', 9010)), app)



##Port Forwarder

""" This is an incredibly simple port forwarder from port 7000 to 22 on
localhost.  It calls a callback function when the socket is closed, to
demonstrate one way that you could start to do interesting things by
starting from a simple framework like this.
"""

import eventlet


def closed_callback():
    print("called back")


def forward(source, dest, cb=lambda: None):
    """Forwards bytes unidirectionally from source to dest"""
    while True:
        d = source.recv(32384)
        if d == '':
            cb()
            break
        dest.sendall(d)

listener = eventlet.listen(('localhost', 7000))
while True:
    client, addr = listener.accept()
    server = eventlet.connect(('localhost', 22))
    # two unidirectional forwarders make a bidirectional one
    eventlet.spawn_n(forward, client, server, closed_callback)
    eventlet.spawn_n(forward, server, client)



##Recursive Web Crawler

"""This is a recursive web crawler.  Don't go pointing this at random sites;
it doesn't respect robots.txt and it is pretty brutal about how quickly it
fetches pages.

The code for this is very short; this is perhaps a good indication
that this is making the most effective use of the primitves at hand.
The fetch function does all the work of making http requests,
searching for new urls, and dispatching new fetches.  The GreenPool
acts as sort of a job coordinator (and concurrency controller of
course).
"""
from __future__ import with_statement

from eventlet.green import urllib2
import eventlet
import re

# http://daringfireball.net/2009/11/liberal_regex_for_matching_urls
url_regex = re.compile(r'\b(([\w-]+://?|www[.])[^\s()<>]+(?:\([\w\d]+\)|([^[:punct:]\s]|/)))')


def fetch(url, seen, pool):
    """Fetch a url, stick any found urls into the seen set, and
    dispatch any new ones to the pool."""
    print("fetching", url)
    data = ''
    with eventlet.Timeout(5, False):
        data = urllib2.urlopen(url).read()
    for url_match in url_regex.finditer(data):
        new_url = url_match.group(0)
        # only send requests to eventlet.net so as not to destroy the internet
        if new_url not in seen and 'eventlet.net' in new_url:
            seen.add(new_url)
            # while this seems stack-recursive, it's actually not:
            # spawned greenthreads start their own stacks
            pool.spawn_n(fetch, new_url, seen, pool)


def crawl(start_url):
    """Recursively crawl starting from *start_url*.  Returns a set of
    urls that were found."""
    pool = eventlet.GreenPool()
    seen = set()
    fetch(start_url, seen, pool)
    pool.waitall()
    return seen

seen = crawl("http://eventlet.net")
print("I saw these urls:")
print("\n".join(seen))



##Producer Consumer Web Crawler

"""This is a recursive web crawler.  Don't go pointing this at random sites;
it doesn't respect robots.txt and it is pretty brutal about how quickly it
fetches pages.

This is a kind of "producer/consumer" example; the fetch function produces
jobs, and the GreenPool itself is the consumer, farming out work concurrently.
It's easier to write it this way rather than writing a standard consumer loop;
GreenPool handles any exceptions raised and arranges so that there's a set
number of "workers", so you don't have to write that tedious management code
yourself.
"""
from __future__ import with_statement

from eventlet.green import urllib2
import eventlet
import re

# http://daringfireball.net/2009/11/liberal_regex_for_matching_urls
url_regex = re.compile(r'\b(([\w-]+://?|www[.])[^\s()<>]+(?:\([\w\d]+\)|([^[:punct:]\s]|/)))')


def fetch(url, outq):
    """Fetch a url and push any urls found into a queue."""
    print("fetching", url)
    data = ''
    with eventlet.Timeout(5, False):
        data = urllib2.urlopen(url).read()
    for url_match in url_regex.finditer(data):
        new_url = url_match.group(0)
        outq.put(new_url)


def producer(start_url):
    """Recursively crawl starting from *start_url*.  Returns a set of
    urls that were found."""
    pool = eventlet.GreenPool()
    seen = set()
    q = eventlet.Queue()
    q.put(start_url)
    # keep looping if there are new urls, or workers that may produce more urls
    while True:
        while not q.empty():
            url = q.get()
            # limit requests to eventlet.net so we don't crash all over the internet
            if url not in seen and 'eventlet.net' in url:
                seen.add(url)
                pool.spawn_n(fetch, url, q)
        pool.waitall()
        if q.empty():
            break

    return seen


seen = producer("http://eventlet.net")
print("I saw these urls:")
print("\n".join(seen))

		
## Brief Introduction to Feedparser 
Universal Feed Parser is a Python module for downloading and parsing syndicated feeds. 
It can handle RSS 0.90, Netscape RSS 0.91, Userland RSS 0.91, RSS 0.92, RSS 0.93, RSS 0.94, RSS 1.0, RSS 2.0, Atom 0.3, Atom 1.0, and CDF feeds.
It also parses several popular extension modules, including Dublin Core and Apple’s iTunes extensions.


#Parsing a feed from a remote URL

>>> import feedparser
>>> d = feedparser.parse('http://feedparser.org/docs/examples/atom10.xml')
>>> d['feed']
u'Sample Feed'


#Parsing a feed from a local file
>>> import feedparser
>>> d = feedparser.parse(r'c:\incoming\atom10.xml')
>>> d['feed']['title']
u'Sample Feed'

#Parsing a feed from a string

>>> import feedparser
>>> rawdata = """<rss version="2.0">
<channel>
<title>Sample Feed</title>
</channel>
</rss>"""
>>> d = feedparser.parse(rawdata)
>>> d['feed']['title']
u'Sample Feed'


		
#Example of Feed Scraper - Dispatch Pattern 
#This is a server that is also a client of some other services
#It uses WSGI to handle POST requests with lot of URLS and uses feedparser module to get title of those 
#Use firefox httprequester to create POST


import eventlet
feedparser = eventlet.import_patched('feedparser')   #To Greenify this lib 

# the pool provides a safety limit on our concurrency
pool = eventlet.GreenPool()


def fetch_title(url):
    d = feedparser.parse(url)
    return d.feed.get('title', '')


def app(environ, start_response):
    if environ['REQUEST_METHOD'] != 'POST':
        start_response('403 Forbidden', [])
        return []

    # the pile collects the result of a concurrent operation -- in this case,
    # the collection of feed titles
    pile = eventlet.GreenPile(pool)
    for line in environ['wsgi.input'].readlines():
        url = line.strip()
        if url:
            pile.spawn(fetch_title, url)
    # since the pile is an iterator over the results,
    # you can use it in all sorts of great Pythonic ways
    titles = '\n'.join(pile)
    start_response('200 OK', [('Content-type', 'text/plain')])
    return [titles]


if __name__ == '__main__':
    from eventlet import wsgi
    wsgi.server(eventlet.listen(('localhost', 9010)), app)

	
	
#db_pool – DBAPI 2 database connection pooling (for MySQLdb and psycopg2)

import eventlet
import MySQLdb
cp = eventlet.ConnectionPool(MySQLdb, host='localhost', user='root', passwd='')
conn = cp.get()  #connect to the database
try:
	result = conn.cursor().execute('SELECT NOW()')
finally:
	cp.put(conn)  #return the connection to the pool
	
	
cp.clear() #at the end, close Connection Pool 

#to connect to multiple databases , a pool of pools, containing a ConnectionPool for every host you connect to.

dc = DatabaseConnector(MySQLdb,
	{'db.internal.example.com': {'user': 'internal', 'passwd': 's33kr1t'},
	'localhost': {'user': 'root', 'passwd': ''}})
pool = dc.get(host, dbname)  #Returns a ConnectionPool to the target host and schema


#WSGI handling - wsgi – WSGI server

from eventlet import wsgi
import eventlet

def hello_world(env, start_response):
    start_response('200 OK', [('Content-Type', 'text/plain')])
    return ['Hello, World!\r\n']

wsgi.server(eventlet.listen(('', 8090)), hello_world)

#SSL - Applications can detect whether they are inside a secure server by the value of the env['wsgi.url_scheme'] environment variable.


wsgi.server(eventlet.wrap_ssl(eventlet.listen(('', 8090)),
                              certfile='cert.crt',
                              keyfile='private.key',
                              server_side=True),    hello_world)



#websocket – Websocket Server


import eventlet
from eventlet import wsgi
from eventlet import websocket
from eventlet.support import six

# demo app
import os
import random


@websocket.WebSocketWSGI
def handle(ws):
    """  This is the websocket handler function.  Note that we
    can dispatch based on path in here, too."""
    if ws.path == '/echo':
        while True:
            m = ws.wait()
            if m is None:
                break
            ws.send(m)

    elif ws.path == '/data':
        for i in six.moves.range(10000):
            ws.send("0 %s %s\n" % (i, random.random()))
            eventlet.sleep(0.1)


def dispatch(environ, start_response):
    """ This resolves to the web page or the websocket depending on
    the path."""
    if environ['PATH_INFO'] == '/data':
        return handle(environ, start_response)
    else:
        start_response('200 OK', [('content-type', 'text/html')])
        return [open(os.path.join(
                     os.path.dirname(__file__),
                     'websocket.html')).read()]

if __name__ == "__main__":
    # run an example app from the command line
    listener = eventlet.listen(('127.0.0.1', 7000))
    print("\nVisit http://localhost:7000/ in your websocket-capable browser.\n")
    wsgi.server(listener, dispatch)
	
	
#timeout – Universal Timeouts - Raises exception in the current greenthread after timeout seconds:

#Example 
timeout = Timeout(seconds, exception)
try:
    ... # execution here is limited by timeout
finally:
    timeout.cancel()


#When exception is omitted or is None, the Timeout instance itself is raised:

>>> Timeout(0.1)
>>> eventlet.sleep(0.2)
Traceback (most recent call last):
 ...
Timeout: 0.1 seconds


#Or with 'with' 

with Timeout(seconds, exception) as timeout:
    pass # ... code block ...


#Example 
data = None
with Timeout(5, False):
    data = mysock.makefile().readline()
if data is None:
    ... # 5 seconds passed without reading a line
else:
    ... # a line was read within 5 seconds

#Methods are 
Timeout.cancel()  											If the timeout is pending, cancel it. 
Timeout.pending  											True if the timeout is scheduled to be raised.
eventlet.with_timeout(seconds, function, *args, **kwds)		Wrap a call to some (yielding) function with a timeout; 
															if the called function fails to return before the timeout, cancel it and return a flag value
data = with_timeout(30, urllib2.open, 'http://www.google.com/', timeout_value="")  #if timeout, returns ""

	

#queue – Queue class

"""This is a recursive web crawler.  
the fetch function produces jobs, and the GreenPool itself is the consumer, farming out work concurrently.
"""


from __future__ import with_statement

from eventlet.green import urllib2
import eventlet
import re

# http://daringfireball.net/2009/11/liberal_regex_for_matching_urls
url_regex = re.compile(r'\b(([\w-]+://?|www[.])[^\s()<>]+(?:\([\w\d]+\)|([^[:punct:]\s]|/)))')


def fetch(url, outq):
    """Fetch a url and push any urls found into a queue."""
    print("fetching", url)
    data = ''
    with eventlet.Timeout(5, False):
        data = urllib2.urlopen(url).read()
    for url_match in url_regex.finditer(data):     #use BeatifulSoup4 for html parsing 
        new_url = url_match.group(0)
        outq.put(new_url)


def producer(start_url, domain):
    """Recursively crawl starting from *start_url*.  Returns a set of
    urls that were found."""
    pool = eventlet.GreenPool()
    seen = set()
    q = eventlet.Queue()
    q.put(start_url)
    # keep looping if there are new urls, or workers that may produce more urls
    while True:
        while not q.empty():
            url = q.get()
            # limit requests to domain  so we don't crash all over the internet
            if url not in seen and domain in url:
                seen.add(url)
                pool.spawn_n(fetch, url, q)
        pool.waitall()
        if q.empty():
            break

    return seen


seen = producer("http://eventlet.net", "eventlet.net")
print("I saw these urls:")
print("\n".join(seen))

    
    
###virtualenv   
#virtualenv is a tool to create isolated Python environments
    
#Used for 
#if you can’t install packages into the global site-packages 
#if you need to have two versions of same library 

$ pip install virtualenv

##Command 
$ virtualenv ENV

#Where ENV is a directory to place the new virtual environment. 
#ENV is created in current directory 

##Few options 

-p PYTHON_EXE, --python=PYTHON_EXE
    The Python interpreter to use, e.g., –python=python2.5 will use the python2.5 interpreter to create the new environment. 
    The default is the interpreter that virtualenv was installed with (like /usr/bin/python)


--system-site-packages
    Give the virtual environment access to the global site-packages.

--always-copy
    Always copy files rather than symlinking.

--relocatable
    Make an EXISTING virtualenv environment relocatable. 
    This fixes up scripts and makes all .pth files relative.

--extra-search-dir=DIR
    Directory to look for setuptools/pip distributions in. This option can be specified multiple times.
   
    
##activate script
#This will change your $PATH so its first entry is the virtualenv’s bin/ directory  

#On Posix systems, this resides in /ENV/bin/
$ source bin/activate

#On Windows, activate script is in the Scripts folder:
> \path\to\env\Scripts\activate

#to undo the changes.
$ deactivate 




##With pycharms 
#https://www.jetbrains.com/help/pycharm/configuring-python-interpreter.html

#In the Settings/Preferences dialog, click Project Interpreter.
#In the drop-down list, choose Add local....
#The Add Local Python Interpreter dialog box opens:
# It has three options - virtualenv environment, Conda environment, System interpreter


#PyCharm provides a dedicated tool for installing, uninstalling, and upgrading Python package
#https://www.jetbrains.com/help/pycharm/installing-uninstalling-and-upgrading-packages.html 

#Open File | Settings | Project Interpreter for Windows and Linux
#In the Project Interpreter page of the project settings, select the desired Python interpreter or virtual environment. 
#Click +
#In the Available Packages dialog box that opens, 
#select the desired package from the list. 




















@@@
###Jinja2 - Modelled on Django templating system 
#http://jinja.pocoo.org/docs
 
$ pip install Jinja2
 
#Example 
from jinja2 import Template
template = Template('Hello {{ name }}!')  #a shared Environment is used 
template.render(name='John Doe') #Hello John Doe!'

##Jinja2 uses a central object called the template Environment
#check options from http://jinja.pocoo.org/docs/2.10/api/#jinja2.Environment

#create an Environment - load templatefile from pkg/templates 
#pkg is python package 
#and enabling autoescaping for HTML and XML files.
from jinja2 import Environment, PackageLoader, select_autoescape
env = Environment(
    loader=PackageLoader('pkg', 'templates'),
    autoescape=select_autoescape(['html', 'xml'])
)
#with filesystem loader 
>>> loader = FileSystemLoader('/path/to/templates')
>>> loader = FileSystemLoader(['/path/to/templates', '/other/path'])

#Usage 
#or use env.from_string(str) to get from string 
template = env.get_template('mytemplate.html') #pkg/templates/mytemplate.html
#To render it with some variables
print template.render(the='variables', go='here')

#Reference of render
render([context])
    This method accepts the same arguments as the dict constructor: 
    A dict, a dict subclass or keyword arguments. 
    If no arguments are given the context will be empty. 
    template.render(knights='that say nih')
    template.render({'knights': 'that say nih'})



#Example - enables autoescaping for templates ending in '.html', '.htm' and '.xml' 
#and disabling it by default for all other extensions. 
from jinja2 import Environment, select_autoescape
env = Environment(autoescape=select_autoescape(['html', 'htm', 'xml']),
                  loader=PackageLoader('mypackage'))


#Example - enable it for all templates created from strings 
#or for all templates with .html and .xml extensions:
from jinja2 import Environment, select_autoescape
env = Environment(autoescape=select_autoescape(
    enabled_extensions=('html', 'xml'),
    default_for_string=True,
))


#Example - to turn it on at all times except if the template ends with .txt:
from jinja2 import Environment, select_autoescape
env = Environment(autoescape=select_autoescape(
    disabled_extensions=('txt',),
    default_for_string=True,
    default=True,
))



##Templating 
•{% ... %} for Statements
•{{ ... }} for Expressions to print to the template output
•{# ... #} for Comments not included in the template output
•#  ... ## for Line Statements

 
#Example 
<!DOCTYPE html>
<html lang="en">
<head>
    <title>{% block title %}{% endblock %}</title>
</head>
<body>
    <ul id="navigation">
    {% for item in navigation %}
        <li><a href="{{ item.href }}">{{ item.caption }}</a></li>
    {% endfor %}
    </ul>

    <h1>My Webpage</h1>
    {{ a_variable }}

    {# a comment #}
</body>
</html>
 
 
##Variables - {{var}}
#Use . or [] to access attribute
#note 'bar' has to be literal string , not a variable 
#[index] can be used to access list, 
#for accessing dictionary, access dict['literalstring'] or iterate like for k,v in dict.items()
{{ foo.bar }}
{{ foo['bar'] }}

# foo.bar in Jinja2 does the following things 
•check for an attribute called bar on foo (getattr(foo, 'bar'))
•if there is not, check for an item 'bar' in foo (foo.__getitem__('bar'))
•if there is not, return an undefined object.

#foo['bar'] works mostly the same with a small difference in sequence:
•check for an item 'bar' in foo. (foo.__getitem__('bar'))
•if there is not, check for an attribute called bar on foo. (getattr(foo, 'bar'))
•if there is not, return an undefined object.

 
##Filters
{{ name|striptags|title }}  #remove all HTML Tags and then capitalize
{{ listx|join(', ') }} 
 
#List of builtin Filters 
#Note first arg is actually used as first_arg|filter_name(remaining_args)
abs(number)
    Return the absolute value of the argument.
attr(obj, name)
    Get an attribute of an object. 
    foo|attr("bar") works like foo.bar just that always an attribute is returned 
    and items are not looked up.
batch(value, linecount, fill_with=None)
    A filter that batches items. 
    It returns a list of lists with the given number of items. 
    If you provide a second parameter this is used to fill up missing items. 
    <table>
    {%- for row in items|batch(3, '&nbsp;') %}
      <tr>
      {%- for column in row %}
        <td>{{ column }}</td>
      {%- endfor %}
      </tr>
    {%- endfor %}
    </table>
capitalize(s)
    Capitalize a value. The first character will be uppercase, all others lowercase.
center(value, width=80)
    Centers the value in a field of a given width.
default(value, default_value=u'', boolean=False)
    Aliases:
        d 
    If the value is undefined it will return the passed default value, 
    otherwise the value of the variable:
    {{ my_variable|default('my_variable is not defined') }}
    This will output the value of my_variable if the variable was defined, 
    otherwise 'my_variable is not defined'. 
    If you want to use default with variables that evaluate to false 
    you have to set the second parameter to true:
    {{ ''|default('the string was empty', true) }}
dictsort(value, case_sensitive=False, by='key', reverse=False)
    Sort a dict and yield (key, value) pairs. 
    {% for item in mydict|dictsort %}
        sort the dict by key, case insensitive
    {% for item in mydict|dictsort(reverse=true) %}
        sort the dict by key, case insensitive, reverse order
    {% for item in mydict|dictsort(true) %}
        sort the dict by key, case sensitive
    {% for item in mydict|dictsort(false, 'value') %}
        sort the dict by value, case insensitive
escape(s)
    Aliases:
        e 
    Convert the characters &, <, >, ‘, and ” in string s to HTML-safe sequences. 
filesizeformat(value, binary=False)
    Format the value like a ‘human-readable’ file size (i.e. 13 kB, 4.1 MB, 102 Bytes, etc). Per default decimal prefixes are used (Mega, Giga, etc.), if the second parameter is set to True the binary prefixes are used (Mebi, Gibi).
first(seq)
    Return the first item of a sequence.
float(value, default=0.0)
    Convert the value into a floating point number. 
    If the conversion doesn’t work it will return 0.0. 
    You can override this default using the first parameter.
forceescape(value)
    Enforce HTML escaping. This will probably double escape variables.
format(value, *args, **kwargs)
    Apply python string formatting on an object:
    {{ "%s - %s"|format("Hello?", "Foo!") }}
        -> Hello? - Foo!
groupby(value, attribute)
    Group a sequence of objects by a common attribute.
    <ul>
    {% for group in persons|groupby('gender') %}
        <li>{{ group.grouper }}<ul>
        {% for person in group.list %}
            <li>{{ person.first_name }} {{ person.last_name }}</li>
        {% endfor %}</ul></li>
    {% endfor %}
    </ul>
    #OR
    <ul>
    {% for grouper, list in persons|groupby('gender') %}
        ...
    {% endfor %}
    </ul>
indent(s, width=4, first=False, blank=False, indentfirst=None)
    Return a copy of the string with each line indented by 4 spaces. 
    The first line and blank lines are not indented by default.
    Parameters:
    •width – Number of spaces to indent by.
    •first – Don’t skip indenting the first line.
    •blank – Don’t skip indenting empty lines. 
int(value, default=0, base=10)
    Convert the value into an integer. 
    If the conversion doesn’t work it will return 0. 
    You can override this default using the first parameter. 
    You can also override the default base (10) in the second parameter, which handles input with prefixes such as 0b, 0o and 0x for bases 2, 8 and 16 respectively. The base is ignored for decimal numbers and non-string values.
join(value, d=u'', attribute=None)
    Return a string which is the concatenation of the strings in the sequence. 
    The separator between elements is an empty string per default, 
    you can define it with the optional parameter:
    {{ [1, 2, 3]|join('|') }}
        -> 1|2|3

    {{ [1, 2, 3]|join }}
        -> 123
    It is also possible to join certain attributes of an object:
    {{ users|join(', ', attribute='username') }}


last(seq)
    Return the last item of a sequence.
length(object)
    Aliases:
        count 
    Return the number of items of a sequence or mapping.
list(value)
    Convert the value into a list. 
    If it was a string the returned list will be a list of characters.
lower(s)
    Convert a value to lowercase.
map()
    Applies a filter on a sequence of objects or looks up an attribute. 
    This is useful when dealing with lists of objects 
    but you are really only interested in a certain value of it.
    The basic usage is mapping on an attribute. 
    Imagine you have a list of users but you are only interested in a list of usernames:
        Users on this page: {{ users|map(attribute='username')|join(', ') }}
    Alternatively you can let it invoke a filter by passing the name of the filter 
    and the arguments afterwards. 
    A good example would be applying a text conversion filter on a sequence:
        Users on this page: {{ titles|map('lower')|join(', ') }}

max(value, case_sensitive=False, attribute=None)
    Return the largest item from the sequence.
    {{ [1, 2, 3]|max }}
        -> 3

min(value, case_sensitive=False, attribute=None)
    Return the smallest item from the sequence.
    {{ [1, 2, 3]|min }}
        -> 1
    Parameters:
    •case_sensitive – Treat upper and lower case strings as distinct.
    •attribute – Get the object with the max value of this attribute.
 
pprint(value, verbose=False)
    Pretty print a variable. Useful for debugging.
random(seq)
    Return a random item from the sequence.
reject()
    Filters a sequence of objects by applying a test to each object, 
    and rejecting the objects with the test succeeding.
    If no test is specified, each object will be evaluated as a boolean.
    Example usage:
    {{ numbers|reject("odd") }}

rejectattr()
    Filters a sequence of objects by applying a test to the specified attribute of each object, 
    and rejecting the objects with the test succeeding.
    If no test is specified, the attribute’s value will be evaluated as a boolean.
    {{ users|rejectattr("is_active") }}
    {{ users|rejectattr("email", "none") }}
replace(s, old, new, count=None)
    Return a copy of the value with all occurrences of a substring replaced with a new one. The first argument is the substring that should be replaced, the second is the replacement string. If the optional third argument count is given, only the first count occurrences are replaced:
    {{ "Hello World"|replace("Hello", "Goodbye") }}
        -> Goodbye World
    {{ "aaaaargh"|replace("a", "d'oh, ", 2) }}
        -> d'oh, d'oh, aaargh
reverse(value)
    Reverse the object or return an iterator that iterates over it the other way round.
round(value, precision=0, method='common')
    Round the number to a given precision. The first parameter specifies the precision (default is 0), the second the rounding method:
    •'common' rounds either up or down
    •'ceil' always rounds up
    •'floor' always rounds down
    If you don’t specify a method 'common' is used.
    {{ 42.55|round }}
        -> 43.0
    {{ 42.55|round(1, 'floor') }}
        -> 42.5
    Note that even if rounded to 0 precision, a float is returned. 
    If you need a real integer, pipe it through int:
    {{ 42.55|round|int }}
        -> 43

safe(value)
    Mark the value as safe which means that in an environment 
    with automatic escaping enabled this variable will not be escaped.
select()
    Filters a sequence of objects by applying a test to each object, 
    and only selecting the objects with the test succeeding.
    If no test is specified, each object will be evaluated as a boolean.
    Example usage:
    {{ numbers|select("odd") }}
    {{ numbers|select("odd") }}
    {{ numbers|select("divisibleby", 3) }}
    {{ numbers|select("lessthan", 42) }}
    {{ strings|select("equalto", "mystring") }}


selectattr()
    Filters a sequence of objects by applying a test to the specified attribute of each object, and only selecting the objects with the test succeeding.
    If no test is specified, the attribute’s value will be evaluated as a boolean.
    Example usage:
    {{ users|selectattr("is_active") }}
    {{ users|selectattr("email", "none") }}


slice(value, slices, fill_with=None)
    Slice an iterator and return a list of lists containing those items. 
    Useful if you want to create a div containing three ul tags that represent columns:
    <div class="columwrapper">
      {%- for column in items|slice(3) %}
        <ul class="column-{{ loop.index }}">
        {%- for item in column %}
          <li>{{ item }}</li>
        {%- endfor %}
        </ul>
      {%- endfor %}
    </div>


sort(value, reverse=False, case_sensitive=False, attribute=None)
    Sort an iterable. Per default it sorts ascending, 
    if you pass it true as first argument it will reverse the sorting.

    If the iterable is made of strings the third parameter can be used to control the case sensitiveness of the comparison which is disabled by default.
    {% for item in iterable|sort %}
        ...
    {% endfor %}
    It is also possible to sort by an attribute 
    (for example to sort by the date of an object) by specifying the attribute parameter:
    {% for item in iterable|sort(attribute='date') %}
        ...
    {% endfor %}

string(object)
    Make a string unicode if it isn’t already. 
    That way a markup string is not converted back to unicode.
striptags(value)
    Strip SGML/XML tags and replace adjacent whitespace by one space.
sum(iterable, attribute=None, start=0)
    Returns the sum of a sequence of numbers plus the value of parameter ‘start’ 
    (which defaults to 0). When the sequence is empty it returns start.
    It is also possible to sum up only certain attributes:
    Total: {{ items|sum(attribute='price') }}

title(s)
    Return a titlecased version of the value. I.e. 
    words will start with uppercase letters, all remaining characters are lowercase.
tojson(value, indent=None)
    Dumps a structure to JSON so that it’s safe to use in <script> tags. 
    It accepts the same arguments and returns a JSON string. 
    Note that this is available in templates through the |tojson filter 
    which will also mark the result as safe. 
    Due to how this function escapes certain characters this is safe even 
    if used outside of <script> tags.
    The following characters are escaped in strings:
    •<
    •>
    •&
    •'

trim(value)
    Strip leading and trailing whitespace.
truncate(s, length=255, killwords=False, end='...', leeway=None)
    Return a truncated copy of the string. 
    The length is specified with the first parameter which defaults to 255. 
    If the second parameter is true the filter will cut the text at length. Otherwise it will discard the last word. If the text was in fact truncated it will append an ellipsis sign ("..."). If you want a different ellipsis sign than "..." you can specify it using the third parameter. Strings that only exceed the length by the tolerance margin given in the fourth parameter will not be truncated.
    {{ "foo bar baz qux"|truncate(9) }}
        -> "foo..."
    {{ "foo bar baz qux"|truncate(9, True) }}
        -> "foo ba..."
    {{ "foo bar baz qux"|truncate(11) }}
        -> "foo bar baz qux"
    {{ "foo bar baz qux"|truncate(11, False, '...', 0) }}
        -> "foo bar..."


unique(value, case_sensitive=False, attribute=None)
    Returns a list of unique items from the the given iterable.
    {{ ['foo', 'bar', 'foobar', 'FooBar']|unique }}
        -> ['foo', 'bar', 'foobar']


upper(s)
    Convert a value to uppercase.
urlencode(value)
    Escape strings for use in URLs (uses UTF-8 encoding). 
    It accepts both dictionaries and regular strings as well as pairwise iterables.

urlize(value, trim_url_limit=None, nofollow=False, target=None, rel=None)
    Converts URLs in plain text into clickable links.
    If you pass the filter an additional integer it will shorten the urls to that number. 
    Also a third argument exists that makes the urls “nofollow”:
    {{ mytext|urlize(40, true) }}
        links are shortened to 40 chars and defined with rel="nofollow"
    If target is specified, the target attribute will be added to the <a> tag:
    {{ mytext|urlize(40, target='_blank') }}



wordcount(s)
    Count the words in that string.
wordwrap(s, width=79, break_long_words=True, wrapstring=None)
    Return a copy of the string passed to the filter wrapped after 79 characters. 
    You can override this default using the first parameter. 
    If you set the second parameter to false Jinja will not split words apart 
    if they are longer than width. 
    By default, the newlines will be the default newlines for the environment, 
    but this can be changed using the wrapstring keyword argument.

xmlattr(d, autospace=True)
    Create an SGML/XML attribute string based on the items in a dict. 
    All values that are neither none nor undefined are automatically escaped:
    <ul{{ {'class': 'my_list', 'missing': none,
            'id': 'list-%d'|format(variable)}|xmlattr }}>
    ...
    </ul>
    Results in something like this:
    <ul class="my_list" id="list-42">
    ...
    </ul>


 
 
##Tests - use with 'is' - 1st arg is used as 1st_arg|test_name(remaining_args)
#Tests can be used to test a variable against a common expression
{% if loop.index is divisibleby 3 %}
{% if loop.index is divisibleby(3) %}

#List of Builtin Tests
callable(object)
    Return whether the object is callable (i.e., some kind of function). 
    Note that classes are callable, as are instances with a __call__() method.
defined(value)
    Return true if the variable is defined:
    {% if variable is defined %}
        value of variable: {{ variable }}
    {% else %}
        variable is not defined
    {% endif %}
divisibleby(value, num)
    Check if a variable is divisible by a number.
eq(a, b)
    Aliases:
        ==, equalto 
escaped(value)
    Check if the value is escaped.
even(value)
    Return true if the variable is even.
ge(a, b)
    Aliases:
        >= 
gt(a, b)
    Aliases:
        >, greaterthan 
in(value, seq)
    Check if value is in seq.
iterable(value)
    Check if it’s possible to iterate over an object.
le(a, b)
    Aliases:
        <= 
lower(value)
    Return true if the variable is lowercased.
lt(a, b)
    Aliases:
        <, lessthan 
mapping(value)
    Return true if the object is a mapping (dict etc.).
ne(a, b)
    Aliases:
        != 
none(value)
    Return true if the variable is none.
number(value)
    Return true if the variable is a number.
odd(value)
    Return true if the variable is odd.
sameas(value, other)
    Check if an object points to the same memory address than another object:
    {% if foo.attribute is sameas false %}
        the foo attribute really is the `False` singleton
    {% endif %}
sequence(value)
    Return true if the variable is a sequence. Sequences are variables that are iterable.
string(value)
    Return true if the object is a string.
undefined(value)
    Like defined() but the other way round.
upper(value)
    Return true if the variable is uppercased.

 
##Whitespace Control
#In the default configuration:
    •a single trailing newline is stripped if present
    •other whitespace (spaces, tabs, newlines etc.) is returned unchanged

#If Line Statements are enabled, 
#they strip leading whitespace automatically up to the beginning of the line.

#By default, Jinja2 also removes trailing newlines. 
#To keep single trailing newlines, configure Jinja to keep_trailing_newline.

#trim_blocks : the first newline after a template tag is removed automatically 
#lstrip_blocks : to strip tabs and spaces from the beginning of a line to the start of a block. 
#(Nothing will be stripped if there are other characters before the start of the block.)

#Example - without the trim_blocks and lstrip_blocks(default) options in Environment 
<div>
    {% if True %}
        yay
    {% endif %}
</div>

#gets rendered with blank lines inside the div:
<div>

        yay

</div>

#with both trim_blocks and lstrip_blocks enabled
env = Environment(
    ...,
    trim_blocks = True, lstrip_blocks = True, ...
)

#the template block lines are removed and other whitespace is preserved:
<div>
        yay
</div>


#To manually disable the lstrip_blocks behavior
#by putting a plus sign (+) at the start of a block:
<div>
        {%+ if something %}yay{% endif %}
</div>


#To  strip whitespace in templates
#add a minus sign (-) to the start or end of a block (e.g. a For tag), 
#a comment, or a variable expression, 
#the whitespaces before or after that block will be removed:
{% for item in seq -%}
    {{ item }}
{%- endfor %}
#If seq was a list of numbers from 1 to 9, the output would be 123456789.


 
 
##Escaping
#to output a literal variable delimiter ({{) 
{{ '{{' }}
#For bigger sections, use 'raw' (raw block is never evaluated)
{% raw %}
    <ul>
    {% for item in seq %}
        <li>{{ item }}</li>
    {% endfor %}
    </ul>
{% endraw %}



##Line Statements
#If line statements are enabled by the application(Environment.line_statement_prefix) 
#it’s possible to mark a line as a statement. 

#For example, if the line statement prefix is configured to # 
#(in Environment.line_statement_prefix='#') 
#Note there is another option line_comment_prefix in Environment
#the following two examples are equivalent:
<ul>
# for item in seq
    <li>{{ item }}</li>
# endfor
</ul>
#OR 
<ul>
# for item in seq:
    <li>{{ item }}</li>
# endfor
</ul>
#OR 
<ul>
{% for item in seq %}
    <li>{{ item }}</li>
{% endfor %}
</ul>


#Line statements can span multiple lines if there are open parentheses, braces or brackets:
<ul>
# for href, caption in [('index.html', 'Index'),
                        ('about.html', 'About')]:
    <li><a href="{{ href }}">{{ caption }}</a></li>
# endfor
</ul>


##Template inheritance 

#base.html- Note endblock tag can mention name for clarity 
<!DOCTYPE html>
<html lang="en">
<head>
    {% block head %}
    <link rel="stylesheet" href="style.css" />
    <title>{% block title %}{% endblock %} - My Webpage</title>
    {% endblock head %}
</head>
<body>
    <div id="content">{% block content %}{% endblock %}</div>
    <div id="footer">
        {% block footer %}
        &copy; Copyright 2008 by <a href="http://domain.invalid/">you</a>.
        {% endblock %}
    </div>
</body>
</html>

#child.html - Note super() call to get parent's value 
{% extends "base.html" %}
{% block title %}Index{% endblock %}
{% block head %}
    {{ super() }}
    <style type="text/css">
        .important { color: #336699; }
    </style>
{% endblock %}
{% block content %}
    <h1>Index</h1>
    <p class="important">
      Welcome to my awesome homepage.
    </p>
{% endblock %}

#the FileSystemLoader allows you to access other templates by giving the filename. 
#You can access templates in subdirectories with a slash:
{% extends "layout/default.html" %}

#to print a block multiple times, 
#use the special self variable and call the block with that name:
<title>{% block title %}{% endblock %}</title>
<h1>{{ self.title() }}</h1>
{% block body %}{% endblock %}

#Block Nesting and Scope
#Blocks can be nested for more complex layouts. 
#However, per default blocks may not access variables from outer scopes:

#This example would output empty <li> items 
#because item is unavailable inside the block.
{% for item in seq %}
    <li>{% block loop_item %}{{ item }}{% endblock %}</li>
{% endfor %}


#specify variables that are available in a block by setting the block to “scoped” 
#When overriding a block, the scoped modifier does not have to be provided.
{% for item in seq %}
    <li>{% block loop_item scoped %}{{ item }}{% endblock %}</li>
{% endfor %}




##Working with Manual Escaping
#ie Environment.autoescape=False(default)
#the escape manually eg 

{{ user.username|e }}


##Working with Automatic Escaping
#ie Environment.autoescape=select_autoescape(...)
#everything is escaped by default except for values explicitly filtered with 'safe'. 

#Jinja2 functions (macros, super, self.BLOCKNAME) always return template data 
#that is marked as safe.


##Control Structures - For 

#for list of object 
<h1>Members</h1>
<ul>
{% for user in users %}
  <li>{{ user.username|e }}</li>
{% endfor %}
</ul>

#for dict 
<dl>
{% for key, value in my_dict.iteritems() %}
    <dt>{{ key|e }}</dt>
    <dd>{{ value|e }}</dd>
{% endfor %}
</dl>


#Inside of a for-loop block, access some special variables:
loop.index          The current iteration of the loop. (1 indexed) 
loop.index0         The current iteration of the loop. (0 indexed) 
loop.revindex       The number of iterations from the end of the loop (1 indexed) 
loop.revindex0      The number of iterations from the end of the loop (0 indexed) 
loop.first          True if first iteration. 
loop.last           True if last iteration. 
loop.length         The number of items in the sequence. 
loop.cycle          A helper function to cycle between a list of sequences. 
loop.depth          Indicates how deep in a recursive loop the rendering currently is. Starts at level 1 
loop.depth0         Indicates how deep in a recursive loop the rendering currently is. Starts at level 0 
loop.previtem       The item from the previous iteration of the loop. Undefined during the first iteration. 
loop.nextitem       The item from the following iteration of the loop. Undefined during the last iteration. 
loop.changed(*val)  True if previously called with a different value (or not called at all). 

#Within a for-loop, to cycle among a list of strings/variables each time through the loop
{% for row in rows %}
    <li class="{{ loop.cycle('odd', 'even') }}">{{ row }}</li>
{% endfor %}

#Unlike in Python, it’s not possible to break or continue in a loop.

#To filter the sequence during iteration, which allows to skip items. 
#the special loop variable will count correctly; thus not counting the users not iterated over.
{% for user in users if not user.hidden %}
    <li>{{ user.username|e }}</li>
{% endfor %}


#If no iteration took place because the sequence was empty 
#or the filtering removed all the items from the sequence, 
#render a default block by using else:

#Note that, in Python, else blocks are executed 
#whenever the corresponding loop did not break. 
#Since Jinja loops cannot break anyway, a slightly different behavior of the else keyword was chosen
<ul>
{% for user in users %}
    <li>{{ user.username|e }}</li>
{% else %}
    <li><em>no users found</em></li>
{% endfor %}
</ul>

#to use loops recursively. 
#This is useful for sitemaps or RDF. 
#add the recursive modifier to the loop definition 
#and call the 'loop' variable with the new iterable where you want to recurse.
#Note {%- means strip whitespace in templates
<ul class="sitemap">
{%- for item in sitemap recursive %}
    <li><a href="{{ item.href|e }}">{{ item.title }}</a>
    {%- if item.children -%}
        <ul class="submenu">{{ loop(item.children) }}</ul>
    {%- endif %}</li>
{%- endfor %}
</ul>

#The 'loop' variable always refers to the closest (innermost) loop. 
#If we have more than one level of loops, 
#rebind the variable loop by writing 
{% set outer_loop = loop %} 
#after the loop that we want to use recursively. 
#Then, we can call it using 
{{ outer_loop(…) }}

#note that assignments in loops will be cleared at the end of the iteration 
#and cannot outlive the loop scope

#Example of useing previtem and nextitem:
{% for value in values %}
    {% if loop.previtem is defined and value > loop.previtem %}
        The value just increased!
    {% endif %}
    {{ value }}
    {% if loop.nextitem is defined and loop.nextitem > value %}
        The value will increase even more!
    {% endif %}
{% endfor %}

#OR If you only care whether the value changed at all
#Use 'changed'
{% for entry in entries %}
    {% if loop.changed(entry.category) %}
        <h2>{{ entry.category }}</h2>
    {% endif %}
    <p>{{ entry.message }}</p>
{% endfor %}



##Control Structures - If
#comparable with the Python if statement. 
#to test if a variable is defined, not empty and not false:
{% if users %}
<ul>
{% for user in users %}
    <li>{{ user.username|e }}</li>
{% endfor %}
</ul>
{% endif %}


#For multiple branches, 
{% if kenny.sick %}
    Kenny is sick.
{% elif kenny.dead %}
    You killed Kenny!  You bastard!!!
{% else %}
    Kenny looks okay --- so far
{% endif %}


##Control Structures - Include
#to include a template 
#and return the rendered contents of that file into the current namespace

#Included templates have access to the variables of the active context by default
#use 'without context' for not to include context 
{% include 'header.html' %}
    Body
{% include 'footer.html' %}


#mark an include with ignore missing; 
#in which case Jinja will ignore the statement 
#if the template to be included does not exist. 
#And can be combined with with or without context
{% include "sidebar.html" ignore missing %}
{% include "sidebar.html" ignore missing with context %}
{% include "sidebar.html" ignore missing without context %}
#with list of templates 
{% include ['page_detailed.html', 'page.html'] %}
{% include ['special_sidebar.html', 'sidebar.html'] ignore missing %}




##Control Structures - Macros
#Macros are comparable with functions in regular programming languages. 

#example of a macro that renders a form element:
{% macro input(name, value='', type='text', size=20) -%}
    <input type="{{ type }}" name="{{ name }}" value="{{
        value|e }}" size="{{ size }}">
{%- endmacro %}

#usage 
<p>{{ input('username') }}</p>
<p>{{ input('password', type='password') }}</p>


#If the macro was defined in a different template, 
#import it first.

#Macros and variables starting with one or more underscores are private 
#and cannot be imported.
#forms.html
{% macro input(name, value='', type='text') -%}
    <input type="{{ type }}" value="{{ value|e }}" name="{{ name }}">
{%- endmacro %}

{%- macro textarea(name, value='', rows=10, cols=40) -%}
    <textarea name="{{ name }}" rows="{{ rows }}" cols="{{ cols
        }}">{{ value|e }}</textarea>
{%- endmacro %}
#Syntax-1 
{% import 'forms.html' as forms %}
<dl>
    <dt>Username</dt>
    <dd>{{ forms.input('username') }}</dd>
    <dt>Password</dt>
    <dd>{{ forms.input('password', type='password') }}</dd>
</dl>
<p>{{ forms.textarea('comment') }}</p>

#Syntax-2
{% from 'forms.html' import input as input_field, textarea %}
<dl>
    <dt>Username</dt>
    <dd>{{ input_field('username') }}</dd>
    <dt>Password</dt>
    <dd>{{ input_field('password', type='password') }}</dd>
</dl>
<p>{{ textarea('comment') }}</p>

#Import Context Behavior
#By default, included templates are passed the current context 
#and imported templates are not
#for example, render_box.html can access box variable 
{% for box in boxes %}
    {% include "render_box.html" %}
{% endfor %}

#to change the default
{% from 'forms.html' import input with context %}
{% include 'header.html' without context %}



#Inside macros, access below special variables:
varargs  - list for more positional args that given in defination 
kwargs   - dict for more keyword args that given in definition 
caller   - If the macro was called from a call tag, 
           the caller is stored in this variable as a callable macro.

#Below attributes are available on Macro object 
name            The name of the macro. {{ input.name }} will print input.
arguments       A tuple of the names of arguments the macro accepts.
defaults        A tuple of default values.
catch_kwargs    This is true if the macro accepts extra keyword arguments (i.e.: accesses the special kwargs variable).
catch_varargs   This is true if the macro accepts extra positional arguments (i.e.: accesses the special varargs variable).
caller          This is true if the macro accesses the special caller variable and may be called from a call tag.


##Control Structures - Call
#a call block works exactly like a macro without a name.
#Note 
caller   - If the macro was called from a call tag, 
           the caller is stored in this variable as a callable macro.
#Example 
{% macro render_dialog(title, class='dialog') -%}
    <div class="{{ class }}">
        <h2>{{ title }}</h2>
        <div class="contents">
            {{ caller() }}    #<-- here comes body of call tag 
        </div>
    </div>
{%- endmacro %}

{% call render_dialog('Hello World') %}
    This is a simple dialog rendered by using a macro and
    a call block.
{% endcall %}


#example of how a call block can be used with arguments:
{% macro dump_users(users) -%}
    <ul>
    {%- for user in users %}
        <li><p>{{ user.username|e }}</p>{{ caller(user) }}</li>
    {%- endfor %}
    </ul>
{%- endmacro %}

{% call(user) dump_users(list_of_user) %}
    <dl>
        <dl>Realname</dl>
        <dd>{{ user.realname|e }}</dd>
        <dl>Description</dl>
        <dd>{{ user.description }}</dd>
    </dl>
{% endcall %}



##Control Structures - filter
#to apply regular Jinja2 filters on a block of template data. 
{% filter upper %}
    This text becomes uppercase
{% endfilter %}



##Control Structures - Assignments
#Inside code blocks, you can also assign values to variables. 

#Assignments at top level (outside of blocks, macros or loops) are exported 
#from the template like top level macros 
#and can be imported by other templates.

#Assignments use the set tag and can have multiple targets:
{% set navigation = [('index.html', 'Index'), ('about.html', 'About')] %}
{% set key, value = call_something() %}


#it is not possible to set variables inside a block 
#and have them show up outside of it. 
#This also applies to loops. 
#The only exception to that rule are 'if' statements which do not introduce a scope. 

#Note inner set 'iterated' is local scope var, does not modify outer 'iterated'
{% set iterated = false %}
{% for item in seq %}
    {{ item }}
    {% set iterated = true %}
{% endfor %}
{% if not iterated %} did not iterate {% endif %}
#above can be written correctly 
{% for item in seq %}
    {{ item }}
{% else %}
    did not iterate
{% endfor %}


#Or use  namespace objects which allow propagating of changes across scopes:
{% set ns = namespace(found=false) %}
{% for item in items %}
    {% if item.check_something() %}
        {% set ns.found = true %}
    {% endif %}
    * {{ item.title }}
{% endfor %}
Found item having something: {{ ns.found }}



##Control Structures - Block Assignments
{% set navigation %}
    <li><a href="/">Index</a>
    <li><a href="/downloads">Downloads</a>
{% endset %}

#with filters.

{% set reply | wordwrap %}
    You wrote:
    {{ message }}
{% endset %}




##Jinja2 - Literals
"Hello World"
42 / 42.23
['list', 'of', 'objects']
('tuple', 'of', 'values')
{'dict': 'of', 'key': 'and', 'value': 'pairs'}
true / false :The special constants true, false, and none are lowercase


##Jinja2 - Operations 
{{ 1 + 1 }} is 2. , can be used with list or tuple concatenation, for string, use ~
{{ 3 - 2 }} is 1.
{{ 1 / 2 }} is {{ 0.5 }}
{{ 20 // 7 }} is 2.
{{ 11 % 7 }} is 4.
{{ 2 * 2 }} would return 4. 
{{ '=' * 80 }} would print a bar of 80 equal signs.
{{ 2**3 }} would return 8.
All comparison operations, ==, !=, ... 
{{ 1 in [1, 2, 3] }} would, return true.
and, or, not, () used for grouping 
~   
    Converts all operands into strings and concatenates them.
    {{ "Hello " ~ name ~ "!" }} would return (assuming name is set to 'John') Hello John!.
inline if  
    The general syntax is <do something> if <something is true> else <do something else>.
    {% extends layout_template if layout_template is defined else 'master.html' %}
    {{ '[%s]' % page.title if page.title }}
    
    
    
##Jinja2 - List of Global Functions
range([start, ]stop[, step])
    Return a list containing an arithmetic progression of integers
    <ul>
    {% for user in users %}
        <li>{{ user.username }}</li>
    {% endfor %}
    {% for number in range(10 - users|count) %}
        <li class="empty"><span>...</span></li>
    {% endfor %}
    </ul>

lipsum(n=5, html=True, min=20, max=100)
    Generates some lorem ipsum for the template. 
    By default, five paragraphs of HTML are generated 
    with each paragraph between 20 and 100 words. 
    If html is False, regular text is returned. 
    This is useful to generate simple contents for layout testing.

dict(**items)
    A convenient alternative to dict literals. 
    {'foo': 'bar'} is the same as dict(foo='bar')
    
class cycler(*items)
    The cycler allows you to cycle among values similar to how loop.cycle works. 
    Unlike loop.cycle, you can use this cycler outside of loops or over multiple loops.
    A cycler has the following attributes and methods:
    reset()
        Resets the cycle to the first item.
    next()
        Goes one item ahead and returns the then-current item.
    current
        Returns the current item.
    {% set row_class = cycler('odd', 'even') %}
    <ul class="browser">
    {% for folder in folders %}
      <li class="folder {{ row_class.next() }}">{{ folder|e }}</li>
    {% endfor %}
    {% for filename in files %}
      <li class="file {{ row_class.next() }}">{{ filename|e }}</li>
    {% endfor %}
    </ul>


class joiner(sep=', ')
    used to "join" multiple sections. 
    A joiner is passed a string and will return that string every time it's called, 
    except the first time (in which case it returns an empty string). 
    {% set pipe = joiner("|") %}
    {% if categories %} {{ pipe() }}
        Categories: {{ categories|join(", ") }}
    {% endif %}
    {% if author %} {{ pipe() }}
        Author: {{ author() }}
    {% endif %}
    {% if can_edit %} {{ pipe() }}
        <a href="?action=edit">Edit</a>
    {% endif %}



class namespace(...)
    Creates a new container that allows attribute assignment using the {% set %} tag:
    {% set ns = namespace() %}
    {% set ns.foo = 'bar' %}
    The main purpose of this is to allow carrying a value 
    from within a loop body to an outer scope. 
    Initial values can be provided as a dict, as keyword arguments, or both 
    {% set ns = namespace(found=false) %}
    {% for item in items %}
        {% if item.check_something() %}
            {% set ns.found = true %}
        {% endif %}
        * {{ item.title }}
    {% endfor %}
    Found item having something: {{ ns.found }}




##jinja2 - Adding Extensions
jinja_env = Environment(extensions=['jinja2.ext.i18n'])

#List of extensions 
i18n Extension          jinja2.ext.i18n
Expression Statement    jinja2.ext.do
Loop Controls           jinja2.ext.loopcontrols
With Statement          jinja2.ext.with_
Autoescape Extension    jinja2.ext.autoescape


##Extension - Expression Statement
#do tag that works exactly like the regular variable expression {{ ... }} 
#except it doesn't print anything. 
#Example - This can be used to modify lists:
{% do navigation.append('a string') %}



##Extension - Loop Controls
#one can use  break and continue in loops. 
#When break is reached, the loop is terminated; 
#if continue is reached, the processing is stopped and continues with the next iteration.
{% for user in users %}
    {%- if loop.index is even %}{% continue %}{% endif %}
    ...
{% endfor %}
#Note that loop.index starts with 1, and loop.index0 starts with 0
{% for user in users %}
    {%- if loop.index >= 10 %}{% break %}{% endif %}
{%- endfor %}




##Extension - With Statement
#The with statement makes it possible to create a new inner scope.
# Variables set within this scope are not visible outside of the scope.


{% with %}
    {% set foo = 42 %}
    {{ foo }}           foo is 42 here
{% endwith %}
#foo is not visible here any longer


#Because it is common to set variables at the beginning of the scope, 
#you can do that within the with statement
{% with foo = 42 %}
    {{ foo }}
{% endwith %}
#equivalent to 
{% with %}
    {% set foo = 42 %}
    {{ foo }}
{% endwith %}


#Accessing earlier variable 
{% with a={}, b=a.attribute %}...{% endwith %}
#equivalent to 
{% with a={} %}
    {% set b = a.attribute %}
{% endwith %}



##jinja2 - Autoescape Overrides
#To activate and deactivate the autoescaping from within the templates.
{% autoescape true %}
    Autoescaping is active within this block
{% endautoescape %}

{% autoescape false %}
    Autoescaping is inactive within this block
{% endautoescape %}



##jinja2 - I18n 
#To mark a section as translatable
#Inside trans tags no statements are allowed, only variable tags and text are allowed
<p>{% trans %}Hello {{ user }}!{% endtrans %}</p>


#To translate a template expression — say, using template filters, 
#or by just accessing an attribute of an object — 
#bind the expression to a name for use within the translation block:
<p>{% trans user=user.username %}Hello {{ user }}!{% endtrans %}</p>


#to bind more than one expression inside a trans tag
{% trans book_title=book.title, author=author.name %}
This is {{ book_title }} by {{ author }}
{% endtrans %}



#To pluralize, specify both the singular and plural forms with the pluralize tag, 
#which appears between trans and endtrans:
{% trans count=list|length %}
There is {{ count }} {{ name }} object.
{% pluralize %}
There are {{ count }} {{ name }} objects.
{% endtrans %}

#By default, the first variable in a block is used to determine the correct singular or plural form. 
#or modify 
{% trans ..., user_count=users|length %}...
{% pluralize user_count %}...{% endtrans %}


#When translating longer blocks of text, 
#whitespace and linebreaks result in rather ugly and error-prone translation strings. 
#To avoid 
{% trans trimmed book_title=book.title %}
    This is {{ book_title }}.
    You should read it!
{% endtrans %}


#If trimming is enabled globally, 
#the 'notrimmed' modifier can be used to disable it for a trans block.



#It's also possible to translate strings in expressions. 
#Use below 
•gettext        translate a single string
•ngettext       translate a pluralizable string
•_              alias for gettext

#To print a translated string like this:
{{ _('Hello World!') }}


#To use placeholders, use the format filter
#For multiple placeholders, always use keyword arguments to format, 
#as other languages may not use the words in the same order.
{{ _('Hello %(user)s!')|format(user=user.username) }}





##Jinja - Custom Filters
#are regular Python functions that take the left side of the filter as first argument 
#and the arguments passed to the filter as extra arguments or keyword arguments.

def datetimeformat(value, format='%H:%M / %d-%m-%Y'):
    return value.strftime(format)

#register it on the template environment 
environment.filters['datetimeformat'] = datetimeformat
#usage 
written on: {{ article.pub_date|datetimeformat }}
publication date: {{ article.pub_date|datetimeformat('%d-%m-%Y') }}


#Filters can also be passed the current template context or environment as first arg
#Use decorators : environmentfilter(), contextfilter() and evalcontextfilter().

#Example 
import re
from jinja2 import evalcontextfilter, Markup, escape

_paragraph_re = re.compile(r'(?:\r\n|\r|\n){2,}')

@evalcontextfilter
def nl2br(eval_ctx, value):
    result = u'\n\n'.join(u'<p>%s</p>' % p.replace('\n', Markup('<br>\n'))
                          for p in _paragraph_re.split(escape(value)))
    if eval_ctx.autoescape:
        result = Markup(result)
    return result

##jinja2 - Evaluation Context
#Currently it is only used to enable and disable the automatic escaping 

#Previous versions:
@environmentfilter
def filter(env, value):
    result = do_something(value)
    if env.autoescape:
        result = Markup(result)
    return result


#In new versions 
@contextfilter
def filter(context, value):
    result = do_something(value)
    if context.eval_ctx.autoescape:
        result = Markup(result)
    return result
#OR 
@evalcontextfilter
def filter(eval_ctx, value):
    result = do_something(value)
    if eval_ctx.autoescape:
        result = Markup(result)
    return result



##jinja2 - Custom Tests
#Tests work like filters just that there is no way for a test 
#to get access to the environment or context and that they can't be chained. The return value of a test should be True or False. The purpose of a test is to give the template designers the possibility to perform type and conformability checks.

import math

def is_prime(n):
    if n == 2:
        return True
    for i in xrange(2, int(math.ceil(math.sqrt(n))) + 1):
        if n % i == 0:
            return False
    return True


#register it on the template environment 
environment.tests['prime'] = is_prime
#usage 
{% if 42 is prime %}
    42 is a prime number
{% else %}
    42 is not a prime number
{% endif %}



##Jinja2- The Global Namespace
#Variables stored in the Environment.globals dict are available for all 
#even for imported templates too, even if they are imported without context. 

#variables in Template.globals are available to a specific template 







###Flask
#http://flask.pocoo.org/docs/0.12/

$ pip install Flask

##Flask- Quick example - Micro blog 
.
├── README.md
├── flaskr.py
├── schema.sql
├── static
│   └── style.css
└── templates
    ├── layout.html
    ├── login.html
    └── show_entries.html

#Execute 
$ python flaskr.py
#Log in with 
Username: admin
Password: default

#schema.sql 
drop table if exists entries;
create table entries (
  id integer primary key autoincrement,
  title text not null,
  text text not null
);

#style.css 
body            { font-family: sans-serif; background: #eee; }
a, h1, h2       { color: #377BA8; }
h1, h2          { font-family: 'Georgia', serif; margin: 0; }
h1              { border-bottom: 2px solid #eee; }
h2              { font-size: 1.2em; }

.page           { margin: 2em auto; width: 35em; border: 5px solid #ccc;
                  padding: 0.8em; background: white; }
.entries        { list-style: none; margin: 0; padding: 0; }
.entries li     { margin: 0.8em 1.2em; }
.entries li h2  { margin-left: -1em; }
.add-entry      { font-size: 0.9em; border-bottom: 1px solid #ccc; }
.add-entry dl   { font-weight: bold; }
.metanav        { text-align: right; font-size: 0.8em; padding: 0.3em;
                  margin-bottom: 1em; background: #fafafa; }
.flash          { background: #CEE5F5; padding: 0.5em;
                  border: 1px solid #AACBE2; }
.error          { background: #F0D6D6; padding: 0.5em; }


#flaskr.py
# -*- coding: utf-8 -*-

from sqlite3 import dbapi2 as sqlite3
from flask import Flask, request, session, g, redirect, url_for, abort, \
     render_template, flash


# create our little application :)
app = Flask(__name__)

# Load default config and Update below with defaults 
app.config.update(dict(
    DATABASE='flaskr.db',
    DEBUG=True,
    SECRET_KEY='development key',
    USERNAME='admin',
    PASSWORD='default'
))
#FLASKR_SETTINGS points to a configuration file 
#containing eg 
#DEBUG = False
#SECRET_KEY = '?\xbf,\xb4\x8d\xa3"<\x9c\xb0@\x0f5\xab,w\xee\x8d$0\x13\x8b83'

app.config.from_envvar('FLASKR_SETTINGS', silent=True)


def connect_db():
    """Connects to the specific database."""
    rv = sqlite3.connect(app.config['DATABASE']) #
    rv.row_factory = sqlite3.Row
    return rv


def init_db():
    """Creates the database tables."""
    #application context is required as there is no request being handled 
    #under request handling, application context is auto created 
    with app.app_context():
        # within this block, current_app points to app
        #get any attributes from app like current_app.attribute 
        db = get_db()
        #open file at root of the project 
        with app.open_resource('schema.sql', mode='r') as f:
            db.cursor().executescript(f.read())
        db.commit()

#flask.g is dict containing global variables 
def get_db():
    """Opens a new database connection if there is none yet for the
    current application context.
    """
    if not hasattr(g, 'sqlite_db'):
        g.sqlite_db = connect_db()
    return g.sqlite_db

#Registers a function to be called when the application context ends
@app.teardown_appcontext
def close_db(error):
    """Closes the database again at the end of the request."""
    if hasattr(g, 'sqlite_db'):
        g.sqlite_db.close()

#Handles "/" for GET (and HEAD ) method  
#must returns a Response object 
@app.route('/')
def show_entries():
    db = get_db()
    cur = db.execute('select title, text from entries order by id desc')
    entries = cur.fetchall()
    #entries can be access in template 
    #by default check from 'templates' dir 
    return render_template('show_entries.html', entries=entries)

#Handles "/add" for POST method  
@app.route('/add', methods=['POST'])
def add_entry():
    #session is a dict containing session data 
    if not session.get('logged_in'):
        abort(401)
    db = get_db()
    #request is DIed when request is getting handled 
    #request.form is dict 
    db.execute('insert into entries (title, text) values (?, ?)',
                 [request.form['title'], request.form['text']])
    db.commit()
    flash('New entry was successfully posted')
    #url_for returns URL given view function name 
    return redirect(url_for('show_entries'))


@app.route('/login', methods=['GET', 'POST'])
def login():
    error = None
    if request.method == 'POST':
        if request.form['username'] != app.config['USERNAME']:
            error = 'Invalid username'
        elif request.form['password'] != app.config['PASSWORD']:
            error = 'Invalid password'
        else:
            session['logged_in'] = True
            flash('You were logged in')
            return redirect(url_for('show_entries'))
    return render_template('login.html', error=error)


@app.route('/logout')
def logout():
    session.pop('logged_in', None)
    flash('You were logged out')
    return redirect(url_for('show_entries'))


if __name__ == '__main__':
    init_db()
    app.run()




#layout.html
#url_for returns URL given view function name 
#or returns URL given directory and filename 
<!doctype html>
<title>Flaskr</title>
<link rel=stylesheet type=text/css href="{{ url_for('static', filename='style.css') }}">
<div class=page>
  <h1>Flaskr</h1>
  <div class=metanav>
  {% if not session.logged_in %}
    <a href="{{ url_for('login') }}">log in</a>
  {% else %}
    <a href="{{ url_for('logout') }}">log out</a>
  {% endif %}
  </div>
  {% for message in get_flashed_messages() %}
    <div class=flash>{{ message }}</div>
  {% endfor %}
  {% block body %}{% endblock %}
</div>

#show_entries.html
#session is dictionary 
{% extends "layout.html" %}
{% block body %}
  {% if session.logged_in %}
    <form action="{{ url_for('add_entry') }}" method=post class=add-entry>
      <dl>
        <dt>Title:
        <dd><input type=text size=30 name=title>
        <dt>Text:
        <dd><textarea name=text rows=5 cols=40></textarea>
        <dd><input type=submit value=Share>
      </dl>
    </form>
  {% endif %}
  <ul class=entries>
  {% for entry in entries %}
    <li><h2>{{ entry.title }}</h2>{{ entry.text|safe }}
  {% else %}
    <li><em>Unbelievable.  No entries here so far</em>
  {% endfor %}
  </ul>
{% endblock %}


#login.html
{% extends "layout.html" %}
{% block body %}
  <h2>Login</h2>
  {% if error %}<p class=error><strong>Error:</strong> {{ error }}{% endif %}
  <form action="{{ url_for('login') }}" method=post>
    <dl>
      <dt>Username:
      <dd><input type=text name=username>
      <dt>Password:
      <dd><input type=password name=password>
      <dd><input type=submit value=Login>
    </dl>
  </form>
{% endblock %}


##Flask- Quick example - Packaging 
#packaging is required to deploy this thorugh external server 
#Once packaging is done, do 'pip install' to install in python path 

#dir structure 
/flaskr
    /flaskr
        __init__.py
        /static
        /templates
        flaskr.py
        schema.sql
    setup.py
    MANIFEST.in

#setup.py
from setuptools import setup

setup(
    name='flaskr',
    packages=['flaskr'],
    include_package_data=True,
    install_requires=[
        'flask',
    ],
)
#MANIFEST.in
#specify any special files that should be included in package 
#In this case, the static and templates directories need to be included, 
#as well as the schema
graft flaskr/templates
graft flaskr/static
include flaskr/schema.sql


#flaskr/__init__.py:
from .flaskr import app

#Update flaskr.py with below 
#such that we can call flask to initialize DB 
@app.cli.command('initdb') 
def initdb_command(): 
    """Creates the database tables.""" 
    init_db() 
    print('Initialized the database.') 



#install in 
#execute at root directory 
$ pip install --editable .

#Execute as 
$ set FLASK_APP=flaskr
$ set FLASK_DEBUG=true
$ flask initdb 
$ flask run 
#or to make it externally visible server 
$ flask run --host=0.0.0.0



##Flask - How to start Development server 

#Option-1 - COmmand Line - recommended
$ export FLASK_APP=my_application
$ export FLASK_DEBUG=1
$ flask run

#then start the server on http://localhost:5000/.
#to disable reloading 
$ flask run --no-reload



#Option-2 -In Code: reloading might not work 
if __name__ == '__main__':
    app.run()
    
    
    


##Flask - Jinja2 as Templating language 
#Jinja2 is configured by Flask as follows:
    • autoescaping is enabled for all templates ending in .html, 
      .htm, .xml ,.xhtml when using render_template().
    • autoescaping is enabled for all strings 
      when using render_template_string().
    • a template has the ability to opt in/out autoescaping 
      with the {% autoescape %} tag.
    • Flask inserts few global functions and helpers into the Jinja2 context


##Flask - Jinja2 - Standard Context
#global variables are available within Jinja2 templates

#these will not show up in the context of imported templates by default 
#use below to get these in imported template 
{% from '_helpers.html' import my_macro with context %}

#list of global variables 
config
    The current configuration object (flask.config)
    Only this is available even in imported templates 
request
    The current request object (flask.request). 
    This variable is unavailable if the template was rendered without an active request context.
session
    The current session object (flask.session). 
    This variable is unavailable if the template was rendered without an active request context.
g
    The request-bound object for global variables (flask.g). 
    This variable is unavailable if the template was rendered without an active request context.
url_for()
    The flask.url_for() function.
get_flashed_messages()
    The flask.get_flashed_messages() function.


##Flask - Jinja2 - additional Standard Filters
tojson()
    This function converts the given object into JSON representation. 
    <script type=text/javascript>
        doSomethingWith({{ user.username|tojson|safe }});
    </script>



##Flask - Jinja2 - Controlling Autoescaping
#by default enabled 
#To disable 
{% autoescape false %}
    <p>autoescaping is disabled here
    <p>{{ will_not_be_escaped }}
{% endautoescape %}


##Flask - Jinja2 - Registering user defined Filters

@app.template_filter('reverse')
def reverse_filter(s):
    return s[::-1]
#OR 
def reverse_filter(s):
    return s[::-1]
app.jinja_env.filters['reverse'] = reverse_filter

#usage 
{% for x in mylist | reverse %}
{% endfor %}



##Flask - Jinja2 - Context Processors
#To inject new variables automatically into the context of a template, 

#A context processor is a function that returns a dictionary. 
#The keys and values of this dictionary are then merged with the template context, for all templates in the app:

#variable 'user' available in the template with the value of g.user. 
@app.context_processor
def inject_user():
    return dict(user=g.user)


#To  make functions available to templates
@app.context_processor
def utility_processor():
    def format_price(amount, currency=u'€'):
        return u'{0:.2f}{1}'.format(amount, currency)
    return dict(format_price=format_price)
#Usage 
{{ format_price(0.33) }}



##Flask - werkzeug.datastructures.MultiDict 
#implements all standard dictionary methods. 
#Internally, it saves all values for a key as a list, 
#but the standard dict access methods will only return the first value for a key.
# http://werkzeug.pocoo.org/docs/0.13/datastructures/#werkzeug.datastructures.MultiDict

>>> d = MultiDict([('a', 'b'), ('a', 'c')])
>>> d
MultiDict([('a', 'b'), ('a', 'c')])
>>> d['a']
'b'
>>> d.getlist('a')
['b', 'c']
>>> 'a' in d
True

##Flask - werkzeug.datastructures.CombinedMultiDict(dicts=None)
#A read only MultiDict
#Created by passing multiple MultiDict instances as sequence 
#and it will combine the return values of all wrapped dicts:

>>> from werkzeug.datastructures import CombinedMultiDict, MultiDict
>>> post = MultiDict([('foo', 'bar')])
>>> get = MultiDict([('blub', 'blah')])
>>> combined = CombinedMultiDict([get, post])
>>> combined['foo']
'bar'
>>> combined['blub']
'blah'

##Flask -  werkzeug.datastructures.FileStorage(stream=None, 
            filename=None, name=None, content_type=None, 
            content_length=None, headers=None)
The FileStorage class is a thin wrapper over incoming files. 
It is used by the request object to represent uploaded files. 
    stream
        The input stream for the uploaded file. 
        This usually points to an open temporary file.        
    filename
        The filename of the file on the client.
    name
        The name of the form field.
    headers
        The multipart headers as Headers object. 
        This usually contains irrelevant information 
        but in combination with custom multipart requests the raw headers might be interesting.
    close()
        Close the underlying file if possible.
    content_length
        The content-length sent in the header. Usually not available
    content_type
        The content-type sent in the header. Usually not available
    mimetype
        Like content_type, but without parameters (eg, without charset, type etc.) and always lowercase. For example if the content type is text/HTML; charset=utf-8 the mimetype would be 'text/html'.
    mimetype_params
        The mimetype parameters as dict. 
        For example if the content type is text/html; charset=utf-8 
        the params would be {'charset': 'utf-8'}    
    read( [size]) 
        Read at most size bytes from the file
    save(dst, buffer_size=16384)
        Save the file to a destination path or file object. 
        If the destination is a file object you have to close it yourself after the call. 
        The buffer size is the number of bytes held in memory during the copy process. 
        It defaults to 16KB.

##Flask -  Application Globals - flask.g - an object 
#Flask provides a special object 
#that ensures it is only valid for the active request 
#and that will return different values for each request. 

flask.g.user = "OK"
user = getattr(flask.g, 'user', None)
user = flask.g.get('user', None)
 

##Flask -  flask.Request(environ, populate_request=True, shallow=False)
#The global flask.request is instance of Request 
#and it's data is valid only for that request 

#attributes 
form
    A MultiDict with the parsed form data from POST or PUT requests. 
args
    A MultiDict with the parsed contents of the query string. 
    (The part in the URL after the question mark).
values
    A CombinedMultiDict with the contents of both form and args.
cookies
    A dict with the contents of all cookies transmitted with the request.
stream
    If the incoming form data was not encoded with a known mimetype 
    the data is stored unmodified in this stream for consumption. 
    Most of the time it is a better idea to use data 
    which will give you that data as a string. 
    The stream only returns the data once.
headers
    The incoming request headers as a dictionary like object.
data
    Contains the incoming request data as string 
    in case it came with a mimetype Flask does not handle.
files
    A MultiDict with files uploaded as part of a POST or PUT request. 
    Each file is stored as FileStorage object. 
environ
    The underlying WSGI environment.
method
    The current request method (POST, GET etc.)
path
full_path
script_root
url
base_url
url_root
    Provides different ways to look at the current IRI.   
    For example http://www.example.com/myapplication/%CF%80/page.html?x=y
        path u'/π/page.html' 
        full_path u'/π/page.html?x=y' 
        script_root u'/myapplication' 
        base_url u'http://www.example.com/myapplication/π/page.html' 
        url u'http://www.example.com/myapplication/π/page.html?x=y' 
        url_root u'http://www.example.com/myapplication/' 
is_xhr
    True if the request was triggered via a JavaScript XMLHttpRequest. 
    This only works with libraries that support the X-Requested-With header 
    and set it to XMLHttpRequest. 
    Libraries that do that are prototype, jQuery and Mochikit 
blueprint
    The name of the current blueprint
endpoint
    The endpoint that matched the request. 
    This in combination with view_args can be used to reconstruct 
    the same or a modified URL. 
    If an exception happened when matching, this will be None.
get_json(force=False, silent=False, cache=True)
    Parses the incoming JSON request data and returns it. 
    By default this function will return None 
    if the mimetype is not application/json 
    but this can be overridden by the force parameter. 
    If parsing fails the on_json_loading_failed() method 
    on the request object will be invoked.
    Parameters:
        •force – if set to True the mimetype is ignored.
        •silent – if set to True this method will fail silently and return None.
        •cache – if set to True the parsed JSON data is remembered on the request.
is_json
    Indicates if this request is JSON or not. 
    By default a request is considered to include JSON data
     if the mimetype is application/json or application/*+json.
json
    If the mimetype is application/json this will contain the parsed JSON data. 
    Otherwise this will be None.
    The get_json() method should be used instead.
max_content_length
    Read-only view of the MAX_CONTENT_LENGTH config key.
module
    The name of the current module if the request was dispatched to an actual module. This is deprecated functionality, use blueprints instead.
on_json_loading_failed(e)
    Called if decoding of the JSON data failed. 
    The return value of this method is used by get_json() when an error occurred. The default implementation just raises a BadRequest exception.
routing_exception = None
    If matching the URL failed, this is the exception that will be raised 
    / was raised as part of the request handling. 
    This is usually a NotFound exception or something similar.
url_rule = None
    The internal URL rule that matched the request.
    This can be useful to inspect which methods are allowed for the URL 
    from a before/after handler (request.url_rule.methods) etc.
view_args = None
    A dict of view arguments that matched the request. 
    If an exception happened when matching, this will be None.


##Flask -  flask.Response(response=None, status=None, headers=None, mimetype=None, content_type=None, direct_passthrough=False)
#call flask.make_response() to create a instance 
#and then modify it before returning from vew function

#Attributes 
headers
    A Headers object representing the response headers.
status
    A string with a response status.
status_code
    The response status as integer.
data
    A descriptor that calls get_data() and set_data(). 
    This should not be used and will eventually get deprecated.
mimetype
    The mimetype (content type without charset etc.)
set_cookie(key, value=”, max_age=None, expires=None, path=’/’, domain=None, secure=None, httponly=False)
    Sets a cookie. The parameters are the same as in the cookie Morsel object in the Python standard library 
    but it accepts unicode data, too.
    Parameters:
        •key – the key (name) of the cookie to be set.
        •value – the value of the cookie.
        •max_age – should be a number of seconds, or None (default) if the cookie should last only as long as the client’s browser session.
        •expires – should be a datetime object or UNIX timestamp.
        •domain – if you want to set a cross-domain cookie. For example, domain=".example.com" will set a cookie that is readable by the domain www.example.com, foo.example.com etc. Otherwise, a cookie will only be readable by the domain that set it.
        •path – limits the cookie to a given path, per default it will span the whole domain.
     


##Flask -  flask.session - a dict 
#set Flask.secret_key(via app.config['SECRET_KEY'] to enable session, uses  a signed cookie.
#A session basically makes it possible to remember information from one request to another. 

#Attributes
new
    True if the session is new, False otherwise.
modified
    True if the session object detected a modification. 
    Be advised that modifications on mutable structures 
    are not picked up automatically, 
    Hence use this attribute 
    # this change is not picked up because a mutable object (here
    # a list) is changed.
    session['objects'].append(42)
    # so mark it as modified yourself
    session.modified = True
permanent
    If set to True the session lives for permanent_session_lifetime seconds. 
    The default is 31 days. 
    If set to False (which is the default) the session will be deleted 
    when the user closes the browser.



##Flask - URL Route Registrations

#three ways to define rules for the routing system:
1.use the flask.Flask.route() decorator.
2.use the flask.Flask.add_url_rule() function.
3.access the underlying Werkzeug routing system 
  which is exposed as flask.Flask.url_map.

#Variable parts in the route can be specified with <> 
#and passed to the view function as keyword arguments.
#eg 
(/user/<username>) # username would be string 
#to convert to other type, use 
<converter:name>.
#converter 
string      accepts any text without a slash (the default) 
int         accepts integers 
float       like int but for floating point values 
path        like the default but also accepts slashes 
any         matches one of the items provided 
uuid        accepts UUID strings 

#Example 
@app.route('/')
def index():
    pass

@app.route('/<username>')
def show_user(username):
    pass

@app.route('/post/<int:post_id>')
def show_post(post_id):
    pass
    
#below are equivalent 
@app.route('/')
def index():
    pass
#OR
def index():
    pass
app.add_url_rule('/', 'index', index)
#OR
app.view_functions['index'] = index




#Custom converters - using flask.Flask.url_map.
from werkzeug.routing import BaseConverter

class ListConverter(BaseConverter):
    def to_python(self, value):
        return value.split(',')
    def to_url(self, values):
        return ','.join(super(ListConverter, self).to_url(value)
                        for value in values)

app = Flask(__name__)
app.url_map.converters['list'] = ListConverter
#usage , where users are user1,user2,...
@app.route('/<list:users>')
def show_post(users):
    pass


#how Flask deals with trailing slashes.
#The idea is to keep each URL unique so the following rules apply:
1.If a rule ends with a slash and is requested without a slash by the user, 
  the user is automatically redirected to the same page with a trailing slash attached.
2.If a rule does not end with a trailing slash 
  and the user requests the page with a trailing slash, 
  a 404 not found is raised.

#To define multiple rules for the same function. 
#They have to be unique however. 
#Defaults can also be specified. 

@app.route('/users/', defaults={'page': 1})
@app.route('/users/page/<int:page>')
def show_users(page):
    pass


#parameters for route(rule, **options) 
#and add_url_rule(rule, endpoint=None, view_func=None, **options) 
rule        the URL rule as string 
endpoint    the endpoint for the registered URL rule. 
            Flask itself assumes that the name of the view function is the name of the endpoint if not explicitly stated. 
view_func   the function to call when serving a request to the provided endpoint. If this is not provided one can specify the function later by storing it in the view_functions dictionary with the endpoint as key. 
defaults    A dictionary with defaults for this rule. 
subdomain   specifies the rule for the subdomain in case subdomain matching is in use. 
            If not specified the default subdomain is assumed. 
options     the options to be forwarded to the underlying Rule object. 
            for example a method options. 
            methods is a list of methods this rule should be limited to (GET, POST etc.). 
            By default a rule just listens for GET (and implicitly HEAD). 
      
      
#To decorate view function 
#such tat view function is enhanced, use python syntax of decorator 

#This example assumes that the login page is called 'login' 
#and that the current user is stored in g.user and is None if there is no-one logged in.

from functools import wraps
from flask import g, request, redirect, url_for

def login_required(f):
    @wraps(f)
    def decorated_function(*args, **kwargs):
        if g.user is None:
            return redirect(url_for('login', next=request.url))
        return f(*args, **kwargs)
    return decorated_function

#The next value will exist in request.args after a GET request for the login page. 
#You’ll have to pass it along when sending the POST request from the login form. 
#You can do this with a hidden input tag, 
#then retrieve it from request.form when logging the user in.
<input type="hidden" name="next" value="{{ request.args.get('next', '') }}"/>

#To use the decorator, apply it as innermost decorator to a view function. 
#remember that the route() decorator is the outermost always 

@app.route('/secret_page')
@login_required
def secret_page():
    pass

          

          
##Flask - what view function can return 
#View function must return a Response or object convertable to Response           
            
#Returning a string 
@app.route('/string')
def string():
    return 'hello world!\n'


#render_template returns a unicode/str string 
@app.route('/template')
def template():
    return render_template('index.html', foo=42)

# Returning a Response
@app.route('/response')
def response():
    response = Response('Hello world!\n', content_type='text/plain')
    response.set_cookie('firstname', 'xurui')
    response.set_cookie('lastname', 'yan')
    return response


# Returning a Response with make_response 
@app.route('/make_response')
def make_response():
    response = flask.make_response(('Hello world!\n'), {'X-My-Header': 'foo'})
    response.set_cookie('username', 'yxr')
    return response



#Returning Json with the application/json mimetype 
@app.route('/jsonify')
def jsonify():
    return flask.jsonify({'username': 'yxr'})



#Return chunked-encoded response
@app.route('/chunked')
def chunked():
    def gen():
        yield 'Hello '
        yield 'World!\n'
    return Response(gen())



#The return value can be a callable WSGI app function, for example: 
def simple_app(environ, start_response):
    status = '200 OK'
    response_headers = [('Content-Length', '13')]
    start_response(status, response_headers)
    return ['Hello World!\n']

@app.route('/wsgi')
def wsgi():
    return simple_app


#Returning HTTPException
#werkzeug.exceptions.HTTPException can be called as WSGI application to render a default error page. 
@app.route('/exception')
def exception():
    # equivalent to raise NotFound
    return NotFound()

#This is equivalent to raise NotFound
#difference between return and raise
#The exception raised will be handled in handle_http_exception. 
#If no handler is registered , it will be return directly like above 


#Returning abort
@app.route('/abort')
def test_abort():
    abort(500, 'something is wrong!\n')



#Returning tuple or list
#Here response can be any form described above 
#3-tuple in (response, status, headers) form
@app.route('/tuple3')
def tuple3():
    return 'hello world!\n', 200, {'X-My-Header': 'foo'}

#2-tuple in (response, status) or (response, headers) form.
#If the the second item is instance of Headers, dict, tuple or list, 
#it is unpacked as the second form. 
#status can be int or any type that can be converted by int constructor. 
@app.route('/tuple2_status')
def tuple2_status():
    return 'hello world!\n', 200

@app.route('/tuple2_header')
def tuple2_header():
    return 'hello world!\n', {'X-My-Header': 'foo'}

@app.route('/tuple2_response')
def tuple2_response():
    response = flask.make_response(('hello world!\n'))
    return response, {'X-My-Header':' foo'}

#response is a callable WSGI app which is rarely used 
@app.route('/tuple2_wsgi')
def tuple2_wsgi():
    return simple_app, {'X-My-Header': 'foo'}


#The status value will override the existing status code 
#and headers will add additional headers. 
@app.route('/tuple2_extend_header')
def tuple2_extend_header():
    response = flask.make_response(('hello world!\n'), {'X-My-Header': 'foo'})
    return response, {'X-My-Header':' bar'}

#None is not allowed
@app.route('/none')
def none():
    # return nothing is equivalent to `return None`
    pass

#response is not allowed to be None even in tuple form 
@app.route('/tuple_none')
def tuple_none():
    return None, 404
    
#Returning a redirect 
@app.route('/redirect')
def redi():
    return redirect(url_for('login', next=request.url))

##Flask - Template Rendering
flask.render_template(template_name_or_list, **context)
    Renders a template from the template folder with the given context.
    Parameters:
    •template_name_or_list – the name of the template to be rendered, 
                            or an iterable with template names the first one existing will be rendered
    •context – the variables that should be available in the context of the template.
 
flask.render_template_string(source, **context)
    Renders a template from the given template source string 
    with the given context. Template variables will be autoescaped.

 
flask.get_template_attribute(template_name, attribute)
    Loads a macro (or variable) a template exports. 
    # _cider.html with the following contents:
    {% macro hello(name) %}Hello {{ name }}!{% endmacro %}
    #You can access this from Python code like this:
    hello = get_template_attribute('_cider.html', 'hello')
    return hello('World')  

    
    
    
    
##Flask - Class-Based Views - using flask.views.View
#Attributes of flask.views.View
classmethod as_view(name, *class_args, **class_kwargs)
    Converts the class into an actual view function 
    The arguments passed to as_view() are forwarded to the constructor of the class.
decorators = ()
dispatch_request()
methods = None

#Example 
class MyView(View):
    methods = ['GET']
    def dispatch_request(self, name):
        return 'Hello %s!' % name

app.add_url_rule('/hello/<name>', view_func=MyView.as_view('myview'))




#To use decorators in View based class 
#The decorators stored in the decorators list are applied one after another 
#when the view function is created. 
class SecretView(View):
    methods = ['GET']
    decorators = [superuser_required]

    def dispatch_request(self):
        ...





##Flask - Class-Based Views - flask.views.MethodView
#Like a regular class-based view 
#but that dispatches requests to particular methods. 

class CounterAPI(MethodView):

    def get(self):  #for get 
        return session.get('counter', 0)

    def post(self): #for post 
        session['counter'] = session.get('counter', 0) + 1
        return 'OK'

app.add_url_rule('/counter', view_func=CounterAPI.as_view('counter'))


##Flask - Useful Functions and Classes
flask.current_app
    Points to the application handling the request. 
    Use with app.app_context() method to get to correct app 

flask.has_request_context()
    To test if a request context is there  or not 

    class User(db.Model):
        def __init__(self, username, remote_addr=None):
            self.username = username
            if remote_addr is None and has_request_context():
                remote_addr = request.remote_addr
            self.remote_addr = remote_addr
    #or just test any of the context bound objects (such as request or g for truthness):
    class User(db.Model):
        def __init__(self, username, remote_addr=None):
            self.username = username
            if remote_addr is None and request:
                remote_addr = request.remote_addr
            self.remote_addr = remote_addr

flask.copy_current_request_context(f)
    A helper function that decorates a function to retain 
    the current request context. 
    This is useful when working with greenlets. 
    The moment the function is decorated a copy of the request context is created 
    and then pushed when the function is called.
    #Example 
    import gevent
    from flask import copy_current_request_context
    @app.route('/')
    def index():
        @copy_current_request_context
        def do_some_work():
            # do some work here, it can access flask.request like in the view function.
            ...
        gevent.spawn(do_some_work)
        return 'Regular response'


flask.has_app_context()
    Works like has_request_context() but for the application context. 


flask.url_for(endpoint, **values)
    Generates a URL to the given endpoint method provided.
    Variable arguments that are unknown to the target endpoint are appended 
    to the generated URL as query arguments. 
    If the value of a query argument is None, the whole pair is skipped. 
    
flask.abort()
    When passed a dict of code -> exception items 
    it can be used as callable that raises exceptions. 
    If the first argument to the callable is an integer 
    it will be looked up in the mapping, 
    if it’s a WSGI application it will be raised in a proxy exception.
    The rest of the arguments are forwarded to the exception constructor.
    
flask.redirect(location, code=302, Response=None)
    Returns a response object (a WSGI application) that, 
    if called, redirects the client to the target location. 
    Supported codes are 301, 302, 303, 305, and 307. 


flask.make_response(*args)
    Used for  setting additional headers in a view.
    This function accepts the very same arguments that can be returned from a view function. This for example creates a response with a 404 error code:
    response = make_response(render_template('not_found.html'), 404)
    If view looked like this and you want to add a new header:
    def index():
        return render_template('index.html', foo=42)
    You can now do something like this:
    def index():
        response = make_response(render_template('index.html', foo=42))
        response.headers['X-Parachutes'] = 'parachutes are cool'
        return response


flask.after_this_request(f)
    Executes a function after this request. 
    This is useful to modify response objects. 
    The function is passed the response object 
    and has to return the same or a new one.
    #Example 
    @app.route('/')
    def index():
        @after_this_request
        def add_header(response):
            response.headers['X-Foo'] = 'Parachute'
            return response
        return 'Hello World!'


flask.send_file(filename_or_fp, mimetype=None, as_attachment=False, 
            attachment_filename=None, add_etags=True, cache_timeout=None, conditional=False, last_modified=None)
    Sends the contents of a file to the client. 
    Please never pass filenames to this function from user sources; 
    you should use send_from_directory() instead.
    Parameters:
        •filename_or_fp – the filename of the file to send in latin-1. This is relative to the root_path if a relative path is specified. Alternatively a file object might be provided in which case X-Sendfile might not work and fall back to the traditional method. Make sure that the file pointer is positioned at the start of data to send before calling send_file().
        •mimetype – the mimetype of the file if provided. If a file path is given, auto detection happens as fallback, otherwise an error will be raised.
        •as_attachment – set to True if you want to send this file with a Content-Disposition: attachment header.
        •attachment_filename – the filename for the attachment if it differs from the file’s filename.
        •add_etags – set to False to disable attaching of etags.
        •conditional – set to True to enable conditional responses.
        •cache_timeout – the timeout in seconds for the headers. When None (default), this value is set by get_send_file_max_age() of current_app.
        •last_modified – set the Last-Modified header to this value, a datetime or timestamp. If a file was passed, this overrides its mtime.
 
flask.send_from_directory(directory, filename, **options)
    Send a file from a given directory with send_file(). 
    This is a secure way to quickly expose static files 
    from an upload folder or something similar.
    #Example usage:
    @app.route('/uploads/<path:filename>')
    def download_file(filename):
        return send_from_directory(app.config['UPLOAD_FOLDER'],
                                   filename, as_attachment=True)

flask.safe_join(directory, *pathnames)
    Safely join directory and zero or more untrusted pathnames components.
    #Example usage:
    @app.route('/wiki/<path:filename>')
    def wiki_page(filename):
        filename = safe_join(app.config['WIKI_FOLDER'], filename)
        with open(filename, 'rb') as fd:
            content = fd.read()  # Read and process the file content...




flask.escape(s) -> markup
    Convert the characters &, <, >, ‘, and ” in string s to HTML-safe sequences. 


class flask.Markup
    Marks a string as being safe for inclusion in HTML/XML output without needing to be escaped. 
    #attributes 
    classmethod escape(s)
        Escape the string.

    striptags()
        Unescape markup into an text_type string and strip all tags. 
        This also resolves known HTML4 and XHTML entities. 
        Whitespace is normalized to one:
        >>> Markup("Main &raquo;  <em>About</em>").striptags()
        u'Main \xbb About'

    unescape()
        Unescape markup again into an text_type string. 
        This also resolves known HTML4 and XHTML entities:
        >>> Markup("Main &raquo; <em>About</em>").unescape()
        u'Main \xbb <em>About</em>'
    #Example 
    >>> Markup("Hello <em>World</em>!")
    Markup(u'Hello <em>World</em>!')
    >>> class Foo(object):
    ...  def __html__(self):
    ...   return '<a href="#">foo</a>'
    ...
    >>> Markup(Foo())
    Markup(u'<a href="#">foo</a>')


    If you want object passed being always treated as unsafe 
    you can use the escape() classmethod to create a Markup object:
    >>> Markup.escape("Hello <em>World</em>!")
    Markup(u'Hello &lt;em&gt;World&lt;/em&gt;!')


    Operations on a markup string are markup aware 
    which means that all arguments are passed through the escape() function:
    >>> em = Markup("<em>%s</em>")
    >>> em % "foo & bar"
    Markup(u'<em>foo &amp; bar</em>')
    >>> strong = Markup("<strong>%(text)s</strong>")
    >>> strong % {'text': '<blink>hacker here</blink>'}
    Markup(u'<strong>&lt;blink&gt;hacker here&lt;/blink&gt;</strong>')
    >>> Markup("<em>Hello</em> ") + "<foo>"
    Markup(u'<em>Hello</em> &lt;foo&gt;')




##Flask - Message Flashing 
flask.flash(message, category=’message’)
    Flashes a message to the next request. 
    In order to remove the flashed message from the session 
    and to display it to the user, the template has to call get_flashed_messages().
    Parameters:
        •message – the message to be flashed.
        •category – the category for the message. 
                    The following values are recommended: 
                    'message' for any kind of message, 
                    'error' for errors, 
                    'info' for information messages 
                    and 'warning' for warnings. 
                    However any kind of string can be used as category.
 
flask.get_flashed_messages(with_categories=False, category_filter=[])
    Pulls all flashed messages from the session and returns them. 
    Further calls in the same request to the function will return the same messages. 
    By default just the messages are returned, 
    but when with_categories is set to True, 
    the return value will be a list of tuples in the form (category, message) instead.



##Flask - JSON Support
#For usage examples, read the json documentation in the standard library
from flask import json


#The following extensions are by default applied to the stdlib’s JSON module:
1.datetime objects are serialized as RFC 822 strings.
2.Any object with an __html__ method (like Markup) will have that method called 
  and then the return value is serialized as string.

#In template 
<script type=text/javascript>
    doSomethingWith({{ user.username|tojson|safe }});
</script>


# Methods 
flask.json.jsonify(*args, **kwargs)
    1.Single argument: Passed straight through to dumps().
    2.Multiple arguments: Converted to an array before being passed to dumps().
    3.Multiple keyword arguments: Converted to a dict before being passed to dumps().
    4.Both args and kwargs: Behavior undefined and will throw an exception.
    #example 
    from flask import jsonify

    @app.route('/_get_current_user')
    def get_current_user():
        return jsonify(username=g.user.username,
                       email=g.user.email,
                       id=g.user.id)


    #This will send a JSON response like this to the browser:
    {
        "username": "admin",
        "email": "admin@localhost",
        "id": 42
    }

flask.json.dumps(obj, **kwargs)
    Serialize obj to a JSON formatted str 

flask.json.dump(obj, fp, **kwargs)
    Like dumps() but writes into a file object.

flask.json.loads(s, **kwargs)
    Unserialize a JSON object from a string s by using the application’s configured decoder (json_decoder) if there is an application on the stack.

flask.json.load(fp, **kwargs)
    Like loads() but reads from a file object.

    
    





##Flask - Configuration - class flask.Config(root_path, defaults=None)
#Works exactly like a dict 
#only uppercase keys are added to the config

#Either fill the config from a config file:
app.config.from_pyfile('yourconfig.cfg')

#Or alternatively , define the configuration options in the module 
DEBUG = True
SECRET_KEY = 'development key'
app.config.from_object(__name__)
#Example 
app.config.from_object('yourapplication.default_config')
#or 
from yourapplication import default_config
app.config.from_object(default_config)

#Or load configurations from an environment variable pointing to a file:
app.config.from_envvar('YOURAPPLICATION_SETTINGS')

#Or load from json file 
app.config.from_json(filename, silent=False)

#or Updates the config like update() ignoring items with non-upper keys.
app.config.from_mapping(*mapping, **kwargs)

#Or update using dict.update 
app.config..update(another_dict, **F)

#Other methods of Config 
get_namespace(namespace, lowercase=True, trim_namespace=True)
    Returns a dictionary containing a subset of configuration options 
    that match the specified namespace/prefix
    app.config['IMAGE_STORE_TYPE'] = 'fs'
    app.config['IMAGE_STORE_PATH'] = '/var/app/images'
    app.config['IMAGE_STORE_BASE_URL'] = 'http://img.website.com'
    image_store_config = app.config.get_namespace('IMAGE_STORE_')

    #The resulting dictionary image_store_config would look like:
    {
        'type': 'fs',
        'path': '/var/app/images',
        'base_url': 'http://img.website.com'
    }


##Flask - Extensions - flask.ext
#to use an extension named “Flask-Foo” 

from flask.ext import foo



##Flask - Stream Helpers
flask.stream_with_context(generator_or_function)
    Request contexts disappear when the response is started on the server. 
    This function however can help you keep the context around for longer
    #Example 
    from flask import stream_with_context, request, Response

    @app.route('/stream')
    def streamed_response():
        @stream_with_context
        def generate():
            yield 'Hello '
            yield request.args['name']
            yield '!'
        return Response(generate())


    #Alternatively it can also be used around a specific generator:
    from flask import stream_with_context, request, Response

    @app.route('/stream')
    def streamed_response():
        def generate():
            yield 'Hello '
            yield request.args['name']
            yield '!'
        return Response(stream_with_context(generate()))



##Flask - Application Object

class flask.Flask(import_name, static_path=None, static_url_path=None, 
        static_folder='static', template_folder='templates', 
        instance_path=None, instance_relative_config=False, root_path=None)
#The flask object implements a WSGI application and acts as the central object.
#Parameters:
    •import_name – the name of the application package
    •static_url_path – can be used to specify a different path for the static files on the web. Defaults to the name of the static_folder folder.
    •static_folder – the folder with static files that should be served at static_url_path. Defaults to the 'static' folder in the root path of the application.
    •template_folder – the folder that contains the templates that should be used by the application. Defaults to 'templates' folder in the root path of the application.
    •instance_path – An alternative instance path for the application. By default the folder 'instance' next to the package or module is assumed to be the instance path.
    •instance_relative_config – if set to True relative filenames for loading the config are assumed to be relative to the instance path instead of the application root.
    •root_path – Flask by default will automatically calculate the path to the root of the application. In certain situations this cannot be achieved (for instance if the package is a Python 3 namespace package) and needs to be manually defined.
 
#Example 
from flask import Flask
app = Flask(__name__)

#For single module, __name__ is always the correct value. 
#For a package, it's usually recommended to hardcode the name of package 
#For example if application is defined in yourapplication/app.py 
app = Flask('yourapplication')
#or
app = Flask(__name__.split('.')[0])


##Flask - Application Object - Attributes
add_template_filter(f, name=None)
    Register a custom template filter. 
    Works exactly like the template_filter() decorator.
    Parameters:
        name – the optional name of the filter, otherwise the function name will be used. 

add_template_global(f, name=None)
    Register a custom template global function. 
    Works exactly like the template_global() decorator.

add_template_test(f, name=None)
    Register a custom template test. 
    Works exactly like the template_test() decorator.

add_url_rule(rule, endpoint=None, view_func=None, **options)
    Connects a URL rule. 
    Works exactly like the route() decorator. 
    If a view_func is provided it will be registered with the endpoint.


after_request(f)
    Register a function to be run after each request.
    function must take one parameter, an instance of response_class 
    and return a new response object or the same

after_request_funcs = None
    A dictionary with lists of functions that should be called after each request. 
    To register a function here, use the after_request() decorator.
    
    
app_context()
    Binds the application only. 
    For as long as the application is bound to the current context 
    the flask.current_app points to that application. 
    An application context is automatically created 
    when a request context is pushed if necessary.
    #Example 
    with app.app_context():
        ...



auto_find_instance_path()
    Tries to locate the instance path if it was not provided 
    to the constructor of the application class. 
    It will basically calculate the path to a folder named instance 
    next to your main file or the package.


before_first_request(f)
    Registers a function to be run before the first request 
    to this instance of the application.
    The function will be called without any arguments 
    and its return value is ignored.


before_first_request_funcs = None
    A lists of functions that should be called at the beginning of the first request to this instance. 
    To register a function here, use the before_first_request() decorator.


before_request(f)
    Registers a function to run before each request.
    The function will be called without any arguments. 
    If the function returns a non-None value, 
    it's handled as if it was the return value from the view 
    and further request handling is stopped.

before_request_funcs = None
    A dictionary with lists of functions that should be called 
    at the beginning of the request. 
    To register a function here, use the before_request() decorator.
    
blueprints = None
    all the attached blueprints in a dictionary by name. 
    Blueprints can be attached multiple times 
    so this dictionary does not tell you how often they got attached.


cli = None
    The click command line context for this application. 
    Commands registered here show up in the flask command 
    once the application has been discovered. 
    The default commands are provided by Flask itself and can be overridden.
    This is an instance of a click.Group object.

config = None
    The configuration dictionary as Config. 
    This behaves exactly like a regular dictionary 
    but supports additional methods to load a config from files.


context_processor(f)
    Registers a template context processor function.
    
    
create_url_adapter(request)
    Creates a URL adapter for the given request. 
    The URL adapter is created at a point 
    This can now also be called without a request object when the URL adapter is created for the application context.

debug
    The debug flag. 
    Set this to True to enable debugging of the application. 

default_config = ImmutableDict({'JSON_AS_ASCII': True, 
        'USE_X_SENDFILE': False, 'SESSION_COOKIE_PATH': None, 'SESSION_COOKIE_DOMAIN': None, 'SESSION_COOKIE_NAME': 'session', 'DEBUG': False, 'LOGGER_HANDLER_POLICY': 'always', 'LOGGER_NAME': None, 'SESSION_COOKIE_SECURE': False, 'SECRET_KEY': None, 'EXPLAIN_TEMPLATE_LOADING': False, 'MAX_CONTENT_LENGTH': None, 'PROPAGATE_EXCEPTIONS': None, 'APPLICATION_ROOT': None, 'SERVER_NAME': None, 'PREFERRED_URL_SCHEME': 'http', 'JSONIFY_PRETTYPRINT_REGULAR': True, 'TESTING': False, 'PERMANENT_SESSION_LIFETIME': datetime.timedelta(31), 'TEMPLATES_AUTO_RELOAD': None, 'TRAP_BAD_REQUEST_ERRORS': False, 'JSON_SORT_KEYS': True, 'JSONIFY_MIMETYPE': 'application/json', 'SESSION_COOKIE_HTTPONLY': True, 'SEND_FILE_MAX_AGE_DEFAULT': datetime.timedelta(0, 43200), 'PRESERVE_CONTEXT_ON_EXCEPTION': None, 'SESSION_REFRESH_EACH_REQUEST': True, 'TRAP_HTTP_EXCEPTIONS': False})
    Default configuration parameters.

dispatch_request()
    Does the request dispatching. 

do_teardown_appcontext(exc=<object object>)
    Called when an application context is popped. 
    This works pretty much the same as do_teardown_request() 
    but for the application context.


do_teardown_request(exc=<object object>)
    Called after the actual request dispatching 
    and will call every  teardown_request() decorated function. 
    This is not actually called by the Flask object itself 
    but is always triggered when the request context is popped. 
  
  
endpoint(endpoint)
    A decorator to register a function as an endpoint. 
    @app.endpoint('example.endpoint')
    def example():
        return "example"




extensions = None
    a place where extensions can store application specific state. 
    For example this is where an extension could store database engines 
    and similar things. 
    #Example 
    if not hasattr(app, 'extensions'):
        app.extensions = {}
    app.extensions['extensionname'] = SomeObject()
    The key must match the name of the extension module. 
    For example in case of a "Flask-Foo" extension in flask_foo, 
    the key would be 'foo'.



full_dispatch_request()
    Dispatches the request and on top of that performs request pre and postprocessing as well as HTTP exception catching and error handling.



get_send_file_max_age(filename)
    Provides default cache_timeout for the send_file() functions.
    By default, this function returns SEND_FILE_MAX_AGE_DEFAULT from the configuration of current_app.
    #Or override in subclass
    class MyFlask(flask.Flask):
        def get_send_file_max_age(self, name):
            if name.lower().endswith('.js'):
                return 60
            return flask.Flask.get_send_file_max_age(self, name)



got_first_request
    This attribute is set to True 
    if the application started handling the first request.



handle_exception(e)
    Default exception handling that kicks in 
    when an exception occurs that is not caught. 
    In debug mode the exception will be re-raised immediately, 
    otherwise it is logged and the handler for a 500 internal server error is used. 
    If no such handler exists, a default 500 internal server error message is displayed.


error_handler_spec = None
    A dictionary of all registered error handlers. 
    To register a error handler, use the errorhandler() decorator.
    
errorhandler(code_or_exception)
    A decorator that is used to register a function given an error code. 
    #Example 
    @app.errorhandler(404)
    def page_not_found(error):
        return 'This page does not exist', 404

    You can also register handlers for arbitrary exceptions:
    @app.errorhandler(DatabaseError)
    def special_exception_handler(error):
        return 'Database connection failed', 500

    #or
    def page_not_found(error):
        return 'This page does not exist', 404
    app.error_handler_spec[None][404] = page_not_found

handle_http_exception(e)
    Handles an HTTP exception. 
    By default this will invoke the registered error handlers 
    and fall back to returning the exception as response.


handle_url_build_error(error, endpoint, values)
    Handle BuildError on url_for().
    
handle_user_exception(e)
    This method is called whenever an exception occurs that should be handled. 
    A special case are HTTPExceptions which are forwarded 
    by this function to the handle_http_exception() method. 
    This function will either return a response value 
    or reraise the exception with the same traceback.



has_static_folder
    This is True if the package bound object's container has a folder for static files.


instance_path = None
    Holds the path to the instance folder.

iter_blueprints()
    Iterates over all blueprints by the order they were registered.

jinja_env
    The Jinja2 environment used to load templates.
    
jinja_environment
    alias of Environment
    
jinja_loader
    The Jinja loader for this package bound object.

jinja_options = ImmutableDict({'extensions': ['jinja2.ext.autoescape', 'jinja2.ext.with_']})
    Options that are passed directly to the Jinja2 environment.

json_decoder
    The JSON decoder class to use. Defaults to JSONDecoder.

json_encoder
    The JSON encoder class to use. Defaults to JSONEncoder.

log_exception(exc_info)
    Logs an exception. 
    This is called by handle_exception() if debugging is disabled 
    and right before the handler is called. 
    The default implementation logs the exception as error on the logger.


logger
    A logging.Logger object for this application. 
    The default configuration is to log to stderr 
    if the application is in debug mode. 
    This logger can be used to log messages. 
    app.logger.debug('A value for debugging')
    app.logger.warning('A warning occurred (%d apples)', 42)
    app.logger.error('An error occurred')



logger_name
    The name of the logger to use. 
    By default the logger name is the package name passed to the constructor.


make_default_options_response()
    This method is called to create the default OPTIONS response. 
    This can be changed through subclassing to change the default behavior of OPTIONS responses.


make_null_session()
    Creates a new instance of a missing session. 
    Instead of overriding this method we recommend replacing the session_interface.


make_response(rv)
    Converts the return value from a view function to a real response object 
    that is an instance of response_class.
    @The following types are allowed for rv:
    response_class  the object is returned unchanged 
    str             a response object is created with the string as body 
    unicode         a response object is created with the string encoded to utf-8 as body 
    WSGI function   the function is called as WSGI application and buffered as response object 
    tuple           A tuple in the form (response, status, headers) or (response, headers) where response is any of the types defined here, status is a string or an integer and headers is a list or a dictionary with header values. 


make_shell_context()
    Returns the shell context for an interactive shell for this application. 
    This runs all the registered shell context processors.


name
    The name of the application. 


open_instance_resource(resource, mode='rb')
    Opens a resource from the application's instance folder (instance_path). 
    Otherwise works like open_resource(). 
    Instance resources can also be opened for writing.
    To access resources within subfolders use forward slashes as separator.


open_resource(resource, mode='rb')
    Opens a resource from the application's resource folder(root dir). 
    

open_session(request)
    Creates or opens a new session. 
    Default implementation stores all session data in a signed cookie. 
    This requires that the secret_key is set. 
    
    
permanent_session_lifetime
    A timedelta which is used to set the expiration date of a permanent session. 
    The default is 31 days which makes a permanent session survive for roughly one month.
    This attribute can also be configured from the config with the PERMANENT_SESSION_LIFETIME configuration key. 
    Defaults to timedelta(days=31)
    
preprocess_request()
    Called before the actual request dispatching 
    and will call each before_request() decorated function, 
    passing no arguments. 
    If any of these functions returns a value, it's handled 
    as if it was the return value from the view 
    and further request handling is stopped.
    This also triggers the url_value_preprocessor() functions 
    before the actual before_request() functions are called.
    
    
preserve_context_on_exception
    Returns the value of the PRESERVE_CONTEXT_ON_EXCEPTION configuration value in case it's set, 
    otherwise a sensible default is returned.



process_response(response)
    Can be overridden in order to modify the response object 
    before it's sent to the WSGI server. 
    By default this will call all the after_request() decorated functions.
    Returns:a new response object or the same, has to be an instance of response_class. 

    
propagate_exceptions
    Returns the value of the PROPAGATE_EXCEPTIONS configuration value in case it's set, otherwise a sensible default is returned.

register_blueprint(blueprint, **options)
    Registers a blueprint on the application.

register_error_handler(code_or_exception, f)
    Alternative error attach function to the errorhandler() decorator 
    that is more straightforward to use for non decorator usage.


request_class
    The class that is used for request objects
    alias of Request
    
request_context(environ)
    Creates a RequestContext from the given environment 
    and binds it to the current context.    
    environ – a WSGI environment 
    #Example usage:
    with app.request_context(environ):
        do_something_with(request)

response_class
    The class that is used for response objects. 
    alias of Response
    
route(rule, **options)
    A decorator that is used to register a view function for a given URL rule. 
    This does the same thing as add_url_rule() but is intended for decorator usage

    
    
run(host=None, port=None, debug=None, **options)
    Runs the application on a local development server.
    Do not use run() in a production setting. 
    
    
save_session(session, response)
    Saves the session if it needs updates. 

 
secret_key
    If a secret key is set, cryptographic components can use this to sign cookies and other things. 
    This attribute can also be configured from the config with the SECRET_KEY configuration key. Defaults to None.
    
    
select_jinja_autoescape(filename)
    Returns True if autoescaping should be active for the given template name. 
    If no template name is given, returns True.


send_file_max_age_default
    A timedelta which is used as default cache_timeout for the send_file() functions. The default is 12 hours.
    This attribute can also be configured from the config with the SEND_FILE_MAX_AGE_DEFAULT configuration key. 
    This configuration variable can also be set with an integer value used as seconds. 
    Defaults to timedelta(hours=12)
    
    
send_static_file(filename)
    Function used internally to send static files 
    from the static folder to the browser.


session_cookie_name
    The secure cookie uses this for the name of the session cookie.
    This attribute can also be configured from the config with the SESSION_COOKIE_NAME configuration key. 
    Defaults to 'session'
    
    
session_interface = <flask.sessions.SecureCookieSessionInterface object>
    the session interface to use. 
    By default an instance of SecureCookieSessionInterface is used here.



shell_context_processor(f)
    Registers a shell context processor function.

shell_context_processors = None
    A list of shell context processor functions that should be run when a shell context is created.


should_ignore_error(error)
    This is called to figure out if an error should be ignored 
    or not as far as the teardown system is concerned. 
    If this function returns True 
    then the teardown handlers will not be passed the error.


static_folder
    The absolute path to the configured static folder.
    

teardown_appcontext(f)
    Registers a function to be called when the application context ends. 
    These functions are typically also called when the request context is popped.
    The return values of teardown functions are ignored.



teardown_appcontext_funcs = None
    A list of functions that are called when the application context is destroyed. 
    

teardown_request(f)
    Register a function to be run at the end of each request, 
    regardless of whether there was an exception or not. 
    These functions are executed when the request context is popped, 
    even if not an actual request was performed.


teardown_request_funcs = None
    A dictionary with lists of functions that are called after each request, 
    even if an exception has occurred. 
    
        
template_context_processors = None
    A dictionary with list of functions that are called without argument 
    to populate the template context. 
    To register a function here, use the context_processor() decorator.
    
    
template_filter(name=None)
    A decorator that is used to register custom template filter.
    You can specify a name for the filter, 
    otherwise the function name will be used. 
    #Example:
    @app.template_filter()
    def reverse(s):
        return s[::-1]



template_global(name=None)
    A decorator that is used to register a custom template global function. 
    You can specify a name for the global function, 
    otherwise the function name will be used. 
    #Example:
    @app.template_global()
    def double(n):
        return 2 * n


template_test(name=None)
    A decorator that is used to register custom template test. 
    You can specify a name for the test, otherwise the function name will be used. 
    #Example:
    @app.template_test()
    def is_prime(n):
        if n == 2:
            return True
        for i in range(2, int(math.ceil(math.sqrt(n))) + 1):
            if n % i == 0:
                return False
        return True


test_client(use_cookies=True, **kwargs)
    Creates a test client for this application. 
    #Example     
    app.testing = True
    client = app.test_client()
    #The test client can be used in a with block to defer the closing down of the context 
    #until the end of the with block. 
    with app.test_client() as c:
        rv = c.get('/?vodka=42')
        assert request.args['vodka'] == '42'


test_request_context(*args, **kwargs)
    Creates a WSGI environment from the given values 
    (see werkzeug.test.EnvironBuilder for more information, 
    class werkzeug.test.EnvironBuilder(path='/', base_url=None, 
        query_string=None, method='GET', input_stream=None, 
        content_type=None, content_length=None, errors_stream=None, 
        multithread=False, multiprocess=False, run_once=False, 
        headers=None, data=None, environ_base=None, 
        environ_overrides=None, charset='utf-8')
    #http://werkzeug.pocoo.org/docs/0.13/test/#werkzeug.test.EnvironBuilder
    #Example 
    >>> from werkzeug.test import EnvironBuilder
    >>> from StringIO import StringIO
    >>> builder = EnvironBuilder(method='POST', data={'foo': 'this is some text',
    ...      'file': (StringIO('my file contents'), 'test.txt')})
    >>> env = builder.get_environ()
    #The resulting environment is a regular WSGI environment 
    #that can be used for further processing:
    >>> from werkzeug.wrappers import Request
    >>> req = Request(env)
    >>> req.form['foo']
    u'this is some text'
    >>> req.files['file']
    <FileStorage: u'test.txt' ('text/plain')>
    >>> req.files['file'].read()
    'my file contents'
    #Another example with form data 
    >>> builder = EnvironBuilder(method='POST', data={'foo': 'bar'})
    >>> builder.content_type
    'application/x-www-form-urlencoded'
    >>> builder.files['foo'] = StringIO('contents')
    >>> builder.content_type
    'multipart/form-data'
    #If a string is provided as data (or an input stream) 
    #you have to specify the content type yourself:
    >>> builder = EnvironBuilder(method='POST', data='{"json": "this is"}')
    >>> builder.content_type
    >>> builder.content_type = 'application/json'



testing
    The testing flag. 
    Set this to True to enable the test mode of Flask extensions 
    This attribute can also be configured from the config with the TESTING configuration key. 
    Defaults to False.
    
    
trap_http_exception(e)
    Checks if an HTTP exception should be trapped or not. 
    By default this will return False for all exceptions 
    except for a bad request key error 
    if TRAP_BAD_REQUEST_ERRORS is set to True. It also returns True if TRAP_HTTP_EXCEPTIONS is set to True.


update_template_context(context)
    Update the template context with some commonly used variables. 
    This injects request, session, config and g into the template context 
    as well as everything template context processors want to inject. 
    
    
url_build_error_handlers = None
    A list of functions that are called when url_for() raises a BuildError. 
    Each function registered here is called with error, endpoint and values. 
    If a function returns None or raises a BuildError the next function is tried.

url_default_functions = None
    A dictionary with lists of functions that can be used as URL value preprocessors. 
   

url_defaults(f)
    Callback function for URL defaults for all view functions of the application. 
    It's called with the endpoint and values 
    and should update the values passed in place.

url_map = None
    The Map for this instance. 
    You can use this to change the routing converters 
    after the class was created but before any routes are connected. 
    #Example 
    from werkzeug.routing import BaseConverter
    class ListConverter(BaseConverter):
        def to_python(self, value):
            return value.split(',')
        def to_url(self, values):
            return ','.join(super(ListConverter, self).to_url(value)
                            for value in values)

    app = Flask(__name__)
    app.url_map.converters['list'] = ListConverter

    
url_rule_class
    The rule object to use for URL rules created. 
    This is used by add_url_rule(). Defaults to werkzeug.routing.Rule.
    alias of Rule
    
    
url_value_preprocessor(f)
    Registers a function as URL value preprocessor 
    for all view functions of the application. 
    It's called before the view functions are called 
    and can modify the url values provided.
    
    
url_value_preprocessors = None
    A dictionary with lists of functions that can be used as URL value processor functions. 
    Whenever a URL is built these functions are called 
    to modify the dictionary of values in place. 
    
    
use_x_sendfile
    Enable this if you want to use the X-Sendfile feature. 
    Keep in mind that the server has to support this. 
    This only affects files sent with the send_file() method.
    This attribute can also be configured from the config with the USE_X_SENDFILE configuration key. 
    Defaults to False.

view_functions = None
    A dictionary of all view functions registered. 
    The keys will be function names which are also used to generate URLs 
    and the values are the function objects themselves. 
    To register a view function, use the route() decorator.

wsgi_app(environ, start_response)
    The actual WSGI application. 
    Note middlewares can be applied without losing a reference to the class. 
    app.wsgi_app = MyMiddleware(app.wsgi_app)


    
    

##Flask - Quickstart- A Minimal Application

#hello.py
from flask import Flask
app = Flask(__name__)

@app.route('/')
def hello_world():
    return 'Hello, World!'

#Execute 

$ export FLASK_APP=hello.py  #name of the module containing app 
$ flask run
 * Running on http://127.0.0.1:5000/

#Externally Visible Server:
$ flask run --host=0.0.0.0

#Debug Mode
$ export FLASK_DEBUG=1
$ flask run



##Flask - Quickstart- Routing

#/ calls index method 
@app.route('/')
def index():
    return 'Index Page'

@app.route('/hello')
def hello():
    return 'Hello, World'

#Variable Rules

@app.route('/user/<username>')
def show_user_profile(username):
    # show the user profile for that user
    return 'User %s' % username

@app.route('/post/<int:post_id>')
def show_post(post_id):
    # show the post with the given id, the id is an integer
    return 'Post %d' % post_id


#Unique URLs / Redirection Behavior:

@app.route('/projects/') #can be accessed as /projects 
def projects():
    return 'The project page'

@app.route('/about')  #can not be accessed as /about/
def about():
    return 'The about page'

#URL Building
#To build a URL to a specific function , use the url_for() function. 
I

>>> from flask import Flask, url_for
>>> app = Flask(__name__)
>>> @app.route('/')
    def index(): pass
...
>>> @app.route('/login')
    def login(): pass
...
>>> @app.route('/user/<username>')
    def profile(username): pass
...
>>> with app.test_request_context():
        print url_for('index')
        print url_for('login')
        print url_for('login', next='/')
        print url_for('profile', username='John Doe')
        
/
/login
/login?next=/
/user/John%20Doe


#HTTP Methods - default is GET and HEAD
from flask import request

@app.route('/login', methods=['GET', 'POST'])
def login():
    if request.method == 'POST':
        do_the_login()
    else:
        show_the_login_form()



##Flask - Quickstart- Static Files

#your web server is configured to serve them for you, 
#but during development Flask can do that as well. 

#create a folder called static in your package or next to  module 
#and it will be available at /static on the application.

#To generate URLs for static files, use the special 'static' endpoint name:
url_for('static', filename='style.css')
#The file has to be stored on the filesystem as static/style.css.


##Flask - Quickstart- Rendering Templates

from flask import render_template

@app.route('/hello/')
@app.route('/hello/<name>')
def hello(name=None):
    return render_template('hello.html', name=name)

#Flask will look for templates in the templates folder.
#Case 1: a module:


/application.py
/templates
    /hello.html

#Case 2: a package:


/application
    /__init__.py
    /templates
        /hello.html


#hello.html: jinja2 template 

<!doctype html>
<title>Hello from Flask</title>
{% if name %}
  <h1>Hello {{ name }}!</h1>
{% else %}
  <h1>Hello, World!</h1>
{% endif %}


#Inside templates , have access to the request, session and g objects ,get_flashed_messages() function.
#Automatic escaping is enabled, so if name contains HTML it will be escaped automatically. 
#or  mark it as safe by using the Markup class or by using the |safe filter 
>>> from flask import Markup
>>> Markup('<strong>Hello %s!</strong>') % '<blink>hacker</blink>'
Markup(u'<strong>Hello &lt;blink&gt;hacker&lt;/blink&gt;!</strong>')
>>> Markup.escape('<blink>hacker</blink>')
Markup(u'&lt;blink&gt;hacker&lt;/blink&gt;')
>>> Markup('<em>Marked up</em> &raquo; HTML').striptags()
u'Marked up \xbb HTML'


##Flask - Quickstart- Accessing Request Data
#for actual request handling, flask.request contains instance of flask.Request 
#for testing 

from flask import request

with app.test_request_context('/hello', method='POST'):
    # now you can do something with the request until the
    # end of the with block, such as basic assertions:
    assert request.path == '/hello'
    assert request.method == 'POST'


#or with a whole WSGI environment to the request_context() method:
from flask import request

with app.request_context(environ):
    assert request.method == 'POST'

#The current request method is available by using the method attribute. 
@app.route('/login', methods=['POST', 'GET'])
def login():
    error = None
    if request.method == 'POST':
        if valid_login(request.form['username'],
                       request.form['password']):
            return log_the_user_in(request.form['username'])
        else:
            error = 'Invalid username/password'
    # the code below is executed if the request method
    # was GET or the credentials were invalid
    return render_template('login.html', error=error)


#if the key does not exist in the form attribute? 
#a special KeyError is raised. 


#To access parameters submitted in the URL (?key=value) 
#use the args attribute:
searchword = request.args.get('key', '')


##Flask - Quickstart- File Uploads
#set the enctype="multipart/form-data" attribute on HTML form, 

#Uploaded files are stored in memory or at a temporary location on the filesystem. 

from flask import request

@app.route('/upload', methods=['GET', 'POST'])
def upload_file():
    if request.method == 'POST':
        f = request.files['the_file']  #<input type="file" name="the_file"
        f.save('/var/www/uploads/uploaded_file.txt')
    ...


#to know how the file was named on the client 
#pass it through the secure_filename() function 

from flask import request
from werkzeug.utils import secure_filename

@app.route('/upload', methods=['GET', 'POST'])
def upload_file():
    if request.method == 'POST':
        f = request.files['the_file']
        f.save('/var/www/uploads/' + secure_filename(f.filename))
    ...


    
##Flask - Quickstart- Cookies
#To access cookies , use the cookies attribute. 
#To set cookies , the set_cookie method of response objects. 

#The cookies attribute of request objects is a dictionary 
#with all the cookies the client transmits. 

#to use sessions, do not use the cookies directly 
#but instead use the Sessions in Flask that add some security on top of cookies 

from flask import request

@app.route('/')
def index():
    username = request.cookies.get('username')
    # use cookies.get(key) instead of cookies[key] to not get a
    # KeyError if the cookie is missing.


from flask import make_response

@app.route('/')
def index():
    resp = make_response(render_template(...))
    resp.set_cookie('username', 'the username')
    return resp


##Flask - Quickstart- Redirects and Errors
#To redirect a user to another endpoint, 
#use the redirect() function; 
#to abort a request early with an error code, use the abort() function:


from flask import abort, redirect, url_for

@app.route('/')
def index():
    return redirect(url_for('login'))

@app.route('/login')
def login():
    abort(401)
    this_is_never_executed()


#By default a black and white error page is shown for each error code. 
#to customize the error page, use the errorhandler() decorator:


from flask import render_template

@app.errorhandler(404)
def page_not_found(error):
    return render_template('page_not_found.html'), 404




##Flask - Quickstart- Sessions
#In order to use sessions , set a secret key
from flask import Flask, session, redirect, url_for, escape, request

app = Flask(__name__)

@app.route('/')
def index():
    if 'username' in session:
        return 'Logged in as %s' % escape(session['username'])
    return 'You are not logged in'

@app.route('/login', methods=['GET', 'POST'])
def login():
    if request.method == 'POST':
        session['username'] = request.form['username']
        return redirect(url_for('index'))
    return '''
        <form method="post">
            <p><input type=text name=username>
            <p><input type=submit value=Login>
        </form>
    '''

@app.route('/logout')
def logout():
    # remove the username from the session if it's there
    session.pop('username', None)
    return redirect(url_for('index'))

# set the secret key.  keep this really secret:
app.secret_key = 'A0Zr98j/3yX R~XHH!jmN]LWX/,?RT'

#How to generate good secret keys:
 
>>> import os
>>> os.urandom(24)
'\xfd{H\xe5<\x95\xf9\xe3\x96.5\xd1\x01O<!\xd5\xa2\xa0\x9fR"\xa1\xa8'
#take that thing and copy/paste it into  code 



##Flask - Quickstart- Message Flashing
#To flash a message use the flash() method, 
#to get hold of the messages , use get_flashed_messages() in templates 



##Flask - Quickstart- Logging


app.logger.debug('A value for debugging')
app.logger.warning('A warning occurred (%d apples)', 42)
app.logger.error('An error occurred')



##Flask - Quickstart- Flask Extensions
#Flask extensions extend the functionality of Flask in various different ways. 

#Flask extensions are listed on the Flask Extension Registry 
#and can be downloaded with easy_install or pip. 

#If you add a Flask extension as dependency to your requirements.txt 
#or setup.py file they are usually installed with a simple command 
#or when your application installs.






##Flask - Configuration Handling

app = Flask(__name__)
app.config['DEBUG'] = True
#Or - Certain configuration values are also forwarded to the Flask object
app.debug = True
#OR To update multiple keys at once you can use the dict.update() method:
app.config.update(
    DEBUG=True,
    SECRET_KEY='...'
)


#Builtin Configuration Values
#http://flask.pocoo.org/docs/0.12/config/#builtin-configuration-values
#The following configuration values are used internally by Flask:
DEBUG                           enable/disable debug mode 
TESTING                         enable/disable testing mode 
SECRET_KEY                      the secret key 
SESSION_COOKIE_NAME             the name of the session cookie 
SESSION_COOKIE_DOMAIN           the domain for the session cookie. If this is not set, the cookie will be valid for all subdomains of SERVER_NAME. 
SESSION_COOKIE_PATH             the path for the session cookie. If this is not set the cookie will be valid for all of APPLICATION_ROOT or if that is not set for '/'. 
SESSION_COOKIE_HTTPONLY         controls if the cookie should be set with the httponly flag. Defaults to True. 
SESSION_COOKIE_SECURE           controls if the cookie should be set with the secure flag. Defaults to False. 
PERMANENT_SESSION_LIFETIME      the lifetime of a permanent session as datetime.timedelta object. Starting with Flask 0.8 this can also be an integer representing seconds. 
SESSION_REFRESH_EACH_REQUEST    this flag controls how permanent sessions are refreshed. If set to True (which is the default) then the cookie is refreshed each request which automatically bumps the lifetime. If set to False a set-cookie header is only sent if the session is modified. Non permanent sessions are not affected by this. 
USE_X_SENDFILE                  enable/disable x-sendfile 
LOGGER_NAME                     the name of the logger 
LOGGER_HANDLER_POLICY           the policy of the default logging handler. The default is 'always' which means that the default logging handler is always active. 'debug' will only activate logging in debug mode, 'production' will only log in production and 'never' disables it entirely. 
SERVER_NAME                     the name and port number of the server. Required for subdomain support (e.g.: 'myapp.dev:5000') Note that localhost does not support subdomains so setting this to "localhost" does not help. Setting a SERVER_NAME also by default enables URL generation without a request context but with an application context. 
APPLICATION_ROOT                If the application does not occupy a whole domain or subdomain this can be set to the path where the application is configured to live. This is for session cookie as path value. If domains are used, this should be None. 
MAX_CONTENT_LENGTH              If set to a value in bytes, Flask will reject incoming requests with a content length greater than this by returning a 413 status code. 
SEND_FILE_MAX_AGE_DEFAULT       Default cache control max age to use with send_static_file() (the default static file handler) and send_file(), as datetime.timedelta or as seconds. Override this value on a per-file basis using the get_send_file_max_age() hook on Flask or Blueprint, respectively. Defaults to 43200 (12 hours). 
TRAP_HTTP_EXCEPTIONS            If this is set to True Flask will not execute the error handlers of HTTP exceptions but instead treat the exception like any other and bubble it through the exception stack. This is helpful for hairy debugging situations where you have to find out where an HTTP exception is coming from. 
TRAP_BAD_REQUEST_ERRORS         Werkzeug's internal data structures that deal with request specific data will raise special key errors that are also bad request exceptions. Likewise many operations can implicitly fail with a BadRequest exception for consistency. Since it's nice for debugging to know why exactly it failed this flag can be used to debug those situations. If this config is set to True you will get a regular traceback instead. 
PREFERRED_URL_SCHEME            The URL scheme that should be used for URL generation if no URL scheme is available. This defaults to http. 
JSON_AS_ASCII                   By default Flask serialize object to ascii-encoded JSON. If this is set to False Flask will not encode to ASCII and output strings as-is and return unicode strings. jsonify will automatically encode it in utf-8 then for transport for instance. 
JSON_SORT_KEYS                  By default Flask will serialize JSON objects in a way that the keys are ordered. This is done in order to ensure that independent of the hash seed of the dictionary the return value will be consistent to not trash external HTTP caches. You can override the default behavior by changing this variable. This is not recommended but might give you a performance improvement on the cost of cacheability. 
JSONIFY_PRETTYPRINT_REGULAR     If this is set to True (the default) jsonify responses will be pretty printed if they are not requested by an XMLHttpRequest object (controlled by the X-Requested-With header) 
JSONIFY_MIMETYPE                MIME type used for jsonify responses. 
TEMPLATES_AUTO_RELOAD           Whether to check for modifications of the template source and reload it automatically. By default the value is None which means that Flask checks original file only in debug mode. 
EXPLAIN_TEMPLATE_LOADING        If this is enabled then every attempt to load a template will write an info message to the logger explaining the attempts to locate the template. This can be useful to figure out why templates cannot be found or wrong templates appear to be loaded. 


#Configuring from Files
app = Flask(__name__)
app.config.from_object('yourapplication.default_settings') #loads the configuration from the yourapplication.default_settings module
app.config.from_envvar('YOURAPPLICATION_SETTINGS')  #then override with the contents of the file the YOURAPPLICATION_SETTINGS environment variable points to


# Example configuration
DEBUG = False
SECRET_KEY = '?\xbf,\xb4\x8d\xa3"<\x9c\xb0@\x0f5\xab,w\xee\x8d$0\x13\x8b83'


#to use classes and inheritance for configuration:

class Config(object):
    DEBUG = False
    TESTING = False
    DATABASE_URI = 'sqlite://:memory:'

class ProductionConfig(Config):
    DATABASE_URI = 'mysql://user@localhost/foo'

class DevelopmentConfig(Config):
    DEBUG = True

class TestingConfig(Config):
    TESTING = True


#To enable such a config you just have to call into from_object():
app.config.from_object('configmodule.ProductionConfig')




##Flask - Instance Folders - Flask.instance_path

#The instance folder is designed to not be under version control 
#and be deployment specific. 
#It's the perfect place to drop things that either change at runtime 
#or configuration files.


#this path must be absolute when provided
app = Flask(__name__, instance_path='/path/to/instance/folder')
#OR auto-configured as below (check from Flask.instance_path)
#For Uninstalled module - /instance 
/myapp.py
/instance
#For Uninstalled package:
/myapp
    /__init__.py
/instance
#For Installed module or package:(check $PREFIX from sys.prefix)
$PREFIX/lib/python2.X/site-packages/myapp
$PREFIX/var/myapp-instance


$PREFIX is the prefix of your Python installation. This can be /usr or the path to your virtualenv. You can print the value of sys.prefix to see what the prefix is set to.


# The behavior of relative paths in config files can be flipped between 
#"relative to the application root" (the default) 
#to "relative to instance folder" via the instance_relative_config 
app = Flask(__name__, instance_relative_config=True)
app.config.from_object('yourapplication.default_settings') 
app.config.from_pyfile('application.cfg', silent=True) #this is from instance folder 


# to open a file from the instance folder with Flask.open_instance_resource().
filename = os.path.join(app.instance_path, 'application.cfg')
with open(filename) as f:
    config = f.read()

# or via open_instance_resource:
with app.open_instance_resource('application.cfg') as f:
    config = f.read()

    
    
    


    

##Flask - Deployment Options
#http://flask.pocoo.org/docs/0.12/deploying/#deployment

#Hosted options
•Deploying Flask on Heroku
•Deploying Flask on OpenShift
•Deploying Flask on Webfaction
•Deploying Flask on Google App Engine
•Deploying Flask on AWS Elastic Beanstalk
•Sharing your Localhost Server with Localtunnel
•Deploying on Azure (IIS)
•Deploying on PythonAnywhere

#Self-hosted options
•mod_wsgi (Apache) 
•Standalone WSGI Containers 
    ◦Gunicorn
    ◦Gevent
    ◦Twisted Web
    ◦Proxy Setups
•uWSGI 
    ◦uwsgi
    ◦nginx
•FastCGI 
•CGI 


#Example - mod_wsgi (Apache)
#There should not be any app.run() calls 
#or  are inside in if __name__ == '__main__': block 

#Step.1 : Install mod_wsgi

#Step.2 : Creating a .wsgi file
#To run , create a yourapplication.wsgi file. 

from yourapplication import app as application

#Store that file somewhere (e.g.: /var/www/yourapplication) 
#and make sure that yourapplication and all the libraries are on the python load path. 
#Eg, Use setuptools and install in the system (or in virtual env)

#OR  patch the path in the .wsgi file before the import:
import sys
sys.path.insert(0, '/path/to/the/application')

#Step.3 : Application configuration file 
#update .wsgi file with 
import os
os.environ['YOURAPPLICATION_CONFIG'] = '/var/www/yourapplication/application.cfg'
from yourapplication import app as application

#then use yourapplication.py 
app = Flask(__name__)
app.config.from_object('yourapplication.default_config')
app.config.from_envvar('YOURAPPLICATION_CONFIG')


#Step.4 : Configuring Apache
#httpd.conf - unix 
<VirtualHost *>
    ServerName example.com

    WSGIDaemonProcess yourapplication user=user1 group=group1 threads=5
    WSGIScriptAlias / /var/www/yourapplication/yourapplication.wsgi

    <Directory /var/www/yourapplication>
        WSGIProcessGroup yourapplication
        WSGIApplicationGroup %{GLOBAL}
        Order deny,allow
        Allow from all
    </Directory>
</VirtualHost>


#httpd.conf - windows 

<VirtualHost *>
        ServerName example.com
        WSGIScriptAlias / C:\yourdir\yourapp.wsgi
        <Directory C:\yourdir>
                Order deny,allow
                Allow from all
        </Directory>
</VirtualHost>

#the syntax for directory permissions has changed from httpd 2.2
Order allow,deny
Allow from all
#to httpd 2.4 syntax
Require all granted



#Working with Virtual Environments

#Add the following lines to the top of your .wsgi file:
activate_this = '/path/to/env/bin/activate_this.py'
execfile(activate_this, dict(__file__=activate_this))

#For Python 3 add the following lines to the top of your .wsgi file:
activate_this = '/path/to/env/bin/activate_this.py'
with open(activate_this) as file_:
    exec(file_.read(), dict(__file__=activate_this))


#Troubleshooting
Problem: application does not run, errorlog shows SystemExit ignored
You have an app.run() call in your application file that is not guarded by an if __name__ == '__main__': condition. 


Problem: application gives permission errors
Probably caused by your application running as the wrong user. 
Make sure the folders the application needs access to have the proper privileges set and the application runs as the correct user 


Problem: application dies with an error on print
Keep in mind that mod_wsgi disallows doing anything with sys.stdout and sys.stderr. 
You can disable this protection from the config by setting the WSGIRestrictStdout to off:
WSGIRestrictStdout Off
#OR Alternatively you can also replace the standard out in the .wsgi file 
#with a different stream:
import sys
sys.stdout = sys.stderr


Problem: accessing resources gives IO errors
Your application probably is a single .py file you symlinked into the site-packages folder. 
Please be aware that this does not work, 
instead you either have to put the folder into the pythonpath 
the file is stored in, or convert your application into a package.



##Flask - Running a  interactive Python shell

$ flask shell

#This will start up an interactive Python shell, 
#setup the correct application context and setup the local variables in the shell. This is done by invoking the Flask.make_shell_context() method of the application. By default you have access to your app and g.


#Custom Commands - to add more commands to the shell script 
#Flask uses click for the command interface 

import click
from flask import Flask

app = Flask(__name__)

@app.cli.command()
def initdb():
    """Initialize the database."""
    click.echo('Init the db')



$ flask initdb
Init the db


##Custom Commands - with Application Context
#@app.cli.command() by default wraps the code with cli.with_appcontext()
#which is not true if a command is added later with add_command() or through other means.

#to disable it 
@app.cli.command(with_appcontext=False)
def example():
    pass




##Flask - Patterns for Flask 
#http://flask.pocoo.org/docs/0.12/patterns/
◦Larger Applications
◦Application Factories
◦Application Dispatching
◦Implementing API Exceptions
◦Using URL Processors
◦Deploying with Setuptools
◦Deploying with Fabric
◦Using SQLite 3 with Flask
◦SQLAlchemy in Flask
◦Uploading Files
◦Caching
◦View Decorators
◦Form Validation with WTForms
◦Template Inheritance
◦Message Flashing
◦AJAX with jQuery
◦Custom Error Pages
◦Lazily Loading Views
◦MongoKit in Flask
◦Adding a favicon
◦Streaming Contents
◦Deferred Request Callbacks
◦Adding HTTP Method Overrides
◦Request Content Checksums
◦Celery Based Background Tasks
◦Subclassing Flask



##Flask - Patterns for Flask  - Using SQLite 3 with Flask


import sqlite3
from flask import g

DATABASE = '/path/to/database.db'

def get_db():
    db = getattr(g, '_database', None)
    if db is None:
        db = g._database = sqlite3.connect(DATABASE)
    return db

@app.teardown_appcontext
def close_connection(exception):
    db = getattr(g, '_database', None)
    if db is not None:
        db.close()


#to use the database, 
#the application must either have an active application context 
#(which is always true if there is a request handling) 
#or create an application context 'with app.app_context()'


@app.route('/')
def index():
    cur = get_db().cursor()
    ...
#or 
with app.app_context():
    # now you can use get_db()



#To simplify working with SQLite, a row factory function is useful.
#It is executed for every result returned from the database to convert the result. 

#For instance, in order to get dictionaries instead of tuples, 
#this could be inserted into the get_db function 

def make_dicts(cursor, row):
    return dict((cursor.description[idx][0], value)
                for idx, value in enumerate(row))

db.row_factory = make_dicts

#OR in get_db , use 
#This would use Row objects rather than dicts to return the results of queries
#so access them either by index or by key(columnName)
db.row_factory = sqlite3.Row


#to provide a query function that combines getting the cursor, 
#executing and fetching the results:

def query_db(query, args=(), one=False):
    cur = get_db().execute(query, args)
    rv = cur.fetchall()
    cur.close()
    return (rv[0] if rv else None) if one else rv


#in combination with a row factory
for user in query_db('select * from users'):
    print user['username'], 'has the id', user['user_id']

#Or if you just want a single result:
user = query_db('select * from users where username = ?',
                [the_username], one=True)
if user is None:
    print 'No such user'
else:
    print the_username, 'has the id', user['user_id']


#Initial Schemas

def init_db():
    with app.app_context():
        db = get_db()
        with app.open_resource('schema.sql', mode='r') as f:
            db.cursor().executescript(f.read())
        db.commit()

#You can then create such a database from the Python shell:
>>> from yourapplication import init_db
>>> init_db()



##Flask - Patterns for Flask  - AJAX with jQuery

#Loading jQuery
<script type=text/javascript src="{{
  url_for('static', filename='jquery.js') }}"></script>

#OR using Google's AJAX Libraries API to load jQuery:

<script src="//ajax.googleapis.com/ajax/libs/jquery/1.9.1/jquery.min.js"></script>
<script>window.jQuery || document.write('<script src="{{
  url_for('static', filename='jquery.js') }}">\x3C/script>')</script>

#To know , script root in client site 

<script type=text/javascript>
  $SCRIPT_ROOT = {{ request.script_root|tojson|safe }};
</script>

#Example 
#jqueryexample.py 
# -*- coding: utf-8 -*-

from flask import Flask, jsonify, render_template, request
app = Flask(__name__)


@app.route('/_add_numbers')
def add_numbers():
    """Add two numbers server side, ridiculous but well..."""
    a = request.args.get('a', 0, type=int)
    b = request.args.get('b', 0, type=int)
    return jsonify(result=a + b)


@app.route('/')
def index():
    return render_template('index.html')

if __name__ == '__main__':
    app.run()

#templates/index.html
{% extends "layout.html" %}
{% block body %}
<script type="text/javascript">
  $(function() {
    var submit_form = function(e) {
      $.getJSON($SCRIPT_ROOT + '/_add_numbers', {
        a: $('input[name="a"]').val(),
        b: $('input[name="b"]').val()
      }, function(data) {
        $('#result').text(data.result);
        $('input[name=a]').focus().select();
      });
      return false;
    };

    $('a#calculate').bind('click', submit_form);

    $('input[type=text]').bind('keydown', function(e) {
      if (e.keyCode == 13) {
        submit_form(e);
      }
    });

    $('input[name=a]').focus();
  });
</script>
<h1>jQuery Example</h1>
<p>
  <input type="text" size="5" name="a"> +
  <input type="text" size="5" name="b"> =
  <span id="result">?</span>
<p><a href=# id="calculate">calculate server side</a>
{% endblock %}

#templates/layout.html
<!doctype html>
<title>jQuery Example</title>
<script type="text/javascript"
  src="http://ajax.googleapis.com/ajax/libs/jquery/1.4.2/jquery.min.js"></script>
<script type="text/javascript">
  var $SCRIPT_ROOT = {{ request.script_root|tojson|safe }};
</script>
{% block body %}{% endblock %}





##Flask - Patterns for Flask  - Uploading Files

1.A <form> tag is marked with enctype=multipart/form-data 
  and an <input type=file> is placed in that form.
2.The application accesses the file from the 'files' dictionary 
  on the request object.
3.use the save() method of the file to save the file permanently somewhere on the filesystem.


#Example 

import os
from flask import Flask, request, redirect, url_for
from werkzeug.utils import secure_filename

UPLOAD_FOLDER = '/path/to/the/uploads'
ALLOWED_EXTENSIONS = set(['txt', 'pdf', 'png', 'jpg', 'jpeg', 'gif'])

app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER


def allowed_file(filename):
    return '.' in filename and \
           filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

@app.route('/', methods=['GET', 'POST'])
def upload_file():
    if request.method == 'POST':
        # check if the post request has the file part
        if 'file' not in request.files:
            flash('No file part')
            return redirect(request.url)
        file = request.files['file']
        # if user does not select file, browser also
        # submit a empty part without filename
        if file.filename == '':
            flash('No selected file')
            return redirect(request.url)
        if file and allowed_file(file.filename):
            filename = secure_filename(file.filename)
            file.save(os.path.join(app.config['UPLOAD_FOLDER'], filename))
            return redirect(url_for('uploaded_file',
                                    filename=filename))
    return '''
    <!doctype html>
    <title>Upload new File</title>
    <h1>Upload new File</h1>
    <form method=post enctype=multipart/form-data>
      <p><input type=file name=file>
         <input type=submit value=Upload>
    </form>
    '''


##Flask - Patterns for Flask  - Message Flashing

from flask import Flask, flash, redirect, render_template, \
     request, url_for

app = Flask(__name__)
app.secret_key = 'some_secret'

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/login', methods=['GET', 'POST'])
def login():
    error = None
    if request.method == 'POST':
        if request.form['username'] != 'admin' or \
                request.form['password'] != 'secret':
            error = 'Invalid credentials'
        else:
            flash('You were successfully logged in')
            return redirect(url_for('index'))
    return render_template('login.html', error=error)


#layout.html 
<!doctype html>
<title>My Application</title>
{% with messages = get_flashed_messages() %}
  {% if messages %}
    <ul class=flashes>
    {% for message in messages %}
      <li>{{ message }}</li>
    {% endfor %}
    </ul>
  {% endif %}
{% endwith %}
{% block body %}{% endblock %}


#index.html 
{% extends "layout.html" %}
{% block body %}
  <h1>Overview</h1>
  <p>Do you want to <a href="{{ url_for('login') }}">log in?</a>
{% endblock %}


#login.html
{% extends "layout.html" %}
{% block body %}
  <h1>Login</h1>
  {% if error %}
    <p class=error><strong>Error:</strong> {{ error }}
  {% endif %}
  <form method=post>
    <dl>
      <dt>Username:
      <dd><input type=text name=username value="{{
          request.form.username }}">
      <dt>Password:
      <dd><input type=password name=password>
    </dl>
    <p><input type=submit value=Login>
  </form>
{% endblock %}



#Flashing With Categories
#setting 
flash(u'Invalid password provided', 'error')

#getting 
{% with messages = get_flashed_messages(with_categories=true) %}
  {% if messages %}
    <ul class=flashes>
    {% for category, message in messages %}
      <li class="{{ category }}">{{ message }}</li>
    {% endfor %}
    </ul>
  {% endif %}
{% endwith %}


#Filtering Flash Messages
#Get only error 
{% with errors = get_flashed_messages(category_filter=["error"]) %}
{% if errors %}
<div class="alert-message block-message error">
  <a class="close" href="#">×</a>
  <ul>
    {%- for msg in errors %}
    <li>{{ msg }}</li>
    {% endfor -%}
  </ul>
</div>
{% endif %}
{% endwith %}



##Flask - Patterns for Flask  - Streaming Contents
#to send an enormous amount of data to the client

from flask import Response

@app.route('/large.csv')
def generate_large_csv():
    def generate():
        for row in iter_all_rows():
            yield ','.join(row) + '\n'
    return Response(generate(), mimetype='text/csv')


#Streaming from Templates


from flask import Response

def stream_template(template_name, **context):
    app.update_template_context(context)
    t = app.jinja_env.get_template(template_name)
    rv = t.stream(context)
    rv.enable_buffering(5)
    return rv

@app.route('/my-large-page.html')
def render_large_template():
    rows = iter_all_rows()
    return Response(stream_template('the_template.html', rows=rows))


#Streaming with Context
#Note that when you stream data, 
#the request context is already gone the moment the function executes.

from flask import stream_with_context, request, Response

@app.route('/stream')
def streamed_response():
    def generate():
        yield 'Hello '
        yield request.args['name']
        yield '!'
    return Response(stream_with_context(generate()))


    
##Flask - Patterns for Flask  - Larger Applications
#use a package instead of a module. 

#a small application looks like this:
/yourapplication
    yourapplication.py
    /static
        style.css
    /templates
        layout.html
        index.html
        login.html
        ...


#convert to large application 
#rename yourapplication.py to __init__.py 
/yourapplication
    setup.py
    /yourapplication
        __init__.py
        /static
            style.css
        /templates
            layout.html
            index.html
            login.html
            ...


#To install this package 
#setup.py 

from setuptools import setup

setup(
    name='yourapplication',
    packages=['yourapplication'],
    include_package_data=True,
    install_requires=[
        'flask',
    ],
)
#then 
$ pip install -e .
$ export FLASK_APP=yourapplication
$ export FLASK_DEBUG=true
$ flask run


#More refactoring 
/yourapplication
    setup.py
    /yourapplication
        __init__.py
        views.py
        /static
            style.css
        /templates
            layout.html
            index.html
            login.html
            ...

            
# __init__.py:
from flask import Flask
app = Flask(__name__)

import yourapplication.views

#views.py

from yourapplication import app

@app.route('/')
def index():
    return 'Hello World!'


